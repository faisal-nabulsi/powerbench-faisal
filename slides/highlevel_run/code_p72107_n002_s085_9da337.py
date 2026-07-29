from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize the Presentation
prs = Presentation()

# Use a blank layout (index 6) for custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Coordinates: Left, Top, Width, Height (in inches)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p_title = title_frame.add_paragraph()
p_title.text = "Types of Blockchain"
p_title.font.size = Pt(36)
p_title.font.bold = True

# --- Add Content ---
# Text box for definitions
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Definitions data
blockchain_types = [
    ("Private Blockchain", "Restricted access network where only invited participants can join and view the ledger. It offers higher privacy and transaction speeds."),
    ("Public Blockchain", "A decentralized network open to everyone. Anyone can read, write, and participate in the consensus process (e.g., Bitcoin, Ethereum)."),
    ("Permissioned Blockchain", "Access is controlled by a permission provider. It is often used in enterprise consortiums to balance decentralization with governance."),
    ("Hybrid Blockchain", "A combination of public and private blockchains. It allows data to be shared selectively, offering flexibility for security and privacy needs.")
]

# Add bullet points
for name, definition in blockchain_types:
    p = content_frame.add_paragraph()
    p.font.size = Pt(20)
    p.space_after = Pt(14)  # Spacing between bullets
    
    # Add a bullet character
    run_bullet = p.add_run()
    run_bullet.text = "• "
    
    # Add the Blockchain Type (Bold)
    run_name = p.add_run()
    run_name.text = f"{name}: "
    run_name.font.bold = True
    
    # Add the Definition (Normal)
    run_def = p.add_run()
    run_def.text = definition
    run_def.font.bold = False

# Save the presentation
prs.save('output.pptx')