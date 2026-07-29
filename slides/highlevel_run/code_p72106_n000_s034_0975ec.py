from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
# Use blank layout (index 6 is standard Blank in default Office templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Slide Title
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
title_tf = title_shape.text_frame
title_tf.text = "Brief highlights of Elon Musk"
title_tf.paragraphs[0].font.size = Pt(36)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add horizontal timeline line
timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(0.05))
timeline.fill.solid()
timeline.fill.fore_color.rgb = RGBColor(0, 51, 102)
timeline.line.fill.background()

# Milestones data: (x_position, year, description)
milestones = [
    (1.5, "1971", "Born in Pretoria,\nSouth Africa"),
    (3.0, "1999", "Co-founded X.com\n(later PayPal)"),
    (4.5, "2002", "Co-founded SpaceX"),
    (6.0, "2004", "Co-founded\nTesla Motors"),
    (7.5, "2020", "Tesla becomes most\nvaluable automaker")
]

# Render milestones on the timeline
for x_pos, year, desc in milestones:
    # Vertical connector line from timeline to node
    conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos - 0.015), Inches(2.55), Inches(0.03), Inches(0.225))
    conn.fill.solid()
    conn.fill.fore_color.rgb = RGBColor(0, 51, 102)
    conn.line.fill.background()

    # Timeline node (circle)
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos - 0.25), Inches(2.275), Inches(0.5), Inches(0.5))
    node.fill.solid()
    node.fill.fore_color.rgb = RGBColor(255, 140, 0)
    node.line.fill.solid()
    node.line.fill.fore_color.rgb = RGBColor(255, 165, 0)
    node.line.width = Pt(2)

    # Year label
    year_tf = slide.shapes.add_textbox(Inches(x_pos - 0.4), Inches(3.1), Inches(0.8), Inches(0.4))
    year_tf.text_frame.word_wrap = True
    year_tf.text_frame.text = year
    year_tf.text_frame.paragraphs[0].font.size = Pt(18)
    year_tf.text_frame.paragraphs[0].font.bold = True
    year_tf.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    year_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description label
    desc_tf = slide.shapes.add_textbox(Inches(x_pos - 0.4), Inches(3.6), Inches(0.8), Inches(1.0))
    desc_tf.text_frame.word_wrap = True
    desc_tf.text_frame.text = desc
    desc_tf.text_frame.paragraphs[0].font.size = Pt(11)
    desc_tf.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    desc_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')