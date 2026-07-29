from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation object
prs = Presentation()

# Get the slide dimensions (default is usually 10x7.5 inches)
slide_width = prs.slide_width
slide_height = prs.slide_height

# Add a blank slide layout
# Index 6 is typically the 'Blank' layout in default templates
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the background image
# We add the picture first so it sits at the back of the stack (background)
# The image fills the entire slide
slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)

# 2. Add the Title Text "ELO's"
# Position: Top center
left_title = Inches(0)
top_title = Inches(0.5)
width_title = slide_width
height_title = Inches(1.5)

txBox_title = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
tf_title = txBox_title.text_frame
tf_title.word_wrap = True

p_title = tf_title.paragraphs[0]
p_title.text = "ELO's"
p_title.font.size = Pt(54)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(255, 255, 255)  # White color for contrast against image
p_title.alignment = PP_ALIGN.CENTER

# 3. Add the Questions
# Position: Centered below the title
left_q = Inches(1.5)
top_q = Inches(2.5)
width_q = Inches(8)
height_q = Inches(3)

txBox_q = slide.shapes.add_textbox(left_q, top_q, width_q, height_q)
tf_q = txBox_q.text_frame
tf_q.word_wrap = True

# Question 1
p1 = tf_q.paragraphs[0]
p1.text = "1. What is market?"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)  # White color
p1.alignment = PP_ALIGN.CENTER

# Question 2
p2 = tf_q.add_paragraph()
p2.text = "2. How product reach to market?"
p2.font.size = Pt(32)
p2.font.bold = True
p2.font.color.rgb = RGBColor(255, 255, 255)  # White color
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)  # Add some vertical spacing between the two questions

# Save the presentation to 'output.pptx'
prs.save('output.pptx')