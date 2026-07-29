from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_tf = title_box.text_frame
title_tf.text = "PART 2"
title_tf.paragraphs[0].font.size = Pt(36)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

# Instruction text
instr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
instr_tf = instr_box.text_frame
instr_tf.text = "Read the provided dialogue excerpts and categorize the verb examples into past simple and past continuous tenses in the table below."
p = instr_tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(51, 51, 51)

# Dialogue excerpts box
dia_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.5))
dia_tf = dia_box.text_frame
dia_tf.word_wrap = True
dia_text = (
    '"A: What did you do yesterday?"  "B: I studied while my sister was preparing dinner."\n'
    '"A: Why was he late?"  "B: He missed the bus because it left while he was running."\n'
    '"A: Did you enjoy the concert?"  "B: Yes, the band played while the crowd was singing."\n'
    '"A: What happened?"  "B: She dropped her coffee because the bag fell while she was walking."\n'
    '"A: Were you busy last night?"  "B: Yes, I worked while my roommate was watching TV."'
)
dia_tf.text = dia_text
for p in dia_tf.paragraphs:
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(70, 70, 70)

# Table for categorization
rows = 6
cols = 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.5), Inches(9), Inches(4.0))
table = table_shape.table

table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(2.25)
table.columns[2].width = Inches(2.25)

headers = ["Dialogue Excerpt", "Past Simple", "Past Continuous"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Populate table rows
dialogues = [
    '"I studied while my sister was preparing dinner."',
    '"He missed the bus because it left while he was running."',
    '"The band played while the crowd was singing."',
    '"She dropped her coffee because the bag fell while she was walking."',
    '"I worked while my roommate was watching TV."'
]

for idx, dialogue in enumerate(dialogues, 1):
    # Dialogue excerpt cell
    cell0 = table.cell(idx, 0)
    cell0.text = dialogue
    for p in cell0.text_frame.paragraphs:
        p.font.size = Pt(10)
    
    # Leave categorization columns empty for user input
    table.cell(idx, 1).text = "________"
    table.cell(idx, 2).text = "________"
    for col in range(1, 3):
        for p in table.cell(idx, col).text_frame.paragraphs:
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

prs.save('output.pptx')