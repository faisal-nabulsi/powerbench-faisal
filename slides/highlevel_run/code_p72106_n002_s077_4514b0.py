from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Add a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Add Quote on the Left ---
# Position and size (left, top, width, height)
left_quote = Inches(1.5)
top_quote = Inches(1.5)
width_quote = Inches(4.5)
height_quote = Inches(5.0)

text_box = slide.shapes.add_textbox(left_quote, top_quote, width_quote, height_quote)
tf = text_box.text_frame
tf.word_wrap = True

# Add the main quote text
p = tf.paragraphs[0]
p.text = '"When something is important enough, you do it even if the odds are not in your favor."'
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.LEFT

# Add the attribution text
p_attr = tf.add_paragraph()
p_attr.text = "- Elon Musk"
p_attr.font.size = Pt(18)
p_attr.font.italic = True
p_attr.space_before = Pt(12)
p_attr.alignment = PP_ALIGN.LEFT

# --- 2. Add Image on the Right ---
# Position and size for the image (using 'image.png' as requested)
left_img = Inches(6.5)
top_img = Inches(1.5)
width_img = Inches(3.5)
height_img = Inches(5.0)

slide.shapes.add_picture('image.png', left_img, top_img, width_img, height_img)

# --- 3. Add "Thank You!" Message at the Bottom ---
# Position and size for the footer text box
left_footer = Inches(1.5)
top_footer = Inches(6.8)
width_footer = Inches(7.0)
height_footer = Inches(0.5)

footer_box = slide.shapes.add_textbox(left_footer, top_footer, width_footer, height_footer)
tf_footer = footer_box.text_frame
tf_footer.word_wrap = True

p_footer = tf_footer.paragraphs[0]
p_footer.text = "Thank You!"
p_footer.font.size = Pt(30)
p_footer.font.bold = True
p_footer.alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')