from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Add a blank slide layout (Index 6 is typically Blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Main Title
# Position: Top of the slide
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "5 Stages of Development of Media"
title_para.font.size = Pt(36)
title_para.font.bold = True

# 2. Add Subtitle (Focus Area)
# Position: Below the main title
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "2. SCRIPTS"
subtitle_para.font.size = Pt(32)
subtitle_para.font.bold = True

# 3. Add Bullet Points (Left Side)
# Position: Left side, below the subtitle
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5), Inches(3))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Bullet 1
p1 = content_frame.paragraphs[0]
p1.text = "• Enhances clarity and organization of the message."
p1.font.size = Pt(20)

# Bullet 2
p2 = content_frame.add_paragraph()
p2.text = "• Allows for precise timing and delivery of content."
p2.font.size = Pt(20)

# 4. Add Image (Right Side)
# Position: Right side, aligned with the text box
# Using the placeholder image 'image.png' as requested
slide.shapes.add_picture('image.png', Inches(6), Inches(2.5), Inches(3.5), Inches(3))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')