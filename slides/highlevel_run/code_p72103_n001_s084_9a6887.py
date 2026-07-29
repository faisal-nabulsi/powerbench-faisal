from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Add a blank slide (Index 6 is typically the 'Blank' layout in standard themes)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Define geometry for the title text box
title_left = Inches(0.5)
title_top = Inches(0.5)
title_width = Inches(9)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf_title = title_box.text_frame
tf_title.text = "Importance of Skills Beyond Academics"

# Format the title paragraph
title_paragraph = tf_title.paragraphs[0]
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.name = 'Arial'
title_paragraph.alignment = PP_ALIGN.LEFT

# --- Add Content ---
# Define geometry for the content text box
content_left = Inches(0.5)
content_top = Inches(1.75)
content_width = Inches(9)
content_height = Inches(5)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
tf_content = content_box.text_frame
tf_content.word_wrap = True

# List of skills and their brief descriptions
skills_data = [
    ("Effective Communication", "The ability to convey ideas clearly and listen actively to others."),
    ("Problem Solving and Critical Thinking", "Analyzing situations logically to identify solutions and make informed decisions."),
    ("Emotional Intelligence", "Understanding and managing one's own emotions and empathizing with others."),
    ("Adaptability and Resilience", "The capacity to adjust to new conditions and recover quickly from difficulties.")
]

# Populate the content text box with skills and descriptions
for skill_name, description in skills_data:
    # Add a new paragraph for each skill entry
    p = tf_content.add_paragraph()
    p.space_after = Pt(14)  # Add space after each item
    p.font.name = 'Arial'
    
    # Add the skill name in bold
    run_skill = p.add_run()
    run_skill.text = f"{skill_name}: "
    run_skill.font.bold = True
    run_skill.font.size = Pt(22)
    
    # Add the description in regular font
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = Pt(18)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')