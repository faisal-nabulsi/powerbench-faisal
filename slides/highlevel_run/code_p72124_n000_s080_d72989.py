from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation object
prs = Presentation()

# Add a blank slide
# In standard templates, index 6 is usually the blank layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Title "QUIZ & REVIEW ACTIVITY"
# Positioning: Left 1", Top 0.5", Width 8", Height 1.5" (Centered horizontally)
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_paragraph = title_frame.paragraphs[0]
title_paragraph.text = "QUIZ & REVIEW ACTIVITY"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.alignment = PP_ALIGN.CENTER

# 2. Add an Icon representing a quiz
# Using 'image.png' as the placeholder image
# Positioning: Centered below the title
icon_left = Inches(3.75)
icon_top = Inches(2.5)
icon_width = Inches(2.5)
icon_height = Inches(2.5)
slide.shapes.add_picture('image.png', icon_left, icon_top, icon_width, icon_height)

# 3. Add the Instruction "Work with a partner"
# Positioning: Below the icon
instruction_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
instruction_frame = instruction_box.text_frame
instruction_frame.word_wrap = True

instruction_paragraph = instruction_frame.paragraphs[0]
instruction_paragraph.text = "Work with a partner"
instruction_paragraph.font.size = Pt(36)
instruction_paragraph.font.italic = True
instruction_paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')