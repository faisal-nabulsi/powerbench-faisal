from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Add a blank slide. Index 6 is typically the blank layout in the default template.
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define the dimensions for the text box
# We use 8 inches width and 2 inches height for a prominent title
box_width = Inches(8)
box_height = Inches(2)

# Calculate the position to center the text box on the slide
# prs.slide_width and prs.slide_height return dimensions in EMUs (English Metric Units)
# We cast to int to ensure we pass integer EMU values to add_textbox
left = int((prs.slide_width - box_width) / 2)
top = int((prs.slide_height - box_height) / 2)

# Add the text box to the slide
text_box = slide.shapes.add_textbox(left, top, box_width, box_height)
text_frame = text_box.text_frame

# Set the text content
text_frame.text = "Thank You!"

# Format the text to be prominent and centered
paragraph = text_frame.paragraphs[0]
paragraph.alignment = PP_ALIGN.CENTER
paragraph.font.size = Pt(60)
paragraph.font.bold = True

# Save the presentation to 'output.pptx'
prs.save('output.pptx')