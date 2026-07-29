from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation instance
prs = Presentation()

# Add a blank slide layout (usually index 6, or access via name if available)
# Using index 6 is standard for 'Blank' in most default presentations.
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Slide Title ---
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.2))
title_frame = title_shape.text_frame
title_frame.word_wrap = True

title_para = title_frame.add_paragraph()
title_para.text = "Academic World and Professional World: Demands and Characteristics"
title_para.font.size = Pt(24)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue title

# --- Helper Function to Create Columns ---
def create_column(slide, left, top, width, height, title, sections_data, title_color):
    # Create text box for the column
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Add Column Header (e.g., "Academic World")
    p_header_col = tf.add_paragraph()
    p_header_col.text = title
    p_header_col.font.size = Pt(18)
    p_header_col.font.bold = True
    p_header_col.font.color.rgb = title_color

    # Iterate through the sections (Goals, Focus, etc.)
    for section_name, points in sections_data.items():
        # Section Sub-header
        p_sub_header = tf.add_paragraph()
        p_sub_header.text = section_name
        p_sub_header.font.size = Pt(14)
        p_sub_header.font.bold = True
        p_sub_header.font.color.rgb = RGBColor(50, 50, 50)  # Dark Grey

        # Bullet points for each section
        for point in points:
            p_bullet = tf.add_paragraph()
            p_bullet.text = point
            p_bullet.font.size = Pt(12)
            p_bullet.level = 0
            p_bullet.space_before = Pt(2)
            p_bullet.space_after = Pt(2)
            # Note: python-pptx renders bullets based on paragraph style. 
            # Explicit indentation or symbols can be added, but clean text is used here.

# --- Data Preparation ---

# Content for Academic World
academic_data = {
    "Goals": [
        "Acquire theoretical knowledge",
        "Develop critical thinking skills",
        "Earn degrees and certifications",
        "Conduct academic research"
    ],
    "Focus": [
        "Understanding 'How' and 'Why'",
        "Individual performance (Grades)",
        "Long-term learning process",
        "Depth of subject matter"
    ],
    "Setting": [
        "Classrooms, Lecture Halls, Seminars",
        "Libraries and Laboratories",
        "Structured academic calendar",
        "Controlled learning environment"
    ],
    "Atmosphere": [
        "Intellectual curiosity",
        "Collaborative peer learning",
        "Safe space for trial and error",
        "Feedback-oriented"
    ]
}

# Content for Professional World
professional_data = {
    "Goals": [
        "Achieve organizational KPIs",
        "Maximize efficiency and profit",
        "Career advancement",
        "Solve real-world business problems"
    ],
    "Focus": [
        "Practical application ('What')",
        "Team output and results",
        "Meeting deadlines and ROI",
        "Continuous improvement"
    ],
    "Setting": [
        "Offices and Corporate campuses",
        "Client sites and Field work",
        "Dynamic/Flexible schedules",
        "High-stakes environment"
    ],
    "Atmosphere": [
        "Performance-driven",
        "Competitive",
        "Accountability for outcomes",
        "Fast-paced decision making"
    ]
}

# --- Add Columns to Slide ---

# Left Column: Academic World
create_column(
    slide, 
    left=Inches(0.5), top=Inches(1.6), width=Inches(4.3), height=Inches(5.5),
    title="Academic World",
    sections_data=academic_data,
    title_color=RGBColor(0, 112, 192)  # Blue
)

# Right Column: Professional World
create_column(
    slide, 
    left=Inches(5.2), top=Inches(1.6), width=Inches(4.3), height=Inches(5.5),
    title="Professional World",
    sections_data=professional_data,
    title_color=RGBColor(192, 0, 0)  # Red
)

# Save the presentation
prs.save('output.pptx')