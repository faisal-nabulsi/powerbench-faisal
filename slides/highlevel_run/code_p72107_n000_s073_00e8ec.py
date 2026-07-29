from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

def set_cell_fill(cell, color_hex):
    """Helper function to set the background color of a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)

# Create a presentation object
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the slide title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "Academic World and Professional World: Demands and Characteristics"
title_paragraph.font.size = Pt(24)
title_paragraph.font.bold = True
title_paragraph.alignment = PP_ALIGN.CENTER

# Define the data for the table
# Structure: (Section Name, Academic Content, Professional Content)
table_data = [
    (
        "Evaluation and Feedback",
        "• Graded exams and assignments\n• Detailed rubrics for grading\n• Peer reviews",
        "• Performance reviews and KPIs\n• Client satisfaction scores\n• 360-degree feedback"
    ),
    (
        "Assessment",
        "• Tests, quizzes, and midterms\n• Final projects and theses\n• Standardized testing",
        "• Project outcomes and deliverables\n• Sales targets and revenue\n• Efficiency and productivity metrics"
    ),
    (
        "Feedback",
        "• Often delayed (weeks)\n• Written comments on papers\n• Focus on learning and correction",
        "• Immediate and ongoing\n• Verbal and written communication\n• Focus on results and improvement"
    ),
    (
        "Timeframe",
        "• Semesters or academic years\n• Fixed schedules and deadlines\n• Long-term progression",
        "• Project deadlines and milestones\n• Fiscal quarters and annual goals\n• Flexible and urgent priorities"
    )
]

# Add a table to the slide
# 5 rows (1 header + 4 data rows), 2 columns
rows = 5
cols = 2
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(4.5)

# Populate and style the header row
header_cell_0 = table.cell(0, 0)
header_cell_0.text = "Academic World"
header_cell_1 = table.cell(0, 1)
header_cell_1.text = "Professional World"

for cell in [header_cell_0, header_cell_1]:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    set_cell_fill(cell, '4472C4')  # Blue background for headers

# Populate and style the data rows
for i, (section_name, academic_text, professional_text) in enumerate(table_data):
    row_idx = i + 1
    
    # Academic World Column
    cell_acad = table.cell(row_idx, 0)
    cell_acad.text = f"{section_name}\n{academic_text}"
    
    # Professional World Column
    cell_prof = table.cell(row_idx, 1)
    cell_prof.text = f"{section_name}\n{professional_text}"
    
    # Style the cells
    for cell in [cell_acad, cell_prof]:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.space_after = Pt(4)
            # Bold the section name (first paragraph)
            if paragraph.text == section_name:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)

# Save the presentation
prs.save('output.pptx')