from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Disadvantages of Blockchain"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# 2. Add Bullet Points (Left Side)
left_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
tf = left_text_box.text_frame
tf.text = "Key Disadvantages:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True

disadvantages = [
    "Scalability Issues: Slow transaction speeds compared to traditional systems.",
    "High Energy Consumption: Proof-of-Work mechanisms require significant power.",
    "Regulatory Uncertainty: Lack of clear legal frameworks in many jurisdictions.",
    "Irreversibility: Transactions cannot be undone once confirmed.",
    "Complexity: Difficult for average users to understand and manage."
]

for item in disadvantages:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.space_after = Pt(10)
    p.level = 0

# 3. Add Diagram Contrasting Disadvantages vs Advantages (Right Side)

# Box 1: Disadvantages (Summary)
box1_left = Inches(5.5)
box1_top = Inches(1.5)
box1_width = Inches(2)
box1_height = Inches(2.5)

shape1 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    box1_left, box1_top, box1_width, box1_height
)
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(220, 53, 69) # Red
shape1.line.color.rgb = RGBColor(255, 255, 255)

tf1 = shape1.text_frame
tf1.word_wrap = True
tf1.text = "Disadvantages"
tf1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tf1.paragraphs[0].font.bold = True
tf1.paragraphs[0].font.size = Pt(14)
tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

p1 = tf1.add_paragraph()
p1.text = "• Slow\n• Expensive\n• Complex"
p1.font.size = Pt(10)
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.alignment = PP_ALIGN.CENTER

# Box 2: Advantages (Summary)
box2_left = Inches(7.8)
box2_top = Inches(1.5)
box2_width = Inches(2)
box2_height = Inches(2.5)

shape2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    box2_left, box2_top, box2_width, box2_height
)
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(40, 167, 69) # Green
shape2.line.color.rgb = RGBColor(255, 255, 255)

tf2 = shape2.text_frame
tf2.word_wrap = True
tf2.text = "Advantages"
tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.size = Pt(14)
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

p2 = tf2.add_paragraph()
p2.text = "• Secure\n• Transparent\n• Decentralized"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.alignment = PP_ALIGN.CENTER

# VS Circle in the middle
vs_left = Inches(7.25)
vs_top = Inches(2.35)
vs_width = Inches(0.8)
vs_height = Inches(0.8)

shape_vs = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    vs_left, vs_top, vs_width, vs_height
)
shape_vs.fill.solid()
shape_vs.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape_vs.line.color.rgb = RGBColor(0, 0, 0)

tf_vs = shape_vs.text_frame
tf_vs.word_wrap = True
tf_vs.text = "VS"
tf_vs.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
tf_vs.paragraphs[0].font.bold = True
tf_vs.paragraphs[0].font.size = Pt(12)
tf_vs.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add a footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
tf_footer = footer_box.text_frame
tf_footer.text = "Comparison of Blockchain Pros and Cons"
tf_footer.paragraphs[0].font.size = Pt(10)
tf_footer.paragraphs[0].font.italic = True
tf_footer.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# Save the presentation
prs.save('output.pptx')