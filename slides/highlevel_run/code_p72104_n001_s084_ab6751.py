from pptx import Presentation
from pptx.util import Inches, Pt

def create_blockchain_slide():
    # Create a Presentation object
    prs = Presentation()

    # Select the layout for a slide with a Title and Content
    # Layout index 1 is typically 'Title and Content' in default templates
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Set the title of the slide
    title = slide.shapes.title
    title.text = "Types of Blockchain"

    # Access the content placeholder
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    # Clear any default text
    tf.clear()

    # Define the blockchain types and their definitions
    blockchain_data = [
        ("Public Blockchain", "A decentralized network where anyone can join, read, write, and audit transactions without permission (e.g., Bitcoin)."),
        ("Private Blockchain", "A restricted network controlled by a single organization. Only invited users can access and validate transactions."),
        ("Permissioned Blockchain", "A hybrid model managed by a consortium of organizations. Access is restricted, but control is shared among a group."),
        ("Hybrid Blockchain", "Combines features of public and private blockchains. It keeps some data public while maintaining privacy for sensitive data.")
    ]

    # Add bullet points for each type and definition
    for term, definition in blockchain_data:
        # Main bullet: The Type of Blockchain
        p = tf.add_paragraph()
        p.text = term
        p.level = 0  # Main level
        # Bold the term for emphasis
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(20)

        # Sub-bullet: The Definition
        p_def = tf.add_paragraph()
        p_def.text = definition
        p_def.level = 1  # Sub-level indent
        # Adjust font size for readability
        p_def.runs[0].font.size = Pt(18)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_blockchain_slide()