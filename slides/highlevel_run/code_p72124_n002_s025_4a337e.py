from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_thank_you_slide():
    # Initialize the Presentation
    prs = Presentation()
    
    # Get slide dimensions (usually 10x7.5 inches for default presentation)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Add a blank slide (Index 6 is typically the Blank layout)
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Set Background to White
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 2. Add a Subtle Green Border
    # We create a rectangle shape covering the entire slide to act as a border frame
    border_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        slide_width, slide_height
    )
    
    # Set fill to transparent (background) so only the line is visible
    border_shape.fill.background()
    
    # Configure the border line
    border_shape.line.color.rgb = RGBColor(0, 128, 0) # Green
    border_shape.line.width = Pt(3) # Subtle thickness
    border_shape.shadow.inherit = False # Remove any default shadows
    
    # 3. Add "THANK YOU" Text in Prominent Green Font
    # Determine text box dimensions and center position
    text_box_width = Inches(8.5)
    text_box_height = Inches(2.5)
    left = (slide_width - text_box_width) / 2
    top = (slide_height - text_box_height) / 2
    
    textbox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # Set text content and formatting
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(80) # Prominent size
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0) # Prominent Green
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    prs.save('output.pptx')
    print("Presentation saved to output.pptx")

if __name__ == "__main__":
    create_thank_you_slide()