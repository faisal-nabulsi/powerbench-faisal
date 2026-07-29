from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "PART 3"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
title_frame = title_box.text_frame
title_frame.text = "PART 3"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 2. Add Instruction Prompt
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
instruction_frame = instruction_box.text_frame
instruction_frame.text = "Choose true (T), false (F), or not given (N) based on video content."
instruction_frame.paragraphs[0].font.size = Pt(18)
instruction_frame.paragraphs[0].font.italic = True
instruction_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 3. Add Conversation Text
conversation_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(5.5), Inches(4.0))
conversation_frame = conversation_box.text_frame
conversation_frame.word_wrap = True

# Use the first paragraph for the header
p_header = conversation_frame.paragraphs[0]
p_header.text = "Conversation:"
p_header.font.size = Pt(20)
p_header.font.bold = True
p_header.space_after = Pt(12)

# Dialogue lines
dialogue = [
    ("Joey:", "How you doin'?"),
    ("Chandler:", "Joey, we're trying to study."),
    ("Joey:", "But I'm hungry. Can we get pizza?"),
    ("Monica:", "Only if you pay for it this time."),
    ("Joey:", "Fine, but I'm getting the extra cheese.")
]

for speaker, text in dialogue:
    p = conversation_frame.add_paragraph()
    p.text = f"{speaker} {text}"
    p.font.size = Pt(16)
    p.space_after = Pt(8)

# 4. Add Image of Joey on the right side
# Using the placeholder 'image.png' as requested
slide.shapes.add_picture('image.png', Inches(6.5), Inches(3.0), Inches(3.5), Inches(4.0))

# Save the presentation
prs.save('output.pptx')