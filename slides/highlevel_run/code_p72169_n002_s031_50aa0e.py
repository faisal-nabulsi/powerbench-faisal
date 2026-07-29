from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create a blank presentation
    prs = Presentation()
    
    # Add a blank slide (Layout index 6 is usually blank)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- 1. Add Slide Title ---
    # Position: Centered top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Academic World and Professional World: Demands and Characteristics"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80) # Dark Slate Blue
    p.alignment = PP_ALIGN.CENTER

    # --- 2. Content Data ---
    # Data for the Academic World column
    academic_data = {
        "Goals": ["Acquire knowledge and skills", "Earn degrees and certifications", "Develop critical thinking"],
        "Focus": ["Theoretical understanding", "Learning process and methodology", "Curriculum-based evaluation"],
        "Setting": ["Classrooms, labs, and libraries", "Structured semesters", "Guided by faculty"],
        "Atmosphere": ["Collaborative and supportive", "Experimentation encouraged", "Low stakes for failure"]
    }

    # Data for the Professional World column
    professional_data = {
        "Goals": ["Solve practical business problems", "Generate value and profit", "Career advancement"],
        "Focus": ["Practical application", "Results, output, and efficiency", "Client and market needs"],
        "Setting": ["Workplace, office, or remote", "Project-based timelines", "Driven by organizational goals"],
        "Atmosphere": ["Competitive and results-oriented", "High accountability", "Adaptability required"]
    }

    # --- 3. Visualization Configuration ---
    col_width = Inches(4.0)
    col1_left = Inches(0.6)
    col2_left = Inches(4.9) # 0.6 + 4.0 + 0.3 gap
    start_top = Inches(1.3)
    section_height = Inches(1.4)
    
    # Colors
    ACADEMIC_COLORS = {
        "header": RGBColor(41, 128, 185),   # Blue
        "text":   RGBColor(40, 40, 40)
    }
    PROFESSIONAL_COLORS = {
        "header": RGBColor(211, 84, 0),     # Dark Orange
        "text":   RGBColor(40, 40, 40)
    }

    # --- 4. Drawing Helpers ---

    def draw_column_header(slide, left, top, width, text, color):
        box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = box.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    def draw_section(slide, left, top, width, section_name, bullets, color_scheme):
        # Draw Section Header
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.25))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = section_name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color_scheme["header"]

        # Draw Bullet Points
        bullet_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.3), width - Inches(0.15), Inches(0.9))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        
        for bullet_text in bullets:
            p = tf.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(12)
            p.font.color.rgb = color_scheme["text"]
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    # --- 5. Render Layout ---
    
    # 5a. Academic World Column
    y_pos = start_top
    draw_column_header(slide, col1_left, y_pos, col_width, "ACADEMIC WORLD", ACADEMIC_COLORS["header"])
    y_pos += 0.5 # Space after header
    
    # Ensure consistent section order
    sections = ["Goals", "Focus", "Setting", "Atmosphere"]
    
    for sec_name in sections:
        draw_section(slide, col1_left, y_pos, col_width, sec_name, academic_data[sec_name], ACADEMIC_COLORS)
        y_pos += section_height

    # 5b. Professional World Column
    y_pos = start_top
    draw_column_header(slide, col2_left, y_pos, col_width, "PROFESSIONAL WORLD", PROFESSIONAL_COLORS["header"])
    y_pos += 0.5 

    for sec_name in sections:
        draw_section(slide, col2_left, y_pos, col_width, sec_name, professional_data[sec_name], PROFESSIONAL_COLORS)
        y_pos += section_height

    # --- 6. Save Presentation ---
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()