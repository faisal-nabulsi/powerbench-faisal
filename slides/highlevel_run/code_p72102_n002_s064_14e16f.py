from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a new presentation
prs = Presentation()

# Add a blank slide to have full control over layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a colorful background
# Using a vibrant blue (RGB: 30, 144, 255) - DodgerBlue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(30, 144, 255)

# 1. Add Title
# Position: Centered at top
left = Inches(0.5)
top = Inches(0.5)
width = Inches(10)
height = Inches(1.2)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_tf = title_box.text_frame
title_tf.word_wrap = True

p_title = title_tf.add_paragraph()
p_title.text = "Innovations (Future Products or Services)"
p_title.font.bold = True
p_title.font.size = Pt(34)
p_title.font.color.rgb = RGBColor(255, 255, 255) # White text
p_title.alignment = PP_ALIGN.CENTER

# 2. Add Description of Airbnb's 2021 Innovations
# Position: Upper middle
desc_left = Inches(0.5)
desc_top = Inches(2.0)
desc_width = Inches(10.5)
desc_height = Inches(2.5)
desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
desc_tf = desc_box.text_frame
desc_tf.word_wrap = True

p_desc_header = desc_tf.add_paragraph()
p_desc_header.text = "Airbnb's 2021 Innovations"
p_desc_header.font.bold = True
p_desc_header.font.size = Pt(22)
p_desc_header.font.color.rgb = RGBColor(255, 255, 255)
p_desc_header.space_after = Pt(10)

p_desc_body = desc_tf.add_paragraph()
# Using explicit bullet characters for safe rendering
p_desc_body.text = (
    "• Airbnb Experiences: Focused on local, in-person activities as travel resumed.\n"
    "• Airbnb Ads: Launched a native advertising tool to help hosts drive bookings.\n"
    "• Reference Links: Enabled users to share specific booking links easily.\n"
    "• Remote Work: Enhanced 'Live there' features for longer stays and digital nomads."
)
p_desc_body.font.size = Pt(16)
p_desc_body.font.color.rgb = RGBColor(255, 255, 255)

# 3. Add Quote from Brian Chesky
# Creating a white rounded rectangle for the quote to ensure readability against the colorful background
quote_left = Inches(1.0)
quote_top = Inches(5.0)
quote_width = Inches(10)
quote_height = Inches(1.8)

quote_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, 
    quote_left, quote_top, quote_width, quote_height
)
quote_bg.fill.solid()
quote_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
quote_bg.line.color.rgb = RGBColor(200, 200, 200)

# Add text to the quote shape
quote_tf = quote_bg.text_frame
quote_tf.word_wrap = True
quote_tf.auto_size = None

p_quote_1 = quote_tf.add_paragraph()
p_quote_1.text = "\"Our mission is to help create a world where anyone can belong anywhere.\""
p_quote_1.font.size = Pt(18)
p_quote_1.font.italic = True
p_quote_1.font.color.rgb = RGBColor(0, 0, 0)
p_quote_1.alignment = PP_ALIGN.CENTER
p_quote_1.space_after = Pt(5)

p_quote_2 = quote_tf.add_paragraph()
p_quote_2.text = "\"We are building the tools and experiences to make that true.\""
p_quote_2.font.size = Pt(18)
p_quote_2.font.italic = True
p_quote_2.font.color.rgb = RGBColor(0, 0, 0)
p_quote_2.alignment = PP_ALIGN.CENTER
p_quote_2.space_before = Pt(0)

p_attr = quote_tf.add_paragraph()
p_attr.text = "— Brian Chesky, CEO & Co-Founder"
p_attr.font.size = Pt(14)
p_attr.font.bold = True
p_attr.font.color.rgb = RGBColor(30, 144, 255) # Matching background color
p_attr.alignment = PP_ALIGN.CENTER
p_attr.space_before = Pt(10)

# Save the presentation
prs.save('output.pptx')