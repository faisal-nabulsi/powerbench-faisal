from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a Presentation object
prs = Presentation()

# Add a blank slide (index 6 is typically Blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title: "PART 4"
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = "PART 4"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# 2. Add Instruction
instruction_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
instruction_frame = instruction_shape.text_frame
instruction_frame.text = "List the events in Joey’s second date with Sarah in order from 2 - 6"
instruction_frame.paragraphs[0].font.size = Pt(18)

# 3. Add Checklist of Events
# We create a list of events. The first one is marked as completed (example).
# Since the specific story events aren't provided, we use generic date events.
events = [
    ("1. Arrived at the restaurant", True),
    ("2. Ordered appetizers", False),
    ("3. Had a conversation", False),
    ("4. Ordered main courses", False),
    ("5. Walked outside", False),
    ("6. Said goodnight", False)
]

y_position = Inches(3.0)
for event_text, is_completed in events:
    # Select checkbox symbol based on completion status
    symbol = "☑" if is_completed else "☐"
    full_text = f"{symbol} {event_text}"
    
    # Create a text box for each item
    item_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(6), Inches(0.5))
    item_frame = item_shape.text_frame
    item_frame.text = full_text
    item_frame.paragraphs[0].font.size = Pt(16)
    item_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Move down for the next item
    y_position += Inches(0.6)

# 4. Add Image of Joey at the bottom right
# Placeholder image 'image.png' is used as specified.
image_left = Inches(7)
image_top = Inches(5.5)
image_width = Inches(3)
image_height = Inches(3)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')