from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation
prs = Presentation()

# Add a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Title "QUIZ & REVIEW ACTIVITY"
# Position: Centered horizontally, near the top
left_title = Inches(1)
top_title = Inches(1)
width_title = Inches(8)
height_title = Inches(1.5)

title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
tf_title = title_box.text_frame
tf_title.word_wrap = True

# Configure title text properties
p_title = tf_title.paragraphs[0]
p_title.text = "QUIZ & REVIEW ACTIVITY"
p_title.font.size = Pt(34)
p_title.font.bold = True
p_title.alignment = PP_ALIGN.CENTER

# 2. Add Icon representing a quiz
# Using the placeholder 'image.png' as requested for images
# Position: Centered horizontally, in the middle of the slide
left_icon = Inches(3.5)
top_icon = Inches(2.8)
width_icon = Inches(3)
height_icon = Inches(3)

slide.shapes.add_picture('image.png', left_icon, top_icon, width_icon, height_icon)

# 3. Add the Instruction "Work with a partner"
# Position: Below the title and icon, centered
left_inst = Inches(1)
top_inst = Inches(6.2)
width_inst = Inches(8)
height_inst = Inches(1)

inst_box = slide.shapes.add_textbox(left_inst, top_inst, width_inst, height_inst)
tf_inst = inst_box.text_frame
tf_inst.word_wrap = True

# Configure instruction text properties
p_inst = tf_inst.paragraphs[0]
p_inst.text = "Work with a partner"
p_inst.font.size = Pt(26)
p_inst.font.italic = True
p_inst.alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')