from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 1. Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 (Standard Widescreen)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Add a Blank Slide
# Index 6 is typically the 'Blank' layout in standard pptx templates
try:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
except IndexError:
    # Fallback to the first available layout if 'Blank' is not found
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

# 3. Create Title Shape
# Position: Top Left
title_box = slide.shapes.add_textbox(
    left=Inches(0.5), 
    top=Inches(0.5), 
    width=Inches(12.333), 
    height=Inches(1.2)
)
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Add Title Text
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "Blockchain in Cryptocurrency"
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_paragraph.alignment = PP_ALIGN.LEFT

# 4. Create Text Box for Content (Bullet Points)
# Position: Middle Left
text_box = slide.shapes.add_textbox(
    left=Inches(0.5), 
    top=Inches(2.0), 
    width=Inches(7.5), 
    height=Inches(5.0)
)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Define the bullet points content
bullet_points = [
    "Definition: A decentralized digital ledger that records transactions across multiple computers securely.",
    "Application: Core technology for cryptocurrencies, smart contracts, and immutable record-keeping.",
    "Governance: Maintained by a network of participants using consensus mechanisms, not a central authority.",
    "Differences from Cryptocurrency: Blockchain is the underlying infrastructure; cryptocurrency is the digital asset.",
    "Decentralization: Distributes data and control across a peer-to-peer network, enhancing security and transparency."
]

# Populate text box with bullets
for point in bullet_points:
    paragraph = text_frame.add_paragraph()
    # Using a unicode bullet character for visual consistency
    paragraph.text = "● " + point
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(50, 50, 50)
    paragraph.space_after = Pt(10)

# 5. Add Image Placeholder
# Position: Middle Right
# Using 'image.png' as the placeholder as requested
# We assume a 4:3 or similar aspect ratio, or just stretch to fit the box
image_shape = slide.shapes.add_picture(
    'image.png', 
    left=Inches(8.5), 
    top=Inches(2.0), 
    width=Inches(4.5), 
    height=Inches(5.0)
)

# 6. Save the presentation
prs.save('output.pptx')

print("Presentation saved to 'output.pptx'")