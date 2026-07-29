from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation object
prs = Presentation()

# Add a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Quote from Elon Musk (Left Side) ---
left = Inches(1)
top = Inches(2)
width = Inches(5)
height = Inches(4)

text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Set the quote text
text_frame.text = '"The beginning of knowledge is the discovery of something we do not understand."'
text_frame.paragraphs[0].font.size = Pt(24)
text_frame.paragraphs[0].font.name = 'Arial'
text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add attribution
attr_para = text_frame.add_paragraph()
attr_para.text = "- Elon Musk"
attr_para.font.size = Pt(18)
attr_para.font.name = 'Arial'
attr_para.alignment = PP_ALIGN.LEFT
attr_para.space_before = Pt(12)

# --- 2. Image of Elon Musk (Right Side) ---
img_left = Inches(6.5)
img_top = Inches(1.5)
img_width = Inches(3)
img_height = Inches(4.5)

# Add the picture using the available placeholder 'image.png'
slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# --- 3. "Thank You!" Message (Bottom) ---
ty_left = Inches(2.5)
ty_top = Inches(6.8)
ty_width = Inches(5)
ty_height = Inches(0.6)

ty_box = slide.shapes.add_textbox(ty_left, ty_top, ty_width, ty_height)
ty_tf = ty_box.text_frame

ty_tf.text = "Thank You!"
ty_tf.paragraphs[0].font.size = Pt(32)
ty_tf.paragraphs[0].font.bold = True
ty_tf.paragraphs[0].font.name = 'Arial'
ty_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')