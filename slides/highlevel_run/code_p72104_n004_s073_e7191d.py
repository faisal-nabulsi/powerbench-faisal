from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_slide():
    # Create a new presentation
    prs = Presentation()

    # Add a blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add the title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Innovations (Future Products or Services)"
    
    # Format the title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.CENTER

    # Add the image placeholder
    # The instruction specifies using 'image.png' for the requested visual
    image_left = Inches(1.0)
    image_top = Inches(2.5)
    image_width = Inches(8.0)
    image_height = Inches(4.5)
    
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()