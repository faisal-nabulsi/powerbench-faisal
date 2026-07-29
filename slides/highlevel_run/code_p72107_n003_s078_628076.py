import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_elon_musk_bio_slide():
    # Create a blank presentation
    prs = Presentation()
    
    # Select a blank layout (usually index 6) to arrange elements freely
    # Standard layouts in a new presentation: 0-Title, 1-Title+Content, ..., 6-Blank
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # --- 1. Add Title ---
    # Position: Top center
    title_left = Inches(0.5)
    title_top = Inches(0.4)
    title_width = Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.clear()  # Remove default empty paragraph
    
    title_para = title_frame.add_paragraph()
    title_para.text = "Early life of Elon Musk"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0) # Black text
    
    # --- 2. Add Portrait Image ---
    # Position: Right side
    img_left = Inches(6.0)
    img_top = Inches(1.6)
    img_width = Inches(3.5)
    img_height = Inches(4.5)
    
    # Add the picture using the required filename
    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
    
    # --- 3. Add Biography Details (Bullet Points) ---
    # Position: Left side
    text_left = Inches(0.5)
    text_top = Inches(1.6)
    text_width = Inches(5.2)
    text_height = Inches(5.0)
    
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Key details content
    biography_details = [
        "Birth: Born on June 28, 1971, in Pretoria, South Africa.",
        "Family: Son of Maye Haldeman (mother) and Errol Musk (father).",
        "Education: Attended Queen's University; earned degrees in Physics and Economics from the University of Pennsylvania.",
        "Early Interests: Created a video game ('Blastar') at age 12; avid reader of sci-fi and encyclopedias."
    ]
    
    # Text frames start with one paragraph. We modify the first one and add the rest.
    # We prepend a bullet character '• ' to ensure visual bullets appear 
    # even in a blank layout without specific theme bullets.
    
    # Process first item using existing paragraph
    p_first = text_frame.paragraphs[0]
    p_first.text = "• " + biography_details[0]
    p_first.font.size = Pt(20)
    p_first.space_after = Pt(10)
    
    # Process remaining items
    for detail in biography_details[1:]:
        p = text_frame.add_paragraph()
        p.text = "• " + detail
        p.font.size = Pt(20)
        p.space_after = Pt(10)
        
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_elon_musk_bio_slide()