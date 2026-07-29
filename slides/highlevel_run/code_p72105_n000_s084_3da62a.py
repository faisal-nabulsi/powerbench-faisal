from pptx import Presentation
from pptx.util import Pt

# Initialize the presentation
prs = Presentation()

# Use the standard Title and Content layout (index 1)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set the slide title
title_shape = slide.shapes.title
title_shape.text = "Importance of Skills Beyond Academics"

# Access the content placeholder
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.clear()

# Define skills and brief descriptions
skills_data = [
    ("Effective Communication", "Conveying ideas clearly and actively listening to foster strong interpersonal relationships and collaboration."),
    ("Problem Solving and Critical Thinking", "Analyzing situations objectively to identify challenges, evaluate options, and implement logical solutions."),
    ("Emotional Intelligence", "Recognizing and managing personal emotions while empathizing with others to navigate social complexities effectively."),
    ("Adaptability and Resilience", "Adjusting swiftly to change and recovering positively from setbacks to maintain progress and achieve goals.")
]

# Populate the content placeholder with formatted runs
for skill, desc in skills_data:
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    
    # Skill name in bold
    run_skill = p.add_run()
    run_skill.text = f"{skill}: "
    run_skill.font.bold = True
    run_skill.font.size = Pt(18)
    
    # Description text
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(18)

# Save the presentation
prs.save('output.pptx')