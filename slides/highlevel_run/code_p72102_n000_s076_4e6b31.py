from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Add a blank slide (Index 6 is typically the 'Blank' layout in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Main Title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_shape.text_frame
title_frame.word_wrap = True
title_para = title_frame.add_paragraph()
title_para.text = "5 Stages of Development of Media"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# 2. Add the Subtitle/Focus
subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
subtitle_frame = subtitle_shape.text_frame
subtitle_frame.word_wrap = True
subtitle_para = subtitle_frame.add_paragraph()
subtitle_para.text = "4. ELECTRONIC MEDIA"
subtitle_para.font.size = Pt(28)
subtitle_para.font.bold = True
subtitle_para.alignment = PP_ALIGN.CENTER

# 3. Add Bullet Points (Left Side)
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.5), Inches(3.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Bullet 1: Components
p1 = text_frame.add_paragraph()
p1.text = "• Components: Includes radio, television, and digital platforms that transmit information via electronic signals."
p1.font.size = Pt(18)
p1.space_after = Pt(12)

# Bullet 2: Significance
p2 = text_frame.add_paragraph()
p2.text = "• Significance: Revolutionized communication by enabling instantaneous, global dissemination of news and entertainment."
p2.font.size = Pt(18)

# 4. Add Image (Right Side)
# The placeholder image 'image.png' is used here
try:
    slide.shapes.add_picture('image.png', Inches(5.5), Inches(2.5), Inches(4), Inches(3.5))
except FileNotFoundError:
    pass

# Save the presentation
prs.save('output.pptx')