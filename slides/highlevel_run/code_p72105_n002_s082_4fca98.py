from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize the presentation
prs = Presentation()

# Use a blank slide layout (typically index 6) to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Slide Title
# Positioned at the top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_run = title_frame.add_paragraph()
title_run.text = "Advantages of Blockchain"
title_run.font.size = Pt(34)
title_run.font.bold = True

# 2. List Key Advantages as Bullet Points
# Positioned on the left side of the slide
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

advantages = [
    "Decentralization: No single point of failure or central control.",
    "Transparency: Transaction history is visible to all participants.",
    "Immutability: Records cannot be altered once written.",
    "Security: Cryptographic techniques ensure data integrity.",
    "Traceability: Provides a complete audit trail of activities."
]

for i, point in enumerate(advantages):
    if i == 0:
        p = left_frame.paragraphs[0]
    else:
        p = left_frame.add_paragraph()
    
    p.text = point
    p.font.size = Pt(16)
    p.space_after = Pt(10)

# 3. Include a Diagram (Placeholder Image)
# Positioned on the right side alongside the text
# The configuration assumes 'image.png' is available in the current directory
try:
    slide.shapes.add_picture(
        'image.png',
        Inches(5.3), Inches(1.8),
        Inches(4.2), Inches(5)
    )
except Exception:
    # In a real scenario, ensure the image file exists.
    pass

# Save the presentation to the specified file name
prs.save('output.pptx')