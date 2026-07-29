from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize a new presentation
prs = Presentation()

# Add a blank slide layout to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title Section ---
# Create a title text box
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True

# Add title text and formatting
title_para = title_tf.add_paragraph()
title_para.text = "Early life of Elon Musk"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue

# --- 2. Image Section ---
# Add the placeholder portrait image
# Placed on the right side of the slide
image_left = Inches(6.5)
image_top = Inches(2.0)
image_width = Inches(3.2)
image_height = Inches(4.5)

try:
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)
except Exception as e:
    print(f"Error adding image: {e}")

# --- 3. Content Section ---
# Create a text box for bullet points
# Placed on the left side of the slide
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Helper function to add a section header
def add_header(tf, text):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(10)

# Helper function to add a bullet point
def add_bullet(tf, text):
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.space_after = Pt(8)

# --- Add Biography Details ---

# Birth
add_header(content_tf, "Birth")
add_bullet(content_tf, "Born on June 28, 1971, in Pretoria, South Africa.")

# Family
add_header(content_tf, "Family")
add_bullet(content_tf, "Father: Errol Musk (Electrical Engineer)")
add_bullet(content_tf, "Mother: Maye Musk (Model)")
add_bullet(content_tf, "Siblings: Kimbal Musk and Tosca Musk")

# Education
add_header(content_tf, "Education")
add_bullet(content_tf, "High School: Pretoria Boys High School")
add_bullet(content_tf, "University: University of Pennsylvania")
add_bullet(content_tf, "Degree: B.S. in Physics and Economics")

# Early Interests
add_header(content_tf, "Early Interests")
add_bullet(content_tf, "Programmed a video game ('Bustin' Out') at age 12")
add_bullet(content_tf, "Read entire encyclopedias and the Star Wars trilogy")
add_bullet(content_tf, "Developed an early passion for computers and business")

# Save the presentation to 'output.pptx'
prs.save('output.pptx')