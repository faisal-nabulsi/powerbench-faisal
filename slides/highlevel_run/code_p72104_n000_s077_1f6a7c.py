from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Add a blank slide layout (index 6 is typically blank in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_shape.text_frame
title_frame.text = "Conclusion"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add Content Text
content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(3.5))
content_frame = content_shape.text_frame
content_frame.word_wrap = True

# Text summarizing importance of top grades
p1 = content_frame.add_paragraph()
p1.text = "Achieving top grades is a significant milestone that validates hard work and opens doors to future academic and professional opportunities."
p1.font.size = Pt(18)
p1.space_after = Pt(14)

# Text emphasizing balance and personal development
p2 = content_frame.add_paragraph()
p2.text = "However, true success requires more than just academic achievement. It is crucial to maintain a healthy balance between studies and personal development."
p2.font.size = Pt(18)
p2.space_after = Pt(14)

# Text on holistic growth
p3 = content_frame.add_paragraph()
p3.text = "Prioritizing well-being, extracurricular engagement, and holistic growth ensures a fulfilling journey and prepares you for the complexities of life beyond the classroom."
p3.font.size = Pt(18)

# Add Image
# Placeholder image 'image.png' representing a person walking towards success
left = Inches(6.5)
top = Inches(1.5)
width = Inches(3)
height = Inches(3)
slide.shapes.add_picture('image.png', left, top, width, height)

# Save the presentation
prs.save('output.pptx')