from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Add a blank slide to the presentation
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "PART 4"
# Position: Top center
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "PART 4"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.alignment = PP_ALIGN.CENTER

# 2. Add Instruction Text
# Position: Below title
instr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(1))
instr_frame = instr_box.text_frame
instr_frame.word_wrap = True
instr_paragraph = instr_frame.add_paragraph()
instr_paragraph.text = "List the events in Joey’s second date with Sarah in order from 2 - 6"
instr_paragraph.font.size = Pt(24)
instr_paragraph.font.color.rgb = RGBColor(50, 50, 50)

# 3. Add Checklist of Events
# The first event is marked as completed (checked) as an example.
# We use dummy events to populate the list.
events = [
    "Met Sarah at the cafe entrance",
    "Ordered drinks",
    "Watched the band perform",
    "Danced on the dancefloor",
    "Had dessert outside",
    "Said goodnight at the car"
]

checklist_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(6), Inches(4))
checklist_frame = checklist_box.text_frame
checklist_frame.word_wrap = True

for index, event_text in enumerate(events):
    paragraph = checklist_frame.add_paragraph()
    
    if index == 0:
        # First event: Completed example
        # Using Unicode check box with check: ☑
        paragraph.text = f"1. ☑ {event_text}"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(100, 100, 100) 
        paragraph.font.bold = False
    else:
        # Remaining events: Unchecked
        # Using Unicode empty check box: ☐
        paragraph.text = f"{index + 1}. ☐ {event_text}"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    paragraph.space_after = Pt(12)

# 4. Add Image of Joey (Placeholder) at the bottom right
# Slide dimensions are typically 10 inches wide by 7.5 inches high.
# We place the image at (7.5, 5.5) with size 2x2 to fit in the bottom right corner.
slide.shapes.add_picture('image.png', Inches(7.5), Inches(5.5), Inches(2), Inches(2))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')