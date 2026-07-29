from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Add a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
tf = title_shape.text_frame
tf.text = "Academic World and Professional World: Demands and Characteristics"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(0, 51, 102)

# 2. Create the Comparison Table
# Dimensions: 4 rows (1 header + 3 sections), 2 columns
rows_count = 4
cols_count = 2
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.0)
height = Inches(5.5)

table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(4.5)

# Helper function to populate and style cells
def style_cell(cell, text, is_header=False):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        if is_header:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            # Set header background color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        else:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

# 3. Populate Table Headers
style_cell(table.cell(0, 0), "Academic World", is_header=True)
style_cell(table.cell(0, 1), "Professional World", is_header=True)

# 4. Populate Table Content

# Row 1: Hierarchy
hierarchy_acad = "Hierarchy\n" \
                 "• Tenure tracks and academic ranks\n" \
                 "• Department chair structure\n" \
                 "• Emphasis on peer recognition"
hierarchy_prof = "Hierarchy\n" \
                 "• Corporate ladder and job titles\n" \
                 "• Managerial reporting lines\n" \
                 "• Emphasis on productivity and results"
style_cell(table.cell(1, 0), hierarchy_acad, is_header=False)
style_cell(table.cell(1, 1), hierarchy_prof, is_header=False)

# Row 2: Learning Approach
learning_acad = "Learning Approach\n" \
                "• Theoretical and conceptual depth\n" \
                "• Research-driven inquiry\n" \
                "• Lifelong education focused"
learning_prof = "Learning Approach\n" \
                "• Practical and applied skills\n" \
                "• Just-in-time learning\n" \
                "• Continuous professional development"
style_cell(table.cell(2, 0), learning_acad, is_header=False)
style_cell(table.cell(2, 1), learning_prof, is_header=False)

# Row 3: Collaboration
collab_acad = "Collaboration\n" \
              "• Co-authoring research papers\n" \
              "• Conference presentations\n" \
              "• Interdisciplinary projects"
collab_prof = "Collaboration\n" \
              "• Cross-functional teams\n" \
              "• Client deliverables and KPIs\n" \
              "• Strategic partnerships"
style_cell(table.cell(3, 0), collab_acad, is_header=False)
style_cell(table.cell(3, 1), collab_prof, is_header=False)

# Save the presentation
prs.save('output.pptx')