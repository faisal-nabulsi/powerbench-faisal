from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen (standard)
slide_width = Inches(13.333)
slide_height = Inches(7.5)
prs.slide_width = slide_width
prs.slide_height = slide_height

# Add a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the background image ('image.png') covering the entire slide
# Adding it first ensures it is placed behind text boxes (z-order)
try:
    slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
except FileNotFoundError:
    # Fallback if image is missing, though instruction guarantees it is available
    pass

# Add the title "ELO's" at the top
# Positioned at the top, centered
title_box = slide.shapes.add_textbox(
    left=Inches(0), 
    top=Inches(0.5), 
    width=slide_width, 
    height=Inches(1.5)
)
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Configure title text styling
title_paragraph = title_frame.paragraphs[0]
title_paragraph.alignment = PP_ALIGN.CENTER
title_paragraph.text = "ELO's"
title_paragraph.font.size = Pt(54)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0) # Black text for readability
title_paragraph.font.name = 'Arial'

# Add the questions text box
# Positioned below the title
questions_box = slide.shapes.add_textbox(
    left=Inches(2.0), 
    top=Inches(2.8), 
    width=Inches(9.333), # Calculated width to center horizontally: 13.333 - 4.0
    height=Inches(4.0)
)
questions_frame = questions_box.text_frame
questions_frame.word_wrap = True

# Add Question 1
p1 = questions_frame.paragraphs[0]
p1.alignment = PP_ALIGN.LEFT
p1.text = "1. What is market?"
p1.font.size = Pt(32)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.font.name = 'Arial'
p1.space_after = Pt(30) # Vertical spacing between questions

# Add Question 2
p2 = questions_frame.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
p2.text = "2. How product reach to market?"
p2.font.size = Pt(32)
p2.font.color.rgb = RGBColor(0, 0, 0)
p2.font.name = 'Arial'

# Save the presentation to 'output.pptx'
prs.save('output.pptx')