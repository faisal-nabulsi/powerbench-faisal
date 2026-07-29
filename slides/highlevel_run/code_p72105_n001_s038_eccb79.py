from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Add a blank slide (Layout 6 is typically Blank in default themes)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Styling variables
TITLE_SIZE = Pt(44)
LABEL_SIZE = Pt(28)
COLOR_TITLE = RGBColor(0, 51, 102) # Dark Blue
COLOR_FILL = RGBColor(0, 128, 0)   # Green
COLOR_DRAIN = RGBColor(204, 0, 0)  # Red

# Slide dimensions for positioning
slide_width = prs.slide_width
slide_height = prs.slide_height

# 1. Add the main question
# Positioned at the top, centered
title_left = Inches(0.5)
title_top = Inches(0.5)
title_width = slide_width - Inches(1)
title_height = Inches(1.5)

txBox_title = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf_title = txBox_title.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "What is filling your bucket today and what’s draining it?"
p_title.font.size = TITLE_SIZE
p_title.font.bold = True
p_title.font.color.rgb = COLOR_TITLE
p_title.alignment = PP_ALIGN.CENTER

# 2. Add the illustration placeholder
# Positioned centrally below the title
img_width = Inches(7)
img_height = Inches(5)
img_top = Inches(2.5)
img_left = (slide_width - img_width) / 2

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# 3. Add "Filling" and "Draining" labels
# Label dimensions
label_width = Inches(1.5)
label_height = Inches(1)
label_top = img_top + (img_height - label_height) / 2

# "Filling" Label - Positioned to the left of the image
fill_label_left = img_left - Inches(2)
txBox_fill = slide.shapes.add_textbox(fill_label_left, label_top, label_width, label_height)
tf_fill = txBox_fill.text_frame
p_fill = tf_fill.paragraphs[0]
p_fill.text = "Filling"
p_fill.font.size = LABEL_SIZE
p_fill.font.bold = True
p_fill.font.color.rgb = COLOR_FILL
p_fill.alignment = PP_ALIGN.CENTER
tf_fill.word_wrap = False

# "Draining" Label - Positioned to the right of the image
drain_label_left = img_left + img_width + Inches(0.5)
txBox_drain = slide.shapes.add_textbox(drain_label_left, label_top, label_width, label_height)
tf_drain = txBox_drain.text_frame
p_drain = tf_drain.paragraphs[0]
p_drain.text = "Draining"
p_drain.font.size = LABEL_SIZE
p_drain.font.bold = True
p_drain.font.color.rgb = COLOR_DRAIN
p_drain.alignment = PP_ALIGN.CENTER
tf_drain.word_wrap = False

# Save the presentation
prs.save('output.pptx')