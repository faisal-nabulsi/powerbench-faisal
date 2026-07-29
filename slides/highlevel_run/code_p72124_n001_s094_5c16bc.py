import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_slide():
    # Initialize presentation
    prs = Presentation()
    
    # Set slide size to 16:9 (Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide (Layout index 6 is typically blank)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # --- 1. Title ---
    # Add a title text box at the top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Advantages of Blockchain"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.alignment = PP_ALIGN.LEFT
    
    # --- 2. Bullet Points (Left Column) ---
    # Create text box for advantages list
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf_txt = text_box.text_frame
    tf_txt.word_wrap = True
    
    # List of advantages with descriptions
    advantages = [
        ("Decentralization", "Removes intermediaries, distributing control across the network for resilience."),
        ("Transparency", "All transactions are visible and verifiable by all participants in the network."),
        ("Immutability", "Data, once recorded, cannot be altered or deleted, ensuring trust and integrity."),
        ("Security", "Advanced cryptography protects data and prevents unauthorized access or tampering."),
        ("Efficiency", "Automates transactions and verification processes, reducing time and costs.")
    ]
    
    # Populate text box with formatted runs
    for i, (title, desc) in enumerate(advantages):
        if i == 0:
            p = tf_txt.paragraphs[0]
        else:
            p = tf_txt.add_paragraph()
        
        p.space_after = Pt(10)
        
        # Title Run (Bold, Blue)
        run_title = p.add_run()
        run_title.text = f"{i+1}. {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(16)
        run_title.font.color.rgb = RGBColor(0x00, 0x5F, 0xA4)
        
        # Description Run (Normal, Dark Grey)
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
    # --- 3. Diagram (Right Column) ---
    # Diagram: A central "Blockchain" node connecting to benefit nodes
    
    # Center Node: "Blockchain"
    # Coordinates: x=8.0, y=3.0 (shifted down slightly to accommodate title)
    # Height 1.4, Width 2.2
    c_x, c_y = Inches(8.0), Inches(3.0)
    c_w, c_h = Inches(2.2), Inches(1.4)
    
    center_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_x, c_y, c_w, c_h)
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = RGBColor(0x00, 0x5F, 0xA4) # Dark Blue
    center_shape.line.color.rgb = RGBColor(255, 255, 255)
    
    tf_c = center_shape.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "Blockchain"
    p_c.font.size = Pt(20)
    p_c.font.color.rgb = RGBColor(255, 255, 255)
    p_c.font.bold = True
    p_c.alignment = PP_ALIGN.CENTER
    p_c.space_before = Pt(15) # Vertical centering adjustment
    
    # Satellite Nodes: Benefits
    node_w, node_h = Inches(2.2), Inches(1.0)
    n_x = Inches(10.8) # To the right of center
    
    # Nodes data: Label, Top Y Position
    # Center Y of the center box is 3.7 (3.0 + 0.7).
    # We align the middle node with this Y.
    nodes_data = [
        ("Decentralization", Inches(2.0)), # Top Node Center Y approx 2.5
        ("Transparency", Inches(3.2)),     # Middle Node Center Y approx 3.7
        ("Immutability", Inches(4.4))      # Bottom Node Center Y approx 4.9
    ]
    
    node_shapes_centers = []
    
    for label, y_pos in nodes_data:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, n_x, y_pos, node_w, node_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xBB, 0xDD, 0xFF) # Light Blue
        shape.line.color.rgb = RGBColor(0x00, 0x5F, 0xA4)
        
        tf_n = shape.text_frame
        tf_n.word_wrap = True
        p_n = tf_n.paragraphs[0]
        p_n.text = label
        p_n.font.size = Pt(14)
        p_n.font.color.rgb = RGBColor(0x00, 0x33, 0x66) # Dark Blue text
        p_n.font.bold = True
        p_n.alignment = PP_ALIGN.CENTER
        
        # Store center Y for connection lines
        node_shapes_centers.append(y_pos + Inches(0.5))
        
    # Add Connectors (Lines)
    # From Center Right Edge to Node Left Edges
    conn_start_x = c_x + c_w
    conn_start_y = c_y + c_h / 2.0
    
    for center_y in node_shapes_centers:
        conn_end_x = n_x
        conn_end_y = center_y
        
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            conn_start_x, conn_start_y, 
            conn_end_x, conn_end_y
        )
        connector.line.color.rgb = RGBColor(0x00, 0x5F, 0xA4)
        connector.line.width = Pt(2)

    # Save the presentation
    prs.save('output.pptx')
    print("Presentation saved successfully as 'output.pptx'.")

if __name__ == '__main__':
    create_slide()