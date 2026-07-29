from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_blockchain_slide():
    # Create a presentation object
    prs = Presentation()

    # Use a blank layout to have full control over positioning
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Add Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Blockchain in Cryptocurrency"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue
    p.alignment = PP_ALIGN.LEFT

    # --- Add Bullet Points ---
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(5.5)
    height = Inches(5.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()

    # Define content
    points = [
        ("Definition", "A decentralized, distributed digital ledger that records transactions across a network of computers, ensuring data cannot be altered retroactively."),
        ("Application", "Primarily known as the backbone for cryptocurrencies like Bitcoin, but also used for smart contracts, supply chain tracking, and secure identity management."),
        ("Governance", "Managed by network participants (nodes) using consensus mechanisms (e.g., Proof of Work or Proof of Stake) rather than a central authority or government."),
        ("Differences", "Cryptocurrency is the digital asset or currency (the 'what'), while blockchain is the underlying technology and protocol (the 'how') that enables it."),
        ("Decentralization", "Distributes control and data across many computers, removing the need for intermediaries (like banks) and significantly enhancing security and transparency.")
    ]

    # Add first point (Title in bold)
    p = tf.add_paragraph()
    p.text = f"{points[0][0]}: {points[0][1]}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(12)
    
    # Bold the category name
    run = p.add_run()
    run.text = f"{points[0][0]}: "
    run.font.bold = True

    # Add remaining points
    for i in range(1, len(points)):
        p = tf.add_paragraph()
        p.text = f"{points[i][0]}: {points[i][1]}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(12)
        
        # Bold the category name
        run = p.add_run()
        run.text = f"{points[i][0]}: "
        run.font.bold = True

    # --- Add Image ---
    # Position image on the right side
    left_img = Inches(6.2)
    top_img = Inches(1.8)
    width_img = Inches(3.5)
    height_img = Inches(5.5)
    
    try:
        slide.shapes.add_picture('image.png', left_img, top_img, width_img, height_img)
    except Exception:
        # Fallback if image is missing, though instructions say it's available
        print("Note: image.png not found, skipping image.")

    # Save the presentation
    prs.save('output.pptx')
    print("Presentation saved to 'output.pptx'")

if __name__ == "__main__":
    create_blockchain_slide()