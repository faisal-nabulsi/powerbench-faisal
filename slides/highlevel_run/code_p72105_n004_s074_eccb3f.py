from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation
prs = Presentation()

# Add a blank slide to allow for custom layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define color palette
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xE7, 0x4C, 0x3C)
TEAL = RGBColor(0x1A, 0xBC, 0x9C)
MUSTARD = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)

# 1. Add Title Bar
title_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(10), Inches(1.5)
)
title_box.fill.solid()
title_box.fill.fore_color.rgb = DARK_BLUE
title_box.line.fill.background()

title_text_frame = title_box.text_frame
title_text_frame.text = "Principles of Graphic Design"
title_text_frame.paragraphs[0].font.size = Pt(36)
title_text_frame.paragraphs[0].font.color.rgb = WHITE
title_text_frame.paragraphs[0].font.bold = True
title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 2. Add Introductory Statement
intro_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.8), Inches(9), Inches(1.0)
)
intro_text_frame = intro_box.text_frame
intro_text_frame.text = "Graphic design is the art of communication and problem-solving through the use of typography, photography, and illustration. These core principles ensure visual harmony and effective message delivery."
intro_text_frame.paragraphs[0].font.size = Pt(14)
intro_text_frame.paragraphs[0].font.color.rgb = DARK_TEXT
intro_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 3. Add Principles List in a Grid Layout
principles = [
    ("Balance", "Visual weight is distributed evenly."),
    ("Contrast", "Differences in elements create interest."),
    ("Emphasis", "Focal points draw the viewer's eye."),
    ("Hierarchy", "Order of importance guides the flow."),
    ("Proximity", "Related items are grouped together."),
    ("Repetition", "Consistent elements create unity."),
    ("Alignment", "Elements connect to create order."),
    ("White Space", "Empty space allows content to breathe.")
]

# Alternating colors for visual contrast
colors = [CORAL, TEAL, MUSTARD, PURPLE, CORAL, TEAL, MUSTARD, PURPLE]

start_top = Inches(3.1)
box_width = Inches(4.2)
box_height = Inches(0.9)
gap = Inches(0.15)
col_width = box_width + gap

for i, (title, desc) in enumerate(principles):
    row = i // 2
    col = i % 2
    
    left = Inches(0.5) + col * col_width
    top = start_top + row * (box_height + gap)
    
    # Add shape for principle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, box_width, box_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors[i]
    shape.line.fill.background()
    
    # Add text to shape
    tf = shape.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT

# Save the presentation
prs.save('output.pptx')