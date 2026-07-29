from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Add a blank slide layout (index 6) to allow custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Title
title_box = slide.shapes.add_textbox(
    left=Inches(0.5), 
    top=Inches(0.5), 
    width=Inches(9), 
    height=Inches(1.2)
)
title_frame = title_box.text_frame
title_frame.text = "Globalization isn’t possible to occur without media"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 2. Add Content (Bullet Points)
text_box = slide.shapes.add_textbox(
    left=Inches(0.5), 
    top=Inches(2.2), 
    width=Inches(4.5), 
    height=Inches(4.5)
)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Bullet Point 1: Social Media
p1 = text_frame.add_paragraph()
p1.text = "• Social Media: Connects individuals globally, facilitating instant communication and cultural exchange across borders."
p1.font.size = Pt(16)
p1.space_before = Pt(6)

# Bullet Point 2: News Media
p2 = text_frame.add_paragraph()
p2.text = "• News Media: Disseminates information worldwide, shaping public perception of international events and fostering political awareness."
p2.font.size = Pt(16)
p2.space_before = Pt(6)

# Bullet Point 3: Streaming Services
p3 = text_frame.add_paragraph()
p3.text = "• Streaming Services: Distributes entertainment globally, enabling the cross-border flow of culture and promoting shared experiences."
p3.font.size = Pt(16)
p3.space_before = Pt(6)

# 3. Add Illustration
# Using 'image.png' as a placeholder for relevant illustrations.
# Positioned on the right side to accompany the text.
slide.shapes.add_picture('image.png', 
                         left=Inches(5.5), 
                         top=Inches(2.2), 
                         width=Inches(4), 
                         height=Inches(4))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')