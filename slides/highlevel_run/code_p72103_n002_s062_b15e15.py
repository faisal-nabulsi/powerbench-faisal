from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Create a presentation
    prs = Presentation()
    
    # Use a blank layout to have full control over positioning
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # --- 1. Title: "PART 3" ---
    # Position at the top center
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PART 3"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204) # A nice blue color
    p.alignment = PP_ALIGN.CENTER

    # --- 2. Prompt: First date scenario ---
    # Position below the title
    prompt_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(2))
    tf = prompt_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Scenario:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "It is your first date at a casual diner. The atmosphere is friendly. Suddenly, the waiter walks by and asks if you'd like to share an order of French fries."
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(70, 70, 70)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    # --- 3. Character Image ---
    # Placeholder image representing the character
    # Position on the left side of the prompt area
    img_left = Inches(1.5)
    img_top = Inches(4.5)
    img_width = Inches(1.5)
    img_height = Inches(2.0)
    
    try:
        slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
    except Exception as e:
        # Fallback if image file is missing, though instructions imply it exists
        print(f"Note: Could not add image. {e}")

    # --- 4. Storytelling Starter (Speech Bubble) ---
    # Position to the right of the character
    bubble_left = Inches(3.5)
    bubble_top = Inches(4.5)
    bubble_width = Inches(6)
    bubble_height = Inches(2.5)
    
    # Add a Rounded Rectangle as the speech bubble
    bubble_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        bubble_left, bubble_top, bubble_width, bubble_height
    )
    
    # Fill the bubble
    fill = bubble_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245) # Light gray background
    
    # Add a border
    line = bubble_shape.line
    line.color.rgb = RGBColor(200, 200, 200)
    line.width = Pt(2)
    
    # Add text to the bubble
    tf = bubble_shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None # Prevent auto-resize shrinking text
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = "Story Starter:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    p2 = tf.add_paragraph()
    p2.text = "\"I usually try to impress people with fancy wine, but honestly, looking at those fries... maybe it's time to let my guard down.\""
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(30, 30, 30)
    p2.space_before = Pt(5)

    # --- Optional: Add a small triangle tail for the bubble ---
    # Approximate the tail using a triangle shape
    tail_left = bubble_left - Inches(0.2)
    tail_top = bubble_top + Inches(0.8)
    tail_width = Inches(0.3)
    tail_height = Inches(0.4)
    
    tail_shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        tail_left, tail_top, tail_width, tail_height
    )
    tail_shape.rotation = 270 # Point to the left
    tail_shape.fill.solid()
    tail_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    tail_shape.line.color.rgb = RGBColor(200, 200, 200)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()