from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide():
    # Create a new presentation instance
    prs = Presentation()
    
    # Add a blank slide to the presentation
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Retrieve slide dimensions to ensure the background covers the entire slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Add the background image (Green gradient with hexagonal patterns)
    # The instructions specify to use 'image.png' for image placeholders.
    try:
        background = slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
    except Exception:
        # Fallback handling if image is not found, though it is expected to be present
        pass

    # --- Add Title ---
    # Positioning the title box slightly above the center
    title_box = slide.shapes.add_textbox(
        left=Inches(1.5), 
        top=Inches(2.0), 
        width=Inches(7.0), 
        height=Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    title_paragraph = title_frame.add_paragraph()
    title_paragraph.text = "Natural Environment"
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    # Using white text for contrast against the green background
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255) 
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # --- Add Details (Speaker, Class, Roll Number) ---
    # Positioning the details box below the title
    details_box = slide.shapes.add_textbox(
        left=Inches(1.5), 
        top=Inches(4.0), 
        width=Inches(7.0), 
        height=Inches(2.5)
    )
    details_frame = details_box.text_frame
    details_frame.word_wrap = True

    # Speaker Name
    p_speaker = details_frame.add_paragraph()
    p_speaker.text = "Speaker: Dr. Emily Green"
    p_speaker.font.size = Pt(24)
    p_speaker.font.color.rgb = RGBColor(255, 255, 255)
    p_speaker.alignment = PP_ALIGN.CENTER

    # Class
    p_class = details_frame.add_paragraph()
    p_class.text = "Class: Environmental Science 101"
    p_class.font.size = Pt(24)
    p_class.font.color.rgb = RGBColor(255, 255, 255)
    p_class.alignment = PP_ALIGN.CENTER

    # Roll Number
    p_roll = details_frame.add_paragraph()
    p_roll.text = "Roll Number: 2023-00-456"
    p_roll.font.size = Pt(24)
    p_roll.font.color.rgb = RGBColor(255, 255, 255)
    p_roll.alignment = PP_ALIGN.CENTER

    # Save the presentation to 'output.pptx'
    prs.save('output.pptx')

if __name__ == "__main__":
    create_title_slide()