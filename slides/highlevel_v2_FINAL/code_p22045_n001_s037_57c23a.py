from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank layout for custom positioning
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "Elon Musk Current Stage"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
p.font.name = "Calibri"

# 2. Portrait Image
slide.shapes.add_picture('image.png', Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.5))

# 3. Roles & Companies List
roles_data = [
    ("CEO & Product Architect", "Tesla"),
    ("CEO & Chief Technology Officer", "SpaceX"),
    ("Owner, CTO & Executive Chairman", "X Corp."),
    ("CEO & Founder", "xAI"),
    ("Founder & Chief Engineer", "Neuralink"),
    ("Founder", "The Boring Company")
]

y_start = Inches(2.0)
height = Inches(0.75)
spacing = Inches(0.15)
x_start = Inches(5.0)
width = Inches(7.8)

for i, (role, company) in enumerate(roles_data):
    y_pos = y_start + i * (height + spacing)
    
    # Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_start, y_pos, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xE8)
    card.line.width = Pt(1)
    
    # Text formatting for card
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.05)
    
    # Role text
    p_role = tf.add_paragraph()
    p_role.text = role
    p_role.font.size = Pt(13)
    p_role.font.italic = True
    p_role.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p_role.font.name = "Calibri"
    
    # Company text
    p_comp = tf.add_paragraph()
    p_comp.text = company
    p_comp.font.size = Pt(20)
    p_comp.font.bold = True
    p_comp.font.color.rgb = RGBColor(0x00, 0x56, 0xB3)
    p_comp.font.name = "Calibri"
    
    # Company Logo Placeholder (using allowed 'image.png' as consistent placeholder)
    logo_box = slide.shapes.add_picture('image.png', 
                                         x_start + Inches(0.3), 
                                         y_pos + Inches(0.12), 
                                         Inches(0.5), 
                                         Inches(0.5))

# Save presentation
prs.save('output.pptx')