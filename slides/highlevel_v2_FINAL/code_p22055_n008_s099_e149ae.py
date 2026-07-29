from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize the presentation
prs = Presentation()

# Set slide width and height to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (layout index 6 is usually blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "PART 2" at the top
# Coordinates: Left, Top, Width, Height
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "PART 2"
title_para.font.size = Pt(48)
title_para.font.bold = True

# 2. Add a prompt to discuss dating
# Positioned below the title
prompt_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
prompt_frame = prompt_box.text_frame
prompt_frame.word_wrap = True

prompt_para = prompt_frame.paragraphs[0]
prompt_para.text = "Discussion Prompt: Let's talk about how we date in the modern era."
prompt_para.font.size = Pt(28)

# 3. List three numbered questions
# Positioned below the prompt
questions_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.333), Inches(3))
questions_frame = questions_box.text_frame
questions_frame.word_wrap = True

questions = [
    "1. What is the most important quality you look for in a partner?",
    "2. How has social media impacted the way you approach dating?",
    "3. What is your ideal scenario for a first date?"
]

# Add questions to the text frame
for i, question in enumerate(questions):
    if i == 0:
        para = questions_frame.paragraphs[0]
    else:
        para = questions_frame.add_paragraph()
    
    para.text = question
    para.font.size = Pt(22)
    # Add space after each question for clarity
    para.space_after = Pt(14)

# Save the presentation
prs.save('output.pptx')