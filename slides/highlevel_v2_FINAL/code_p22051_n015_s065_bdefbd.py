from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the image (placeholder)
# Centered horizontally and vertically positioned to leave room for text
image_width = Inches(5.5)
image_height = Inches(5.5)
image_left = (prs.slide_width - image_width) / 2
image_top = Inches(0.5)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Add the text phrase
# Positioned below the image
text_top = image_top + image_height + Inches(0.3)
text_left = Inches(1.5)
text_width = Inches(10.333)
text_height = Inches(1.5)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Create a paragraph for the phrase
paragraph = text_frame.paragraphs[0]
paragraph.text = "global media cultures"

# Format the text
paragraph.font.size = Pt(44)
paragraph.font.bold = True
paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Gray
paragraph.alignment = 1 # Center alignment (1 = PP_ALIGN.CENTER)

# Save the presentation
prs.save('output.pptx')