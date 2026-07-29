from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation and set slide size to 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set background to a dark color to simulate a dark market theme
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A) # Dark Grey/Black

# Add a red section on the top right
# Position: Right side, Top edge
# Size: covers approx 1/3 width and top portion
red_shape = slide.shapes.add_shape(
    1, # Rectangle shape ID
    left=Inches(9.0),
    top=Inches(0),
    width=Inches(4.333),
    height=Inches(3.5)
)
red_shape.fill.solid()
red_shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00) # Red
red_shape.line.fill.background() # Remove border

# Add text "WHAT IS MARKET?"
# Positioned to the left of the red section for balance
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4.5))
tf = text_box.text_frame
tf.word_wrap = True

# Configure the paragraph
p = tf.paragraphs[0]
p.text = "WHAT IS MARKET?"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text for contrast

# Save the presentation
prs.save('output.pptx')