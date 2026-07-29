from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the Blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Dynamics of Local and Global Culture"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79) # Dark Blue
title_para.alignment = PP_ALIGN.CENTER

# --- Overview Section ---
overview_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.5))
overview_frame = overview_box.text_frame
overview_frame.word_wrap = True

# Overview Header
ov_header_para = overview_frame.paragraphs[0]
ov_header_run = ov_header_para.add_run()
ov_header_run.text = "Overview of Global Cultural Flows"
ov_header_run.font.size = Pt(24)
ov_header_run.font.bold = True
ov_header_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Overview Body
ov_body_para = overview_frame.add_paragraph()
ov_body_run = ov_body_para.add_run()
ov_body_run.text = "Global cultural flows refer to the movement of ideas, values, and practices across borders. These flows influence local identities, creating complex interactions and tensions between maintaining local distinctiveness and adapting to global influences."
ov_body_run.font.size = Pt(18)
ov_body_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# --- Perspectives Section ---
perspectives_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.333), Inches(3.0))
perspectives_frame = perspectives_box.text_frame
perspectives_frame.word_wrap = True

# Perspectives Header
pers_header_para = perspectives_frame.paragraphs[0]
pers_header_run = pers_header_para.add_run()
pers_header_run.text = "Three Key Perspectives"
pers_header_run.font.size = Pt(24)
pers_header_run.font.bold = True
pers_header_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Perspective 1: Cultural Differentialism
p1_para = perspectives_frame.add_paragraph()
p1_para.space_before = Pt(14)
p1_title_run = p1_para.add_run()
p1_title_run.text = "1. Cultural Differentialism: "
p1_title_run.font.size = Pt(18)
p1_title_run.font.bold = True
p1_title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p1_desc_run = p1_para.add_run()
p1_desc_run.text = "Emphasizes the distinctiveness and boundaries of cultures, often resisting global homogenization to preserve local identity."
p1_desc_run.font.size = Pt(18)

# Perspective 2: Hybridization
p2_para = perspectives_frame.add_paragraph()
p2_para.space_before = Pt(14)
p2_title_run = p2_para.add_run()
p2_title_run.text = "2. Hybridization: "
p2_title_run.font.size = Pt(18)
p2_title_run.font.bold = True
p2_title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p2_desc_run = p2_para.add_run()
p2_desc_run.text = "Describes the blending of local and global elements to create new, mixed cultural forms and identities."
p2_desc_run.font.size = Pt(18)

# Perspective 3: Convergence
p3_para = perspectives_frame.add_paragraph()
p3_para.space_before = Pt(14)
p3_title_run = p3_para.add_run()
p3_title_run.text = "3. Convergence: "
p3_title_run.font.size = Pt(18)
p3_title_run.font.bold = True
p3_title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p3_desc_run = p3_para.add_run()
p3_desc_run.text = "Suggests that global forces lead to a uniformity in culture, where local differences diminish over time."
p3_desc_run.font.size = Pt(18)

# Save the presentation
prs.save('output.pptx')