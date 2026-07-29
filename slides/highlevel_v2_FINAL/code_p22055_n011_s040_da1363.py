from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set slide background to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "Principles of Graphic Design"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

# --- Introductory Statement ---
intro_box = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9.333), Inches(1.0))
intro_frame = intro_box.text_frame
intro_frame.word_wrap = True
intro_para = intro_frame.paragraphs[0]
intro_para.alignment = PP_ALIGN.CENTER
intro_run = intro_para.add_run()
intro_run.text = "Effective design relies on a set of core principles that guide the arrangement of visual elements to communicate messages clearly and aesthetically."
intro_run.font.size = Pt(18)
intro_run.font.color.rgb = RGBColor(51, 51, 51) # Dark Grey
intro_run.font.italic = True

# --- List of Principles Section ---
# Background Rectangle for the list area to provide contrast
list_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.5), Inches(9.333), Inches(3.5))
list_bg.fill.solid()
list_bg.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light Grey
list_bg.line.fill.background() # Remove border

# Textbox for the list content
list_box = slide.shapes.add_textbox(Inches(2.2), Inches(3.7), Inches(8.933), Inches(3.1))
list_frame = list_box.text_frame
list_frame.word_wrap = True

principles = [
    ("Balance", "Distributing visual weight evenly across the composition."),
    ("Contrast", "Highlighting differences to create visual interest and focus."),
    ("Emphasis", "Focusing attention on a specific element or area."),
    ("Hierarchy", "Organizing elements to show order of importance."),
    ("Proximity", "Grouping related items together to create organization."),
    ("Repetition", "Reusing visual elements to create unity and consistency."),
    ("White Space", "Using empty space to reduce clutter and improve readability.")
]

for i, (title, desc) in enumerate(principles):
    if i > 0:
        p = list_frame.add_paragraph()
    else:
        p = list_frame.paragraphs[0]
    
    p.space_after = Pt(6)
    
    # Principle Name
    run_title = p.add_run()
    run_title.text = f"{title}: "
    run_title.font.bold = True
    run_title.font.size = Pt(16)
    run_title.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
    
    # Description
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(14)
    run_desc.font.color.rgb = RGBColor(51, 51, 51) # Dark Grey

# Save the presentation
prs.save('output.pptx')