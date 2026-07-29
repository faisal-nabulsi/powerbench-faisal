from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
title_text = "Innovations (Future Products or Services)"
# Positioning: Centered at the top
title_left = Inches(0.5)
title_top = Inches(0.3)
title_width = Inches(12.333)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Style the title
title_paragraph = title_frame.paragraphs[0]
title_paragraph.text = title_text
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.alignment = PP_ALIGN.CENTER

# --- Add Illustration (Placeholder Image) ---
# Representing "colorful illustration of diverse Airbnb accommodations"
# Positioning: Centered horizontally, taking up the lower portion of the slide
image_width = Inches(10)
image_height = Inches(4.2)
image_left = (prs.slide_width - image_width) / 2
image_top = Inches(1.6)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')