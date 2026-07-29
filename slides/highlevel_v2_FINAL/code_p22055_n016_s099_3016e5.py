from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation and set dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.5), Inches(1.5))
title_tf = title_shape.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Elon's Future Plans"
title_run.font.size = Pt(44)
title_run.font.bold = True

# Add Bullet Points
content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(7.5), Inches(5))
content_tf = content_shape.text_frame
content_tf.word_wrap = True

bullets = [
    "Development of Starship for Mars colonization",
    "Advancing Neuralink brain-computer interfaces",
    "Expanding The Boring Company tunnel network",
    "Deploying Tesla Optimus humanoid robots",
    "Scaling X (formerly Twitter) into an 'everything app'"
]

for i, bullet in enumerate(bullets):
    if i == 0:
        para = content_tf.paragraphs[0]
    else:
        para = content_tf.add_paragraph()
    
    run = para.add_run()
    run.text = bullet
    run.font.size = Pt(20)
    
    # Add bullet character manually if needed, or rely on text structure
    # simple approach: just text with newlines, but pptx doesn't auto-bullet without XML manipulation
    # or we can use a bullet character.
    # Let's format properly.
    para.level = 0

# Add Image on the right
try:
    slide.shapes.add_picture('image.png', Inches(8.5), Inches(2.2), Inches(4.5), Inches(4.5))
except:
    pass

# Save the presentation
prs.save('output.pptx')