from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Layout 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define dimensions and position for the text box to center it on the slide
# Center X: (13.333 width - 10 box_width) / 2 = 1.6665
# Center Y: (7.5 height - 2.5 box_height) / 2 = 2.5
left = Inches(1.6665)
top = Inches(2.5)
width = Inches(10)
height = Inches(2.5)

# Add the text box shape
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

# Clear any existing paragraphs (though it's a new box, good practice)
text_frame.clear()

# Access the first paragraph
paragraph = text_frame.paragraphs[0]

# Center align the text within the text box
paragraph.alignment = PP_ALIGN.CENTER

# Create a run for the text and apply formatting
run = paragraph.add_run()
run.text = "THANK YOU!"
run.font.size = Pt(72)   # Large font size
run.font.bold = True     # Bold font
run.font.color.rgb = RGBColor(0, 0, 0) # Black text color for contrast

# Save the presentation
prs.save('output.pptx')