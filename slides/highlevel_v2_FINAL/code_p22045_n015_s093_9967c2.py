from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Global and Local Cultural Products"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add the definition of cultural products
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0))
text_tf = text_box.text_frame
text_tf.word_wrap = True
text_para = text_tf.paragraphs[0]
text_run = text_para.add_run()
text_run.text = (
    "Definition:\n\n"
    "Cultural products are goods and services that embody cultural meanings, values, and identities. "
    "They include items like art, music, literature, fashion, food, and traditional crafts, which reflect "
    "the heritage and creativity of a community or society. These products serve as a bridge between "
    "global markets and local traditions, preserving cultural diversity while fostering exchange and understanding."
)
text_run.font.size = Pt(18)

# Add placeholder images to showcase various cultural items
# Image 1
slide.shapes.add_picture('image.png', Inches(7.0), Inches(1.8), Inches(3.0), Inches(2.5))
# Image 2
slide.shapes.add_picture('image.png', Inches(10.2), Inches(1.8), Inches(3.0), Inches(2.5))
# Image 3
slide.shapes.add_picture('image.png', Inches(8.6), Inches(4.5), Inches(3.0), Inches(2.5))

# Save the presentation
prs.save('output.pptx')