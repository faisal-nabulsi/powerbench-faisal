from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title "PART 2"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "PART 2"
title_run.font.size = Pt(40)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

# Add Content Textbox
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(4.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Instruction 1
para_intro = content_frame.paragraphs[0]
run_intro = para_intro.add_run()
run_intro.text = "Instructions:"
run_intro.font.size = Pt(22)
run_intro.font.bold = True

para_verb = content_frame.add_paragraph()
para_verb.space_after = Pt(10)
run_verb = para_verb.add_run()
run_verb.text = "Please fill in the correct verb forms in the provided dialogue."
run_verb.font.size = Pt(18)

# Part 3 Section
para_part3_title = content_frame.add_paragraph()
para_part3_title.space_before = Pt(20)
run_part3_title = para_part3_title.add_run()
run_part3_title.text = "PART 3"
run_part3_title.font.size = Pt(28)
run_part3_title.font.bold = True
run_part3_title.font.color.rgb = RGBColor(153, 0, 0) # Dark Red

para_part3_desc = content_frame.add_paragraph()
run_part3_desc = para_part3_desc.add_run()
run_part3_desc.text = "Retell Joey's part with emphasis on stress and intonation."
run_part3_desc.font.size = Pt(18)

# Save the presentation
prs.save('output.pptx')