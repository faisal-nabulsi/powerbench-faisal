from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title "PART 3" at the top
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_para.text = "PART 3"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Add Prompt Text Box
# Positioned below the title on the left side
prompt_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.5), Inches(3.5))
prompt_tf = prompt_box.text_frame
prompt_tf.word_wrap = True

prompt_para = prompt_tf.paragraphs[0]
prompt_para.text = "Prompt: What is your go-to strategy for making a good first impression on a date?"
prompt_para.font.size = Pt(24)

# Add Speech Bubble Image (Placeholder)
# Positioned on the right side
image_left = Inches(9.5)
image_top = Inches(2.5)
image_width = Inches(3.5)
image_height = Inches(3.5)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')