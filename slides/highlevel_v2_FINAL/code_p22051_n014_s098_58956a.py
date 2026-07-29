from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_p = title_tf.paragraphs[0]
title_run = title_p.add_run()
title_run.text = "PRONUNCIATION ACTIVITY"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add Question
question_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
question_tf = question_box.text_frame
question_tf.word_wrap = True
question_p = question_tf.paragraphs[0]
question_run = question_p.add_run()
question_run.text = "What is the general rule for the pronunciation scheme?"
question_run.font.size = Pt(24)
question_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Add Option A
opt_a_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.8))
opt_a_tf = opt_a_box.text_frame
opt_a_tf.word_wrap = True
opt_a_p = opt_a_tf.paragraphs[0]
opt_a_run = opt_a_p.add_run()
opt_a_run.text = "a) Vowels are always short."
opt_a_run.font.size = Pt(20)
opt_a_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add Option B
opt_b_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.333), Inches(0.8))
opt_b_tf = opt_b_box.text_frame
opt_b_tf.word_wrap = True
opt_b_p = opt_b_tf.paragraphs[0]
opt_b_run = opt_b_p.add_run()
opt_b_run.text = "b) Stress falls on the first syllable."
opt_b_run.font.size = Pt(20)
opt_b_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add Option C
opt_c_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.333), Inches(0.8))
opt_c_tf = opt_c_box.text_frame
opt_c_tf.word_wrap = True
opt_c_p = opt_c_tf.paragraphs[0]
opt_c_run = opt_c_p.add_run()
opt_c_run.text = "c) Consonants are silent at the end."
opt_c_run.font.size = Pt(20)
opt_c_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Save the presentation
prs.save('output.pptx')