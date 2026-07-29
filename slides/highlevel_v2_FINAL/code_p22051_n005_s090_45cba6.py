from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Title: "The Marketing Mix"
# Top center area
title_left = Inches(1.5)
title_top = Inches(0.3)
title_width = Inches(10)
title_height = Inches(0.8)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "The Marketing Mix"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_para.alignment = PP_ALIGN.CENTER

# 2. "Marketing" Text: Prominently on the right
mk_left = Inches(9.5)
mk_top = Inches(1.2)
mk_width = Inches(3.5)
mk_height = Inches(5.0)

mk_box = slide.shapes.add_textbox(mk_left, mk_top, mk_width, mk_height)
mk_frame = mk_box.text_frame
mk_para = mk_frame.paragraphs[0]
mk_para.text = "MARKETING"
mk_para.font.size = Pt(60) # Very large
mk_para.font.bold = True
mk_para.font.color.rgb = RGBColor(0, 51, 102)
mk_para.alignment = PP_ALIGN.CENTER
# Maybe vertical flip? No, standard center in box is fine. 
# To make it look like a sidebar title, maybe rotate? 
# Let's keep it simple horizontal.

# 3. Central Target Market (Circle)
# Positioned slightly left of center to balance the right text
center_x = Inches(5.0)
center_y = Inches(3.5)
center_size = Inches(1.0)

target_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x, center_y, center_size, center_size)
target_frame = target_shape.text_frame
target_frame.word_wrap = True
target_para = target_frame.paragraphs[0]
target_para.text = "Target\nMarket"
target_para.font.size = Pt(12)
target_para.font.color.rgb = RGBColor(255, 255, 255) # White
target_para.alignment = PP_ALIGN.CENTER
target_shape.fill.solid()
target_shape.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue match title
target_shape.line.fill.background() # No line

# 4. Hexagons (The 4 Ps)
# Arranged around the center (5.5, 4.0)
# Center of diagram shifted slightly right to match visual balance? 
# Let's rigidly place them around the Target shape center.
# Target Shape Center is at (5.0 + 0.5, 3.5 + 0.5) = (5.5, 4.0).

hex_width = Inches(1.4)
hex_height = Inches(1.2)

# Positions (Left, Top)
# Product: Top of circle. Circle top Y = 3.5. Hex bottom needs to be near there? 
# Gap of 0.4 inches.
# Center Y = 4.0. 
# Top Hex Center Y = 4.0 - 1.8 = 2.2.
# Top Hex Top Y = 2.2 - 0.6 = 1.6.
# Left X needed to center on 5.5 -> 5.5 - 0.7 = 4.8.

pos_product = (Inches(4.8), Inches(1.6))
pos_price = (Inches(6.8), Inches(3.4))   # Right. Center X = 5.5 + 1.8 = 7.3. Left X = 7.3 - 0.7 = 6.6. Let's say 6.8.
pos_place = (Inches(4.8), Inches(5.4))   # Bottom. Center Y = 4.0 + 1.8 = 5.8. Top Y = 5.8 - 0.6 = 5.2.
pos_promotion = (Inches(2.8), Inches(3.4)) # Left. Center X = 5.5 - 1.8 = 3.7. Left X = 3.7 - 0.7 = 3.0.

# Adjustments for exact symmetry around (5.5, 4.0)
# Hex Width 1.4 (Half 0.7), Height 1.2 (Half 0.6)
# Center of Hex needs to be:
# Top: (5.5, 2.2) -> Left 4.8, Top 1.6
# Bottom: (5.5, 5.8) -> Left 4.8, Top 5.2
# Left: (3.7, 4.0) -> Left 3.0, Top 3.4
# Right: (7.3, 4.0) -> Left 6.6, Top 3.4

final_positions = {
    "Product": (Inches(4.8), Inches(1.6)),
    "Price": (Inches(6.6), Inches(3.4)),
    "Place": (Inches(4.8), Inches(5.2)),
    "Promotion": (Inches(3.0), Inches(3.4))
}

# Colors for hexagons
hex_colors = {
    "Product": RGBColor(0, 128, 128),    # Teal
    "Price": RGBColor(255, 127, 80),     # Coral
    "Place": RGBColor(106, 90, 205),     # SlateBlue
    "Promotion": RGBColor(255, 140, 0)   # DarkOrange
}

for label, (left, top) in final_positions.items():
    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, hex_width, hex_height)
    shape.text_frame.text = label
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_colors[label]
    shape.line.fill.background() # Remove outline for cleaner look

# Save the presentation
prs.save('output.pptx')