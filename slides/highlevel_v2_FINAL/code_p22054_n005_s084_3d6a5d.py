from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the Presentation object
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout to have full control over elements
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# -------------------------
# 1. Add Title: "Graphic Design"
# -------------------------
# Define dimensions and position (Centered at the top)
title_width = Inches(10.5)
title_height = Inches(1.5)
title_left = Inches((13.333 - 10.5) / 2) # Center calculation
title_top = Inches(0.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Configure text properties
p = title_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Graphic Design"
run.font.size = Pt(54)
run.font.bold = True
run.font.color.rgb = RGBColor(33, 33, 33) # Dark grey for a modern look

# -------------------------
# 2. Add Image: Laptop with Paint Splashes
# -------------------------
# Use the provided placeholder image 'image.png'
# Position centrally below the title to be the focal point
img_width = Inches(8.5)
img_height = Inches(4.5)
img_left = Inches((13.333 - 8.5) / 2) # Center calculation
img_top = Inches(2.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# -------------------------
# Save the Presentation
# -------------------------
prs.save('output.pptx')