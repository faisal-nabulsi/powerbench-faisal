from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6] # 6 is usually 'Blank'
slide = prs.slides.add_slide(slide_layout)

# --- 1. Add Title ---
# Title dimensions and position
title_left = Inches(0.5)
title_top = Inches(0.5)
title_width = Inches(12.333)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
text = "Globalization isn’t possible to occur without media"
run = p.add_run()
run.text = text
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey
p.alignment = 1 # 1 is Center

# --- 2. Add Content Text (Bullet Points) ---
# Text box dimensions and position (Left side)
text_left = Inches(0.5)
text_top = Inches(2.0)
text_width = Inches(6.0)
text_height = Inches(4.5)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
tf_contents = text_box.text_frame
tf_contents.word_wrap = True

# Clear default paragraph
tf_contents.clear()

# Bullet Point 1
p1 = tf_contents.add_paragraph()
p1.space_after = Pt(14)
run1 = p1.add_run()
run1.text = "•  Print Media: "
run1.font.bold = True
run1.font.size = Pt(16)
run1.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) # Blue accent
run1_desc = p1.add_run()
run1_desc.text = "Disseminates knowledge across borders, laying the groundwork for global awareness."
run1_desc.font.size = Pt(16)
run1_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Bullet Point 2
p2 = tf_contents.add_paragraph()
p2.space_after = Pt(14)
run2 = p2.add_run()
run2.text = "•  Broadcast Media: "
run2.font.bold = True
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
run2_desc = p2.add_run()
run2_desc.text = "Enables real-time cultural exchange, fostering a shared global consciousness."
run2_desc.font.size = Pt(16)
run2_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Bullet Point 3
p3 = tf_contents.add_paragraph()
run3 = p3.add_run()
run3.text = "•  Digital Media: "
run3.font.bold = True
run3.font.size = Pt(16)
run3.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
run3_desc = p3.add_run()
run3_desc.text = "Facilitates instant connectivity and economic integration worldwide."
run3_desc.font.size = Pt(16)
run3_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# --- 3. Add Illustrations (Images) ---
# Images arranged vertically on the right side
img_width = Inches(5.5)
img_height = Inches(1.3)
img_left = Inches(7.0) # Right side alignment
img_spacer = Inches(0.2)

# Image 1 (Print)
img_top_1 = Inches(2.0)
slide.shapes.add_picture('image.png', img_left, img_top_1, img_width, img_height)

# Image 2 (Broadcast)
img_top_2 = Inches(2.0) + img_height + img_spacer
slide.shapes.add_picture('image.png', img_left, img_top_2, img_width, img_height)

# Image 3 (Digital)
img_top_3 = Inches(2.0) + (img_height + img_spacer) * 2
slide.shapes.add_picture('image.png', img_left, img_top_3, img_width, img_height)

# Save the presentation
prs.save('output.pptx')