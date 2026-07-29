from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Central Graphic
img_w = Inches(5.5)
img_h = Inches(5.5)
img_left = (prs.slide_width - img_w) / 2
img_top = (prs.slide_height - img_h) / 2
slide.shapes.add_picture('image.png', img_left, img_top, img_w, img_h)

# Text
text = "global media cultures"
t_w = Inches(10)
t_h = Inches(0.8)
t_left = (prs.slide_width - t_w) / 2
t_top = Inches(0.3) # Slightly indented from top

tb = slide.shapes.add_textbox(t_left, t_top, t_w, t_h)
tb.text_frame.paragraphs[0].text = text
tb.text_frame.paragraphs[0].font.size = Pt(42)
tb.text_frame.paragraphs[0].font.bold = True
tb.text_frame.paragraphs[0].font.name = "Arial"
tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

prs.save('output.pptx')