from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set the slide size to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (index 6 is typically the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the Fluentize logo (placeholder 'image.png') at the top left corner
# Position: 0.5 inches from left, 0.5 inches from top
# Size: 2.0 inches width, 1.0 inches height
slide.shapes.add_picture('image.png', left=Inches(0.5), top=Inches(0.5), width=Inches(2.0), height=Inches(1.0))

# Define text box properties for "THANK YOU!"
text_width = Inches(8.0)
text_height = Inches(2.0)

# Calculate position to center horizontally and place at the bottom
# Horizontal Center: (Slide Width - Text Width) / 2
text_left = Inches((13.333 - 8.0) / 2)
# Vertical Position: Bottom Margin (1.0 inch) above bottom edge, adjusted for text height
text_top = Inches(7.5 - text_height.inches - 1.0)

# Add the text box
text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Style the text
paragraph = text_frame.paragraphs[0]
paragraph.text = "THANK YOU!"
paragraph.font.size = Pt(60)
paragraph.font.bold = True
paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey
paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')