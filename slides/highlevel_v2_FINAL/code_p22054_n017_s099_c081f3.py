from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
title_left = Inches(1)
title_top = Inches(0.5)
title_width = Inches(11.333)
title_height = Inches(1.2)

title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_shape.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "PART 1"
title_run.font.size = Pt(40)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

# --- Content Section ---
content_left = Inches(1)
content_top = Inches(1.8)
content_width = Inches(7)
content_height = Inches(5.2)

content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
content_tf = content_shape.text_frame
content_tf.word_wrap = True

# Part 1: Discussion Prompt
p_prompt_header = content_tf.paragraphs[0]
run_prompt_header = p_prompt_header.add_run()
run_prompt_header.text = "Discussion Prompt"
run_prompt_header.font.size = Pt(20)
run_prompt_header.font.bold = True

p_prompt_question = content_tf.add_paragraph()
p_prompt_question.text = "How did Joey's date go?"
p_prompt_question.font.size = Pt(16)

content_tf.add_paragraph() # Spacer

p_prompt_answers = content_tf.add_paragraph()
p_prompt_answers.text = "Please provide short answers:"
p_prompt_answers.font.size = Pt(14)
p_prompt_answers.font.italic = True

# Placeholders for answers
for i in range(1, 4):
    p_ans = content_tf.add_paragraph()
    p_ans.text = f"Short Answer {i}:"
    p_ans.font.size = Pt(12)
    p_ans.level = 1

content_tf.add_paragraph() # Spacer

# Part 2: Instruction
p_part2_header = content_tf.add_paragraph()
run_part2_header = p_part2_header.add_run()
run_part2_header.text = "PART 2"
run_part2_header.font.size = Pt(22)
run_part2_header.font.bold = True
run_part2_header.font.color.rgb = RGBColor(180, 0, 0)

p_part2_text = content_tf.add_paragraph()
p_part2_text.text = "Refer back to the previous section to compare Joey's expectations with the actual outcome."
p_part2_text.font.size = Pt(14)

# --- Image Section ---
image_left = Inches(8.5)
image_top = content_top
image_width = Inches(4.5)
image_height = Inches(4.5)

try:
    # Use the provided placeholder image for Joey
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)
except FileNotFoundError:
    # Fallback if image is missing, though instructions imply it exists
    pass

# Save the presentation
prs.save('output.pptx')