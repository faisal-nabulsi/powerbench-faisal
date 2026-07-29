from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize presentation
prs = Presentation()

# Set 16:9 dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a slide with 'Title and Content' layout
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set Title
title = slide.shapes.title
title.text = "Dynamics of Local and Global Culture"

# Set Content
body = slide.placeholders[1]
tf = body.text_frame
tf.clear()

# Function to add paragraphs
def add_p(text, level=0, bold=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    run = p.runs[0]
    run.font.size = Pt(24)
    if bold:
        run.font.bold = True
    p.space_after = Pt(10)

# 1. Cultural Differentialism
add_p("1. Cultural Differentialism", level=0, bold=True)
add_p("Emphasizes the distinctness and preservation of unique cultural identities.", level=1)

# 2. Cultural Hybridization
add_p("2. Cultural Hybridization", level=0, bold=True)
add_p("Definition: The merging of global and local cultures to form new cultural expressions.", level=1)
add_p("Integration of Cultures: The blending of diverse traditions and global influences.", level=1)
add_p("Glocalization: Adapting global trends to fit specific local cultural contexts.", level=1)

prs.save('output.pptx')