from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p_title = title_frame.paragraphs[0]
p_title.text = "DEFINITION OF TERMS"
p_title.alignment = PP_ALIGN.CENTER
r_title = p_title.runs[0]
r_title.font.size = Pt(40)
r_title.font.bold = True
r_title.font.color.rgb = RGBColor(0x2E, 0x40, 0x57) # Dark Navy

# --- Helper Function to add definition blocks ---
def add_definition_block(slide, left, top, width, height, term, definition):
    # Create textbox
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Term (Heading)
    p_term = tf.paragraphs[0]
    p_term.text = term
    r_term = p_term.runs[0]
    r_term.font.size = Pt(28)
    r_term.font.bold = True
    r_term.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) # Bright Blue
    
    # Definition (Body)
    p_def = tf.add_paragraph()
    p_def.text = definition
    p_def.alignment = PP_ALIGN.LEFT
    r_def = p_def.runs[0]
    r_def.font.size = Pt(20)
    r_def.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey

# --- Layout Constants ---
margin_left = Inches(1)
margin_top = Inches(2.5)
col_width = Inches(5.5)
col_height = Inches(4.0)
gap = Inches(1.333) # Adjusts spacing between the two columns

# --- Media Section ---
media_def = ("Media refers to the channels of mass communication, such as television, radio, movies, "
             "newspapers, magazines, the Internet, and mobile devices. It is the primary means for "
             "disseminating information and entertainment to a broad audience, significantly "
             "influencing public opinion and social trends.")
add_definition_block(slide, margin_left, margin_top, col_width, col_height, "MEDIA", media_def)

# --- Culture Section ---
culture_def = ("Culture is the collective programming of the mind that distinguishes the members of "
               "one group or category of people from another. It encompasses the customary beliefs, "
               "social forms, material traits, language, religion, clothing, food, and arts that "
               "define a specific society or community.")
add_definition_block(slide, margin_left + col_width + gap, margin_top, col_width, col_height, "CULTURE", culture_def)

# Save the presentation
prs.save('output.pptx')