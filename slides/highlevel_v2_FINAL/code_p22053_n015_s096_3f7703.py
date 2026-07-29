from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_elon_musk_slide():
    # Create a new presentation
    prs = Presentation()

    # Set slide dimensions to 16:9
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    title_tf = title_box.text_frame
    title_para = title_tf.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "Lessons that we can learn from Elon Musk's life"
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00) # Black
    title_para.alignment = PP_ALIGN.LEFT

    # --- Bullet Points (Left Side) ---
    bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5))
    bullets_tf = bullets_box.text_frame
    bullets_tf.word_wrap = True

    lessons = [
        "Embrace First Principles: Break down complex problems into fundamental truths rather than reasoning by analogy.",
        "Resilience in Failure: View failure as feedback. Musk cites failure as essential for innovation.",
        "Relentless Work Ethic: Dedication and hard work are often required to achieve extraordinary goals.",
        "Take Calculated Risks: Opportunities for massive impact often come with significant risk.",
        "Follow Your Passion: Pursue missions that inspire you, such as sustainable energy and space exploration."
    ]

    for lesson in lessons:
        para = bullets_tf.add_paragraph()
        para.space_after = Pt(10)
        run = para.add_run()
        run.text = lesson
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Gray

    # --- Image (Right Top) ---
    # Using placeholder 'image.png'
    try:
        slide.shapes.add_picture('image.png', Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.2))
    except FileNotFoundError:
        # Fallback if image is missing, though instructions imply it exists
        placeholder_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.2))
        placeholder_tf = placeholder_box.text_frame
        placeholder_para = placeholder_tf.paragraphs[0]
        placeholder_run = placeholder_para.add_run()
        placeholder_run.text = "[Image: Elon Musk]"
        placeholder_run.font.size = Pt(12)
        placeholder_para.alignment = PP_ALIGN.CENTER
        placeholder_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Quote (Right Bottom) ---
    quote_box = slide.shapes.add_textbox(Inches(7.0), Inches(5.0), Inches(5.5), Inches(1.8))
    quote_tf = quote_box.text_frame
    quote_tf.word_wrap = True
    
    quote_text = "\"Failure is an option here. If things are not failing, you are not innovating enough.\""
    quote_para = quote_tf.paragraphs[0]
    quote_run = quote_para.add_run()
    quote_run.text = quote_text
    quote_run.font.size = Pt(14)
    quote_run.font.italic = True
    quote_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    quote_para.alignment = PP_ALIGN.LEFT

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_elon_musk_slide()