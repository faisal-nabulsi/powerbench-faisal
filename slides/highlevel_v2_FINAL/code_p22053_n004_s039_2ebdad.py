import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_branding_slide():
    # Initialize Presentation
    prs = Presentation()

    # Set Slide Dimensions (16:9 Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Define Airbnb-like colors
    AIRBNB_CORAL = RGBColor(255, 90, 95)  # #FF5A5F
    TEXT_DARK = RGBColor(49, 49, 49)
    TEXT_GRAY = RGBColor(111, 111, 111)
    LINE_GRAY = RGBColor(200, 200, 200)

    # --- Title Section ---
    # Title: "Branding"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.0), Inches(1.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Branding"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    p_title.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(8.0), Inches(0.5))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Airbnb's Brand Evolution"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_GRAY
    p_sub.alignment = PP_ALIGN.LEFT

    # --- Visuals (Image Placeholder) ---
    # Placed on the right side to represent the Bélo symbol visually
    # Coordinates: Left=9.0, Top=0.8, Width=3.5, Height=2.2
    # This places it next to the title area, above the timeline
    img_shape = slide.shapes.add_picture('image.png', Inches(9.0), Inches(0.8), Inches(3.5), Inches(2.2))

    # Caption for the image
    caption_box = slide.shapes.add_textbox(Inches(9.0), Inches(3.1), Inches(3.5), Inches(0.4))
    p_cap = caption_box.text_frame.paragraphs[0]
    p_cap.text = "The Bélo Symbol & Color Palette"
    p_cap.font.size = Pt(14)
    p_cap.font.bold = True
    p_cap.font.color.rgb = AIRBNB_CORAL
    p_cap.alignment = PP_ALIGN.CENTER

    # --- Timeline Section ---
    # Timeline horizontal line
    timeline_y = Inches(4.5)
    timeline_x_start = Inches(0.5)
    timeline_x_end = Inches(12.5)
    
    line_shape = slide.shapes.add_shape(1, timeline_x_start, timeline_y, timeline_x_end - timeline_x_start, Inches(0.1))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = LINE_GRAY
    line_shape.line.fill.background()

    # Function to add a timeline node
    def add_timeline_node(x_pos, year_text, description_text):
        # Node Circle
        circle_shape = slide.shapes.add_shape(20, x_pos, timeline_y - Inches(0.25), Inches(0.5), Inches(0.5))
        circle_shape.fill.solid()
        circle_shape.fill.fore_color.rgb = AIRBNB_CORAL
        circle_shape.line.fill.background()

        # Year Label (Below Circle)
        year_box = slide.shapes.add_textbox(x_pos - Inches(0.5), timeline_y + Inches(0.5), Inches(1.0), Inches(0.4))
        p_year = year_box.text_frame.paragraphs[0]
        p_year.text = year_text
        p_year.font.size = Pt(16)
        p_year.font.bold = True
        p_year.font.color.rgb = TEXT_DARK
        p_year.alignment = PP_ALIGN.CENTER

        # Description Box (Below Year)
        desc_box = slide.shapes.add_textbox(x_pos - Inches(1.0), timeline_y + Inches(1.1), Inches(2.5), Inches(1.5))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = description_text
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_GRAY
        p_desc.alignment = PP_ALIGN.CENTER

    # Add Nodes for Airbnb's Evolution
    # 2008: The "Bosh" era
    add_timeline_node(Inches(2.5), "2008", "\"Bosh\" Logo\nPlayful, hand-drawn aesthetic\nrepresenting startup roots.")

    # 2014: Introduction of Bélo
    add_timeline_node(Inches(6.5), "2014", "Launch of Bélo\nA custom icon encapsulating\nBelonging, Love, and Life.")

    # Present: Refined Identity
    add_timeline_node(Inches(10.5), "Present", "Refined Identity\nSimplified Bélo mark.\nFocus on Coral Red (#FF5A5F).")

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_branding_slide()