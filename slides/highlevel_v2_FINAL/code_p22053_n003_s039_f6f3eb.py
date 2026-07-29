from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide (Layout 6 is typically blank)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Color definitions
    DARK_BLUE = RGBColor(0, 51, 102)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)

    # --- Title ---
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 4"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT

    # --- Instruction ---
    instr_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(10), Inches(0.8))
    tf = instr_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "List the events in Joey’s second date with Sarah in order from 2 - 6,"
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    # --- Checklist ---
    # Generic events for the date since specific context isn't provided
    events = [
        "Met at the cafe at 2:00 PM",
        "Went to the cinema",
        "Had dinner at an Italian restaurant",
        "Took a walk by the river",
        "Bought ice cream",
        "Waved goodbye at the bus stop"
    ]

    check_start_y = 2.8
    check_spacing = 0.8
    check_margin_x = 1.0

    for i, event in enumerate(events):
        y = check_start_y + (i * check_spacing)
        
        # 1. Checkbox
        box_size = 0.4
        box_left = Inches(check_margin_x)
        box_top = Inches(y + 0.05) # Slight vertical adjustment to align with text
        box_w = Inches(box_size)
        box_h = Inches(box_size)
        
        # Add rectangle shape
        shape = slide.shapes.add_shape(1, box_left, box_top, box_w, box_h)
        
        if i == 0:
            # First item completed (Example)
            shape.fill.solid()
            shape.fill.fore_color.rgb = DARK_BLUE
            shape.line.fill.background()
            
            # Add checkmark text inside the box
            # Using a text box overlaid is safer for positioning
            check_text_box = slide.shapes.add_textbox(box_left, box_top, box_w, box_h)
            p = check_text_box.text_frame.paragraphs[0]
            p.text = "✓"
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            check_text_box.line.fill.background()
            
            # Text content: Numbered 1 to show example
            text_content = f"1. {event}"
        else:
            # Unchecked item
            shape.fill.background()
            shape.line.color.rgb = DARK_BLUE
            shape.line.width = Pt(2)
            
            # Text content: Blank for user to number
            text_content = f"__. {event}"
        
        # 2. Event Text
        text_left = Inches(check_margin_x + box_size + 0.2)
        text_top = Inches(y)
        text_width = Inches(8.5)
        text_height = Inches(0.6)
        
        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text_content
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.LEFT

    # --- Image ---
    # Placeholder image 'image.png' at bottom right
    img_w = Inches(2.5)
    img_h = Inches(2.5)
    # 0.5 inch margins
    img_left = prs.slide_width - img_w - Inches(0.5)
    img_top = prs.slide_height - img_h - Inches(0.5)
    
    try:
        slide.shapes.add_picture('image.png', img_left, img_top, img_w, img_h)
    except FileNotFoundError:
        # If image is missing, draw a placeholder shape
        shape = slide.shapes.add_shape(1, img_left, img_top, img_w, img_h)
        shape.fill.background()
        shape.line.color.rgb = RGBColor(150, 150, 150)
        p = shape.text_frame.paragraphs[0]
        p.text = "[ Image of Joey ]"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)

    # Save
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()