from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background color to white (clean look)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245) # Light gray background

# 1. Add the main question text prominently
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "What is filling your bucket today and what’s draining it?"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(51, 51, 51) # Dark gray
p.alignment = PP_ALIGN.CENTER

# 2. Add the illustration (placeholder image)
# Placing it in the center. 
# Slide width 13.333, Image width 6 -> Left = (13.333 - 6) / 2 = 3.6665
# Slide height 7.5, Image height 4 -> Top = (7.5 - 4) / 2 + 1 (offset for title) = 1.75 (approx)
left_img = Inches(3.666)
top_img = Inches(2.2)
width_img = Inches(6)
height_img = Inches(4)

slide.shapes.add_picture('image.png', left_img, top_img, width_img, height_img)

# 3. Add "Filling" section label
# Positioned to the left of the image
left_fill = Inches(0.5)
top_fill = Inches(3.5)
width_fill = Inches(2.5)
height_fill = Inches(1.5)

fill_box = slide.shapes.add_textbox(left_fill, top_fill, width_fill, height_fill)
tf_fill = fill_box.text_frame
tf_fill.word_wrap = True

p_fill = tf_fill.paragraphs[0]
p_fill.text = "Filling"
p_fill.font.size = Pt(32)
p_fill.font.bold = True
p_fill.font.color.rgb = RGBColor(0, 128, 0) # Green
p_fill.alignment = PP_ALIGN.CENTER

# Optional: Add a small arrow or shape pointing to the image for "Filling"
# Simple rectangle arrow pointing right
arrow_x = Inches(2.2)
arrow_y = Inches(3.8)
arrow_w = Inches(1.2)
arrow_h = Inches(0.5)
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, arrow_w, arrow_h)
arrow.fill.solid()
arrow.fill.fore_color.rgb = RGBColor(0, 128, 0)
arrow.line.color.rgb = RGBColor(0, 128, 0)

# 4. Add "Draining" section label
# Positioned to the right of the image
left_drain = Inches(9.5)
top_drain = Inches(3.5)
width_drain = Inches(2.5)
height_drain = Inches(1.5)

drain_box = slide.shapes.add_textbox(left_drain, top_drain, width_drain, height_drain)
tf_drain = drain_box.text_frame
tf_drain.word_wrap = True

p_drain = tf_drain.paragraphs[0]
p_drain.text = "Draining"
p_drain.font.size = Pt(32)
p_drain.font.bold = True
p_drain.font.color.rgb = RGBColor(200, 0, 0) # Red
p_drain.alignment = PP_ALIGN.CENTER

# Optional: Add a small arrow or shape pointing to the image for "Draining"
# Simple rectangle arrow pointing left
arrow2_x = Inches(10)
arrow2_y = Inches(3.8)
arrow2_w = Inches(1.2)
arrow2_h = Inches(0.5)
arrow2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, arrow2_x, arrow2_y, arrow2_w, arrow2_h)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(200, 0, 0)
arrow2.line.color.rgb = RGBColor(200, 0, 0)

# Save the presentation
prs.save('output.pptx')