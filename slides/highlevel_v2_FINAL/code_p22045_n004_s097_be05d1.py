from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Try to find a blank layout, fallback to index 6 (standard for Blank in default themes)
blank_layout = None
for layout in prs.slide_layouts:
    if layout.name.lower() == 'blank':
        blank_layout = layout
        break
if not blank_layout:
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[0] # Fallback to title slide if blank not found

slide = prs.slides.add_slide(blank_layout)

# --- Background Configuration ---
# Setting a dark background for contrast (Dark Slate Blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(7.5), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "ZENEROM UAE"
title_para.font.size = Pt(48)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
title_para.font.name = "Calibri"

# --- Company Description ---
desc_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(7.5), Inches(1.2))
desc_tf = desc_box.text_frame
desc_tf.word_wrap = True
desc_para = desc_tf.paragraphs[0]
desc_para.text = "Zenerom is a pioneer in renewable energy solutions across the UAE. We focus on sustainable power generation, energy storage technologies, and smart grid solutions to empower a greener future for businesses and communities."
desc_para.font.size = Pt(18)
desc_para.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0) # Light Grey
desc_para.font.name = "Calibri"

# --- Services Section ---
# Header
svc_header_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(7.5), Inches(0.6))
svc_header_tf = svc_header_box.text_frame
svc_header_para = svc_header_tf.paragraphs[0]
svc_header_para.text = "Our Services"
svc_header_para.font.size = Pt(28)
svc_header_para.font.bold = True
svc_header_para.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00) # Deep Orange Accent
svc_header_para.font.name = "Calibri"

# List
svc_list_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(7.5), Inches(1.5))
svc_list_tf = svc_list_box.text_frame
svc_list_tf.word_wrap = True

services = [
    "☀️ Solar Energy Systems (PV)",
    "🔋 Battery Energy Storage (BESS)",
    "⚡ Energy Management",
    "🛠️ Site Services & Maintenance"
]

for i, svc in enumerate(services):
    if i == 0:
        para = svc_list_tf.paragraphs[0]
    else:
        para = svc_list_tf.add_paragraph()
    para.text = svc
    para.font.size = Pt(18)
    para.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    para.font.name = "Calibri"
    para.space_after = Pt(6)

# --- Contact Information ---
contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(7.5), Inches(1.0))
contact_tf = contact_box.text_frame
contact_tf.word_wrap = True

# Line 1
c_para1 = contact_tf.paragraphs[0]
c_para1.text = "📍 Dubai, UAE    |    📞 +971 50 123 4567"
c_para1.font.size = Pt(16)
c_para1.font.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
c_para1.font.name = "Calibri"
c_para1.space_after = Pt(4)

# Line 2
c_para2 = contact_tf.add_paragraph()
c_para2.text = "✉️ info@zenerom.com    |    🌐 www.zenerom.com"
c_para2.font.size = Pt(16)
c_para2.font.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
c_para2.font.name = "Calibri"

# --- Imagery ---
# Add placeholder image on the right side
slide.shapes.add_picture('image.png', Inches(9.0), Inches(1.0), Inches(4.0), Inches(5.5))

# Save the presentation
prs.save('output.pptx')