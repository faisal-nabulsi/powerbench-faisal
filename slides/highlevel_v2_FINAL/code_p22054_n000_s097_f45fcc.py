from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # Initialize Presentation
    prs = Presentation()

    # Set Canvas Size to 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # 6 is typically the index for a Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # --- Add Title ---
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_title = tf_title.paragraphs[0]
    p_title.text = "Importance of Skills Beyond Academics"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue
    
    # --- Define Content ---
    skills_data = [
        {
            "title": "Effective Communication",
            "desc": "The ability to convey ideas clearly and listen actively to ensure mutual understanding in professional and personal contexts."
        },
        {
            "title": "Problem Solving and Critical Thinking",
            "desc": "Analyzing complex situations objectively to identify root causes and develop logical, effective solutions."
        },
        {
            "title": "Emotional Intelligence",
            "desc": "Recognizing, understanding, and managing one's own emotions, as well as empathizing with the emotions of others."
        },
        {
            "title": "Adaptability and Resilience",
            "desc": "Flexibility in dealing with unexpected changes and the ability to recover quickly from setbacks or challenges."
        }
    ]

    # --- Layout Configuration (2x2 Grid) ---
    # Margins
    left_margin = Inches(1)
    top_margin_content = Inches(2.2)
    
    # Dimensions
    box_width = Inches(5.6)
    box_height = Inches(2.2)
    gap = Inches(0.5)

    # Positions calculation
    col1_left = left_margin
    col2_left = left_margin + box_width + gap
    
    row1_top = top_margin_content
    row2_top = top_margin_content + box_height + (gap / 2)

    positions = [
        (col1_left, row1_top),
        (col2_left, row1_top),
        (col1_left, row2_top),
        (col2_left, row2_top)
    ]

    # --- Add Content Boxes ---
    for i, skill in enumerate(skills_data):
        left, top = positions[i]
        
        # Add Textbox
        txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
        tf_box = txBox.text_frame
        tf_box.word_wrap = True
        
        # Title of the skill
        p_skill = tf_box.paragraphs[0]
        p_skill.text = skill["title"]
        p_skill.font.size = Pt(24)
        p_skill.font.bold = True
        p_skill.font.color.rgb = RGBColor(0, 84, 136) # Slightly lighter blue
        p_skill.space_after = Pt(12)
        
        # Description of the skill
        p_desc = tf_box.add_paragraph()
        p_desc.text = skill["desc"]
        p_desc.font.size = Pt(18)
        p_desc.font.color.rgb = RGBColor(50, 50, 50) # Dark Grey

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()