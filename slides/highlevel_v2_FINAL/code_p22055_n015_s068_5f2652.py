from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide():
    prs = Presentation()

    # Set slide width and height to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Create a dark market-themed background
    # Since no specific image is requested for the background itself, we use a dark solid color
    # to represent the "dark market" theme.
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2B, 0x2B, 0x2B) # Dark Charcoal

    # 2. Overlay a red section on the top right
    # Dimensions: Width 4 inches, Height 1.5 inches
    # Position: Top-Right corner
    red_section_width = Inches(4.0)
    red_section_height = Inches(1.5)
    red_section_left = prs.slide_width - red_section_width
    red_section_top = Inches(0)

    red_box = slide.shapes.add_shape(
        1, # ID for Rectangle
        red_section_left,
        red_section_top,
        red_section_width,
        red_section_height
    )

    # Style the red section
    red_box.fill.solid()
    red_box.fill.fore_color.rgb = RGBColor(0xE6, 0x00, 0x00) # Bright Red
    red_box.line.fill.background() # Remove border

    # 3. Add the text "WHAT IS MARKET?"
    # Position: Center-Left area to balance the red section
    # Dimensions: Width 6 inches, Height 1.5 inches
    text_left = Inches(1.0)
    text_top = Inches(0.75)
    text_width = Inches(7.5)
    text_height = Inches(1.5)

    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = text_box.text_frame
    tf.auto_size = None

    # Add paragraph and formatting
    paragraph = tf.paragraphs[0]
    run = paragraph.add_run()
    run.text = "WHAT IS MARKET?"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text
    run.font.name = 'Arial' # Clear font
    
    # Align text
    paragraph.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True

    # Save the presentation
    prs.save('output.pptx')

if __name__ == '__main__':
    create_slide()