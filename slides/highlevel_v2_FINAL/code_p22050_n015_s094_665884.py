from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_supply_chain_slide():
    # Create presentation and set slide size
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
    title_tf = title_box.text_frame
    title_para = title_tf.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "From Factory to the Market"
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # Define steps data
    steps = [
        {
            "title": "1. Producer",
            "desc": "Source of raw materials and initial resources.",
            "color": RGBColor(0x00, 0x70, 0xC0), # Blue
            "left": Inches(0.9)
        },
        {
            "title": "2. Manufacturer",
            "desc": "Factory processing transforms materials into goods.",
            "color": RGBColor(0x00, 0xB0, 0x50), # Green
            "left": Inches(3.9)
        },
        {
            "title": "3. Distributor",
            "desc": "Manages logistics, storage, and bulk transport.",
            "color": RGBColor(0xFF, 0xA5, 0x00), # Orange
            "left": Inches(6.9)
        },
        {
            "title": "4. Retailer",
            "desc": "Sells final products directly to consumers.",
            "color": RGBColor(0xE0, 0x1B, 0x20), # Red
            "left": Inches(9.9)
        }
    ]
    
    box_width = Inches(2.5)
    box_height = Inches(2.5)
    box_top = Inches(2.0)
    
    for step in steps:
        # Add Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, step["left"], box_top, box_width, box_height)
        
        # Set Fill Color
        shape.fill.solid()
        shape.fill.fore_color.rgb = step["color"]
        
        # Remove Lines/Outline
        shape.line.fill.background()
        
        # Add Text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # Vertical alignment logic isn't direct on tf, so we add paragraphs
        
        # Title Paragraph
        p_title = tf.paragraphs[0]
        r_title = p_title.add_run()
        r_title.text = step["title"]
        r_title.font.size = Pt(20)
        r_title.font.bold = True
        r_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p_title.space_after = Pt(6)
        
        # Description Paragraph
        p_desc = tf.add_paragraph()
        r_desc = p_desc.add_run()
        r_desc.text = step["desc"]
        r_desc.font.size = Pt(14)
        r_desc.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Center text horizontally roughly by adjusting margins/padding if needed, 
        # but simple text frame usually centers left-aligned text fine enough for boxes.
        # To center text inside the shape, we can't easily do it with python-pptx without 
        # calculating runs, but left alignment is standard. 
        # To mimic centering, we can set the paragraph alignment to CENTER.
        p_title.alignment = pp.PP_ALIGN.CENTER # Need to import PP_ALIGN? No, just use int or enum if available.
        # Actually PP_ALIGN is in pptx.enum.text. Let's import it.
        # But to keep imports minimal and standard, I will skip centering or use standard module.
        # Let's stick to left alignment which is safe, or add import.
        
    # Re-running loop to fix alignment just for cleanliness
    # I'll add the import for PP_ALIGN to ensure nice centering
    from pptx.enum.text import PP_ALIGN
    
    # Clear previous shapes logic for re-implementation? 
    # No, I'll just rewrite the script block properly inside the function to be robust.
    pass

# Corrected Script Implementation
def generate_slide():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
    title_tf = title_shape.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_r = title_p.add_run()
    title_r.text = "From Factory to the Market"
    title_r.font.size = Pt(36)
    title_r.font.bold = True
    title_r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    steps = [
        ("1. Producer", "Generates raw materials for the supply chain.", RGBColor(0x00, 0x70, 0xC0), Inches(0.9)),
        ("2. Manufacturer", "Processes materials into finished products.", RGBColor(0x00, 0xB0, 0x50), Inches(3.9)),
        ("3. Distributor", "Handles logistics, storage, and nationwide transport.", RGBColor(0xFF, 0xA5, 0x00), Inches(6.9)),
        ("4. Retailer", "Sells products directly to the end consumer.", RGBColor(0xE0, 0x1B, 0x20), Inches(9.9)),
    ]
    
    box_w = Inches(2.5)
    box_h = Inches(2.5)
    box_y = Inches(2.0)
    
    for title_text, desc_text, color, left_pos in steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, box_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title_text
        r1.font.size = Pt(20)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p1.space_after = Pt(8)
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc_text
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save('output.pptx')

if __name__ == "__main__":
    generate_slide()