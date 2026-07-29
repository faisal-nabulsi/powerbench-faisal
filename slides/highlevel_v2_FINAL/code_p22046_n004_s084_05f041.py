from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a slide using the 'Title and Content' layout (index 1)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set the title of the slide
title_shape = slide.shapes.title
title_shape.text = "What is Blockchain?"

# Access the content placeholder (index 1 usually refers to the content box in Title and Content layout)
content_placeholder = slide.placeholders[1]
text_frame = content_placeholder.text_frame

# Clear existing paragraphs to ensure clean insertion
text_frame.clear()

# Define the content for the bullet points
# We use a list of dictionaries to structure headers and sub-points
slide_content = [
    {
        "title": "Definition",
        "details": "A distributed, decentralized digital ledger that records transactions across many computers. It ensures data integrity and security without a central authority."
    },
    {
        "title": "Structure",
        "details": [
            "Blocks: Containers storing transaction data, timestamps, and the cryptographic hash of the previous block.",
            "Chain: Blocks are linked chronologically; altering one block changes its hash, invalidating all subsequent blocks."
        ]
    },
    {
        "title": "Characteristics",
        "details": [
            "Decentralized: No single point of control over the network.",
            "Immutable: Records cannot be altered retroactively without consensus.",
            "Transparent: Visible to all network participants.",
            "Secure: Protected by advanced cryptographic algorithms."
        ]
    }
]

# Function to add a paragraph with specific text and level
def add_paragraph(text_frame, text, level, font_size):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    # Apply font size to the paragraph's runs
    for run in p.runs:
        run.font.size = font_size

# Populate the text frame with content
for item in slide_content:
    # Add the main bullet point (Header)
    add_paragraph(text_frame, item["title"], level=0, font_size=Pt(22))
    
    # Add sub-points or description
    if isinstance(item["details"], list):
        for point in item["details"]:
            add_paragraph(text_frame, point, level=1, font_size=Pt(18))
    else:
        add_paragraph(text_frame, item["details"], level=1, font_size=Pt(18))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')