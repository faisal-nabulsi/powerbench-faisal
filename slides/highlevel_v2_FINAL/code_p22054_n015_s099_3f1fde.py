from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_biography_slide():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # --- Title Section ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
    title_tf = title_box.text_frame
    title_para = title_tf.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "Early life of Elon Musk"
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    
    # --- Content Section ---
    
    # 1. Text Box (Left side)
    # Position: Left 0.5, Top 2.0, Width 6.5, Height 5.0
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.5), Inches(5.0))
    text_tf = text_box.text_frame
    text_tf.word_wrap = True
    
    bullet_points = [
        ("Birth", "June 28, 1971 in Pretoria, South Africa."),
        ("Family", "Parents Errol Musk (electromechanical engineer) and Maye Musk (dietitian/model). He has two older brothers, Kim and Tosca."),
        ("Education", "Graduated from the University of Pennsylvania with degrees in Physics and Economics. Briefly attended Stanford University for a PhD before dropping out."),
        ("Early Interests", "Developed a passion for computing at a young age. Learned to program in BASIC at age 10 and sold his first video game, 'Bust-a-Move', at age 12.")
    ]
    
    for label, detail in bullet_points:
        p = text_tf.add_paragraph()
        p.space_after = Pt(14)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        # Bold label
        run_label = p.add_run()
        run_label.text = f"{label}: "
        run_label.font.bold = True
        run_label.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) # Blue-ish header color
        
        # Detail text
        run_detail = p.add_run()
        run_detail.text = detail
        run_detail.font.bold = False
        
        # Add bullet symbol manually if desired, or rely on paragraph style. 
        # Using a simple hyphen for safety across themes.
        p.text = f"• {p.text}"
        
    # 2. Image (Right side)
    # Position: Left 7.2, Top 2.0, Width 5.6, Height 5.0
    # Using the required placeholder image
    try:
        image = slide.shapes.add_picture('image.png', Inches(7.2), Inches(2.0), Inches(5.6), Inches(5.0))
    except FileNotFoundError:
        # Fallback if image is missing, add a note or skip
        print("Warning: image.png not found. Skipping image addition.")
        
    # Save the presentation
    prs.save('output.pptx')
    print("Presentation saved to 'output.pptx'")

if __name__ == "__main__":
    create_biography_slide()