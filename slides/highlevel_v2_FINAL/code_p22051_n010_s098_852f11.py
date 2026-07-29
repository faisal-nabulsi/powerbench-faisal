from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Layout index 6 is typically Blank in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title ---
# Add a text box for the title
# Positioned at the top, spanning most of the width
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = False
title_para = title_tf.paragraphs[0]
title_para.text = "Look Around"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50) # Dark Blue-Grey
title_para.alignment = PP_ALIGN.CENTER

# --- Left Content: Prompt for Producers ---
# Add text box on the left side
prompt_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.8), Inches(4.5))
prompt_tf = prompt_box.text_frame
prompt_tf.word_wrap = True

# Header
p_header = prompt_tf.paragraphs[0]
p_header.text = "Prompt for Producers:"
p_header.font.size = Pt(24)
p_header.font.bold = True
p_header.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
p_header.space_after = Pt(15)

# Body Text
p_body = prompt_tf.add_paragraph()
p_body.text = (
    "1. OBSERVE: Scan the environment and the people present.\n\n"
    "2. IDENTIFY: Note body language, energy levels, and key needs.\n\n"
    "3. ADAPT: Adjust your production approach to resonate with the immediate context."
)
p_body.font.size = Pt(18)
p_body.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p_body.line_spacing = Pt(28)

# --- Right Content: Illustrative Icon ---
# Add the placeholder image 'image.png' on the right side
# Positioned to align vertically with the text block
slide.shapes.add_picture('image.png', Inches(7), Inches(2.5), Inches(5.5), Inches(4.5))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')