from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation and set slide size to 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to have full control over layout
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- 1. Teal Header ---
# Create a rectangle shape spanning the width of the slide for the header
header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.5))

# Set the fill color to Teal
header_fill = header_shape.fill
header_fill.solid()
header_fill.fore_color.rgb = RGBColor(0, 128, 128)  # Teal RGB

# Add the title "Values" to the header
header_text_frame = header_shape.text_frame
header_text_frame.word_wrap = True
p_header = header_text_frame.paragraphs[0]
p_header.text = "Values"
p_header.font.size = Pt(48)
p_header.font.bold = True
p_header.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast
p_header.alignment = PP_ALIGN.CENTER

# --- 2. Content: List of Values ---
# Define the six values and their descriptions
values_data = [
    ("Integrity", "We act with honesty and transparency in all our dealings."),
    ("Innovation", "We embrace creativity and continuous improvement."),
    ("Collaboration", "We succeed by working together and supporting one another."),
    ("Respect", "We value diverse perspectives and treat everyone with dignity."),
    ("Excellence", "We strive for the highest quality in everything we do."),
    ("Sustainability", "We are committed to a sustainable future for our community.")
]

# Add a text box on the left side for the list
text_left = Inches(0.5)
text_top = Inches(2.0)
text_width = Inches(6.5)
text_height = Inches(5.0)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
tf = text_box.text_frame
tf.word_wrap = True

for i, (title, description) in enumerate(values_data):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    
    p.space_after = Pt(15)  # Spacing between list items

    # Add Value Title (Bold, Teal color)
    run_title = p.add_run()
    run_title.text = f"{i+1}. {title}\n"
    run_title.font.bold = True
    run_title.font.size = Pt(18)
    run_title.font.color.rgb = RGBColor(0, 128, 128)

    # Add Value Description (Regular, Black color)
    run_desc = p.add_run()
    run_desc.text = description
    run_desc.font.size = Pt(14)
    run_desc.font.color.rgb = RGBColor(0, 0, 0)

# --- 3. Images (Right Side) ---
# Place the placeholder image on the right side of the slide
img_left = Inches(7.5)
img_top = Inches(2.0)
img_width = Inches(5.5)
img_height = Inches(5.0)

# Add the picture using the available placeholder file
try:
    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
except FileNotFoundError:
    pass  # Handle case where image is missing, though instructions say it is available

# Save the presentation to 'output.pptx'
prs.save('output.pptx')