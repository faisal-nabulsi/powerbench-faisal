from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # Initialize the presentation
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Add Title "PART 3" ---
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "PART 3"
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x00, 0x6B, 0x3F) # Dark Green

    # --- Add Instructions ---
    instr_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    instr_tf = instr_box.text_frame
    instr_tf.word_wrap = True
    instr_p = instr_tf.paragraphs[0]
    instr_run = instr_p.add_run()
    instr_run.text = "Choose true (T), false (F), or not given (N) based on video content."
    instr_run.font.size = Pt(18)
    instr_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Add Conversation Content ---
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6), Inches(4))
    text_tf = text_box.text_frame
    text_tf.word_wrap = True

    # List of conversation lines and statements
    lines = [
        ("Joey:", "I just got a waitress job at Central Perk."),
        ("Friends:", "That's awesome, Joey! Do you have an outfit yet?"),
        ("Joey:", "I'm thinking of wearing my 'How You Doin'?' shirt."),
        ("Friends:", "Joey, you can't wear that to work!"),
        ("Joey:", "Why not? It's my lucky shirt."),
        ("Friends:", "Because the boss said no distracting customers."),
        ("", ""), # Empty line for spacing
        ("Task:", "Based on the conversation, determine if the statements are True, False, or Not Given."),
        ("1.", "Joey works at a restaurant."),
        ("2.", "Joey owns a 'How You Doin'?' shirt."),
        ("3.", "The boss likes the shirt."),
    ]

    for speaker, text_content in lines:
        p = text_tf.add_paragraph()
        p.space_after = Pt(8)
        
        if speaker == "":
            p.space_before = Pt(12)
            continue
            
        run_speaker = p.add_run()
        run_speaker.text = speaker + " "
        run_speaker.font.bold = True
        run_speaker.font.size = Pt(16)
        
        if text_content:
            run_text = p.add_run()
            run_text.text = text_content
            run_text.font.size = Pt(16)

    # --- Add Image ---
    # Placing image on the right side
    try:
        slide.shapes.add_picture('image.png', 
                                 left=Inches(7.5), 
                                 top=Inches(1.2), 
                                 width=Inches(5), 
                                 height=Inches(4))
    except FileNotFoundError:
        print("Warning: 'image.png' not found. Image skipped.")

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()