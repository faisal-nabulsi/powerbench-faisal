from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Main Question
# Add a text box for the question
question_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.333), Inches(1.8))
tf = question_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "What is filling your bucket today and what’s draining it?"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102) # Navy Blue

# 2. Illustration (Image Placeholder)
# Center the image below the title
# Image dimensions: 5.5 inches wide, 4.5 inches tall
img_left = Inches(3.9165) # (13.333 - 5.5) / 2
img_top = Inches(2.5)
img_width = Inches(5.5)
img_height = Inches(4.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# 3. Labeled Sections
# We want "Filling" on the left and "Draining" on the right of the image/bucket concept.
# However, since the image is central, let's place the labels near the image to associate them.
# Or, strictly "labeled sections" implies areas of the slide or distinct groups.
# Let's place them above the image on the sides to point to the image.

# Filling Label
filling_left = Inches(1)
filling_top = Inches(2.5)
filling_width = Inches(2.5)
filling_height = Inches(1.5)

filling_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, filling_left, filling_top, filling_width, filling_height)
filling_shape.fill.solid()
filling_shape.fill.fore_color.rgb = RGBColor(200, 255, 200) # Light Green
filling_shape.line.color.rgb = RGBColor(0, 128, 0)
filling_shape.line.width = Pt(2)

filling_tf = filling_shape.text_frame
filling_tf.word_wrap = True
filling_p = filling_tf.paragraphs[0]
filling_p.text = "Filling"
filling_p.alignment = PP_ALIGN.CENTER
filling_p.font.size = Pt(28)
filling_p.font.bold = True
filling_p.font.color.rgb = RGBColor(0, 100, 0)

# Draining Label
draining_left = Inches(9.833) # 13.333 - 1 - 2.5
draining_top = Inches(2.5)
draining_width = Inches(2.5)
draining_height = Inches(1.5)

draining_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, draining_left, draining_top, draining_width, draining_height)
draining_shape.fill.solid()
draining_shape.fill.fore_color.rgb = RGBColor(255, 200, 200) # Light Red
draining_shape.line.color.rgb = RGBColor(200, 0, 0)
draining_shape.line.width = Pt(2)

draining_tf = draining_shape.text_frame
draining_tf.word_wrap = True
draining_p = draining_tf.paragraphs[0]
draining_p.text = "Draining"
draining_p.alignment = PP_ALIGN.CENTER
draining_p.font.size = Pt(28)
draining_p.font.bold = True
draining_p.font.color.rgb = RGBColor(150, 0, 0)

# Add arrows pointing to the center image
# Left Arrow pointing right (from Filling box)
# Right Arrow pointing left (from Draining box)

# Arrow 1 (Left side)
# Start at right edge of Filling box, point to left edge of Image
arrow_left_x = filling_left + filling_width # 3.5
arrow_left_y = filling_top + (filling_height / 2)
arrow_left_w = Inches(0.5)
arrow_left_h = Inches(0.3)

# Using MSO_SHAPE.RIGHT_ARROW is not in MSO_SHAPE enum directly by name in some versions? 
# MSO_SHAPE is an enum in pptx.enum.shapes.
# Common shapes: MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL.
# Arrows: MSO_SHAPE.RIGHT_ARROW (ID 143 usually? No, let's check docs or use standard integer if unsure, but enum is safer).
# In python-pptx, MSO_SHAPE.RIGHT_ARROW exists.

try:
    slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, 
        Inches(3.6), 
        Inches(3.0), 
        Inches(0.4), 
        Inches(0.3)
    )
    slide.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW, 
        Inches(9.3), 
        Inches(3.0), 
        Inches(0.4), 
        Inches(0.3)
    )
except:
    pass # Fallback if attribute error, though it should exist.

# Save
prs.save('output.pptx')