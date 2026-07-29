from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KEY reasons why academic success is important in society"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Function to add a section with icon and text
def add_section(slide, left, top, width, height, title, description):
    # Background box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.fill.background()
    
    # Icon (using placeholder image.png)
    img_w = Inches(1.5)
    img_h = Inches(1.5)
    img_l = left + (width - img_w) / 2
    img_t = top + Inches(0.3)
    slide.shapes.add_picture('image.png', img_l, img_t, img_w, img_h)
    
    # Section Title
    t_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(2.0), width - Inches(0.4), Inches(0.5))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_p = t_tf.paragraphs[0]
    t_p.text = title
    t_p.font.size = Pt(18)
    t_p.font.bold = True
    t_p.font.color.rgb = RGBColor(0, 0, 0)
    t_p.alignment = PP_ALIGN.CENTER
    
    # Description
    d_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(2.6), width - Inches(0.4), Inches(1.5))
    d_tf = d_box.text_frame
    d_tf.word_wrap = True
    d_p = d_tf.paragraphs[0]
    d_p.text = description
    d_p.font.size = Pt(14)
    d_p.font.color.rgb = RGBColor(50, 50, 50)
    d_p.alignment = PP_ALIGN.CENTER

# Section 1
add_section(slide, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), 
            "Personal Growth and Development", 
            "Fosters critical thinking, self-discipline, and a lifelong love for learning.")

# Section 2
add_section(slide, Inches(4.8), Inches(2.0), Inches(3.8), Inches(4.5), 
            "Employability and Career Opportunities", 
            "Opens doors to diverse career paths, increasing job security and earning potential.")

# Section 3
add_section(slide, Inches(9.1), Inches(2.0), Inches(3.8), Inches(4.5), 
            "Economic Impact", 
            "Drives innovation and productivity, leading to economic growth and stability.")

prs.save('output.pptx')