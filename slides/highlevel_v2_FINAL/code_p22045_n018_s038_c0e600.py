from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a slide with a title layout (Layout 5 is Title Only)
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Set the slide title
slide.shapes.title.text = "PART 1"

# Part 1: Discussion Prompt
box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(2))
tf1 = box1.text_frame
p1 = tf1.paragraphs[0]
run1 = p1.add_run()
run1.text = "Discussion Prompt:\nHow did Joey's date go? Please provide short answers."
run1.font.size = Pt(18)

# Part 2: Instruction
box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(6), Inches(2))
tf2 = box2.text_frame

# Part 2 Heading
p2_head = tf2.add_paragraph()
run2_head = p2_head.add_run()
run2_head.text = "PART 2"
run2_head.font.size = Pt(24)
run2_head.font.bold = True

# Part 2 Body
p2_body = tf2.add_paragraph()
run2_body = p2_body.add_run()
run2_body.text = "Refer back to the previous section for comparison."
run2_body.font.size = Pt(18)

# Image of Joey
slide.shapes.add_picture('image.png', Inches(7), Inches(1.2), Inches(6), Inches(4.5))

prs.save('output.pptx')