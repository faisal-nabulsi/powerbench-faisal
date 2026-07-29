from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # Add Title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    p_title = title_frame.paragraphs[0]
    p_title.text = "Dynamics of Local and Global Culture"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(33, 33, 33)

    # Add Content
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(4))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    
    # Remove default empty paragraph to avoid extra space if adding new ones
    p0 = content_frame.paragraphs[0]
    p0.text = ""

    # Bullet 1
    p1 = content_frame.add_paragraph()
    p1.text = "• Cultural Differentialism"
    p1.font.size = Pt(26)
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.space_after = Pt(10)

    # Bullet 2
    p2 = content_frame.add_paragraph()
    p2.text = "• Cultural Hybridization"
    p2.font.size = Pt(26)
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.space_after = Pt(6)

    # Sub-bullets for Hybridization
    # We manually indent and use a smaller bullet or dash
    sub1 = content_frame.add_paragraph()
    sub1.text = "    - Definition: The blending of distinct cultural elements to form new, mixed cultural forms."
    sub1.font.size = Pt(20)
    sub1.font.color.rgb = RGBColor(60, 60, 60)

    sub2 = content_frame.add_paragraph()
    sub2.text = "    - Integration of cultures: The process where local and global traits merge and coexist."
    sub2.font.size = Pt(20)
    sub2.font.color.rgb = RGBColor(60, 60, 60)

    sub3 = content_frame.add_paragraph()
    sub3.text = '    - "glocalization": The adaptation of global phenomena to fit local contexts and preferences.'
    sub3.font.size = Pt(20)
    sub3.font.color.rgb = RGBColor(60, 60, 60)

    prs.save('output.pptx')

if __name__ == "__main__":
    create_pptx()