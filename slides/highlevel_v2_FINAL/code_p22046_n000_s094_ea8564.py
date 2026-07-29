from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add title "PART 2"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.0))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_run = title_frame.add_paragraph()
title_run.text = "PART 2"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 0, 0)
title_run.alignment = PP_ALIGN.LEFT

# Add instruction text
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(1.2))
instruction_frame = instruction_box.text_frame
instruction_frame.word_wrap = True
instruction_para = instruction_frame.add_paragraph()
instruction_para.text = "Read the provided sentences below and categorize the examples of Past Simple and Past Continuous in the table. Ensure you include the dialogue excerpts from the conversation."
instruction_para.font.size = Pt(16)
instruction_para.font.color.rgb = RGBColor(50, 50, 50)

# Add dialogue excerpts section
dialogue_header = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(0.5))
dh_frame = dialogue_header.text_frame
dh_frame.word_wrap = True
dh_para = dh_frame.add_paragraph()
dh_para.text = "Dialogue Excerpts:"
dh_para.font.size = Pt(18)
dh_para.font.bold = True
dh_para.font.color.rgb = RGBColor(0, 51, 102)

dialogue_text = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1.5))
dt_frame = dialogue_text.text_frame
dt_frame.word_wrap = True

dialogue_examples = [
    '"I was studying when my friend called me."',
    '"She walked to the park and saw an old friend."',
    '"What were you doing at 8 PM yesterday?"',
    '"We visited the museum last weekend."'
]

for ex in dialogue_examples:
    p = dt_frame.add_paragraph()
    p.text = ex
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.space_before = Pt(6)

# Create table for categorization
rows = 1 + len(dialogue_examples)  # Header row + example rows
cols = 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.9), Inches(12.333), Inches(2.2))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(5.0)
table.columns[1].width = Inches(3.667)
table.columns[2].width = Inches(3.666)

# Add table headers
headers = ["Dialogue Excerpt", "Past Simple", "Past Continuous"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    # Set cell background color to dark blue
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '003366')

# Populate table rows with dialogue excerpts
for row_idx, example in enumerate(dialogue_examples):
    # First column: full dialogue excerpt
    cell = table.cell(row_idx + 1, 0)
    cell.text = example
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_ALIGN.LEFT
    
    # Second and third columns: empty for user to fill in
    for col_idx in [1, 2]:
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = ""
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)

prs.save('output.pptx')