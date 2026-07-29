from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the Fluentize logo (using 'image.png' as placeholder) at the top left corner
# Coordinates: x=0.5, y=0.5
# Size: 2.5 inches wide, 1 inch high
logo_left = Inches(0.5)
logo_top = Inches(0.5)
logo_width = Inches(2.5)
logo_height = Inches(1.0)
slide.shapes.add_picture('image.png', logo_left, logo_top, logo_width, logo_height)

# Add the "THANK YOU!" title
# Requirements: Centered horizontally, positioned at the bottom
# Text box dimensions
tb_width = Inches(10.0)
tb_height = Inches(1.5)

# Calculate coordinates
# Center horizontally: (Slide Width - Box Width) / 2
tb_left = (prs.slide_width - tb_width) / 2

# Position at the bottom: Slide Height - Box Height - Padding
padding_bottom = Inches(0.5)
tb_top = prs.slide_height - tb_height - padding_bottom

# Create text box
text_box = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Add text and format
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = "THANK YOU!"
run.font.size = Pt(48)
run.font.bold = True
paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')