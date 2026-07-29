from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation and set dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title "VIEWING ACTIVITY"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.LEFT
title_run = title_para.add_run()
title_run.text = "VIEWING ACTIVITY"
title_run.font.size = Pt(32)
title_run.font.bold = True

# Add Video Camera Icon (using placeholder image.png)
# Positioned on the right side, vertically centered with the title
slide.shapes.add_picture('image.png', Inches(11.2), Inches(0.6), Inches(1.8), Inches(1.2))

# Add Subtitle "Friends | Joey Doesn't Share Food!"
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(10), Inches(0.8))
sub_frame = sub_box.text_frame
sub_frame.word_wrap = True
sub_para = sub_frame.paragraphs[0]
sub_para.alignment = PP_ALIGN.LEFT
sub_run = sub_para.add_run()
sub_run.text = "Friends | Joey Doesn't Share Food!"
sub_run.font.size = Pt(20)

# Save the presentation
prs.save('output.pptx')