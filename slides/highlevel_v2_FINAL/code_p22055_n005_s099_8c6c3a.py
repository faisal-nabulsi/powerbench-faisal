from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Initialize the presentation
prs = Presentation()

# 2. Set slide dimensions to 16:9 Widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 3. Add a blank slide (Layout index 6 is typically 'Blank')
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 4. Set the background color to Green
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(34, 139, 34)  # Forest Green

# 5. Add the Title
# Position: Top Center
title_left = Inches(1.0)
title_top = Inches(0.5)
title_width = Inches(11.333)
title_height = Inches(1.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Format Title Text
p = title_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Human Impact or Influence on Environment"
run.font.size = Pt(42)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)  # White for contrast
run.font.name = "Arial"  # Clear, readable font

# 6. Add Explanatory Text
# Position: Left side below title
text_left = Inches(0.5)
text_top = Inches(2.2)
text_width = Inches(5.5)
text_height = Inches(4.5)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Add Text Content
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = (
    "Human activities significantly influence the natural world through industrialization, "
    "deforestation, and urbanization. These actions can degrade ecosystems and contribute "
    "to climate change. However, sustainable practices and conservation efforts are essential "
    "to mitigate negative impacts and preserve the environment for future generations."
)
run.font.size = Pt(22)
run.font.color.rgb = RGBColor(255, 255, 255)  # White text
run.font.name = "Arial"

# 7. Add Image
# Position: Right side
# Using the provided placeholder 'image.png'
slide.shapes.add_picture(
    'image.png', 
    Inches(6.5), 
    Inches(2.0), 
    Inches(6.3), 
    Inches(4.8)
)

# 8. Save the presentation
prs.save('output.pptx')