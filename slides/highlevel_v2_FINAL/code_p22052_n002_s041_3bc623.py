from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()

# Set 16:9 Widescreen Dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Slide Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KEY Reasons Why Academic Success Is Important in Society"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# --- Section Content Definition ---
sections = [
    {
        "title": "Innovation and Technological Advancement",
        "text": "Academic success fuels research and development, driving technological breakthroughs that solve complex societal problems and improve efficiency.",
        "x": 0.5
    },
    {
        "title": "Social Mobility and Equality",
        "text": "Education acts as a powerful equalizer, providing individuals with the skills to overcome socioeconomic barriers and fostering a more inclusive society.",
        "x": 4.84
    },
    {
        "title": "Civic Engagement and Informed Citizenship",
        "text": "Higher education cultivates critical thinking and civic awareness, enabling citizens to participate actively in democracy and contribute to community well-being.",
        "x": 9.18
    }
]

# Layout constants
col_width = 3.64
icon_size = 1.5
section_top_y = 1.6
section_height = 4.8

# --- Processing Sections ---
for data in sections:
    x = data["x"]
    
    # 1. Background Card (Rounded Rectangle)
    left = Inches(x)
    top = Inches(section_top_y)
    width = Inches(col_width)
    height = Inches(section_height)
    
    try:
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.color.rgb = RGBColor(200, 200, 200)
        bg_shape.line.width = Pt(1)
    except Exception:
        # Fallback if shape addition fails (unlikely)
        pass
    
    # 2. Icon (Using placeholder image)
    icon_x = x + (col_width - icon_size) / 2
    icon_y = section_top_y + 0.3
    try:
        slide.shapes.add_picture('image.png', Inches(icon_x), Inches(icon_y), Inches(icon_size), Inches(icon_size))
    except Exception:
        pass
    
    # 3. Section Title
    title_y = icon_y + icon_size + 0.4
    title_h = 0.6
    tb_title = slide.shapes.add_textbox(Inches(x + 0.2), Inches(title_y), Inches(col_width - 0.4), Inches(title_h))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = data["title"]
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 102, 153)
    p_title.alignment = PP_ALIGN.CENTER
    
    # 4. Explanatory Text
    text_y = title_y + title_h + 0.2
    text_h = section_height - (text_y - section_top_y)
    tb_text = slide.shapes.add_textbox(Inches(x + 0.3), Inches(text_y), Inches(col_width - 0.6), Inches(text_h))
    tf_text = tb_text.text_frame
    tf_text.word_wrap = True
    p_text = tf_text.paragraphs[0]
    p_text.text = data["text"]
    p_text.font.size = Pt(14)
    p_text.font.color.rgb = RGBColor(80, 80, 80)
    p_text.alignment = PP_ALIGN.LEFT

# Save the presentation
prs.save('output.pptx')