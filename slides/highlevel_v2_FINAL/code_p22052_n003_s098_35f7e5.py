from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Find a blank layout or default to index 6
slide_layout = None
for layout in prs.slide_layouts:
    if layout.name == 'Blank':
        slide_layout = layout
        break
if slide_layout is None:
    slide_layout = prs.slide_layouts[6]

slide = prs.slides.add_slide(slide_layout)

# Set a clear white background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# --- Add Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "PRONUNCIATION ACTIVITY"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_para.alignment = PP_ALIGN.CENTER

# --- Add Question ---
question_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1))
question_frame = question_box.text_frame
question_frame.word_wrap = True
question_para = question_frame.paragraphs[0]
question_para.text = "What is the general rule for the pronunciation of the letter 'y' at the end of a word?"
question_para.font.size = Pt(26)
question_para.font.color.rgb = RGBColor(40, 40, 40)
question_para.alignment = PP_ALIGN.LEFT

# --- Add Options ---
options_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(3))
options_frame = options_box.text_frame
options_frame.word_wrap = True

# Option a
p_a = options_frame.paragraphs[0]
p_a.text = "a. It is always pronounced as a long 'i' sound."
p_a.font.size = Pt(22)
p_a.font.color.rgb = RGBColor(60, 60, 60)
p_a.space_after = Pt(15)

# Option b
p_b = options_frame.add_paragraph()
p_b.text = "b. It is always pronounced as a short 'e' sound."
p_b.font.size = Pt(22)
p_b.font.color.rgb = RGBColor(60, 60, 60)
p_b.space_after = Pt(15)

# Option c
p_c = options_frame.add_paragraph()
p_c.text = "c. It can be pronounced as a long 'i' or a long 'e' sound depending on the word."
p_c.font.size = Pt(22)
p_c.font.color.rgb = RGBColor(60, 60, 60)

# Save the presentation
prs.save('output.pptx')