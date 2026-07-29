from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set the slide size to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "Content"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Add the Body Text (Bullet points on the left side)
# Positioning: Left aligned relative to the slide, taking up about half the width
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(4.8))
body_frame = body_box.text_frame
body_frame.word_wrap = True

# List of key topics about Elon Musk
topics = [
    "Early Life: Born in Pretoria, South Africa, in 1971.",
    "Early Ventures: Co-founded Zip2 and PayPal (formerly X.com).",
    "SpaceX: Founder and CEO, developed reusable rocket technology (Falcon 9, Starship).",
    "Tesla: CEO and Product Architect, accelerating the world's transition to sustainable energy.",
    "Neuralink: Developing brain-computer interface technology.",
    "The Boring Company: Aiming to reduce traffic congestion through tunnel networks.",
    "X (formerly Twitter): Acquired the social media platform in 2022.",
    "Vision: Long-term goals include colonizing Mars and sustainable energy for Earth."
]

# Add each topic as a bullet point
for i, topic in enumerate(topics):
    if i == 0:
        para = body_frame.paragraphs[0]
    else:
        para = body_frame.add_paragraph()
    
    para.text = topic
    para.level = 0
    para.space_after = Pt(8)
    
    # Formatting the run
    run = para.runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')