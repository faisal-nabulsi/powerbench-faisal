from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set the slide canvas size to 16:9 widescreen (13.333" x 7.5")
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background color to solid blue
background = slide.background
fill = background.fill
fill.solid()
# Using a vibrant solid blue color (RGB)
fill.fore_color.rgb = RGBColor(0, 112, 192)

# Add the image (light bulb icon) on the left side
# Note: Using 'image.png' as the placeholder for the icon as per instructions
left = Inches(1.5)
top = Inches(1.5)
width = Inches(3.5)
height = Inches(4.5)
slide.shapes.add_picture('image.png', left, top, width, height)

# Add the title "VIEWING FOLLOW-UP" on the right side
text_left = Inches(5.5)
text_top = Inches(3.0)
text_width = Inches(7)
text_height = Inches(1.5)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Add the text content and formatting
paragraph = text_frame.paragraphs[0]
paragraph.text = "VIEWING FOLLOW-UP"
paragraph.font.size = Pt(44)
paragraph.font.bold = True
# Set text color to white for contrast against the blue background
paragraph.font.color.rgb = RGBColor(255, 255, 255)
paragraph.alignment = PP_ALIGN.LEFT

# Save the presentation to 'output.pptx'
prs.save('output.pptx')