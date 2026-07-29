from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "KEY reasons why academic success is important in society"
title_run.font.size = Pt(28)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79) # Dark Blue

# --- Content Sections ---
# Data for the three sections
sections_data = [
    {
        "title": "Personal Growth and Development",
        "text": "Fosters critical thinking, self-discipline, and a lifelong love for learning, shaping well-rounded individuals.",
        "x_pos": 1.0
    },
    {
        "title": "Employability and Career Opportunities",
        "text": "Opens doors to diverse career paths, enhances job security, and leads to higher earning potential.",
        "x_pos": 5.0
    },
    {
        "title": "Economic Impact",
        "text": "Drives innovation, increases productivity, and contributes significantly to the overall economic growth of the nation.",
        "x_pos": 9.0
    }
]

# Layout constants
icon_top = Inches(1.8)
icon_size = Inches(1.5)
section_title_top = Inches(3.5)
section_title_height = Inches(0.6)
text_top = Inches(4.2)
text_height = Inches(2.5)
column_width = Inches(3.0)

# Loop to create sections
for section in sections_data:
    x = Inches(section["x_pos"])
    
    # 1. Add Icon (Placeholder Image)
    slide.shapes.add_picture('image.png', x, icon_top, icon_size, icon_size)
    
    # 2. Add Section Title
    title_box = slide.shapes.add_textbox(x, section_title_top, column_width, section_title_height)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_para = title_tf.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = section["title"]
    title_run.font.size = Pt(16)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    
    # 3. Add Explanatory Text
    text_box = slide.shapes.add_textbox(x, text_top, column_width, text_height)
    text_tf = text_box.text_frame
    text_tf.word_wrap = True
    text_para = text_tf.paragraphs[0]
    text_run = text_para.add_run()
    text_run.text = section["text"]
    text_run.font.size = Pt(12)
    text_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')