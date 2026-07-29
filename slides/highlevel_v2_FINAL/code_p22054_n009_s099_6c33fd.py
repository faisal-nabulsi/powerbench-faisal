from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "PRONUNCIATION ACTIVITY"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 0, 0)
title_para.alignment = 1 # Center

# --- Question ---
question_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
question_tf = question_box.text_frame
question_tf.word_wrap = True
question_para = question_tf.paragraphs[0]
question_para.text = "What is the general rule for the pronunciation of the 'th' digraph?"
question_para.font.size = Pt(24)
question_para.font.color.rgb = RGBColor(0, 0, 0)
question_para.alignment = 1 # Center

# --- Options ---
options_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2.5))
options_tf = options_box.text_frame
options_tf.word_wrap = True

# Option a
para_a = options_tf.paragraphs[0]
para_a.text = "a) It is always pronounced as a voiced sound /ð/."
para_a.font.size = Pt(20)
para_a.font.color.rgb = RGBColor(0, 0, 0)

# Option b
para_b = options_tf.add_paragraph()
para_b.text = "b) It is always pronounced as an unvoiced sound /θ/."
para_b.font.size = Pt(20)
para_b.font.color.rgb = RGBColor(0, 0, 0)

# Option c
para_c = options_tf.add_paragraph()
para_c.text = "c) It can be pronounced as either voiced /ð/ or unvoiced /θ/ depending on the word."
para_c.font.size = Pt(20)
para_c.font.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
prs.save('output.pptx')