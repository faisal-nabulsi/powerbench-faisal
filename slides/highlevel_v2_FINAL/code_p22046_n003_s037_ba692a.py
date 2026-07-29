from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize the presentation
prs = Presentation()

# Set the slide size to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Gradient Background ---
# Create a rectangle shape that covers the entire slide
# Adding it first ensures it is placed behind other elements
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    0, 0, prs.slide_width, prs.slide_height
)

# Apply a gradient fill to the background shape
try:
    gradient = bg_shape.fill.gradient
    # Set gradient stops: 0% (start) and 100% (end)
    gradient.add_stop(0, RGBColor(41, 128, 185))   # Light Blue
    gradient.add_stop(1, RGBColor(44, 62, 80))     # Dark Blue/Grey
except AttributeError:
    # Fallback to solid color if gradient API is not supported
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)

# --- 2. Title ---
# Add a text box for the title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Format the title
title_para = title_frame.paragraphs[0]
title_para.text = "Home Fun:"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.name = "Arial"  # Clear font
title_para.font.color.rgb = RGBColor(255, 255, 255) # White text
title_para.alignment = PP_ALIGN.LEFT

# --- 3. Bullet Points ---
# Add a text box for the content
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(3))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Bullet Point 1
p1 = content_frame.paragraphs[0]
p1.text = "Market opportunities"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.name = "Arial"
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.level = 0

# Bullet Point 2
p2 = content_frame.add_paragraph()
p2.text = "Inequality"
p2.font.size = Pt(32)
p2.font.bold = True
p2.font.name = "Arial"
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.level = 0

# Save the presentation
prs.save('output.pptx')