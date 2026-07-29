from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
# Positioning the title to be prominent and centered on the left/center area
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(2.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "Friends | Joey Doesn't Share Food!"
title_run.font.size = Pt(44)
title_run.font.bold = True

# Add Image (Character)
# Positioning the image on the right
left = Inches(9.5)
top = Inches(1.5)
width = Inches(3.5)
height = Inches(5)

slide.shapes.add_picture('image.png', left, top, width, height)

# Add Label "B1" at the bottom
label_box = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(1.5), Inches(0.8))
label_frame = label_box.text_frame
label_para = label_frame.paragraphs[0]
label_para.alignment = PP_ALIGN.CENTER
label_run = label_para.add_run()
label_run.text = "B1"
label_run.font.size = Pt(24)

# Save the presentation
prs.save('output.pptx')