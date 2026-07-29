from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Initialize Presentation
prs = Presentation()

# 2. Set Slide Dimensions (16:9 Widescreen)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 3. Add a Blank Slide
# Layout index 6 is standard for blank slides in default templates
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- SECTION 1: Quote on the Left ---
# Add a text box for the quote
left_quote = Inches(1.5)
top_quote = Inches(2.0)
width_quote = Inches(5.5)
height_quote = Inches(3.5)

txBox_quote = slide.shapes.add_textbox(left_quote, top_quote, width_quote, height_quote)

# Modified the first paragraph instead of adding a new one to avoid empty space
p_quote = txBox_quote.text_frame.paragraphs[0]
p_quote.text = "\"If things seem under control, you're not going fast enough.\""
p_quote.font.size = Pt(26)
p_quote.font.name = "Arial"
p_quote.font.color.rgb = RGBColor(50, 50, 50)
p_quote.alignment = PP_ALIGN.LEFT

# Add attribution in a new paragraph
p_attrib = txBox_quote.text_frame.add_paragraph()
p_attrib.text = "- Elon Musk"
p_attrib.font.size = Pt(20)
p_attrib.font.name = "Arial"
p_attrib.font.color.rgb = RGBColor(100, 100, 100)
p_attrib.font.italic = True
p_attrib.alignment = PP_ALIGN.LEFT

# --- SECTION 2: Image on the Right ---
# Add the placeholder image
img_left = Inches(8.2) # Positioned to the right of the quote
img_top = Inches(1.5)
img_width = Inches(4.5)
img_height = Inches(4.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# --- SECTION 3: Thank You Message at the Bottom ---
footer_left = Inches(1.667) # Centered horizontally: (13.333 - 10) / 2
footer_top = Inches(6.3)    # Positioned at the bottom
footer_width = Inches(10)
footer_height = Inches(1.0)

txBox_footer = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)

p_footer = txBox_footer.text_frame.paragraphs[0]
p_footer.text = "Thank You!"
p_footer.font.size = Pt(40)
p_footer.font.bold = True
p_footer.font.name = "Arial"
p_footer.font.color.rgb = RGBColor(0, 112, 192) # Blue
p_footer.alignment = PP_ALIGN.CENTER

# 4. Save the Presentation
prs.save('output.pptx')