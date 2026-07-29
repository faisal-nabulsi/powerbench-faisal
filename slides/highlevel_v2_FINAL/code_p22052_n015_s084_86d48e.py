from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set the slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define image properties
image_width = Inches(4)
image_height = Inches(4)
# Center the image horizontally and place it in the upper-middle area
image_left = (prs.slide_width - image_width) / 2
image_top = Inches(1.5)

# Add the placeholder image (representing the colorful globe with headphones)
slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Define text properties
phrase = "global media cultures"
text_box_width = Inches(5)
text_box_height = Inches(1)
# Center text horizontally and place it just below the image
text_left = (prs.slide_width - text_box_width) / 2
text_top = image_top + image_height + Inches(0.2)

# Add the text box
text_box = slide.shapes.add_textbox(text_left, text_top, text_box_width, text_box_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Configure the paragraph
paragraph = text_frame.paragraphs[0]
paragraph.text = phrase
paragraph.font.size = Pt(28)
paragraph.font.bold = True
paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')