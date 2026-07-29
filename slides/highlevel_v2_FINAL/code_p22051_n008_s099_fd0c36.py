from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to allow for custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Add Title ---
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1.5))
title_frame = title_shape.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "Top Grades Are Not the Sole Determinant of Success"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x2F, 0x54, 0x96) # Professional Blue

# --- 2. Add Body Content ---
# Text box for the explanation and key terms
body_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(7.5), Inches(4.5))
body_frame = body_shape.text_frame
body_frame.word_wrap = True

# Clear the default empty paragraph
body_frame.paragraphs[0].clear()

# Introductory text
intro_run = body_frame.paragraphs[0].add_run()
intro_run.text = "Success is a multifaceted journey that extends far beyond the classroom. While academic performance is valuable, it is only one piece of the puzzle. True achievement is often defined by:"
intro_run.font.size = Pt(18)
intro_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Key Term 1: Resilience
p1 = body_frame.add_paragraph()
p1.space_before = Pt(12)
run1_label = p1.add_run()
run1_label.text = "Resilience: "
run1_label.font.size = Pt(20)
run1_label.font.bold = True
run1_label.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) # Blue
run1_desc = p1.add_run()
run1_desc.text = "The capacity to recover quickly from difficulties; toughness."
run1_desc.font.size = Pt(18)

# Key Term 2: Creativity
p2 = body_frame.add_paragraph()
p2.space_before = Pt(12)
run2_label = p2.add_run()
run2_label.text = "Creativity: "
run2_label.font.size = Pt(20)
run2_label.font.bold = True
run2_label.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00) # Orange
run2_desc = p2.add_run()
run2_desc.text = "The use of imagination or original ideas to create something."
run2_desc.font.size = Pt(18)

# Key Term 3: Emotional Intelligence
p3 = body_frame.add_paragraph()
p3.space_before = Pt(12)
run3_label = p3.add_run()
run3_label.text = "Emotional Intelligence: "
run3_label.font.size = Pt(20)
run3_label.font.bold = True
run3_label.font.color.rgb = RGBColor(0x10, 0x7C, 0x10) # Green
run3_desc = p3.add_run()
run3_desc.text = "The ability to understand, use, and manage your own emotions in positive ways."
run3_desc.font.size = Pt(18)

# --- 3. Add Image ---
# Adding the placeholder image representing a trophy and target
# Positioned on the right side of the slide
image_left = Inches(9)
image_top = Inches(2)
image_width = Inches(4)
image_height = Inches(4)

try:
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)
except Exception:
    pass

# Save the presentation
prs.save('output.pptx')