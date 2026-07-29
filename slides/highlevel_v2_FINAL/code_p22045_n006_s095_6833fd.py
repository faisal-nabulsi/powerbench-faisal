from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout (Index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title ---
# Title: "Think Line:"
# Positioned at the top center/left
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_paragraph = title_frame.paragraphs[0]
title_paragraph.alignment = PP_ALIGN.LEFT

title_run = title_paragraph.add_run()
title_run.text = "Think Line:"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x00, 0x5C, 0xB8) # Professional Blue

# --- Image ---
# Graphic on the left using 'image.png'
img_left = Inches(0.5)
img_top = Inches(1.8)
img_width = Inches(4.5)
img_height = Inches(4.2)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# --- Question Text ---
# Text: "Do you think market create opportunity. How?"
# Positioned on the right
q_left = Inches(6)
q_top = Inches(2.5)
q_width = Inches(6.5)
q_height = Inches(3.5)

q_box = slide.shapes.add_textbox(q_left, q_top, q_width, q_height)
q_frame = q_box.text_frame
q_frame.word_wrap = True

q_paragraph = q_frame.paragraphs[0]
q_paragraph.alignment = PP_ALIGN.LEFT

q_run = q_paragraph.add_run()
q_run.text = "Do you think market create opportunity. How?"
q_run.font.size = Pt(24)
q_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# --- Background ---
# Note: python-pptx does not natively support gradient fills for slide backgrounds 
# without complex XML manipulation. A solid color or default background is used here 
# to ensure script stability. For a visual enhancement, a shape could be added, 
# but the focus is on the layout and content requested.

# Save the presentation
prs.save('output.pptx')