from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background color to Blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0) # A professional blue shade

# Add the Title "PRONUNCIATION ACTIVITY"
# Position: Top center
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(10.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True

p = title_tf.paragraphs[0]
p.text = "PRONUNCIATION ACTIVITY"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text for contrast
p.alignment = PP_ALIGN.CENTER

# Add the Icon (using the provided placeholder image)
# Position: Centered below the title
img_width = Inches(4.5)
img_height = Inches(4.5)
img_left = (prs.slide_width - img_width) / 2
img_top = Inches(3.0)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation
prs.save('output.pptx')