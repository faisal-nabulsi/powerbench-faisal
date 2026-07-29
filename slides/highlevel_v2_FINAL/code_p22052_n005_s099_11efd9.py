from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title ---
# Add a text box for the title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "PRONUNCIATION ACTIVITY"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_para.alignment = PP_ALIGN.CENTER

# --- Question ---
# Add a text box for the question
question_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
question_tf = question_box.text_frame
question_tf.word_wrap = True
question_para = question_tf.paragraphs[0]
question_para.text = "What is the general rule for the pronunciation of unstressed vowels in English?"
question_para.font.size = Pt(28)
question_para.font.color.rgb = RGBColor(50, 50, 50) # Dark Gray
question_para.alignment = PP_ALIGN.CENTER

# --- Options ---
# Add a text box for the answer options
options_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2.5))
options_tf = options_box.text_frame
options_tf.word_wrap = True

# Option a
p_a = options_tf.paragraphs[0]
p_a.text = "a) They are pronounced with their full vowel sound."
p_a.font.size = Pt(24)
p_a.font.color.rgb = RGBColor(0, 0, 0) # Black
p_a.space_after = Pt(15)
p_a.alignment = PP_ALIGN.CENTER

# Option b
p_b = options_tf.add_paragraph()
p_b.text = "b) They are often reduced to a schwa sound (/ə/)."
p_b.font.size = Pt(24)
p_b.font.color.rgb = RGBColor(0, 0, 0) # Black
p_b.space_after = Pt(15)
p_b.alignment = PP_ALIGN.CENTER

# Option c
p_c = options_tf.add_paragraph()
p_c.text = "c) They are completely silent."
p_c.font.size = Pt(24)
p_c.font.color.rgb = RGBColor(0, 0, 0) # Black
p_c.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')