from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to the presentation
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Title: "PART 2" at the top
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_tf = title_shape.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "PART 2"
title_para.font.size = Pt(50)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x29, 0x80, 0xB9) # Professional Blue
title_para.alignment = PP_ALIGN.CENTER

# 2. Prompt to discuss dating
prompt_shape = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(1.0))
prompt_tf = prompt_shape.text_frame
prompt_tf.word_wrap = True
prompt_para = prompt_tf.paragraphs[0]
prompt_para.text = "Discussion Topic: Dating"
prompt_para.font.size = Pt(34)
prompt_para.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey
prompt_para.alignment = PP_ALIGN.CENTER

# 3. List of three numbered questions
list_shape = slide.shapes.add_textbox(Inches(2), Inches(3.3), Inches(9.333), Inches(3.5))
list_tf = list_shape.text_frame
list_tf.word_wrap = True

# Clear the default empty paragraph
list_tf.paragraphs[0].clear()

questions = [
    "1. What are your 'non-negotiables' or deal-breakers in a relationship?",
    "2. How do you prefer to meet new people: through apps, mutual friends, or hobbies?",
    "3. What do you think is the biggest challenge in dating today?"
]

for q in questions:
    p = list_tf.add_paragraph()
    p.text = q
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(20) # Add spacing between questions for clarity

# Save the presentation
prs.save('output.pptx')