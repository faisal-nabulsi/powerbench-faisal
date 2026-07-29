from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 2. Define Content - Title and Milestones
title_text = "Brief highlights of Elon Musk"
milestones = [
    ("1971", "Born in Pretoria, South Africa."),
    ("1995", "Co-founded Zip2, sold for $307M."),
    ("2002", "Founded SpaceX, aiming to reduce space transport costs."),
    ("2004", "Joined Tesla Motors as Product Architect and CEO."),
    ("2008", "First successful Falcon 1 launch & Model S prototype."),
    ("2012", "SpaceX Dragon became first commercial spacecraft to dock with ISS."),
    ("2020", "NASA astronauts launch to ISS via SpaceX Crew Dragon.")
]

# 3. Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = title_text
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(51, 51, 51)  # Dark grey
p.alignment = PP_ALIGN.LEFT

# 4. Add Image (Placeholder)
# Using a simple rectangle background for the image to make it look intentional if image.png is missing/generic
image_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(1), Inches(3.5), Inches(5.5))
image_box.fill.solid()
image_box.fill.fore_color.rgb = RGBColor(230, 230, 230)
image_box.line.fill.background()

# Insert the actual image
try:
    # Attempt to insert the image into the slide
    img = slide.shapes.add_picture('image.png', Inches(9.7), Inches(1.2), Inches(3.1), Inches(5.1), image_from_memory=False)
    # Remove the background box if image exists
    container = slide.shapes._spTree
    # Note: Removing a shape from tree is complex in python-pptx, easier to just overlay or ignore.
    # Let's just overlay the image over the grey box logic above but keep it simple.
    # Actually, let's just replace the grey box logic:
    # We will assume image.png is available as per instructions. 
    # We'll add it directly at the position.
    slide.shapes.add_picture('image.png', Inches(9.5), Inches(1), Inches(3.5), Inches(5.5))
except Exception:
    pass # Fallback to grey box if image fails (though prompt guarantees existence)

# 5. Draw Timeline Line
# Vertical position for the timeline line
timeline_y = Inches(2.5)
line_start_x = Inches(0.5)
line_end_x = Inches(9.0)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    line_start_x, timeline_y, line_end_x - line_start_x, Inches(0.05)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(45, 45, 45) # Black/Dark Grey
line.line.fill.background()

# 6. Draw Milestones
# Spread milestones evenly across the line
margin_left = 0.5
margin_right = 9.0
available_width = margin_right - margin_left
num_items = len(milestones)
step_x = available_width / (num_items - 1)

for i, (year, desc) in enumerate(milestones):
    x_pos = margin_left + (i * step_x)
    
    # A. Add Year Label (Above the line)
    year_box = slide.shapes.add_textbox(Inches(x_pos - 0.3), Inches(1.8), Inches(0.6), Inches(0.5))
    tf = year_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0) # Orange/Red accent
    p.alignment = PP_ALIGN.CENTER

    # B. Add Circle Node (On the line)
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos - 0.125), Inches(2.375), Inches(0.25), Inches(0.25))
    node.fill.solid()
    node.fill.fore_color.rgb = RGBColor(255, 102, 0) # Orange/Red
    node.line.fill.background()

    # C. Add Description Label (Below the line)
    desc_box_width = Inches(1.8)
    desc_box = slide.shapes.add_textbox(Inches(x_pos - 0.9), Inches(2.8), desc_box_width, Inches(2.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER
    
    # Adjust spacing for nodes at ends if necessary, but this loop covers center alignment logic for text.

# 7. Save Presentation
prs.save('output.pptx')