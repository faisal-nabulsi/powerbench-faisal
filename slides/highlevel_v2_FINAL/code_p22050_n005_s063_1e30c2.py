from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout for custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Set Background Color (Dark Charcoal to make visuals pop)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(30, 40, 50) # Dark Blue-Gray

# 2. Add Title "Graphic Design"
# Position: Top center
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "Graphic Design"
title_para.font.size = Pt(64)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255) # White text
title_para.alignment = PP_ALIGN.CENTER

# 3. Add Image (Laptop with Paint Splashes)
# Dimensions: Wide to utilize the widescreen format
image_width = Inches(9.0)
image_height = Inches(5.5)

# Calculate left position to center the image
image_left = Inches((13.333 - 9.0) / 2)
image_top = Inches(1.8) # Positioned below the title

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')