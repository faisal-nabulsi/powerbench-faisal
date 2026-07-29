from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the Presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 Widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title Section ---
# Add a text box for the title "PRONUNCIATION ACTIVITY"
title_left = Inches(1)
title_top = Inches(0.5)
title_width = Inches(11.333)
title_height = Inches(1.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_para.text = "PRONUNCIATION ACTIVITY"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_para.alignment = PP_ALIGN.CENTER

# --- 2. Question Section ---
# Add a text box for the question about the general rule
question_left = Inches(1.5)
question_top = Inches(2.5)
question_width = Inches(10.333)
question_height = Inches(1.5)

question_box = slide.shapes.add_textbox(question_left, question_top, question_width, question_height)
question_tf = question_box.text_frame
question_tf.word_wrap = True

question_para = question_tf.paragraphs[0]
# Using a placeholder question text as requested
question_para.text = "What is the general rule for this pronunciation scheme?"
question_para.font.size = Pt(26)
question_para.font.color.rgb = RGBColor(50, 50, 50) # Dark Gray
question_para.alignment = PP_ALIGN.CENTER

# --- 3. Options Section ---
# Add three text boxes for answer options a, b, and c
# Centered horizontally on the slide
option_text_width = Inches(9.0)
option_left_pos = (13.333 - 9.0) / 2 

# Option A
opt_a_top = Inches(4.5)
box_a = slide.shapes.add_textbox(option_left_pos, opt_a_top, option_text_width, Inches(0.8))
tf_a = box_a.text_frame
p_a = tf_a.paragraphs[0]
p_a.text = "a) Vowels in unstressed syllables are reduced."
p_a.font.size = Pt(22)
p_a.font.color.rgb = RGBColor(80, 80, 80) # Medium Gray
p_a.alignment = PP_ALIGN.CENTER

# Option B
opt_b_top = Inches(5.4)
box_b = slide.shapes.add_textbox(option_left_pos, opt_b_top, option_text_width, Inches(0.8))
tf_b = box_b.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "b) Stress is always applied to the root syllable."
p_b.font.size = Pt(22)
p_b.font.color.rgb = RGBColor(80, 80, 80) # Medium Gray
p_b.alignment = PP_ALIGN.CENTER

# Option C
opt_c_top = Inches(6.3)
box_c = slide.shapes.add_textbox(option_left_pos, opt_c_top, option_text_width, Inches(0.8))
tf_c = box_c.text_frame
p_c = tf_c.paragraphs[0]
p_c.text = "c) Consonant clusters are simplified at the end of words."
p_c.font.size = Pt(22)
p_c.font.color.rgb = RGBColor(80, 80, 80) # Medium Gray
p_c.alignment = PP_ALIGN.CENTER

# Save the presentation to 'output.pptx'
prs.save('output.pptx')