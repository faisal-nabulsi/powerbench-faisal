from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (index 6 is standard for 'Blank' layout in default themes)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Title "PART 2"
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.3), Inches(11.333), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_para.text = "PART 2"
title_para.font.size = Pt(42)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# 2. Add Dialogue Excerpts
# Adding sample dialogue relevant to the grammar topic as a placeholder
dialogue_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(10.333), Inches(3.0))
dialogue_tf = dialogue_box.text_frame
dialogue_tf.word_wrap = True

# Dialogue Header
header_para = dialogue_tf.paragraphs[0]
header_para.text = "Dialogue Excerpts for Discussion"
header_para.font.size = Pt(22)
header_para.font.bold = True
header_para.space_after = Pt(12)

# Sample Dialogue Lines
dialogue_lines = [
    "Speaker A: \"What were you doing when the phone rang?\"",
    "Speaker B: \"I was cooking dinner.\"",
    "Speaker A: \"Did you answer it?\"",
    "Speaker B: \"No, I was too busy!\"",
    "Speaker C: \"While I was watching TV, it started to rain.\""
]

for line in dialogue_lines:
    p = dialogue_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.space_after = Pt(6)

# 3. Add the two Questions
questions_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.333), Inches(2.0))
questions_tf = questions_box.text_frame
questions_tf.word_wrap = True

# Question 1
q1_para = questions_tf.paragraphs[0]
q1_para.text = "When is the past continuous used in comparison with the past simple?"
q1_para.font.size = Pt(20)
q1_para.font.bold = True
q1_para.space_after = Pt(12)

# Question 2
q2_para = questions_tf.add_paragraph()
q2_para.text = "How is the structure for the past simple different than the past continuous?"
q2_para.font.size = Pt(20)
q2_para.font.bold = True

# Save the presentation
prs.save('output.pptx')