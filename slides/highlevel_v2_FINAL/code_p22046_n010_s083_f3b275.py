from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set the slide size to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout in default themes)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Define the dimensions of the text box
text_box_width = Inches(10)
text_box_height = Inches(2)

# Calculate the position to center the text box on the slide
left = (prs.slide_width - text_box_width) / 2
top = (prs.slide_height - text_box_height) / 2

# Add the text box
text_box = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
tf = text_box.text_frame

# Add the text "THANK YOU!"
p = tf.paragraphs[0]
p.text = "THANK YOU!"
p.alignment = PP_ALIGN.CENTER

# Apply formatting to the text
run = p.runs[0]
run.font.size = Pt(72)
run.font.bold = True
run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
prs.save('output.pptx')