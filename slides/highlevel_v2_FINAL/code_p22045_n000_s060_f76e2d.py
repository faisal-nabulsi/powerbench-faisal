from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title = title_frame.add_paragraph()
title.text = "'DATE GONE WRONG' STORY"
title.font.size = Pt(36)
title.font.bold = True
title.font.color.rgb = RGBColor(2, 77, 144)
title.alignment = PP_ALIGN.LEFT

# Add light bulb icon placeholder
slide.shapes.add_picture('image.png', Inches(10.5), Inches(1.5), Inches(2.5), Inches(2.5))

# Add short story for pronunciation practice
story_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.5), Inches(4.5))
story_frame = story_box.text_frame
story_frame.word_wrap = True

story_paragraph = story_frame.add_paragraph()
story_paragraph.text = (
    "My best friend Sam insisted I meet his cousin at the café downtown. "
    "I arrived early, nervously adjusting my tie in the reflection of the window. "
    "When she walked in, wearing an emerald dress and laughing with her friends, my heart skipped a beat. "
    "I waved enthusiastically, but she glanced past me, completely unfazed. "
    "Minutes ticked by. Finally, I overheard her say, 'Where's Sam? I came to meet his brother, not some random guy.' "
    "I froze, face burning crimson, as Sam rushed in, apologizing profusely for mixing up the names. "
    "The date was a disaster before it even began!"
)
story_paragraph.font.size = Pt(18)
story_paragraph.font.color.rgb = RGBColor(51, 51, 51)
story_paragraph.space_after = Pt(10)

# Add pronunciation practice instructions
instructions_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
instructions_frame = instructions_box.text_frame
instructions_frame.word_wrap = True

instructions = instructions_frame.add_paragraph()
instructions.text = (
    "Speech Bubble: Read the story aloud with clear pronunciation and proper intonation. "
    "Then listen to the accompanying audio clip to compare your delivery and receive feedback on areas for improvement."
)
instructions.font.size = Pt(16)
instructions.font.color.rgb = RGBColor(100, 100, 100)
instructions.font.italic = True

# Save presentation
prs.save('output.pptx')