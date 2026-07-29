from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "Importance of Skills Beyond Academics"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
title_para.alignment = PP_ALIGN.CENTER

# --- Intro Text ---
intro_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(10.333), Inches(0.5))
intro_tf = intro_box.text_frame
intro_tf.word_wrap = True
intro_para = intro_tf.paragraphs[0]
intro_para.text = "Mastering these essential skills ensures readiness for the modern workplace."
intro_para.font.size = Pt(18)
intro_para.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
intro_para.alignment = PP_ALIGN.CENTER

# --- Key Skills List ---
skills = [
    ("Critical Thinking", "Enables logical analysis and problem-solving in complex situations."),
    ("Communication", "Facilitates clear idea exchange and effective collaboration within teams."),
    ("Adaptability", "Crucial for navigating rapid changes and learning new technologies."),
    ("Emotional Intelligence", "Enhances interpersonal relationships and conflict resolution skills."),
    ("Leadership", "Inspires direction, motivates others, and drives project completion."),
    ("Digital Literacy", "Ensures proficiency with tools essential for modern productivity.")
]

content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True
content_tf.clear()

for skill, desc in skills:
    para = content_tf.add_paragraph()
    para.space_after = Pt(12)
    
    # Skill Name (Bold and Blue)
    run_skill = para.add_run()
    run_skill.text = f"• {skill}: "
    run_skill.font.bold = True
    run_skill.font.size = Pt(20)
    run_skill.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    
    # Description (Regular and Dark Grey)
    run_desc = para.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(18)
    run_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')