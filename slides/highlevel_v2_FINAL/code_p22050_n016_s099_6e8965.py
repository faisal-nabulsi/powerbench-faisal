from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create the presentation object
prs = Presentation()

# Set the slide size to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Importance of Skills Beyond Academics"
title_run.font.size = Pt(36)
title_run.font.bold = True

# Add the Content (Skills and Descriptions)
content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(4.5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

skills_data = [
    ("Effective Communication", "Ability to convey information clearly and listen actively."),
    ("Problem Solving and Critical Thinking", "Analyzing complex issues to find logical solutions."),
    ("Emotional Intelligence", "Understanding and managing one's own emotions and empathizing with others."),
    ("Adaptability and Resilience", "Adjusting to change and recovering quickly from setbacks.")
]

for i, (skill, description) in enumerate(skills_data):
    if i > 0:
        para = content_tf.add_paragraph()
    else:
        para = content_tf.paragraphs[0]
    
    para.space_after = Pt(16)
    
    # Skill Name (Bold)
    skill_run = para.add_run()
    skill_run.text = f"{skill}: "
    skill_run.font.bold = True
    skill_run.font.size = Pt(18)
    
    # Description
    desc_run = para.add_run()
    desc_run.text = description
    desc_run.font.size = Pt(18)

# Save the presentation
prs.save('output.pptx')