from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title "PART 1"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.0))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "PART 1"
title_p.font.size = Pt(36)
title_p.font.bold = True

# Add Discussion Prompt for Joey's Date
discussion_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(3.5))
discussion_tf = discussion_box.text_frame
discussion_tf.word_wrap = True

p1 = discussion_tf.paragraphs[0]
p1.text = "Discussion Prompt:"
p1.font.size = Pt(20)
p1.font.bold = True

p2 = discussion_tf.add_paragraph()
p2.text = "Think about Joey's recent date."
p2.font.size = Pt(18)

p3 = discussion_tf.add_paragraph()
p3.text = "1. What was the most memorable moment?"
p3.font.size = Pt(18)

p4 = discussion_tf.add_paragraph()
p4.text = "2. How did Joey react to the surprise gift?"
p4.font.size = Pt(18)

p5 = discussion_tf.add_paragraph()
p5.text = "3. Would you recommend this restaurant to a friend? Why or why not?"
p5.font.size = Pt(18)

p6 = discussion_tf.add_paragraph()
p6.text = "(Write your short answers below)"
p6.font.size = Pt(18)
p6.font.italic = True

# Add "PART 2" Instruction
part2_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.5))
part2_tf = part2_box.text_frame

p_part2 = part2_tf.paragraphs[0]
p_part2.text = "PART 2"
p_part2.font.size = Pt(24)
p_part2.font.bold = True

p_inst = part2_tf.add_paragraph()
p_inst.text = "Instruction: Please refer back to the previous section for comparison with Joey's date experience."
p_inst.font.size = Pt(16)

# Add Image of Joey
# Assuming 'image.png' exists in the working directory
try:
    slide.shapes.add_picture('image.png', Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.0))
except FileNotFoundError:
    # Add a placeholder text if image is missing, though instructions say it's available
    placeholder = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.0))
    placeholder.text = "[Image of Joey]"

# Save the presentation
prs.save('output.pptx')