from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Title text box: Left, Top, Width, Height
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "5 Stages of Development of Media"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# --- Add Content (Left Side) ---
# Text box for "3. PRINTING PRESS" and the bullet point
content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.5), Inches(4.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Use the first paragraph for the heading
heading_para = content_tf.paragraphs[0]
heading_run = heading_para.add_run()
heading_run.text = "3. PRINTING PRESS"
heading_run.font.size = Pt(28)
heading_run.font.bold = True

# Add a new paragraph for the bullet point
bullet_para = content_tf.add_paragraph()
bullet_para.space_after = Pt(12)
bullet_run = bullet_para.add_run()
bullet_run.text = (
    "Significance: The invention of the printing press in the 15th century democratized "
    "access to information. By mechanizing book production, it drastically reduced costs, "
    "accelerated the spread of literacy, and played a pivotal role in the Reformation, "
    "the Renaissance, and the Scientific Revolution."
)
bullet_run.font.size = Pt(16)

# --- Add Image (Right Side) ---
# Positioning the image on the right side of the slide
# Left margin starts at 1.0, Text width is 5.5 (ends at 6.5).
# Image starts at 7.0 to leave a small gap, width 5.333 to stay within 13.333 canvas.
image_left = Inches(7.0)
image_top = Inches(2.5)
image_width = Inches(5.333)
image_height = Inches(4.0)

# Add the placeholder image
slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')