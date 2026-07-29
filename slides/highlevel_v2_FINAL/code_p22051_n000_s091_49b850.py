from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Create a presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use the 'Title and Content' layout (Index 1)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add and format the Title
title_shape = slide.shapes.title
title_shape.text = "Types of Blockchain"
# Center the title for better aesthetics
title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Access the content placeholder
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.clear() # Remove default placeholder text

# Define the data to be added
blockchain_types = [
    ("Private Blockchain", "A restricted network where participation is controlled by a single organization or group. It prioritizes privacy, security, and transaction speed over decentralization."),
    ("Public Blockchain", "An open, decentralized network accessible to anyone via the internet. Anyone can read, write data, or audit the network, ensuring maximum transparency and censorship resistance."),
    ("Permissioned Blockchain", "A network with restricted access where users must be authorized to join as nodes. It allows for consensus mechanisms and access controls defined by the network members."),
    ("Hybrid Blockchain", "A composite architecture that combines elements of public and private blockchains. It allows data to be kept private while utilizing public blockchain's benefits like smart contracts.")
]

# Add bullet points
for term, definition in blockchain_types:
    # Add a new paragraph
    p = tf.add_paragraph()
    p.level = 0
    p.space_before = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    # Add the Term (Bold)
    run_term = p.add_run()
    run_term.text = term + ": "
    run_term.bold = True
    run_term.font.size = Pt(18)
    
    # Add the Definition
    run_def = p.add_run()
    run_def.text = definition
    run_def.font.size = Pt(18)

# Save the presentation
prs.save('output.pptx')