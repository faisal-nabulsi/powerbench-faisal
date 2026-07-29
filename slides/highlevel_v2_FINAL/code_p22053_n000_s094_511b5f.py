from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]  # Blank layout index is usually 6
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.2))
title_frame = title_shape.text_frame
title_frame.word_wrap = True
p_title = title_frame.add_paragraph()
p_title.text = "Types of Blockchain"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0, 51, 102)
p_title.alignment = PP_ALIGN.LEFT

# --- Add Content (Bullet Points) ---
# Create a text box for the content below the title
left = Inches(0.5)
top = Inches(2.0)
width = Inches(12)
height = Inches(5)
content_shape = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_shape.text_frame
content_frame.word_wrap = True

# Define the blockchain types and their definitions
blockchain_data = [
    {
        "term": "Private Blockchain Networks",
        "definition": "A network where access is restricted to authorized users only. Typically used within a single organization to maintain high security, control, and transaction speed."
    },
    {
        "term": "Public Blockchain Networks",
        "definition": "Decentralized and open-source networks accessible to the general public. Anyone can join, participate in consensus, and read/write to the ledger (e.g., Bitcoin, Ethereum)."
    },
    {
        "term": "Permissioned Blockchain Networks",
        "definition": "A type where consensus is restricted to a specific set of approved nodes. It balances decentralization with the privacy and scalability needed for enterprise applications."
    },
    {
        "term": "Hybrid Blockchain Networks",
        "definition": "Combines features of public and private blockchains. It allows organizations to keep certain data private while interacting with the public network when necessary for transparency."
    }
]

# Populate the content box with bullet points
for item in blockchain_data:
    p = content_frame.add_paragraph()
    p.space_before = Pt(15)
    
    # Add bullet point symbol and term (Bold)
    run_bold = p.add_run()
    run_bold.text = f"• {item['term']}"
    run_bold.font.bold = True
    run_bold.font.size = Pt(22)
    run_bold.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add definition (Regular, Indented)
    run_desc = p.add_run()
    run_desc.text = f"\n    {item['definition']}"
    run_desc.font.bold = False
    run_desc.font.size = Pt(20)
    run_desc.font.color.rgb = RGBColor(50, 50, 50)

# Save the presentation
prs.save('output.pptx')