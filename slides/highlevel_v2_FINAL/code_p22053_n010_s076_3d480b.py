from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Title: "Branding"
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_shape.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "Branding"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.alignment = PP_ALIGN.CENTER

# 2. Timeline Line
line_left = Inches(1.5)
line_top = Inches(3.5)
line_width = Inches(10.333)
line_height = Inches(0.05)
line_color = RGBColor(50, 50, 50)

line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height)
line_shape.line.fill.background()
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = line_color

# Helper to calculate node Y position (centered on line)
def get_node_y(line_top, size):
    return line_top - (size / 2) + (line_height / 2)

# 3. Node 1: 2008 - AirBed&Breakfast
node1_size = Inches(0.4)
node1_x = line_left - (node1_size / 2)
node1_y = get_node_y(line_top, node1_size)

node1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, node1_x, node1_y, node1_size, node1_size)
node1.fill.solid()
node1.fill.fore_color.rgb = RGBColor(255, 90, 95) # Airbnb Brand Red
node1.line.fill.background()

label1 = slide.shapes.add_textbox(node1_x, node1_y + node1_size + Inches(0.2), Inches(1.0), Inches(1.0))
label1_tf = label1.text_frame
label1_tf.word_wrap = True
p1 = label1_tf.paragraphs[0]
p1.text = "2008\nAirBed&Breakfast"
p1.font.size = Pt(12)
p1.alignment = PP_ALIGN.CENTER

# 4. Node 2: 2009 - Airbnb
mid_x = line_left + (line_width / 2)
node2_size = Inches(0.4)
node2_y = get_node_y(line_top, node2_size)

node2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, mid_x - (node2_size / 2), node2_y, node2_size, node2_size)
node2.fill.solid()
node2.fill.fore_color.rgb = RGBColor(255, 90, 95)
node2.line.fill.background()

label2 = slide.shapes.add_textbox(mid_x - Inches(0.5), node2_y - Inches(1.2), Inches(1.0), Inches(1.0))
label2_tf = label2.text_frame
p2 = label2_tf.paragraphs[0]
p2.text = "2009\nAirbnb"
p2.font.size = Pt(12)
p2.alignment = PP_ALIGN.CENTER

# 5. Node 3: 2014 - Bélo Symbol
end_x = line_left + line_width
img_size = Inches(1.0)
img_x = end_x - (img_size / 2)
img_y = get_node_y(line_top, img_size)

# Add the placeholder image representing the Bélo symbol
slide.shapes.add_picture('image.png', img_x, img_y, img_size, img_size)

label3 = slide.shapes.add_textbox(img_x, img_y + img_size + Inches(0.2), Inches(1.0), Inches(1.0))
label3_tf = label3.text_frame
label3_tf.word_wrap = True
p3 = label3_tf.paragraphs[0]
p3.text = "2014\nThe Bélo"
p3.font.size = Pt(12)
p3.alignment = PP_ALIGN.CENTER

# 6. Quote about significance
quote_box = slide.shapes.add_textbox(Inches(3.0), Inches(5.2), Inches(7.333), Inches(1.5))
quote_tf = quote_box.text_frame
quote_tf.word_wrap = True
quote_p = quote_tf.paragraphs[0]
quote_p.text = "\"The Bélo represents belonging anywhere and is a symbol of the sharing economy.\""
quote_p.font.size = Pt(16)
quote_p.font.italic = True
quote_p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')