from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 wide screen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Add the dark market-themed background image
# We add this first so it remains at the back of the visual stack
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Add a red section on the top right
# Calculate position to anchor to the top-right corner
red_width = Inches(4.5)
red_height = Inches(3.0)
red_left = prs.slide_width - red_width
red_top = Inches(0)

red_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, red_left, red_top, red_width, red_height)
red_shape.fill.solid()
red_shape.fill.fore_color.rgb = RGBColor(255, 0, 0) # Red color
red_shape.line.fill.background() # Remove border/line

# 3. Add the text "WHAT IS MARKET?"
# Position text inside the red section with padding
text_left = red_left + Inches(0.4)
text_top = red_top + Inches(0.4)
text_width = red_width - Inches(0.8)
text_height = red_height - Inches(0.8)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Configure paragraph and font
paragraph = text_frame.paragraphs[0]
paragraph.text = "WHAT IS MARKET?"
paragraph.font.bold = True
paragraph.font.size = Pt(42) # Clear, large font size
paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text for contrast

# Save the presentation
prs.save('output.pptx')