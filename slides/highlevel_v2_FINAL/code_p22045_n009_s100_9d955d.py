from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Add Main Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_p = title_tf.paragraphs[0]
title_run = title_p.add_run()
title_run.text = "KEY Reasons Why Academic Success Is Important in Society"
title_run.font.size = Pt(28)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79) # Dark Blue
title_p.alignment = PP_ALIGN.CENTER

# 2. Define Section Content
sections = [
    {
        "title": "Global Competitiveness",
        "desc": "Builds a skilled workforce that drives economic growth and enhances international standing."
    },
    {
        "title": "Scientific Advancement and Research",
        "desc": "Encourages innovation and critical thinking, leading to breakthroughs in science and technology."
    },
    {
        "title": "Cultural Enrichment and Social Cohesion",
        "desc": "Fosters understanding, empathy, and shared values, building stronger and more unified communities."
    }
]

# 3. Layout Configuration (in inches)
margin_left = 0.5
start_y = 1.6
col_width = 3.78
gap = 0.5
img_width = 1.5
img_height = 1.5

# Calculate X positions for columns and centered images
col_lefts = [
    margin_left,
    margin_left + col_width + gap,
    margin_left + 2 * (col_width + gap)
]
# Image X position = Column Left + (Column Width - Image Width) / 2
img_offset = (col_width - img_width) / 2
img_lefts = [c + img_offset for c in col_lefts]

# 4. Add Icons and Text for each section
for i, section in enumerate(sections):
    # Add Placeholder Image (Icon)
    slide.shapes.add_picture('image.png', Inches(img_lefts[i]), Inches(start_y), Inches(img_width), Inches(img_height))
    
    # Add Text Box below the icon
    text_top = start_y + img_height + 0.2
    text_box = slide.shapes.add_textbox(Inches(col_lefts[i]), Inches(text_top), Inches(col_width), Inches(3.5))
    text_tf = text_box.text_frame
    text_tf.word_wrap = True
    
    # Section Title
    p_title = text_tf.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = section["title"]
    r_title.font.size = Pt(16)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    p_title.alignment = PP_ALIGN.CENTER
    
    # Explanatory Text
    p_desc = text_tf.add_paragraph()
    p_desc.space_before = Pt(6)
    r_desc = p_desc.add_run()
    r_desc.text = section["desc"]
    r_desc.font.size = Pt(12)
    r_desc.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p_desc.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')