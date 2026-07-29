from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout for custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title Section ---
title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "How he came up with his ideas?"
title_para.font.size = Pt(42)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
title_para.alignment = PP_ALIGN.LEFT

# --- 2. Image Section (Right Side) ---
# Assuming 'image.png' exists in the working directory as per instructions
try:
    slide.shapes.add_picture('image.png', Inches(9.0), Inches(1.8), Inches(3.733), Inches(5.2))
except FileNotFoundError:
    pass # Image handling would go here if required, but prompt guarantees file

# --- 3. Narrative Section (Left Side, Top) ---
narrative_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(7.8), Inches(2.2))
narrative_frame = narrative_box.text_frame
narrative_frame.word_wrap = True

narrative_para = narrative_frame.paragraphs[0]
narrative_para.text = (
    "During his time at the University of Pennsylvania, Elon Musk didn't just focus on getting good grades; "
    "he spent his free time pondering the long-term future of civilization. He realized that many grand challenges "
    "were ignored because they seemed too difficult or distant. This reflection led him to adopt 'first-principles' thinking, "
    "breaking down complex problems into their fundamental truths to build revolutionary solutions from the ground up."
)
narrative_para.font.size = Pt(18)
narrative_para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
narrative_para.alignment = PP_ALIGN.JUSTIFY
narrative_para.space_after = Pt(6)

# --- 4. Bullet Points Section (Left Side, Bottom) ---
bullet_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(7.8), Inches(2.8))
bullet_frame = bullet_box.text_frame
bullet_frame.word_wrap = True

key_ideas = [
    "Sustainable Energy Systems",
    "Multi-planetary Civilization",
    "Neural Interface Technology",
    "Hyperloop Transportation Network",
    "Beneficial Artificial Intelligence"
]

for i, idea in enumerate(key_ideas):
    if i == 0:
        para = bullet_frame.paragraphs[0]
    else:
        para = bullet_frame.add_paragraph()
    
    para.text = idea
    para.font.size = Pt(20)
    para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    para.level = 0
    para.space_after = Pt(8)

# Save the presentation
prs.save('output.pptx')