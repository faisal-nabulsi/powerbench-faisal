from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen before adding slides
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the Blank layout)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Fluentize logo at the top left corner
# Using 'image.png' as the placeholder for the logo
# Coordinates: Left 0.5", Top 0.5", Width 2.5", Height 1.25"
slide.shapes.add_picture('image.png', Inches(0.5), Inches(0.5), Inches(2.5), Inches(1.25))

# Add "THANK YOU!" title centered at the bottom
# Slide Width: 13.333 inches
# Text Box Width: 8 inches
# Left Position (Centered): (13.333 - 8) / 2 = 2.6665 inches
text_box_width = Inches(8.0)
left_pos = Inches((13.333 - 8.0) / 2)

# Slide Height: 7.5 inches
# Text Box Height: 1.2 inches
# Top Position: Slide Height - Margin - Text Box Height
# Setting a 0.5 inch margin from the bottom
text_box_height = Inches(1.2)
top_pos = Inches(7.5 - 0.5 - 1.2)

text_box = slide.shapes.add_textbox(left_pos, top_pos, text_box_width, text_box_height)
tf = text_box.text_frame
tf.word_wrap = True

# Configure the text paragraph
p = tf.paragraphs[0]
p.text = "THANK YOU!"
p.font.size = Pt(48)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')