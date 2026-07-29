from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Initialize Presentation
prs = Presentation()

# Set 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Layout index 6 is usually blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Background Styling ---
# Set a contrasting dark blue background for a professional look
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0A, 0x19, 0x2F) # Dark Navy Blue

# --- 2. Imagery ---
# Add the placeholder image on the left side
# Dimensions: 5" width, 6.5" height, positioned near the center vertically
try:
    slide.shapes.add_picture(
        'image.png',
        left=Inches(0.5),
        top=Inches(0.5),
        width=Inches(5.0),
        height=Inches(6.5)
    )
except Exception:
    # Fallback if image is missing, creates a gray box
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(5.0), Inches(6.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)

# --- 3. Text Content ---
# Text block positions (Right side)
text_left = Inches(6.0)
text_width = Inches(7.0)

# Title: "ZENEROM UAE"
title_top = Inches(1.0)
title_height = Inches(1.2)
title_box = slide.shapes.add_textbox(text_left, title_top, text_width, title_height)
tf_title = title_box.text_frame
p_title = tf_title.paragraphs[0]
p_title.text = "ZENEROM UAE"
p_title.font.size = Pt(46)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text

# Description and Services Section
desc_top = Inches(2.5)
desc_height = Inches(3.5)
desc_box = slide.shapes.add_textbox(text_left, desc_top, text_width, desc_height)
tf_desc = desc_box.text_frame
tf_desc.word_wrap = True

# Clear default paragraph
tf_desc.paragraphs[0].text = ""

# Company Description
p_desc = tf_desc.add_paragraph()
p_desc.text = "Zenerom UAE is a pioneering firm dedicated to transforming the energy landscape. We provide innovative, sustainable solutions tailored to the unique needs of businesses operating in the Middle East, focusing on efficiency and carbon reduction."
p_desc.font.size = Pt(18)
p_desc.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0) # Light Grey
p_desc.space_after = Pt(15)

# Services Header
p_services_header = tf_desc.add_paragraph()
p_services_header.text = "Our Services"
p_services_header.font.size = Pt(22)
p_services_header.font.bold = True
p_services_header.font.color.rgb = RGBColor(0xFF, 0xC1, 0x07) # Gold/Yellow Accent
p_services_header.space_after = Pt(8)

# Services List
services_list = [
    "Energy Efficiency Consulting",
    "Smart Grid & IoT Solutions",
    "Renewable Energy Integration",
    "Sustainable Infrastructure Planning"
]

for svc in services_list:
    p_item = tf_desc.add_paragraph()
    p_item.text = f"• {svc}"
    p_item.font.size = Pt(18)
    p_item.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p_item.space_after = Pt(4)

# Contact Information
contact_top = Inches(6.3)
contact_box = slide.shapes.add_textbox(text_left, contact_top, text_width, Inches(0.8))
tf_contact = contact_box.text_frame
p_contact = tf_contact.paragraphs[0]
p_contact.text = "Contact: info@zenerom.ae  |  +971 4 123 4567  |  Dubai, UAE"
p_contact.font.size = Pt(14)
p_contact.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
p_contact.font.italic = True

# Save the presentation
prs.save('output.pptx')