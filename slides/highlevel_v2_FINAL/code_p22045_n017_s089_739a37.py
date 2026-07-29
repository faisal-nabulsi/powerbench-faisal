from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_slide():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.0))
    title_tf = title_box.text_frame
    title_para = title_tf.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = "Academic World and Professional World: Demands and Characteristics"
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    
    # Define content for Academic World
    academic_content = {
        "Title": "Academic World",
        "Sections": {
            "Goals": [
                "Knowledge creation and dissemination",
                "Mastery of specific disciplines",
                "Critical thinking and research skills"
            ],
            "Focus": [
                "Theoretical understanding",
                "Long-term learning and development",
                "Individual academic performance"
            ],
            "Setting": [
                "Universities, research institutes",
                "Classrooms, labs, and libraries",
                "Mentorship-based hierarchy"
            ],
            "Atmosphere": [
                "Intellectual curiosity and debate",
                "Slower, deliberate pace",
                "Grant-dependent funding models"
            ]
        }
    }
    
    # Define content for Professional World
    professional_content = {
        "Title": "Professional World",
        "Sections": {
            "Goals": [
                "Business growth and profitability",
                "Problem-solving and efficiency",
                "Career progression and value creation"
            ],
            "Focus": [
                "Practical application and results",
                "Teamwork and collaboration",
                "Adaptability and market relevance"
            ],
            "Setting": [
                "Corporations, startups, agencies",
                "Offices, remote work, client sites",
                "Matrix or functional structures"
            ],
            "Atmosphere": [
                "Fast-paced and dynamic",
                "Results-oriented and competitive",
                "Networking and continuous change"
            ]
        }
    }
    
    # Function to populate text box
    def populate_text_box(shape, content):
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Add Title
        title_para = tf.paragraphs[0]
        title_run = title_para.add_run()
        title_run.text = content["Title"]
        title_run.font.size = Pt(24)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        title_para.space_after = Pt(6)
        
        # Add Sections
        for section_name, bullets in content["Sections"].items():
            section_para = tf.add_paragraph()
            section_run = section_para.add_run()
            section_run.text = section_name
            section_run.font.size = Pt(18)
            section_run.font.bold = True
            section_run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
            section_para.space_after = Pt(2)
            
            for bullet in bullets:
                bullet_para = tf.add_paragraph()
                bullet_para.level = 0
                bullet_run = bullet_para.add_run()
                bullet_run.text = f"• {bullet}"
                bullet_run.font.size = Pt(14)
                bullet_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                bullet_para.space_after = Pt(4)

    # Add Left Column (Academic World)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(6.0))
    populate_text_box(left_box, academic_content)
    
    # Add Right Column (Professional World)
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6.0), Inches(6.0))
    populate_text_box(right_box, professional_content)
    
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()