from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation with 16:9 Widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Title: "QUIZ & REVIEW ACTIVITY"
# Centered near the top
title_left = Inches(1.666)
title_top = Inches(1)
title_width = Inches(10)
title_height = Inches(1)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "QUIZ & REVIEW ACTIVITY"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

# 2. Instruction: "Work with a partner"
# Below the title
instr_top = Inches(2.2)
instr_height = Inches(0.8)

instr_box = slide.shapes.add_textbox(title_left, instr_top, title_width, instr_height)
tf = instr_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Work with a partner"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(80, 80, 80)
p.alignment = PP_ALIGN.CENTER

# 3. Icon: Placeholder image 'image.png'
# Centered below the text elements
# Dimensions: 4x4 inches
img_width = Inches(4)
img_height = Inches(4)
# Horizontal Centering: (13.333 - 4) / 2 = 4.6665
img_left = Inches(4.667)
# Vertical Positioning: Centered in the remaining space below text (approx 3.5 to 7.5)
# Available height ~5 inches. Image 4 inches. Top margin ~2.5?
# Let's just place it at 3.5 to be safe and balanced.
img_top = Inches(3.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation
prs.save('output.pptx')