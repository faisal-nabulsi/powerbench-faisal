from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Importance of Skills Beyond Academics"
title_run.font.size = Pt(28)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

# Add Content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Helper to add items
items = [
    ("Effective Communication", "The ability to convey ideas clearly and listen actively is crucial for collaboration, leadership, and success in any workplace."),
    ("Problem Solving and Critical Thinking", "Analyzing complex situations, identifying root causes, and developing innovative solutions are essential for navigating challenges."),
    ("Emotional Intelligence", "Understanding and managing one's own emotions, as well as empathizing with others, fosters stronger relationships and better decision-making."),
    ("Adaptability and Resilience", "The capacity to adjust to new conditions and bounce back from setbacks is vital in a rapidly changing world.")
]

for i, (header, desc) in enumerate(items):
    # Add Header
    if i > 0:
        content_frame.add_paragraph() # Add spacing
    header_para = content_frame.add_paragraph()
    header_para.space_before = Pt(6)
    header_run = header_para.add_run()
    header_run.text = f"{header}"
    header_run.font.size = Pt(16)
    header_run.font.bold = True
    header_run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
    
    # Add Description
    desc_para = content_frame.add_paragraph()
    desc_para.space_before = Pt(6)
    desc_para.space_after = Pt(12)
    desc_run = desc_para.add_run()
    desc_run.text = f"{desc}"
    desc_run.font.size = Pt(12)
    desc_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')