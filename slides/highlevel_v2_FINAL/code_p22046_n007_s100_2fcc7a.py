from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set the slide size to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title: "PART 2" at the top ---
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_text_frame = title_shape.text_frame
title_text_frame.word_wrap = True
title_paragraph = title_text_frame.paragraphs[0]
title_paragraph.alignment = PP_ALIGN.CENTER

title_run = title_paragraph.add_run()
title_run.text = "PART 2"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(30, 30, 30) # Dark Black/Grey

# --- 2. Prompt to discuss dating ---
prompt_shape = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9.333), Inches(1.5))
prompt_text_frame = prompt_shape.text_frame
prompt_text_frame.word_wrap = True
prompt_paragraph = prompt_text_frame.paragraphs[0]
prompt_paragraph.alignment = PP_ALIGN.CENTER

prompt_run = prompt_paragraph.add_run()
prompt_run.text = "Let's open the floor to discuss dating."
prompt_run.font.size = Pt(28)
prompt_run.font.color.rgb = RGBColor(60, 60, 60)

# --- 3. Three numbered questions for discussion ---
questions_shape = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(3))
questions_text_frame = questions_shape.text_frame
questions_text_frame.word_wrap = True

# Question 1
p1 = questions_text_frame.paragraphs[0]
p1.space_after = Pt(12)
run_q1_num = p1.add_run()
run_q1_num.text = "1. "
run_q1_num.font.bold = True
run_q1_num.font.size = Pt(22)
run_q1_text = p1.add_run()
run_q1_text.text = "What qualities do you value most in a romantic partner?"
run_q1_text.font.size = Pt(22)

# Question 2
p2 = questions_text_frame.add_paragraph()
p2.space_after = Pt(12)
run_q2_num = p2.add_run()
run_q2_num.text = "2. "
run_q2_num.font.bold = True
run_q2_num.font.size = Pt(22)
run_q2_text = p2.add_run()
run_q2_text.text = "How has your approach to dating changed over the last five years?"
run_q2_text.font.size = Pt(22)

# Question 3
p3 = questions_text_frame.add_paragraph()
run_q3_num = p3.add_run()
run_q3_num.text = "3. "
run_q3_num.font.bold = True
run_q3_num.font.size = Pt(22)
run_q3_text = p3.add_run()
run_q3_text.text = "What is a common misconception people have about modern dating?"
run_q3_text.font.size = Pt(22)

# Save the presentation
prs.save('output.pptx')