from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Select a blank slide layout (Index 6 is typically Blank, fallback to 0)
try:
    slide_layout = prs.slide_layouts[6]
except IndexError:
    slide_layout = prs.slide_layouts[0]

# Add the slide
slide = prs.slides.add_slide(slide_layout)

# --- 1. Image on the Right ---
# Positioning the image placeholder 'image.png' on the right side
# Coordinates: Left, Top, Width, Height
img_left = Inches(7.5)
img_top = Inches(1.0)
img_width = Inches(5.5)
img_height = Inches(5.0)
slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# --- 2. Quote from Elon Musk on the Left ---
# Textbox positioning
quote_left = Inches(0.5)
quote_top = Inches(1.5)
quote_width = Inches(6.5)
quote_height = Inches(4.0)

textbox_quote = slide.shapes.add_textbox(quote_left, quote_top, quote_width, quote_height)
tf_quote = textbox_quote.text_frame
tf_quote.word_wrap = True

# Add the quote text
p_quote = tf_quote.paragraphs[0]
p_quote.text = "“When something is important enough, you do it even if the odds are not in your favor.”"
p_quote.font.size = Pt(28)
p_quote.font.color.rgb = RGBColor(30, 30, 30) # Dark Grey
p_quote.font.name = 'Arial'

# Add the attribution run
run_attribution = p_quote.add_run()
run_attribution.text = "\n\n— Elon Musk"
run_attribution.font.size = Pt(20)
run_attribution.font.color.rgb = RGBColor(80, 80, 80) # Medium Grey
run_attribution.font.name = 'Arial'
run_attribution.font.italic = True

# --- 3. "Thank You!" Message at the Bottom ---
# Textbox positioned at the bottom, centered horizontally
ty_width = Inches(5.0)
ty_height = Inches(0.8)
ty_top = Inches(6.5)
# Calculate left position for centering: (SlideWidth - BoxWidth) / 2
ty_left = Inches((13.333 - 5.0) / 2)

textbox_ty = slide.shapes.add_textbox(ty_left, ty_top, ty_width, ty_height)
tf_ty = textbox_ty.text_frame
p_ty = tf_ty.paragraphs[0]
p_ty.text = "Thank You!"
p_ty.alignment = PP_ALIGN.CENTER
p_ty.font.size = Pt(36)
p_ty.font.color.rgb = RGBColor(0, 0, 0) # Black
p_ty.font.name = 'Arial'
p_ty.font.bold = True

# Save the presentation to 'output.pptx'
prs.save('output.pptx')