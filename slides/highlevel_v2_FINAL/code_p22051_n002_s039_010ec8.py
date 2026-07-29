from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Colorful Background
# Using the placeholder image to create a colorful background
try:
    # Add image covering the whole slide
    slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)
except Exception:
    # Fallback to a solid color if image is missing
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(100, 149, 237) # Cornflower Blue

# 2. Title
# Positioning title at top
left_title = Inches(1.5)
top_title = Inches(1.0)
width_title = Inches(10.333)
height_title = Inches(1.0)

title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_paragraph = title_frame.paragraphs[0]

title_paragraph.text = "Innovations (Future Products or Services)"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text
title_paragraph.font.name = "Arial"
title_paragraph.space_after = Pt(10)

# 3. Description
# Positioning description below title
left_desc = Inches(1.5)
top_desc = Inches(2.5)
width_desc = Inches(10.333)
height_desc = Inches(2.5)

desc_box = slide.shapes.add_textbox(left_desc, top_desc, width_desc, height_desc)
desc_frame = desc_box.text_frame
desc_frame.word_wrap = True
desc_paragraph = desc_frame.paragraphs[0]

desc_paragraph.text = "Brief Description of Airbnb's 2021 Innovations:"
desc_paragraph.font.size = Pt(24)
desc_paragraph.font.bold = True
desc_paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text
desc_paragraph.font.name = "Arial"
desc_paragraph.space_after = Pt(10)

details_paragraph = desc_frame.add_paragraph()
details_paragraph.text = (
    "In 2021, Airbnb launched significant innovations focused on enhancing travel experiences. "
    "Key updates included the expansion of 'Airbnb Experiences' to include local classes and workshops, "
    "introducing 'Airbnb for Work' features, and rolling out new safety and verification tools for hosts and guests."
)
details_paragraph.font.size = Pt(18)
details_paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text
details_paragraph.font.name = "Arial"
details_paragraph.space_after = Pt(20)

# 4. Quote
# Positioning quote at the bottom
left_quote = Inches(1.5)
top_quote = Inches(5.5)
width_quote = Inches(10.333)
height_quote = Inches(1.5)

quote_box = slide.shapes.add_textbox(left_quote, top_quote, width_quote, height_quote)
quote_frame = quote_box.text_frame
quote_frame.word_wrap = True
quote_paragraph = quote_frame.paragraphs[0]

quote_paragraph.text = "\"We are building a platform for the future of travel, where connection and belonging are at the heart of every stay.\""
quote_paragraph.font.size = Pt(22)
quote_paragraph.font.italic = True
quote_paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text
quote_paragraph.font.name = "Arial"
quote_paragraph.alignment = PP_ALIGN.LEFT
quote_paragraph.space_after = Pt(5)

attribution_paragraph = quote_frame.add_paragraph()
attribution_paragraph.text = "— Brian Chesky, Co-founder and CEO of Airbnb"
attribution_paragraph.font.size = Pt(16)
attribution_paragraph.font.color.rgb = RGBColor(200, 200, 200)
attribution_paragraph.font.name = "Arial"
attribution_paragraph.alignment = PP_ALIGN.LEFT

# Save to output.pptx
prs.save('output.pptx')