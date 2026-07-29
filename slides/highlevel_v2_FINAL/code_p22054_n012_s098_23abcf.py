from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title "PART 3" at the top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_run = title_p.add_run()
title_run.text = "PART 3"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 0, 0)

# Add prompt about making a good first impression on a date
prompt_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.5), Inches(4.5))
prompt_tf = prompt_box.text_frame
prompt_tf.word_wrap = True
prompt_p = prompt_tf.paragraphs[0]
prompt_run = prompt_p.add_run()
prompt_run.text = "Imagine you are going on a first date.\n\nWhat is the key to making a good first impression?"
prompt_run.font.size = Pt(22)
prompt_run.font.color.rgb = RGBColor(50, 50, 50)

# Add a speech bubble graphic using the placeholder image
# Positioned on the right side of the slide
try:
    slide.shapes.add_picture('image.png', Inches(6.5), Inches(2.5), Inches(6.5), Inches(4.5))
except Exception:
    # In case the image file is missing, skip silently or handle error
    pass

# Save the presentation
prs.save('output.pptx')