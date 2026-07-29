from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set dark background
# Note: python-pptx does not support gradient fills on slide backgrounds via the public API.
# A solid dark color is used to achieve the dark theme.
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

# Add Image (Light Bulb)
# Positioned centrally
img_width = Inches(3)
img_height = Inches(3)
img_left = (prs.slide_width - img_width) / 2
img_top = Inches(1.5)
slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Add Text
# Positioned below the image
text_left = Inches(2)
text_top = Inches(5)
text_width = Inches(9.333)
text_height = Inches(1.5)

textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
tf = textbox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "THINK LINE: Q. How products reach to us?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

prs.save('output.pptx')