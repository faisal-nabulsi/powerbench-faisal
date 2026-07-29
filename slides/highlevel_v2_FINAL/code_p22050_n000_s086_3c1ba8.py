from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Create presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank layout and add a slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# ── Title ──────────────────────────────────────────────
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "PART 2"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# ── Instruction text ───────────────────────────────────
instr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.7))
tf2 = instr_box.text_frame
tf2.word_wrap = True
p2 = tf2.add_paragraph()
p2.text = (
    "Read the dialogue excerpts below and categorize the verb forms into "
    "Past Simple or Past Continuous in the table."
)
p2.font.size = Pt(16)
p2.font.italic = True
p2.font.color.rgb = RGBColor(51, 51, 51)
p2.alignment = PP_ALIGN.CENTER

# ── Table ──────────────────────────────────────────────
rows, cols = 8, 3
table_shape = slide.shapes.add_table(
    rows, cols,
    left=Inches(0.5),
    top=Inches(2.0),
    width=Inches(12.333),
    height=Inches(4.8),
)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(6.5)
table.columns[1].width = Inches(2.9165)
table.columns[2].width = Inches(2.9165)

# Turn on row height formatting
for row in table.rows:
    row.height_format = True

# Header row
headers = ["Dialogue Excerpt", "Past Simple", "Past Continuous"]
for c_idx, header_text in enumerate(headers):
    cell = table.cell(0, c_idx)
    cell.text = header_text
    cell.fill.background()  # clear fill
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)

    para = cell.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.font.bold = True
    para.font.size = Pt(14)
    para.font.color.rgb = RGBColor(255, 255, 255)

# Dialogue excerpts to populate the first column
examples = [
    '"Actually, while Sarah was cooking , the phone rang."',
    '"He opened the window because it was getting hot."',
    '"We were playing football when it started to rain."',
    '"I was reading a book as she was preparing tea."',
    '"We watched the movie while he was cleaning the house."',
    '"She arrived home as they were having dinner."',
]

HIGHLIGHT_RGB = RGBColor(235, 241, 250)
NORMAL_RGB = RGBColor(255, 255, 255)

for row_idx, example in enumerate(examples, start=1):
    # Dialogue excerpt cell
    cell0 = table.cell(row_idx, 0)
    cell0.text = example
    cell0.text_frame.paragraphs[0].font.size = Pt(11)
    cell0.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    cell0.text_frame.word_wrap = True
    cell0.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill colour (alternating stripes)
    bg_color = HIGHLIGHT_RGB if row_idx % 2 == 0 else NORMAL_RGB
    for col_idx in range(3):
        cell = table.cell(row_idx, col_idx)
        cell.fill.background()
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    # Border colour for header
    for col_idx in [1, 2]:
        cell = table.cell(row_idx, col_idx)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.vertical_anchor = MSO_ANCHOR.MIDDLE

# ── Small footer text ─────────────────────────────────
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
ft_tf = footer_box.text_frame
p3 = ft_tf.add_paragraph()
p3.text = "Analysing verb tenses in context — Part 2"
p3.font.size = Pt(10)
p3.font.color.rgb = RGBColor(128, 128, 128)
p3.alignment = PP_ALIGN.CENTER

# Save
prs.save("output.pptx")