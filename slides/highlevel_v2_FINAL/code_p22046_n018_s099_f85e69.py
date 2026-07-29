from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "GLOBAL MEDIA CULTURES"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

# Add Bullet Points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6.5), Inches(5.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

bullets = [
    "Globalization enables the seamless flow of media, allowing cultural stories and values to transcend geographical boundaries.",
    "It encourages cultural hybridization, merging local traditions with global trends to form unique, evolving identities.",
    "Digital connectivity amplifies this spread, empowering individuals to share and consume diverse cultural content instantly."
]

for bullet in bullets:
    para = content_tf.add_paragraph()
    para.text = bullet
    para.font.size = Pt(18)
    para.space_after = Pt(12)
    para.level = 0

# Add Image
# Positioning on the right side
# Left: 7.2 (0.5 + 6.5 + 0.2 gap)
# Top: 2.2
# Width: 5.633 (13.333 - 0.5 - 6.5 - 0.5 right margin)
# Height: 5.0
try:
    slide.shapes.add_picture('image.png', Inches(7.2), Inches(2.2), Inches(5.633), Inches(5.0))
except FileNotFoundError:
    # Fallback if image is missing, though instructions say it's available
    pass

# Save the presentation
prs.save('output.pptx')