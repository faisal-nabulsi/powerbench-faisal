from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 1. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 1"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # 2. Instructions
    instr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    tf = instr_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Instructions: Describe and rank the six date ideas (A-F) below."
    p.font.size = Pt(18)
    p.font.italic = True

    # 3. Image Grid & Labels
    # Grid params
    img_w, img_h = Inches(3.0), Inches(2.0)
    label_w, label_h = Inches(0.6), Inches(0.6)
    
    # Positions
    row1_y = Inches(1.7)
    row2_y = Inches(4.1) # 1.7 + 2.0 + 0.4
    col1_x = Inches(0.8)
    col2_x = Inches(4.6)
    col3_x = Inches(8.4)

    images_data = [
        (col1_x, row1_y, "A"),
        (col2_x, row1_y, "B"),
        (col3_x, row1_y, "C"),
        (col1_x, row2_y, "D"),
        (col2_x, row2_y, "E"),
        (col3_x, row2_y, "F"),
    ]

    for x, y, char in images_data:
        # Add Image
        # Try wrapping in try-except just in case image missing, though prompt guarantees it.
        try:
            slide.shapes.add_picture('image.png', x, y, img_w, img_h)
        except Exception:
            # Fallback placeholder text if image fails (should not happen based on prompt)
            pass

        # Add Label
        # Center on image
        center_x = x + Inches(1.5)
        center_y = y + Inches(1.0)
        lbl_x = center_x - Inches(0.3)
        lbl_y = center_y - Inches(0.3)
        
        lbl_box = slide.shapes.add_textbox(lbl_x, lbl_y, label_w, label_h)
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = char
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # Optional: White background for label box to make it pop? 
        # pptx doesn't easily support background color on textbox, usually requires shape.
        # But text is black, image might be busy. 
        # Let's stick to simple text.

    # 4. Ranking Section
    # Header
    rank_hdr_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(3), Inches(0.4))
    tf = rank_hdr_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ranking (1-6):"
    p.font.size = Pt(20)
    p.font.bold = True

    # Inputs
    rank_input_top = Inches(6.4)
    # Spacing: (13.333 - 1.0) / 6 approx? No, 6 items.
    # Let's fit them comfortably.
    # Total width ~ 12. Items width 1.5. 6 * 1.5 = 9.0. Remaining gap 3.0.
    # Gap ~ 0.5 each.
    
    current_x = Inches(0.5)
    step = Inches(2.0) 
    
    for i in range(1, 7):
        box_w = Inches(1.8)
        box_h = Inches(0.5)
        
        rank_box = slide.shapes.add_textbox(current_x, rank_input_top, box_w, box_h)
        tf = rank_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{i}: _______"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.LEFT
        
        current_x += step

    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()