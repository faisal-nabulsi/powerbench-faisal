from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random

# Initialize presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_run = title_frame.add_paragraph()
title_run.text = "What is Graphic Design?"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50) # Dark Blue-Grey
title_run.alignment = PP_ALIGN.LEFT

# --- 2. Bullet Points ---
bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(5))
bullet_frame = bullet_box.text_frame
bullet_frame.word_wrap = True

points = [
    "Visual communication using text and images.",
    "Combines art and technology to solve problems.",
    "Essential for branding and marketing.",
    "Focuses on aesthetics and functionality.",
    "Uses elements like typography, color, and layout."
]

for i, point in enumerate(points):
    if i == 0:
        p = bullet_frame.paragraphs[0]
    else:
        p = bullet_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.space_after = Pt(12)
    p.level = 0

# --- 3. Word Cloud ---
# Terms with relative importance (font size)
terms = [
    ("Typography", 38), ("Layout", 30), ("Color", 34), ("Branding", 32),
    ("Illustration", 26), ("UI/UX", 28), ("Print", 22), ("Digital", 24),
    ("Creativity", 30), ("Composition", 24), ("Balance", 20), ("Contrast", 22),
    ("Hierarchy", 26), ("Vector", 20), ("Raster", 18), ("Adobe", 24),
    ("Photoshop", 22), ("Illustrator", 24), ("InDesign", 22), ("Figma", 26),
    ("Sketch", 20), ("Canva", 22), ("Pixel", 18), ("Resolution", 20),
    ("Grid", 22), ("Whitespace", 24), ("Serif", 18), ("Sans", 18),
    ("Kerning", 20), ("Leading", 18), ("Palette", 22), ("Mockup", 20)
]

# Cloud area parameters
cloud_left = Inches(6.5)
cloud_top = Inches(1.5)
cloud_width = Inches(6)
cloud_height = Inches(5)

random.seed(100) # Consistent randomness

# Colors for word cloud
cloud_colors = [
    RGBColor(0xE7, 0x4C, 0x3C), # Red
    RGBColor(0x34, 0x98, 0xDB), # Blue
    RGBColor(0xF3, 0x9C, 0x12), # Orange
    RGBColor(0x1A, 0xB8, 0x1A), # Green
    RGBColor(0x9B, 0x59, 0xB6), # Purple
    RGBColor(0x34, 0x49, 0x5E), # Dark Blue
    RGBColor(0x95, 0xA5, 0xA6), # Grey
]

for term, size in terms:
    # Random position within the cloud bounding box
    margin_x = 1.5 # inches
    margin_y = 0.5 # inches
    
    x_pos = cloud_left + Inches(random.uniform(0, cloud_width.inches - margin_x))
    y_pos = cloud_top + Inches(random.uniform(0, cloud_height.inches - margin_y))
    
    # Add text box
    box_width = Inches(2.5) 
    box_height = Inches(0.6)
    
    shape = slide.shapes.add_textbox(x_pos, y_pos, box_width, box_height)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = term
    p.font.size = Pt(size)
    p.font.color.rgb = random.choice(cloud_colors)
    p.alignment = PP_ALIGN.CENTER
    
    # Random rotation for some terms to make it look more organic
    if random.random() > 0.7:
        shape.rotation = random.uniform(-15, 15)

# --- 4. Geometric Shapes ---
# Shape 1: Large Circle Top Right
shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.2), Inches(2.5), Inches(2.5))
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
shape1.line.fill.background()
shape1.rotation = 15

# Shape 2: Rectangle Bottom Left
shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(5.5), Inches(2.2), Inches(1.8))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0x4E, 0xC9, 0xB0)
shape2.line.fill.background()
shape2.rotation = -10

# Shape 3: Triangle Top Left
shape3 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.2), Inches(0.5), Inches(1.8), Inches(1.8))
shape3.fill.solid()
shape3.fill.fore_color.rgb = RGBColor(0x6B, 0x5B, 0x95)
shape3.line.fill.background()
shape3.rotation = 45

# Shape 4: Circle Bottom Right
shape4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(5.2), Inches(1.8), Inches(1.8))
shape4.fill.solid()
shape4.fill.fore_color.rgb = RGBColor(0xF7, 0xCA, 0x18)
shape4.line.fill.background()

# Shape 5: Small Square near title
shape5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(0.4), Inches(0.6), Inches(0.6))
shape5.fill.solid()
shape5.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
shape5.line.fill.background()

# Save the presentation
prs.save('output.pptx')