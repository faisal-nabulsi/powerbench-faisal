from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set the slide size to 16:9 Widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the Blank layout)
try:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
except IndexError:
    # Fallback if blank layout isn't at index 6
    slide = prs.slides.add_slide(prs.slide_layouts[0])

# Set the background to a solid blue color
background_fill = slide.background.fill
background_fill.solid()
background_fill.fore_color.rgb = RGBColor(0, 0, 255)  # Pure Blue

# Add the light bulb icon (placeholder image) on the left side
# Position: Left x=1.5, Top y=2.5, Size: 3x3 inches
left_image = Inches(1.5)
top_image = Inches(2.5)
width_image = Inches(3.0)
height_image = Inches(3.0)
slide.shapes.add_picture('image.png', left_image, top_image, width_image, height_image)

# Add the Title text
# Position: Spanning the width, top centered
left_title = Inches(0.5)
top_title = Inches(0.5)
width_title = Inches(12.333)
height_title = Inches(1.5)

title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "VIEWING FOLLOW-UP"
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(255, 255, 255)  # White text
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')