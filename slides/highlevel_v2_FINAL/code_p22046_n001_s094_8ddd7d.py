from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Title "PART 3" at the top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.0))
title_frame = title_box.text_frame
title_frame.paragraphs[0].text = "PART 3"
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 2. Add the prompt about a first date involving French fries
# Positioned below the title
prompt_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9.333), Inches(1.0))
prompt_frame = prompt_box.text_frame
prompt_frame.word_wrap = True
prompt_frame.paragraphs[0].text = "Prompt: Describe a first date scenario where a basket of French fries becomes the center of attention."
prompt_frame.paragraphs[0].font.size = Pt(24)
prompt_frame.paragraphs[0].font.bold = True
prompt_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
prompt_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 3. Character Image (using 'image.png')
# Positioned on the left side
image_left = Inches(1)
image_top = Inches(3.5)
image_width = Inches(3)
image_height = Inches(3.5)
slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# 4. Storytelling Starter with a Speech Bubble
# Create a rounded rectangle shape to act as a speech bubble
bubble_left = Inches(4.5)
bubble_top = Inches(4)
bubble_width = Inches(7)
bubble_height = Inches(2)
bubble_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bubble_left, bubble_top, bubble_width, bubble_height)

# Style the speech bubble (white background, black border)
bubble_shape.fill.solid()
bubble_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
bubble_shape.line.color.rgb = RGBColor(0, 0, 0)
bubble_shape.line.width = Pt(2)

# Add text to the speech bubble
bubble_tf = bubble_shape.text_frame
bubble_tf.word_wrap = True
bubble_p = bubble_tf.paragraphs[0]
bubble_p.text = "\"I didn't mean to spill ketchup \non the dinner table!\""
bubble_p.font.size = Pt(22)
bubble_p.font.color.rgb = RGBColor(0, 0, 0)
bubble_p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')