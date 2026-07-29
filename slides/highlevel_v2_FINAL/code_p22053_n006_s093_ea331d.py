from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Add Title
# Positioned at top center
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "More About the Market"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey
p_title.alignment = PP_ALIGN.CENTER

# 2. Define Flowchart Configuration
steps = [
    "Marketing", 
    "Storage", 
    "Sales", 
    "Transportation", 
    "Distribution", 
    "After-sales services"
]

# Dimensions
BOX_WIDTH = Inches(1.3)
BOX_HEIGHT = Inches(1.0)
ARROW_WIDTH = Inches(0.4)
ARROW_HEIGHT = Inches(0.4)

# Positioning
# Calculate total width to center the flowchart horizontally
# 6 boxes * 1.3 + 5 arrows * 0.4 = 7.8 + 2.0 = 9.8 inches
# Start X = (13.333 - 9.8) / 2 approx 1.76. Let's use 1.75 for padding.
START_X = Inches(1.75)
# Y Positions
# Vertical center of flowchart roughly at 4.25 inches (leaving room for title)
# Box Y = 3.75 (Height 1.0 -> Center 4.25)
# Arrow Y = 4.05 (Height 0.4 -> Center 4.25)
Y_BOX = Inches(3.75)
Y_ARROW = Inches(4.05)

# 3. Draw Flowchart
current_x = START_X

for i, step_text in enumerate(steps):
    # --- Draw Step Box ---
    left_box = current_x
    top_box = Y_BOX
    
    box_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_box, top_box, BOX_WIDTH, BOX_HEIGHT)
    
    # Style Box
    box_shape.fill.solid()
    box_shape.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0) # Office Blue
    box_shape.line.fill.background() # No border line
    
    # Box Text
    tf_box = box_shape.text_frame
    tf_box.word_wrap = True
    try:
        p_box = tf_box.paragraphs[0]
    except IndexError:
        p_box = tf_box.add_paragraph()
        
    p_box.text = step_text
    p_box.font.size = Pt(14)
    p_box.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
    p_box.alignment = PP_ALIGN.CENTER
    
    # --- Draw Arrow Connector ---
    # Only draw arrow if it's not the last step
    if i < len(steps) - 1:
        left_arrow = current_x + BOX_WIDTH
        top_arrow = Y_ARROW
        
        arrow_shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_arrow, top_arrow, ARROW_WIDTH, ARROW_HEIGHT)
        
        # Style Arrow
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88) # Grey
        arrow_shape.line.fill.background() # No border
        
        # Advance x position for the next box
        current_x += BOX_WIDTH + ARROW_WIDTH

# Save the presentation
prs.save('output.pptx')