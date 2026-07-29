from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set Background Color (Solid Green to represent the theme robustly)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(34, 139, 34) # Forest Green

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_p = title_tf.paragraphs[0]
title_p.text = "Environment"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.LEFT

# Add Content Box (Bullets)
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(8), Inches(4.5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Bullet 1
p1 = content_tf.paragraphs[0]
p1.text = "Definition: The sum total of all that surrounds living organisms."
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.space_after = Pt(12)

# Bullet 2
p2 = content_tf.add_paragraph()
p2.text = "Components: Includes biotic (living) and abiotic (non-living) factors."
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.space_after = Pt(12)

# Bullet 3
p3 = content_tf.add_paragraph()
p3.text = "Significance: Plays a critical role in sustaining life on Earth."
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(255, 255, 255)

# Add Image
# Position: Right side
# Dimensions: Fit nicely on the right
image_left = Inches(9.2)
image_top = Inches(1)
image_width = Inches(3.8)
image_height = Inches(6)

try:
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)
except FileNotFoundError:
    # Fallback if image.png is missing, though instructions say it's available
    pass

# Save
prs.save('output.pptx')