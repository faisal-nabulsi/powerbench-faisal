from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the background image
# The image covers the entire slide area
slide.shapes.add_picture('image.png', Inches(0), Inches(0), prs.slide_width, prs.slide_height)

# Add the Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = False

title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Now: I, CAN"
title_run.font.size = Pt(40)
title_run.font.bold = True
# Using White for high contrast against a typical photo background
title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Add the Content (Bullet Points)
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(5.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Bullet Point 1
p1 = content_tf.add_paragraph()
p1.text = "Differentiate between different types of market"
p1.space_after = Pt(18)
run1 = p1.runs[0]
run1.font.size = Pt(24)
run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White

# Bullet Point 2
p2 = content_tf.add_paragraph()
p2.text = "Explain why society needs different types of market,"
p2.space_after = Pt(18)
run2 = p2.runs[0]
run2.font.size = Pt(24)
run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White

# Save the presentation
prs.save('output.pptx')