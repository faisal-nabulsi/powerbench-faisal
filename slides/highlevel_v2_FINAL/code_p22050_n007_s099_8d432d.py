from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation instance
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.text = "Dynamics of Local and Global Culture"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

# --- Add Content ---
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(4.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Clear the default empty paragraph
if content_frame.paragraphs[0].text == '':
    content_frame.paragraphs[0].clear()

# Helper function to add formatted bullet points
def add_bullet(text_frame, level, text, is_bold=False):
    """Adds a bullet point to the text frame."""
    # Use the first paragraph if it's empty, otherwise add a new one
    if len(text_frame.paragraphs) == 1 and text_frame.paragraphs[0].text == '':
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = is_bold
    run.font.color.rgb = RGBColor(51, 51, 51) # Dark Grey
    return p

# 1. Cultural Differentialism (Starting point)
add_bullet(content_frame, 0, "Cultural Differentialism", is_bold=True)

# 2. Cultural Hybridization (With key points)
add_bullet(content_frame, 0, "Cultural Hybridization", is_bold=True)

# Sub-bullets for Hybridization
add_bullet(content_frame, 1, "Definition: The blending of cultural elements from different societies to create new cultural forms.")
add_bullet(content_frame, 1, "Integration of Cultures: The merging of distinct social practices, beliefs, and identities into a cohesive whole.")
add_bullet(content_frame, 1, "Glocalization: The adaptation of global products and practices to fit specific local cultural contexts.")

# 3. Third Perspective (Added to fulfill "outlines the three perspectives" requirement)
# Commonly "Cultural Homogenization" or "Cultural Imperialism" completes this triad.
add_bullet(content_frame, 0, "Cultural Homogenization", is_bold=True)

# Save the presentation
prs.save('output.pptx')