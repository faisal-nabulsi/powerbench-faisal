from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()
# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Position: Left 1", Top 0.5", Width ~11.3", Height 1.2"
title_shape = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.2))
title_tf = title_shape.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "The Technology used in Blockchain"
title_run.font.size = Pt(36)
title_run.font.bold = True

# --- Add Bullet Points ---
# Position: Left 1", Top 2.2", Width 5.5", Height 4"
content_shape = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.5), Inches(4.0))
content_tf = content_shape.text_frame
content_tf.word_wrap = True

bullet_items = [
    "Cryptographic Keys",
    "Network Protocol",
    "Distributed Ledger Technology",
    "Hashing"
]

# Set the first paragraph (default empty paragraph)
p = content_tf.paragraphs[0]
p.text = f"• {bullet_items[0]}"
p.font.size = Pt(24)
p.space_after = Pt(10)

# Add remaining items
for item in bullet_items[1:]:
    new_p = content_tf.add_paragraph()
    new_p.text = f"• {item}"
    new_p.font.size = Pt(24)
    new_p.space_after = Pt(10)

# --- Add Image ---
# Position: Left 7.2", Top 2.2", Width 5", Height 3.3"
# Placing it to the right of the text content
slide.shapes.add_picture('image.png', Inches(7.2), Inches(2.2), Inches(5.0), Inches(3.3))

# Save the presentation
prs.save('output.pptx')