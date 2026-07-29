from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the background image
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# Define dimensions and position for the red box
box_left = Inches(2.1665)
box_top = Inches(2.0)
box_width = Inches(9)
box_height = Inches(3.5)

# Add the red box shape
red_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height)
red_box.fill.solid()
red_box.fill.fore_color.rgb = RGBColor(255, 0, 0)
red_box.line.fill.background() # Remove border

# Define dimensions and position for the text box
text_left = Inches(2.4165)
text_top = Inches(2.25)
text_width = Inches(8.5)
text_height = Inches(3.0)

# Add the text box
text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
tf = text_box.text_frame
tf.word_wrap = True

# Add the text content
p = tf.paragraphs[0]
p.text = "Which type of market is beneficial for the customers? Why?"
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')