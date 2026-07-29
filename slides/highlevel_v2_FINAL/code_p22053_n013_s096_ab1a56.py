from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_p = title_tf.paragraphs[0]
title_run = title_p.add_run()
title_run.text = "Global and Local Cultural Products"
title_run.font.size = Pt(32)
title_run.font.bold = True
title_p.alignment = 1 # CENTER

# --- Section Header ---
header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
header_tf = header_box.text_frame
header_tf.word_wrap = True
header_p = header_tf.paragraphs[0]
header_run = header_p.add_run()
header_run.text = "Global Product"
header_run.font.size = Pt(24)
header_run.font.bold = True

# --- Definition ---
def_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.8))
def_tf = def_box.text_frame
def_tf.word_wrap = True
def_p = def_tf.paragraphs[0]
def_run = def_p.add_run()
def_run.text = "A global product is a tangible good or intangible service marketed and sold internationally. It often transcends cultural boundaries by maintaining a consistent brand identity, packaging, and core value proposition across diverse markets worldwide."
def_run.font.size = Pt(14)

# --- Brand Logos ---
# Layout parameters
num_logos = 4
logo_width = Inches(2.0)
logo_height = Inches(2.0)
gap = Inches(0.4)
total_width = num_logos * logo_width + (num_logos - 1) * gap
start_x = (Inches(13.333) - total_width) / 2
start_y = Inches(4.2)

brands = ["Coca-Cola", "McDonald's", "Apple", "Adidas"]

for i in range(num_logos):
    x = start_x + i * (logo_width + gap)
    
    # Add Placeholder Image
    try:
        slide.shapes.add_picture('image.png', x, start_y, logo_width, logo_height)
    except FileNotFoundError:
        # In case the image is missing, we skip adding it but keep the label
        pass
    
    # Add Brand Label
    label_box = slide.shapes.add_textbox(x, start_y + logo_height + Inches(0.1), logo_width, Inches(0.4))
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_p = label_tf.paragraphs[0]
    label_run = label_p.add_run()
    label_run.text = brands[i]
    label_p.alignment = 1
    label_run.font.size = Pt(12)
    label_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')