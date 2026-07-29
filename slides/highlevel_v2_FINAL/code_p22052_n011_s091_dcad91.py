from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize Presentation
prs = Presentation()

# Set 16:9 Widescreen Dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a slide using the 'Title and Content' layout (Index 1)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set the Slide Title
title_shape = slide.shapes.title
title_shape.text = "Origin of Blockchain Technology"

# Get the Content Placeholder
# In standard layouts, the content placeholder is typically at index 1
content_placeholder = slide.shapes.placeholders[1]
text_frame = content_placeholder.text_frame
text_frame.word_wrap = True

# Add Bullet Point 1: Creation
p1 = text_frame.paragraphs[0]
p1.text = "Creation: Blockchain technology was created in 2008 by an anonymous entity known as Satoshi Nakamoto to enable decentralized, trustless digital transactions."
p1.font.size = Pt(18)
p1.font.name = 'Calibri'

# Add Bullet Point 2: Implementation in Bitcoin
p2 = text_frame.add_paragraph()
p2.text = "Implementation in Bitcoin: First implemented as the underlying database technology for the Bitcoin cryptocurrency, launching in 2009 as a public, immutable ledger."
p2.font.size = Pt(18)
p2.font.name = 'Calibri'

# Add Bullet Point 3: Proposal of Private Blockchains
p3 = text_frame.add_paragraph()
p3.text = "Proposal of Private Blockchains: Later proposals adapted the technology for business use, introducing private or permissioned blockchains to offer enterprise-level security, privacy, and control."
p3.font.size = Pt(18)
p3.font.name = 'Calibri'

# Save the presentation
prs.save('output.pptx')