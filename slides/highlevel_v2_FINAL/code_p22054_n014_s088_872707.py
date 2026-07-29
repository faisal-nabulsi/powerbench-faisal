from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set dark background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)

# Add textured image on the right side
# Image starts at x=7.5, covering the right ~5.8 inches
slide.shapes.add_picture(
    'image.png',
    left=Inches(7.5),
    top=Inches(0),
    width=Inches(5.833),
    height=Inches(7.5)
)

# Add Title
title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
title_tf = title_shape.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Now: I, CAN,"
title_run.font.size = Pt(40)
title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
title_run.font.bold = True

# Add Content (Bullet Points)
content_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(6), Inches(3.5))
content_tf = content_shape.text_frame
content_tf.word_wrap = True

# Bullet 1
p1 = content_tf.add_paragraph()
p1.text = "• Define the term Market"
p1.space_after = Pt(14)
run1 = p1.runs[0]
run1.font.size = Pt(24)
run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Bullet 2
p2 = content_tf.add_paragraph()
p2.text = "• Explain how products reach to market"
p2.space_after = Pt(14)
run2 = p2.runs[0]
run2.font.size = Pt(24)
run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Save the presentation
prs.save('output.pptx')