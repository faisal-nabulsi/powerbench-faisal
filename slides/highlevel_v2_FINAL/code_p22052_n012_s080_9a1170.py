from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()

# Set 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.2))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "Successful Individuals with Top Academic Records"
title_p.font.size = Pt(36)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79) # Dark Slate Blue

# --- Content Data ---
profiles = [
    {
        "name": "Neil deGrasse Tyson",
        "title": "Astrophysicist & Author",
        "education": "Harvard University (BA), Columbia University (PhD)",
    },
    {
        "name": "Mae Jemison",
        "title": "Engineer & NASA Astronaut",
        "education": "Stanford University (BS), Cornell University (MD)",
    },
    {
        "name": "Carl Sagan",
        "title": "Astronomer & Science Popularizer",
        "education": "University of Chicago (BS/MS), UC Berkeley (PhD)",
    }
]

# --- Layout Constants ---
IMAGE_WIDTH = Inches(3.0)
IMAGE_HEIGHT = Inches(3.0)
IMAGE_TOP = Inches(1.8)
TEXT_BOX_HEIGHT = Inches(2.5)
LEFT_MARGINS = [Inches(1.15), Inches(5.25), Inches(9.35)] 

for i in range(len(profiles)):
    left_pos = LEFT_MARGINS[i]
    
    # Add Placeholder Image
    picture = slide.shapes.add_picture('image.png', left_pos, IMAGE_TOP, IMAGE_WIDTH, IMAGE_HEIGHT)

    # Add Text Box for Name, Title, Education
    text_top = IMAGE_TOP + IMAGE_HEIGHT + Inches(0.3)
    text_box = slide.shapes.add_textbox(left_pos, text_top, IMAGE_WIDTH, TEXT_BOX_HEIGHT)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Name
    name_para = text_frame.paragraphs[0]
    name_run = name_para.add_run()
    name_run.text = profiles[i]["name"]
    name_run.font.bold = True
    name_run.font.size = Pt(18)
    name_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    
    # Title
    title_para = text_frame.add_paragraph()
    title_para.space_before = Pt(6)
    title_run = title_para.add_run()
    title_run.text = profiles[i]["title"]
    title_run.font.italic = True
    title_run.font.size = Pt(14)
    title_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    
    # Education
    edu_para = text_frame.add_paragraph()
    edu_para.space_before = Pt(6)
    edu_run = edu_para.add_run()
    edu_run.text = profiles[i]["education"]
    edu_run.font.size = Pt(13)
    edu_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the file
prs.save('output.pptx')