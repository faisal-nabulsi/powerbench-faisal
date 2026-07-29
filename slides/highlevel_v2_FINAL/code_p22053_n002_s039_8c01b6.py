from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide. Index 6 is typically the Blank layout in default presentations.
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the background image
# 'image.png' is specified to be in the working directory.
# We stretch it to fill the entire slide canvas.
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Add the Title Text "ELO's"
# Positioned at the top, centered.
left_title = Inches(0)
top_title = Inches(0.5)
width_title = Inches(13.333)
height_title = Inches(1.5)

title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
tf_title = title_box.text_frame
tf_title.word_wrap = True

# Style the title paragraph
p_title = tf_title.paragraphs[0]
p_title.text = "ELO's"
p_title.font.size = Pt(60)          # Large font size
p_title.font.bold = True            # Bold for emphasis
p_title.font.color.rgb = RGBColor(255, 255, 255) # White color for contrast
p_title.font.name = "Arial"         # Clear, readable font
p_title.alignment = PP_ALIGN.CENTER # Center alignment

# 3. Add the Questions
# Positioned centrally below the title.
left_q = Inches(1.5)
top_q = Inches(3.0)
width_q = Inches(10.333)
height_q = Inches(3.0)

questions_box = slide.shapes.add_textbox(left_q, top_q, width_q, height_q)
tf_q = questions_box.text_frame
tf_q.word_wrap = True

# First Question
p1 = tf_q.paragraphs[0]
p1.text = "1. What is market?"
p1.font.size = Pt(36)
p1.font.color.rgb = RGBColor(255, 255, 255) # White
p1.font.name = "Arial"
p1.alignment = PP_ALIGN.CENTER

# Second Question
p2 = tf_q.add_paragraph()
p2.text = "2. How product reach to market?"
p2.font.size = Pt(36)
p2.font.color.rgb = RGBColor(255, 255, 255) # White
p2.font.name = "Arial"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(25) # Add vertical spacing between the two questions

# Save the final presentation
prs.save('output.pptx')