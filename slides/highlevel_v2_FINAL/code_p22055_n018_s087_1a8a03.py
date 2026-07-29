from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Blockchain Technology"
title_run.font.size = Pt(40)
title_run.font.bold = True
title_para.alignment = 1  # Center alignment

# Add the first image placeholder (Blockchain illustration)
# Using 'image.png' as the placeholder
slide.shapes.add_picture('image.png', Inches(0.5), Inches(2.2), Inches(5.5), Inches(4.8))

# Add the second image placeholder (BLOCKCHAIN with keyboard and smartphone)
# Using 'image.png' as the placeholder
slide.shapes.add_picture('image.png', Inches(7.3), Inches(2.2), Inches(5.5), Inches(4.8))

# Save the presentation
prs.save('output.pptx')