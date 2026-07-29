from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set the slide width and height to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
title_text = "VIEWING ACTIVITY"
title_width = Inches(6.0)
title_left = Inches((13.333 - 6.0) / 2)
title_top = Inches(0.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER

title_run = title_para.add_run()
title_run.text = title_text
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# --- Add Video Camera Icon ---
# Using 'image.png' as the placeholder for the icon
icon_width = Inches(4.0)
icon_height = Inches(3.0)
icon_left = Inches((13.333 - 4.0) / 2)
icon_top = Inches(1.8)

slide.shapes.add_picture('image.png', icon_left, icon_top, icon_width, icon_height)

# --- Add Subtitle ---
subtitle_text = "Friends | Joey Doesn't Share Food!"
sub_width = Inches(10.0)
sub_left = Inches((13.333 - 10.0) / 2)
sub_top = Inches(5.2)

sub_box = slide.shapes.add_textbox(sub_left, sub_top, sub_width, Inches(1.0))
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_para = sub_tf.paragraphs[0]
sub_para.alignment = PP_ALIGN.CENTER

sub_run = sub_para.add_run()
sub_run.text = subtitle_text
sub_run.font.size = Pt(24)
sub_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# Save the presentation
prs.save('output.pptx')