from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Top Grades Are Not the Sole Determinant of Success"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
title_para.alignment = PP_ALIGN.CENTER

# Add Content Text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7.5), Inches(4.5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Intro text
intro_para = content_tf.paragraphs[0]
intro_run = intro_para.add_run()
intro_run.text = "True success is multifaceted. It relies on:"
intro_run.font.size = Pt(20)
intro_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Key terms with colors
terms = [
    ("Resilience", RGBColor(0xE6, 0x00, 0x00)), # Red
    ("Creativity", RGBColor(0x00, 0x71, 0xC5)), # Blue
    ("Empathy", RGBColor(0x00, 0x99, 0x00)),    # Green
    ("Adaptability", RGBColor(0xFF, 0x80, 0x00)) # Orange
]

for term, color in terms:
    para = content_tf.add_paragraph()
    run = para.add_run()
    run.text = f"• {term}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = color

# Closing text
closing_para = content_tf.add_paragraph()
closing_run = closing_para.add_run()
closing_run.text = "\nand many other soft skills."
closing_run.font.size = Pt(20)
closing_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Add Image
# Using placeholder image.png
image = slide.shapes.add_picture('image.png', Inches(8.5), Inches(2.5), Inches(4.5), Inches(4.5))

# Save the presentation
prs.save('output.pptx')