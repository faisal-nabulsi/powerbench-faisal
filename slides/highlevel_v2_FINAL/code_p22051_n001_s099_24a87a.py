from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Create a text box for the title at the top
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p_title = title_frame.paragraphs[0]
p_title.text = "Advantages of Blockchain"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.alignment = PP_ALIGN.CENTER

# --- Add Bullet Points ---
# List of key advantages
advantages = [
    "Decentralization: No central authority; distributed control.",
    "Transparency: Shared, viewable ledger of transactions.",
    "Security: Cryptographic hashing protects data integrity.",
    "Immutability: Data cannot be altered once recorded.",
    "Efficiency: Streamlines processes and reduces costs.",
    "Traceability: Complete audit trail of all activities."
]

# Define position and size for the bullet points text box (Left side)
txt_left = Inches(0.5)
txt_top = Inches(1.8)
txt_width = Inches(6.0)
txt_height = Inches(5.2)

textbox = slide.shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
tf = textbox.text_frame
tf.word_wrap = True

# Populate text box with bullet points
for i, item in enumerate(advantages):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    # Add a bullet character manually for consistency
    p.text = f"• {item}"
    p.font.size = Pt(22)
    p.space_after = Pt(12)

# --- Add Diagram (Image) ---
# Define position and size for the diagram (Right side)
img_left = Inches(7.2)
img_top = Inches(1.8)
img_width = Inches(5.6)
img_height = Inches(5.2)

# Add the placeholder image 'image.png'
slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')