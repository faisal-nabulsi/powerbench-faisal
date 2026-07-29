from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
# Add a text box for the title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_para = title_tf.paragraphs[0]
title_para.text = "LESSON REFLECTION:"
title_para.font.size = Pt(36)
title_para.font.bold = True

# --- Questions Section ---
# Add a text box for the questions
questions_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(4.5))
questions_tf = questions_box.text_frame
questions_tf.word_wrap = True

# Question 1
p1 = questions_tf.paragraphs[0]
p1.text = "1. What was the most important thing you learned today?"
p1.font.size = Pt(24)
p1.space_after = Pt(24)

# Question 2
p2 = questions_tf.add_paragraph()
p2.text = "2. What part of the lesson was challenging for you?"
p2.font.size = Pt(24)
p2.space_after = Pt(24)

# Question 3
p3 = questions_tf.add_paragraph()
p3.text = "3. How can you apply what you learned to real-life situations?"
p3.font.size = Pt(24)
p3.space_after = Pt(24)

# Save the presentation
prs.save('output.pptx')