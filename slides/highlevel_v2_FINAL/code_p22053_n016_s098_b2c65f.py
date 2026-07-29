from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "5 Stages of Development of Media"
title_run.font.size = Pt(32)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 51, 102)

# Add Subtitle/Heading
heading_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.8))
heading_tf = heading_box.text_frame
heading_para = heading_tf.paragraphs[0]
heading_run = heading_para.add_run()
heading_run.text = "4. ELECTRONIC MEDIA"
heading_run.font.size = Pt(24)
heading_run.font.bold = True
heading_run.font.color.rgb = RGBColor(0, 102, 0)

# Add Bullet Points
bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.5))
bullets_tf = bullets_box.text_frame
bullets_tf.word_wrap = True

# Bullet 1
bullet1_para = bullets_tf.add_paragraph()
bullet1_para.level = 0
bullet1_run = bullet1_para.add_run()
bullet1_run.text = "Components: Includes radio, television, and digital platforms, enabling real-time transmission of audio and visual content."
bullet1_run.font.size = Pt(14)

# Bullet 2
bullet2_para = bullets_tf.add_paragraph()
bullet2_para.level = 0
bullet2_run = bullet2_para.add_run()
bullet2_run.text = "Significance: Revolutionized communication by breaking geographical barriers and fostering a globally connected society."
bullet2_run.font.size = Pt(14)

# Add Image
# Positioning on the right side
image_left = Inches(7.0)
image_top = Inches(1.5)
image_width = Inches(6.0)
image_height = Inches(5.5)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')