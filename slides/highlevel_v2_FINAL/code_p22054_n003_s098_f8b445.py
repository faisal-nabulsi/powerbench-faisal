from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Create Presentation
prs = Presentation()

# 2. Set Slide Dimensions (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 3. Add a Blank Slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 4. Add Title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_frame = title_shape.text_frame
title_frame.word_wrap = True
p_title = title_frame.paragraphs[0]
p_title.alignment = PP_ALIGN.LEFT
# Using the exact text from instruction (note the curly apostrophe if present, using standard for safety or copy-paste)
p_title.text = "Globalization isn’t possible to occur without media"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0, 51, 102)

# 5. Add Content Text
# Left side text box
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(5))
tf = text_box.text_frame
tf.word_wrap = True

# Bullet 1
p1 = tf.paragraphs[0]
r1 = p1.add_run()
r1.text = "1. Digital Media"
r1.font.bold = True
r1.font.size = Pt(22)
r1.font.color.rgb = RGBColor(51, 51, 51)
p1_desc = tf.add_paragraph()
r1d = p1_desc.add_run()
r1d.text = "Internet and social media break down geographical barriers, enabling real-time global communication and connectivity."
r1d.font.size = Pt(16)
r1d.font.color.rgb = RGBColor(80, 80, 80)
p1_desc.space_after = Pt(20)

# Bullet 2
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "2. Mass Media (Broadcast)"
r2.font.bold = True
r2.font.size = Pt(22)
r2.font.color.rgb = RGBColor(51, 51, 51)
p2_desc = tf.add_paragraph()
r2d = p2_desc.add_run()
r2d.text = "Television and news networks disseminate information globally, unifying audiences and spreading cultural content."
r2d.font.size = Pt(16)
r2d.font.color.rgb = RGBColor(80, 80, 80)
p2_desc.space_after = Pt(20)

# Bullet 3
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = "3. Print Media"
r3.font.bold = True
r3.font.size = Pt(22)
r3.font.color.rgb = RGBColor(51, 51, 51)
p3_desc = tf.add_paragraph()
r3d = p3_desc.add_run()
r3d.text = "Books and newspapers facilitate the spread of literacy, language, and complex ideas across borders."
r3d.font.size = Pt(16)
r3d.font.color.rgb = RGBColor(80, 80, 80)

# 6. Add Images
# Using 'image.png' for illustrations
slide.shapes.add_picture('image.png', Inches(7), Inches(1.8), Inches(5.5), Inches(2.5))
slide.shapes.add_picture('image.png', Inches(7), Inches(4.5), Inches(5.5), Inches(2.5))

# 7. Save
prs.save('output.pptx')