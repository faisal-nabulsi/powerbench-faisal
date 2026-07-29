from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout in default templates)
try:
    slide_layout = prs.slide_layouts[6]
except IndexError:
    # Fallback to first layout if blank layout index is unavailable
    slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Create a text box for the title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "DEFINITION OF TERMS"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.alignment = 2  # Center alignment (0=Left, 1=Right, 2=Center)

# --- Add Media Section (Left Side) ---
media_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.5), Inches(4.0))
media_tf = media_box.text_frame
media_tf.word_wrap = True

# Media Heading (Bold)
media_heading = media_tf.paragraphs[0]
media_run = media_heading.add_run()
media_run.text = "MEDIA"
media_run.font.size = Pt(28)
media_run.font.bold = True

# Media Description
media_desc = media_tf.add_paragraph()
media_desc.space_before = Pt(12)
media_desc_run = media_desc.add_run()
media_desc_run.text = "The collective means of mass communication, including print, broadcast, and digital platforms, used to inform, educate, or entertain the public."
media_desc_run.font.size = Pt(16)

# --- Add Culture Section (Right Side) ---
culture_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(5.5), Inches(4.0))
culture_tf = culture_box.text_frame
culture_tf.word_wrap = True

# Culture Heading (Bold)
culture_heading = culture_tf.paragraphs[0]
culture_run = culture_heading.add_run()
culture_run.text = "CULTURE"
culture_run.font.size = Pt(28)
culture_run.font.bold = True

# Culture Description
culture_desc = culture_tf.add_paragraph()
culture_desc.space_before = Pt(12)
culture_desc_run = culture_desc.add_run()
culture_desc_run.text = "The shared patterns of behaviors, beliefs, values, and artifacts that the members of a society use to cope with their world and transmit from generation to generation."
culture_desc_run.font.size = Pt(16)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')