from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Create Red Background
# We add a full-slide rectangle shape to set the background color reliably.
# Shape ID 1 corresponds to a Rectangle.
bg_shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(255, 0, 0) # Red
bg_shape.line.fill.background() # Remove border/line
# Note: Since this shape is added first, subsequent textboxes will render on top of it.

# 2. Add "BIG QUESTION" Text
# Requirements: Smaller blue text, positioned above the main question, centered.
# Horizontal Centering: (13.333 height - width) / 2. Let's use width 4 inches.
small_text_left = Inches(4.6665) 
small_text_top = Inches(2.0)
small_text_width = Inches(4.0)
small_text_height = Inches(1.0)

txBox_small = slide.shapes.add_textbox(small_text_left, small_text_top, small_text_width, small_text_height)
tf_small = txBox_small.text_frame
tf_small.word_wrap = True

p_small = tf_small.paragraphs[0]
p_small.text = "BIG QUESTION"
p_small.alignment = PP_ALIGN.CENTER

run_small = p_small.runs[0]
run_small.font.size = Pt(32) # Smaller font
run_small.font.color.rgb = RGBColor(0, 0, 255) # Blue
run_small.font.bold = True

# 3. Add "How do markets help us?" Text
# Requirements: Large white text, centered.
# Horizontal Centering: Let's use width 10 inches.
main_text_left = Inches(1.6665)
main_text_top = Inches(3.5) # Positioned below the first text box
main_text_width = Inches(10.0)
main_text_height = Inches(1.5)

txBox_main = slide.shapes.add_textbox(main_text_left, main_text_top, main_text_width, main_text_height)
tf_main = txBox_main.text_frame
tf_main.word_wrap = True

p_main = tf_main.paragraphs[0]
p_main.text = "How do markets help us?"
p_main.alignment = PP_ALIGN.CENTER

run_main = p_main.runs[0]
run_main.font.size = Pt(54) # Large font
run_main.font.color.rgb = RGBColor(255, 255, 255) # White
run_main.font.bold = True

# Save the presentation
prs.save('output.pptx')