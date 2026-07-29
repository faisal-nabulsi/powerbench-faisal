from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_nyc_slide():
    # Initialize the presentation
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add a blank slide to allow for custom layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Add Title "NYC"
    # Positioning: Top center
    title_left = Inches(1.5)
    title_top = Inches(0.3)
    title_width = Inches(10.333)
    title_height = Inches(1.0)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "NYC"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0) # Black text for visibility
    p.alignment = PP_ALIGN.CENTER

    # 2. Add Prominent Image of Brooklyn Bridge
    # Using 'image.png' as the placeholder for the bridge image.
    # Positioning: Center of the slide, below the title.
    img_left = Inches(1.5)
    img_top = Inches(1.5)
    img_width = Inches(10.333)
    img_height = Inches(5.0)

    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

    # 3. Add Airbnb Logo at the bottom
    # Using 'image.png' as the placeholder for the logo, as per instructions.
    # Positioning: Bottom right corner.
    logo_left = Inches(11.0)
    logo_top = Inches(6.5)
    logo_width = Inches(2.0)
    logo_height = Inches(0.8)

    slide.shapes.add_picture('image.png', logo_left, logo_top, logo_width, logo_height)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_nyc_slide()