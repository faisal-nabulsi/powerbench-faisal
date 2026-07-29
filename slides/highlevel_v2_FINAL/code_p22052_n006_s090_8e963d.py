from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Select a blank slide layout
blank_layout = None
for layout in prs.slide_layouts:
    if layout.name == "Blank":
        blank_layout = layout
        break
if not blank_layout:
    # Fallback to the last layout if 'Blank' is not explicitly named
    blank_layout = prs.slide_layouts[-1]

slide = prs.slides.add_slide(blank_layout)

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p_title = title_frame.paragraphs[0]
p_title.text = "Obstacles that Elon faced"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(40, 40, 40)

# --- Content ---
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(4.8))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Data: (Challenge Name, Brief Description)
challenges = [
    (
        "2008 Bankruptcy Crisis",
        "Both Tesla and SpaceX were simultaneously days away from total financial collapse before securing critical last-minute funding."
    ),
    (
        "SpaceX Launch Failures",
        "The Falcon 1 rocket failed three times in a row, nearly ending the company before achieving orbit on the fourth attempt."
    ),
    (
        "Model 3 Production Hell",
        "Scaling manufacturing to meet massive demand caused severe bottlenecks, leading Musk to sleep on the factory floor for weeks."
    ),
    (
        "Regulatory & SEC Scrutiny",
        "Faced intense investigations regarding public statements on social media and complex compliance issues with federal regulations."
    ),
    (
        "Cybertruck Engineering",
        "Developing a viable manufacturing process for an ultra-hard stainless steel exoskeleton presented unique industrial hurdles."
    )
]

# Styles
color_blue = RGBColor(0, 85, 164)
color_gray = RGBColor(60, 60, 60)
size_title = Pt(24)
size_desc = Pt(18)

# Populate Content
for i, (title, desc) in enumerate(challenges):
    # Add Challenge Title (Bullet)
    if i == 0:
        p_t = content_frame.paragraphs[0]
    else:
        p_t = content_frame.add_paragraph()
    
    p_t.text = "• " + title
    p_t.font.size = size_title
    p_t.font.bold = True
    p_t.font.color.rgb = color_blue
    p_t.space_after = Pt(2)  # Minimal space between title and description
    
    # Add Description
    p_d = content_frame.add_paragraph()
    p_d.text = "  " + desc
    p_d.font.size = size_desc
    p_d.font.color.rgb = color_gray
    p_d.space_after = Pt(12)  # Space before the next challenge

prs.save('output.pptx')