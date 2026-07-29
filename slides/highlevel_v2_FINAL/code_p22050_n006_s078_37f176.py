from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen as requested
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to allow free positioning of elements
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "PART 1"
# Positioned at the top center
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "PART 1"
title_para.font.size = Pt(48)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# 2. Add Discussion Prompt about Joey's date
# Positioned in the left column
prompt_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5.5), Inches(3.0))
prompt_frame = prompt_box.text_frame
prompt_frame.word_wrap = True

p_prompt_header = prompt_frame.paragraphs[0]
p_prompt_header.text = "Discussion & Short Answers"
p_prompt_header.font.size = Pt(24)
p_prompt_header.font.bold = True

p_prompt_body = prompt_frame.add_paragraph()
p_prompt_body.text = "Reflect on Joey's date:\n\n• What were the key moments that went wrong?\n• How did his personality traits impact the outcome?\n• Provide short answers on how to improve."
p_prompt_body.font.size = Pt(18)
p_prompt_body.space_before = Pt(12)

# 3. Add "PART 2" Instruction
# Positioned in the left column, below the discussion prompt
part2_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(5.5), Inches(2.0))
part2_frame = part2_box.text_frame
part2_frame.word_wrap = True

p_part2_header = part2_frame.paragraphs[0]
p_part2_header.text = "PART 2"
p_part2_header.font.size = Pt(24)
p_part2_header.font.bold = True

p_part2_body = part2_frame.add_paragraph()
p_part2_body.text = "Refer back to the previous section for a comparison of the date's outcome versus expectations."
p_part2_body.font.size = Pt(18)
p_part2_body.space_before = Pt(12)

# 4. Add Image of Joey
# Positioned in the right column
# Coordinates: Left=7 (aligned next to text), Top=1.5, Width=5.5, Height=5.0
slide.shapes.add_picture('image.png', Inches(7), Inches(1.5), Inches(5.5), Inches(5.0))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')