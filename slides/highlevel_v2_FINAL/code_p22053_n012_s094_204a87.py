from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background color to black for high contrast with yellow text
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add image on the left side
# Using the provided placeholder 'image.png'
image = slide.shapes.add_picture('image.png', Inches(0.5), Inches(1.0), Inches(4.5), Inches(5.5))

# Add text box on the right side
left = Inches(5.5)
top = Inches(2.0)
width = Inches(7.5)
height = Inches(3.5)
text_box = slide.shapes.add_textbox(left, top, width, height)

# Configure text properties
tf = text_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "As an individual how we can ensure equality in society?"
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 0) # Yellow
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')