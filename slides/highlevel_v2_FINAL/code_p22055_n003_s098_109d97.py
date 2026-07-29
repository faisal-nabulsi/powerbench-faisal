from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "Disadvantages of Blockchain"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.LEFT

# --- Left Column: Bullet Points ---
# Sub-header
header_left = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.5))
header_frame = header_left.text_frame
header_para = header_frame.paragraphs[0]
header_para.text = "Key Disadvantages"
header_para.font.size = Pt(24)
header_para.font.bold = True
header_para.font.color.rgb = RGBColor(0, 102, 153)

# Bullet points content
disadvantages = [
    "Scalability: Limited transaction throughput compared to centralized systems.",
    "Energy Consumption: High electricity usage, particularly in Proof-of-Work.",
    "Security Risks: Vulnerabilities in smart contracts and private key management.",
    "Regulatory Uncertainty: Evolving legal frameworks create compliance challenges.",
    "Irreversibility: Transactions cannot be undone once confirmed.",
    "Complexity: Steep learning curve for development and integration."
]

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(6.0), Inches(4.5))
tf = text_box.text_frame
tf.word_wrap = True

for i, item in enumerate(disadvantages):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(10)

# --- Right Column: Diagram ---
# Sub-header
header_right = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(6.0), Inches(0.5))
header_frame_r = header_right.text_frame
header_para_r = header_frame_r.paragraphs[0]
header_para_r.text = "Advantages vs. Disadvantages"
header_para_r.font.size = Pt(24)
header_para_r.font.bold = True
header_para_r.font.color.rgb = RGBColor(0, 102, 153)

# Diagram Image
# Using the placeholder image 'image.png' as requested for visual elements
slide.shapes.add_picture('image.png', Inches(7.0), Inches(2.7), Inches(6.0), Inches(4.5))

# Save the presentation
prs.save('output.pptx')