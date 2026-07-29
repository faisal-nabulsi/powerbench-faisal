from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Find the blank slide layout to add a clean slide
slide_layout = None
for layout in prs.slide_layouts:
    if layout.name == "Blank":
        slide_layout = layout
        break

# Fallback to index 6 (usually Blank) if not found by name
if slide_layout is None:
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(slide_layout)

# Set the background color to Coral (RGB: 255, 127, 80)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0x7F, 0x50)

# Calculate positions to center the text box on the slide
# Box width ~8.333, Box height ~2.5
# Slide Width ~13.333, Slide Height ~7.5
left = Inches(2.5)
top = Inches(2.5)
width = Inches(8.333)
height = Inches(2.5)

# Add the text box
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

# Configure the paragraph alignment
paragraph = text_frame.paragraphs[0]
paragraph.alignment = PP_ALIGN.CENTER

# Add text and formatting
run = paragraph.add_run()
run.text = "This is Bélo"
run.font.bold = True
run.font.size = Pt(48)
run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# Save the presentation
prs.save('output.pptx')