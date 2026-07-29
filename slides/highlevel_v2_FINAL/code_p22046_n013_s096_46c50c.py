from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()

# Set Slide Dimensions (16:9 Widescreen)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "PART 3"
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.2))
title_frame = title_shape.text_frame
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "PART 3"
title_run.font.size = Pt(54)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5) # Blue

# 2. Add Instruction Text
instr_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.8))
instr_frame = instr_shape.text_frame
instr_para = instr_frame.paragraphs[0]
instr_run = instr_para.add_run()
instr_run.text = "Instruction: Fill in the blanks with the correct Past Simple or Past Continuous form of the verbs in parentheses."
instr_run.font.size = Pt(18)
instr_run.font.bold = True
instr_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey

# 3. Add Story Text
# Since no specific story was provided, here is a sample story for the exercise.
story_text = """Last weekend, the weather was beautiful. I ______ (decide) to go for a hike in the mountains. 
When I ______ (arrive) at the trail, I ______ (meet) some friends. 
We ______ (talk) about our weekend plans while we ______ (hike) up the steep path. 

Suddenly, one of us ______ (slip) on some loose rocks, but luckily, we ______ (not/hurt) ourselves. 
After the scrape, we ______ (continue) the hike and ______ (reach) the summit by noon. 
From there, we ______ (see) a breathtaking view of the valley. 

We ______ (take) many photos while the sun ______ (set), before we ______ (head) back down."""

story_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(4.5))
story_frame = story_shape.text_frame
story_frame.word_wrap = True

story_para = story_frame.paragraphs[0]
story_run = story_para.add_run()
story_run.text = story_text
story_run.font.size = Pt(20)
story_run.font.line_spacing = Pt(32)
story_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00) # Black

# Save the presentation
prs.save('output.pptx')