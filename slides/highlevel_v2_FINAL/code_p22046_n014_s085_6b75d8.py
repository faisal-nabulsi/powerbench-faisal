from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "PART 4"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

# Add Instruction
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(10), Inches(1.0))
instruction_tf = instruction_box.text_frame
instruction_para = instruction_tf.paragraphs[0]
instruction_run = instruction_para.add_run()
instruction_run.text = "List the events in Joey’s second date with Sarah in order from 2 - 6:"
instruction_run.font.size = Pt(18)
instruction_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Add Checklist
checklist_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(10), Inches(4.0))
checklist_tf = checklist_box.text_frame
checklist_tf.word_wrap = True

# Define checklist items
items = [
    "[x] Met Sarah at the cinema",
    "[ ] Watched the movie",
    "[ ] Ate dinner at the Italian restaurant",
    "[ ] Walked along the river",
    "[ ] Said goodnight",
    "[ ] Texted her later that night"
]

for i, item in enumerate(items):
    if i == 0:
        para = checklist_tf.paragraphs[0]
    else:
        para = checklist_tf.add_paragraph()
    
    run = para.add_run()
    run.text = item
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    
    # Add some spacing between items
    para.space_after = Pt(10)

# Add Image of Joey in the bottom right corner
# Using placeholder 'image.png'
image_width = Inches(2.5)
image_height = Inches(2.5)
image_left = prs.slide_width - image_width - Inches(0.5)
image_top = prs.slide_height - image_height - Inches(0.5)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')