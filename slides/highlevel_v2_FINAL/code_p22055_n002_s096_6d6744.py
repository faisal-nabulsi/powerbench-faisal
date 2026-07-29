from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set slide width and height to 16:9 (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
# Index 6 is typically the "Blank" layout in a default presentation
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Add Quote from Elon Musk on the Left ---
# Positioning: Left=0.5, Top=1.5, Width=6.0, Height=4.0
quote_box = slide.shapes.add_textbox(
    left=Inches(0.5),
    top=Inches(1.5),
    width=Inches(6.0),
    height=Inches(4.0)
)

quote_tf = quote_box.text_frame
quote_tf.word_wrap = True

# Format the paragraph
p_quote = quote_tf.paragraphs[0]
p_quote.text = '"When something is important enough, you do it even if the odds are not in your favor."'
p_quote.font.size = Pt(32)
p_quote.font.bold = True
p_quote.font.name = 'Calibri'

# --- 2. Add Image on the Right ---
# Positioning: Left=7.0 (0.5 start + 6.0 width + 0.5 gap), Top=1.5
# Width=5.333, Height=4.0
# Using the placeholder image 'image.png'
slide.shapes.add_picture(
    'image.png',
    left=Inches(7.0),
    top=Inches(1.5),
    width=Inches(5.333),
    height=Inches(4.0)
)

# --- 3. Add "Thank You!" message at the Bottom ---
# Positioning: Left=0.5, Top=6.2, Width=12.333, Height=1.0
thank_you_box = slide.shapes.add_textbox(
    left=Inches(0.5),
    top=Inches(6.2),
    width=Inches(12.333),
    height=Inches(1.0)
)

ty_tf = thank_you_box.text_frame
ty_tf.word_wrap = True

p_ty = ty_tf.paragraphs[0]
p_ty.text = 'Thank You!'
p_ty.font.size = Pt(48)
p_ty.font.bold = True
p_ty.font.name = 'Calibri'
p_ty.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')