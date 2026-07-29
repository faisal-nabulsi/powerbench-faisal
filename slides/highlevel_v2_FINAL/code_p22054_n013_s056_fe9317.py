from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set dark background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)

# Add placeholder images on the left side to represent various markets
# Note: python-pptx does not support cropping images to circles natively, 
# so standard picture shapes are used.
img_width = Inches(2.2)
img_height = Inches(2.2)
slide.shapes.add_picture('image.png', Inches(0.5), Inches(1.0), img_width, img_height)
slide.shapes.add_picture('image.png', Inches(0.5), Inches(3.5), img_width, img_height)
slide.shapes.add_picture('image.png', Inches(0.5), Inches(6.0), img_width, img_height)

# Add title "Types of Markets" at the bottom right
title_box = slide.shapes.add_textbox(Inches(8.5), Inches(5.5), Inches(4.5), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True

p = title_tf.paragraphs[0]
p.text = "Types of Markets"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

prs.save('output.pptx')