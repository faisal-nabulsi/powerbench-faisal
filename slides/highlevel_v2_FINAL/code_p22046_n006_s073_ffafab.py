from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Academic World and Professional World: Demands and Characteristics"
title_run.font.size = Pt(32)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
title_para.alignment = PP_ALIGN.CENTER

# --- Data ---
academic_data = [
    ("Evaluation and Feedback", [
        "Focus on grades and theoretical understanding.",
        "Feedback is often delayed and formal.",
        "Based on standardized rubrics."
    ]),
    ("Assessment", [
        "Exams, essays, and research papers.",
        "Standardized criteria for success.",
        "Emphasis on knowledge retention."
    ]),
    ("Feedback", [
        "Corrective and developmental nature.",
        "Usually provided by instructors.",
        "Written comments on assignments."
    ]),
    ("Timeframe", [
        "Semester-based or academic year.",
        "Fixed deadlines for submissions.",
        "Structured academic calendar."
    ])
]

professional_data = [
    ("Evaluation and Feedback", [
        "Focus on performance, results, and ROI.",
        "Feedback is continuous and often informal.",
        "Tied to business objectives."
    ]),
    ("Assessment", [
        "KPIs, project outcomes, and peer reviews.",
        "Dynamic and goal-oriented criteria.",
        "Emphasis on skill application."
    ]),
    ("Feedback", [
        "Action-oriented and immediate.",
        "Provided by managers, peers, and clients.",
        "360-degree feedback mechanisms."
    ]),
    ("Timeframe", [
        "Project-based or fiscal year.",
        "Flexible but deadline-driven.",
        "Quarterly or annual performance reviews."
    ])
]

# --- Helper Function ---
def add_column(slide, left, top, width, height, header, data, header_color):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Header
    p_header = tf.paragraphs[0]
    run_header = p_header.add_run()
    run_header.text = header
    run_header.font.size = Pt(24)
    run_header.font.bold = True
    run_header.font.color.rgb = header_color
    p_header.alignment = PP_ALIGN.CENTER
    
    # Content
    for section_title, bullets in data:
        p_section = tf.add_paragraph()
        p_section.space_before = Pt(14)
        run_section = p_section.add_run()
        run_section.text = section_title
        run_section.font.size = Pt(18)
        run_section.font.bold = True
        run_section.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        for bullet in bullets:
            p_bullet = tf.add_paragraph()
            p_bullet.level = 1
            p_bullet.space_before = Pt(4)
            run_bullet = p_bullet.add_run()
            run_bullet.text = bullet
            run_bullet.font.size = Pt(14)
            run_bullet.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# --- Add Columns ---
# Academic World (Left)
add_column(
    slide, 
    Inches(1), Inches(1.5), Inches(5.5), Inches(5.5), 
    "Academic World", 
    academic_data, 
    RGBColor(0x2E, 0x74, 0xB5)
)

# Professional World (Right)
add_column(
    slide, 
    Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.5), 
    "Professional World", 
    professional_data, 
    RGBColor(0xC0, 0x39, 0x2B)
)

# Save
prs.save('output.pptx')