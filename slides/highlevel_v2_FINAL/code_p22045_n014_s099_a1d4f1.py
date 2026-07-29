from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Title Section ---
# Define title text box properties
title_left = Inches(0.5)
title_top = Inches(0.5)
title_width = Inches(12.333)
title_height = Inches(1.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_box.text_frame
title_tf.word_wrap = True

# Set title text
title_paragraph = title_tf.paragraphs[0]
title_paragraph.text = "Origin of Blockchain Technology"
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50) # Dark Slate Gray

# --- Content Section ---
# Define content text box properties
content_left = Inches(0.5)
content_top = Inches(2.5)
content_width = Inches(12.333)
content_height = Inches(4.5)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Helper function to format bullet points
def add_bullet(text_frame, text, level=0):
    if level == 0:
        paragraph = text_frame.paragraphs[0]
    else:
        paragraph = text_frame.add_paragraph()
    
    paragraph.text = text
    paragraph.font.size = Pt(20)
    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Gray
    paragraph.space_after = Pt(12)
    return paragraph

# Add Bullet Points
# 1. Creation
bullet1 = add_bullet(content_tf, "Creation: Conceptualized in 2008 by Satoshi Nakamoto through a whitepaper titled 'Bitcoin: A Peer-to-Peer Electronic Cash System', introducing a decentralized ledger technology independent of central authorities.")

# 2. Implementation in Bitcoin
bullet2 = add_bullet(content_tf, "Implementation in Bitcoin: Launched in January 2009, Bitcoin marked the first practical application of blockchain technology, using it to secure transactions and manage the issuance of digital currency without intermediaries.", level=1)

# 3. Proposal of Private Blockchains
bullet3 = add_bullet(content_tf, "Proposal of Private Blockchains for Business Use: As the technology matured, proposals emerged for private or permissioned blockchains (e.g., Hyperledger, Ethereum Enterprise). These aimed to leverage blockchain's security and transparency for internal business operations, supply chain management, and smart contracts without requiring public cryptocurrency mining.", level=1)

# Save the presentation
prs.save('output.pptx')

print("Presentation created successfully as 'output.pptx'")