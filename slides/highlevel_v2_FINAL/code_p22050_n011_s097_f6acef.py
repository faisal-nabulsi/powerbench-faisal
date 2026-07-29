from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Title "VIEWING ACTIVITY" at the top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_p = title_tf.paragraphs[0]
title_p.text = "VIEWING ACTIVITY"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(0, 0, 0)
title_p.alignment = PP_ALIGN.CENTER

# 2. Video Camera Icon
# Using a Unicode character for the icon, placed centrally below the title
icon_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(2.3), Inches(1.0))
icon_tf = icon_box.text_frame
icon_p = icon_tf.paragraphs[0]
icon_p.text = "📹"
icon_p.font.size = Pt(60)
icon_p.alignment = PP_ALIGN.CENTER

# 3. Image 'image.png'
# Placed on the left side
# Left=1.0, Top=2.4, Width=5.5, Height=4.0
slide.shapes.add_picture('image.png', Inches(1.0), Inches(2.4), Inches(5.5), Inches(4.0))

# 4. Subtitle "Friends | Joey Doesn't Share Food!"
# Placed to the right of the image.
# Image ends at x = 1.0 + 5.5 = 6.5. Subtitle starts at x = 7.0.
subtitle_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.4), Inches(5.8), Inches(4.0))
subtitle_tf = subtitle_box.text_frame
subtitle_tf.word_wrap = True
subtitle_p = subtitle_tf.paragraphs[0]
subtitle_p.text = "Friends | Joey Doesn't Share Food!"
subtitle_p.font.size = Pt(32)
subtitle_p.font.bold = True
subtitle_p.font.color.rgb = RGBColor(139, 0, 0) # Dark Red
subtitle_p.alignment = PP_ALIGN.LEFT

# 5. Instruction "Watch the first 2 minutes of the video"
# Placed below the image and subtitle content.
# Content ends at y = 2.4 + 4.0 = 6.4. Instruction starts at y = 6.6.
instruction_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.333), Inches(0.6))
instruction_tf = instruction_box.text_frame
instruction_p = instruction_tf.paragraphs[0]
instruction_p.text = "Watch the first 2 minutes of the video"
instruction_p.font.size = Pt(24)
instruction_p.font.italic = True
instruction_p.font.color.rgb = RGBColor(50, 50, 50)
instruction_p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')