from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (layout index 6) to allow custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a contrasting background color (Dark Blue) for high contrast with text
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

# Add Title "ELO's" at the top
# Coordinates: Left 1, Top 0.5, Width 11.333, Height 1.5
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "ELO's"
title_run.font.size = Pt(48)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White

# Add Bullet Points section
# Coordinates: Left 1, Top 2.5, Width 11.333, Height 4.5
content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(4.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

points = [
    "Market and Opportunity",
    "Market and Equality"
]

for i, point in enumerate(points):
    if i == 0:
        para = content_frame.paragraphs[0]
    else:
        para = content_frame.add_paragraph()
    
    run = para.add_run()
    # Add a bullet character manually for guaranteed visual rendering
    run.text = "• " + point
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White
    run.font.name = "Arial"

# Save the presentation
prs.save('output.pptx')