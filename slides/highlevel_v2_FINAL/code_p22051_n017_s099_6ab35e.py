from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Create a presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a subtle background color (Light Grey/Blue)
background = slide.background
fill = background.fill
fill.type_name = 'SOLID'
fill.solid()
fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

# Add Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Homework"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

# Add Question Text
question_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(4))
question_tf = question_box.text_frame
question_tf.word_wrap = True

# Paragraph for the question
q_para = question_tf.paragraphs[0]
q_run = q_para.add_run()
q_run.text = "Q1. Explain how a chain of market is formed? What purpose does it serve. (3+2=5 MARKS)"
q_run.font.size = Pt(24)
q_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')