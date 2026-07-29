from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to the presentation
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the Main Title
# Position: Top center/left, spanning most width
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "5 Stages of Development of Media"
title_run.font.size = Pt(36)
title_run.font.bold = True

# 2. Add the Content Text Box
# Position: Left side, below title
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Subtitle focusing on Digital Media
p_header = content_frame.paragraphs[0]
run_header = p_header.add_run()
run_header.text = "5. DIGITAL MEDIA"
run_header.font.size = Pt(28)
run_header.font.bold = True

# Bullet Point 1: Features
p1 = content_frame.add_paragraph()
p1.space_before = Pt(14)
run1 = p1.add_run()
run1.text = "\u2022 Features: Ubiquitous access to real-time information, multimedia integration, and high interactivity through internet platforms."
run1.font.size = Pt(18)

# Bullet Point 2: Importance
p2 = content_frame.add_paragraph()
p2.space_before = Pt(14)
run2 = p2.add_run()
run2.text = "\u2022 Importance: Democratizes content creation, breaking down barriers to entry and enabling instant global connectivity and social engagement."
run2.font.size = Pt(18)

# 3. Add the Image
# Position: Right side
# Using the placeholder 'image.png' as requested
slide.shapes.add_picture('image.png', Inches(9), Inches(1.8), Inches(4), Inches(4))

# Save the presentation to 'output.pptx'
prs.save('output.pptx')