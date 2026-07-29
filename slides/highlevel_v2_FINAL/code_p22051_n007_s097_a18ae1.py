from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize the presentation
prs = Presentation()

# Set the slide dimensions to 16:9 Widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
# Note: Index 6 is typically the blank layout in the default theme
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Position: Left=1, Top=0.5, Width=11.333, Height=1.5
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_paragraph = title_frame.paragraphs[0]
title_paragraph.text = "Blockchain in Bitcoin"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True

# --- Add Introduction Text ---
# Position: Left=1, Top=2.2, Width=6.5, Height=4.5
intro_text = (
    "Bitcoin is a decentralized digital currency, based on blockchain technology - a distributed ledger "
    "enforced by a diverse network of computers around the world. It was introduced in 2009 by an unknown "
    "person or group of people using the name Satoshi Nakamoto. Transactions are verified by network nodes "
    "through cryptography and recorded in a public distributed ledger called a blockchain."
)

text_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(6.5), Inches(4.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

text_paragraph = text_frame.paragraphs[0]
text_paragraph.text = intro_text
text_paragraph.font.size = Pt(18)

# --- Add Image ---
# Position: Left=8, Top=2, Width=5, Height=4.5
# Using the placeholder image provided in the working directory
image_path = 'image.png'
slide.shapes.add_picture(image_path, Inches(8), Inches(2), Inches(5), Inches(4.5))

# Save the presentation
prs.save('output.pptx')