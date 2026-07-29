from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize the presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout in standard templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- PART 1 CONTENT ---

# Positioning constants for the left column
left_margin = Inches(0.8)
column_width = Inches(5.5)

# PART 1 Title
part1_top = Inches(0.8)
part1_height = Inches(0.6)

title_box1 = slide.shapes.add_textbox(left_margin, part1_top, column_width, part1_height)
tf1 = title_box1.text_frame
p1 = tf1.paragraphs[0]
p1.text = "PART 1"
p1.font.size = Pt(32)
p1.font.bold = True

# PART 1 Discussion Prompt & Short Answers
part1_content_top = Inches(1.5)
part1_content_height = Inches(2.8)

content_box1 = slide.shapes.add_textbox(left_margin, part1_content_top, column_width, part1_content_height)
tf_content1 = content_box1.text_frame
tf_content1.word_wrap = True
p_content1 = tf_content1.paragraphs[0]
p_content1.text = (
    "Discussion Prompt:\n\n"
    "Review the details regarding Joey's date.\n\n"
    "Short Answers:\n"
    "1. Describe the atmosphere of the date.\n"
    "2. How did Joey handle the conversation?\n"
    "3. What was the outcome of the evening?"
)
p_content1.font.size = Pt(14)

# --- PART 2 CONTENT ---

# PART 2 Title
part2_top = Inches(4.8)
part2_height = Inches(0.6)

title_box2 = slide.shapes.add_textbox(left_margin, part2_top, column_width, part2_height)
tf2 = title_box2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "PART 2"
p2.font.size = Pt(32)
p2.font.bold = True

# PART 2 Comparison Instruction
part2_content_top = Inches(5.5)
part2_content_height = Inches(1.5)

content_box2 = slide.shapes.add_textbox(left_margin, part2_content_top, column_width, part2_content_height)
tf_content2 = content_box2.text_frame
tf_content2.word_wrap = True
p_content2 = tf_content2.paragraphs[0]
p_content2.text = (
    "Comparison Task:\n\n"
    "Refer back to the previous section of this presentation "
    "to compare and contrast the events of Joey's date with "
    "the scenarios discussed earlier."
)
p_content2.font.size = Pt(14)

# --- IMAGE ---

# Add the placeholder image to the right side of the slide
img_left = Inches(7.0)
img_top = Inches(1.0)
img_width = Inches(5.8)
img_height = Inches(5.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')