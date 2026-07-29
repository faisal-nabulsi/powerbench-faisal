from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Main Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "5 Stages of Development of Media"
title_run.font.size = Pt(36)
title_run.font.bold = True

# Add Subtitle focusing on "2. SCRIPTS"
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.8))
subtitle_tf = subtitle_box.text_frame
subtitle_tf.word_wrap = True
subtitle_para = subtitle_tf.paragraphs[0]
subtitle_run = subtitle_para.add_run()
subtitle_run.text = "2. SCRIPTS"
subtitle_run.font.size = Pt(28)
subtitle_run.font.bold = True

# Add Bullet Points about benefits
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6.0), Inches(4.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Bullet 1
p1 = content_tf.paragraphs[0]
p1.text = "Scripts provide a structured framework for communication, ensuring clarity and consistency."
p1.level = 0
p1.font.size = Pt(18)

# Bullet 2
p2 = content_tf.add_paragraph()
p2.text = "They allow for precise timing and pacing, enhancing the overall delivery and impact of the message."
p2.level = 0
p2.font.size = Pt(18)

# Add Image on the right side
# Using 'image.png' as requested
slide.shapes.add_picture('image.png', Inches(7.0), Inches(2.5), Inches(5.833), Inches(4.5))

# Save the presentation
prs.save('output.pptx')