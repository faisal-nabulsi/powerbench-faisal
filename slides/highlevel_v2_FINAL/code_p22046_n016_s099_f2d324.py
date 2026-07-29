from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "Elon's Future Plans"
title_run.font.size = Pt(40)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0, 0, 0)

# Add the Bullet-point list
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(7.5), Inches(5.0))
content_tf = content_box.text_frame
content_tf.word_wrap = True

projects = [
    "Mars Colony Construction Initiatives",
    "Neuralink Clinical Trials Expansion",
    "Optimus Humanoid Robot Mass Production",
    "Starship Interplanetary Transport System",
    "The Boring Company Urban Loop Networks"
]

for i, project in enumerate(projects):
    if i == 0:
        para = content_tf.paragraphs[0]
    else:
        para = content_tf.add_paragraph()
    
    run = para.add_run()
    run.text = project
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(51, 51, 51)
    para.space_after = Pt(12)

# Add the Image on the right side
# Position: x=8.5, y=2.2, Width=4.5, Height=5.0
image_shape = slide.shapes.add_picture('image.png', Inches(8.5), Inches(2.2), Inches(4.5), Inches(5.0))

# Save the presentation
prs.save('output.pptx')