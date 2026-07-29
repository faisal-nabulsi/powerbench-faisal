from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize the presentation
prs = Presentation()

# Set slide dimensions to 16:9 (13.333 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide to allow full customization
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the background image ('image.png') covering the entire slide
# We assume 'image.png' is a visually engaging design-related image
slide.shapes.add_picture('image.png', 0, 0, Inches(13.333), Inches(7.5))

# Create a semi-transparent dark container for the text to ensure readability
# positioned on the left side of the slide
left, top, width, height = Inches(1.0), Inches(1.0), Inches(8.0), Inches(5.5)
container_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

# Style the container: Dark grey background
container_shape.fill.solid()
container_shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
container_shape.line.fill.background() # Remove border

# Add Title
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(7.0), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p_title = title_frame.paragraphs[0]
p_title.text = "Components of Visual Communication"
p_title.font.size = Pt(42)
p_title.font.color.rgb = RGBColor(255, 255, 255) # White text
p_title.font.bold = True
p_title.font.name = 'Calibri'

# Add Content (Bullet Points)
content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7.0), Inches(3.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Define a helper function to add formatted bullets
def add_bullet_point(frame, heading, description):
    # Heading
    p_head = frame.add_paragraph()
    p_head.text = heading
    p_head.level = 0
    p_head.font.size = Pt(28)
    p_head.font.bold = True
    p_head.font.color.rgb = RGBColor(255, 200, 100) # Gold/Orange accent
    
    # Description
    p_desc = frame.add_paragraph()
    p_desc.text = description
    p_desc.level = 1
    p_desc.font.size = Pt(20)
    p_desc.font.color.rgb = RGBColor(240, 240, 240) # Light grey
    p_desc.space_before = Pt(5)
    p_desc.space_after = Pt(12)

# Add specific bullet points requested
add_bullet_point(content_frame, "Craftsmanship",
                 "The art of creating visual elements with precision, skill, and high attention to detail. It reflects the mastery of tools and the quality of execution in the final design.")

add_bullet_point(content_frame, "Visual Components",
                 "The fundamental building blocks of design, including color theory, typography, lines, shapes, and textures, which are combined to construct meaningful visual messages.")

add_bullet_point(content_frame, "Composition & Layout",
                 "The strategic arrangement of visual elements within a space to guide the viewer's eye, create balance, and ensure the communication of ideas is clear and impactful.")

# Save the presentation
prs.save('output.pptx')