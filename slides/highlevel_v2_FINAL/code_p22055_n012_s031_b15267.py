from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# Create a new Presentation
prs = Presentation()

# Set slide dimensions to 16:9 Widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Add the background image (Forklift and shipping containers)
# Using the available placeholder image file
try:
    background_image = slide.shapes.add_picture(
        'image.png',
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.333),
        height=Inches(7.5)
    )
except FileNotFoundError:
    # Fallback in case the image file is missing
    pass

# 2. Add a semi-transparent overlay shape
# Positioned centrally to create a banner for the title
overlay_left = Inches(2.5)
overlay_top = Inches(2.25)
overlay_width = Inches(8.5)
overlay_height = Inches(3.0)

overlay_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    overlay_left, overlay_top, overlay_width, overlay_height
)

# Configure overlay appearance: Black fill, no border
overlay_shape.fill.solid()
overlay_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
overlay_shape.line.fill.background()

# Apply transparency (Alpha) to the overlay fill
# 50000 represents 50% transparency (scale 00000 to 100000)
spPr = overlay_shape._element.find(qn('p:spPr'))
solidFill = spPr.find(qn('a:solidFill'))
if solidFill is not None:
    alpha = etree.SubElement(solidFill, qn('a:alpha'))
    alpha.set('val', '50000')

# 3. Add the Title Textbox
# The textbox is placed on top of the overlay
text_box_height = Inches(1.0)
# Calculate vertical centering within the overlay
text_box_top = overlay_top + (overlay_height - text_box_height) / 2
text_box_left = overlay_left + Inches(0.5)
text_box_width = overlay_width - Inches(1.0)

title_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.auto_size = None

# Clear default paragraph
for p in title_frame.paragraphs:
    p.clear()

# Add title paragraph
title_paragraph = title_frame.paragraphs[0]
title_run = title_paragraph.add_run()
title_run.text = "Transportation and Storage"

# Format title text
title_run.font.size = Pt(42)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(255, 255, 255) # White color for contrast
title_run.font.name = "Calibri"

# Center text horizontally
title_paragraph.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')