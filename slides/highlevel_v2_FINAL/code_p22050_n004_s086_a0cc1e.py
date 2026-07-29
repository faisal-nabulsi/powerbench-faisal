from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize the presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Index 6 is typically the blank layout in default templates)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Create Teal Header
# Rectangle shape for the header background
header_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(1.2)
)
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = RGBColor(0, 128, 128) # Teal color
header_shape.line.fill.background() # Remove border

# Add "Values" text to the header
tf_header = header_shape.text_frame
tf_header.word_wrap = True
p_header = tf_header.paragraphs[0]
p_header.text = "Values"
p_header.font.size = Pt(36)
p_header.font.bold = True
p_header.font.color.rgb = RGBColor(255, 255, 255) # White text
p_header.alignment = PP_ALIGN.CENTER
p_header.space_before = Pt(6) # Adjust vertical alignment

# 2. List Six Values with Descriptions
# Position the text box on the left side
left_margin = Inches(0.5)
content_top = Inches(1.5)
content_width = Inches(8.5)
content_height = Inches(5.5)

values_box = slide.shapes.add_textbox(left_margin, content_top, content_width, content_height)
tf_body = values_box.text_frame
tf_body.word_wrap = True

# Data: List of six values with descriptions
values_data = [
    ("Integrity", "Acting with honesty and transparency."),
    ("Innovation", "Embracing creative solutions."),
    ("Excellence", "Committing to high quality standards."),
    ("Collaboration", "Working together for shared success."),
    ("Accountability", "Taking ownership of results."),
    ("Customer Focus", "Prioritizing client satisfaction.")
]

# Populate the text box
for i, (title, description) in enumerate(values_data):
    if i == 0:
        p = tf_body.paragraphs[0]
    else:
        p = tf_body.add_paragraph()
    
    # Add bold title
    run_title = p.add_run()
    run_title.text = f"{i+1}. {title}"
    run_title.font.bold = True
    run_title.font.size = Pt(20)
    run_title.font.color.rgb = RGBColor(0, 100, 100) # Dark Teal for contrast
    
    # Add description
    run_desc = p.add_run()
    run_desc.text = f" {description}"
    run_desc.font.size = Pt(16)
    
    p.space_after = Pt(12)

# 3. Add Relevant Images on the Right Side
# Define area for images on the right
img_x = Inches(9.5)
img_y = Inches(1.8)
img_width = Inches(3.5)
img_height = Inches(1.5)
img_gap = Inches(0.3)

# Add 3 placeholder images stacked vertically
for _ in range(3):
    slide.shapes.add_picture('image.png', img_x, img_y, img_width, img_height)
    img_y += img_height + img_gap

# Save the presentation
prs.save('output.pptx')