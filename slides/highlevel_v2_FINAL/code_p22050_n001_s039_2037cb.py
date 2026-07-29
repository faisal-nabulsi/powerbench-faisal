from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Title ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_text_range = title_frame.add_paragraph()
title_text_range.text = "Advantages of Blockchain"
title_text_range.font.size = Pt(40)
title_text_range.font.bold = True
title_text_range.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_text_range.alignment = PP_ALIGN.LEFT

# --- 2. Bullet Points ---
# Define bullet points content
bullets = [
    "Decentralization: No single point of failure or control.",
    "Security: Cryptography protects data integrity and prevents tampering.",
    "Transparency: All participants can view the ledger in real-time.",
    "Traceability: Full history of transactions is recorded and auditable.",
    "Efficiency: Faster settlement times and reduced intermediary costs."
]

# Add text box for bullets
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(5.0))
text_frame = text_box.text_frame
text_frame.word_wrap = True

for i, bullet in enumerate(bullets):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    
    p.text = bullet
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(15)
    p.level = 0 # Bullet level

    # Add a small circle bullet shape manually or rely on default. 
    # python-pptx default bullets can be tricky, so we just use text formatting.
    # To make it look like a bullet list, we can prepend a bullet char or use formatting.
    # Let's use a simple character bullet for visual clarity.
    p.text = "• " + bullet

# --- 3. Diagram (Illustrating a Blockchain Chain) ---
# We will draw a visual diagram using shapes since no specific image was provided 
# and drawing a diagram programmatically is more precise than a generic placeholder.
# The diagram will show blocks connected by hash links.

# Diagram Area Coordinates
margin_left = 7.0
top_y = 2.5
height = 4.5
width = 5.8

# Background for the diagram area to make it stand out
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(margin_left), 
    Inches(top_y), 
    Inches(width), 
    Inches(height)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(240, 245, 250) # Light grey-blue
bg_shape.line.color.theme_color = 1 # Black/Dark
bg_shape.line.width = Pt(1)

# Function to add a block
def add_block(left, top, width, height, text, color, slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0, 0, 0)
    
    # Add text to shape
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape

# Block 1
b1_color = RGBColor(0, 112, 192) # Blue
block1 = add_block(margin_left + 0.5, top_y + 1.5, 1.5, 1.5, "Block 1", b1_color, slide)

# Block 2
b2_color = RGBColor(0, 176, 240) # Light Blue
block2 = add_block(margin_left + 2.5, top_y + 1.5, 1.5, 1.5, "Block 2", b2_color, slide)

# Block 3
b3_color = RGBColor(31, 119, 180) # Medium Blue
block3 = add_block(margin_left + 4.5, top_y + 1.5, 1.5, 1.5, "Block 3", b3_color, slide)

# Connecting Arrows/Lines
# Arrow 1 (Block 1 to Block 2)
arrow1 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, 
    Inches(margin_left + 2.0), 
    Inches(top_y + 1.85), 
    Inches(0.4), 
    Inches(0.3)
)
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(100, 100, 100)
arrow1.line.fill.background()

# Arrow 2 (Block 2 to Block 3)
arrow2 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, 
    Inches(margin_left + 4.0), 
    Inches(top_y + 1.85), 
    Inches(0.4), 
    Inches(0.3)
)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(100, 100, 100)
arrow2.line.fill.background()

# Labels for arrows (Hash)
hash1_text = slide.shapes.add_textbox(Inches(margin_left + 2.0), Inches(top_y + 1.2), Inches(0.4), Inches(0.4))
t = hash1_text.text_frame
t.paragraphs[0].text = "Hash"
t.paragraphs[0].font.size = Pt(10)
t.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
t.paragraphs[0].alignment = PP_ALIGN.CENTER

hash2_text = slide.shapes.add_textbox(Inches(margin_left + 4.0), Inches(top_y + 1.2), Inches(0.4), Inches(0.4))
t = hash2_text.text_frame
t.paragraphs[0].text = "Hash"
t.paragraphs[0].font.size = Pt(10)
t.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
t.paragraphs[0].alignment = PP_ALIGN.CENTER

# Title for the diagram
diagram_title = slide.shapes.add_textbox(Inches(margin_left), Inches(top_y + 0.2), Inches(width), Inches(0.5))
tf = diagram_title.text_frame
p = tf.paragraphs[0]
p.text = "Blockchain Structure"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.LEFT

# Save the presentation
prs.save('output.pptx')