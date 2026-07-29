from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # Set the slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add a slide with a title and content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title
    title = slide.shapes.title
    title.text = "Importance of Skills Beyond Academics"
    
    # Get the content placeholder
    # Typically, index 1 is the content placeholder in the Title and Content layout
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    
    # Define the skills and their descriptions
    skills = [
        "Effective Communication: The ability to convey ideas clearly, both verbally and in writing, and to listen actively.",
        "Problem Solving and Critical Thinking: The capacity to analyze complex situations, evaluate information, and develop logical solutions.",
        "Emotional Intelligence: The ability to recognize, understand, and manage one's own emotions, as well as empathize with others.",
        "Adaptability and Resilience: The flexibility to adjust to new conditions and the strength to recover from challenges or setbacks."
    ]
    
    # Add each skill as a bullet point
    # The placeholder usually comes with one empty paragraph, so we use that first
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
        # Clear any existing text
        p.text = ""
        
    p.text = skills[0]
    p.level = 0
    
    for skill in skills[1:]:
        p = tf.add_paragraph()
        p.text = skill
        p.level = 0
        
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()