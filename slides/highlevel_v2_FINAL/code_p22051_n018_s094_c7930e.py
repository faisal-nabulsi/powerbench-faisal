from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Set the slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add the title "PART 4"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "PART 4"
title_run.font.size = Pt(36)
title_run.font.bold = True

# Add the instruction text
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1))
instruction_tf = instruction_box.text_frame
instruction_para = instruction_tf.paragraphs[0]
instruction_run = instruction_para.add_run()
instruction_run.text = "List the events in Joey’s second date with Sarah in order from 2 - 6"
instruction_run.font.size = Pt(18)

# Add the checklist
checklist_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(3.5))
checklist_tf = checklist_box.text_frame
checklist_tf.word_wrap = True

# Define the checklist items
# The first item is marked as completed [x] as an example
checklist_items = [
    "1. [x] Event 1 (Example)",
    "2. [ ]",
    "3. [ ]",
    "4. [ ]",
    "5. [ ]",
    "6. [ ]"
]

for i, item in enumerate(checklist_items):
    if i == 0:
        para = checklist_tf.paragraphs[0]
    else:
        para = checklist_tf.add_paragraph()
    
    run = para.add_run()
    run.text = item
    run.font.size = Pt(16)
    
    # Make the example item bold
    if i == 0:
        run.font.bold = True

# Add the image of Joey in the bottom right corner
# Image size: 2x2 inches
# Margins: 0.5 inches from right and bottom
img_width = Inches(2)
img_height = Inches(2)
img_left = prs.slide_width - img_width - Inches(0.5)
img_top = prs.slide_height - img_height - Inches(0.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation
prs.save('output.pptx')