from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set the background to a dark color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x2c, 0x2c, 0x2c)

# Add the title shape
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = False

title_para = title_tf.paragraphs[0]
title_para.text = "Now: I, CAN,"
title_para.font.size = Pt(48)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
title_para.alignment = PP_ALIGN.LEFT

# Add the content text box for bullet points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(8.5), Inches(4))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# First bullet point
p1 = content_tf.add_paragraph()
p1.text = "Define the term Market"
p1.font.size = Pt(24)
p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p1.space_after = Pt(18)

# Second bullet point
p2 = content_tf.add_paragraph()
p2.text = "Explain how products reach to market"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2.space_after = Pt(18)

# Add the textured image on the right side
slide.shapes.add_picture('image.png', Inches(9), Inches(0.5), Inches(4), Inches(6.5))

# Save the presentation
prs.save('output.pptx')