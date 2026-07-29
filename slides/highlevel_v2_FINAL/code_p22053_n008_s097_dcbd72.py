from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (Layout index 6)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Add Title ---
# Position: Top of the slide, spanning most of the width
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
title_frame = title_shape.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "The Technology used in Blockchain"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) # Professional Blue

# --- Add Bullet Points ---
# Position: Left side of the slide
text_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(4.5))
text_frame = text_shape.text_frame
text_frame.word_wrap = True

bullets = [
    "Cryptographic Keys",
    "Network Protocol",
    "Distributed Ledger Technology",
    "Hashing"
]

for i, bullet in enumerate(bullets):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    
    run = p.add_run()
    run.text = "• " + bullet
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Grey
    p.space_after = Pt(14)

# --- Add Image ---
# Position: Right side of the slide, aligned with text
img_left = Inches(7)
img_top = Inches(2)
img_width = Inches(5.833)
img_height = Inches(4.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# Save the presentation
prs.save('output.pptx')