from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Title at the top
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
title_tf = title_shape.text_frame
title_tf.word_wrap = True
p_title = title_tf.paragraphs[0]
p_title.alignment = PP_ALIGN.CENTER
title_run = p_title.add_run()
title_run.text = "VIEWING ACTIVITY"
title_run.font.size = Pt(40)
title_run.font.bold = True

# 2. Image (represents the video camera icon / visual content)
# Placed on the left side of the slide
im_left = Inches(1)
im_top = Inches(2)
im_width = Inches(5)
im_height = Inches(4.5)
slide.shapes.add_picture('image.png', im_left, im_top, im_width, im_height)

# 3. Text content positioned to the right of the image
text_left = Inches(6.5)
text_top = Inches(2)
text_width = Inches(5.833)
text_height = Inches(4.5)

text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
text_tf = text_box.text_frame
text_tf.word_wrap = True

# Subtitle: "Friends | Joey Doesn't Share Food!"
p1 = text_tf.paragraphs[0]
run1 = p1.add_run()
run1.text = "Friends | Joey Doesn't Share Food!"
run1.font.size = Pt(30)
run1.font.bold = True
p1.space_after = Pt(15)

# Instruction: "Watch the first 2 minutes of the video" below the subtitle
p2 = text_tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Watch the first 2 minutes of the video"
run2.font.size = Pt(24)
p2.space_before = Pt(10)

# Save the file
prs.save('output.pptx')