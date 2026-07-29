from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Set the slide width and height for 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add Title: "PART 3"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
p = title_tf.paragraphs[0]
p.text = "PART 3"
p.font.size = Pt(32)
p.font.bold = True

# Add Instruction Text
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1))
instruction_tf = instruction_box.text_frame
instruction_tf.word_wrap = True
p = instruction_tf.paragraphs[0]
p.text = "Instructions: Fill in the correct past simple or past continuous forms of the verbs in parentheses in the story below."
p.font.size = Pt(18)

# Add Story Text with blanks
story_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(4.0))
story_tf = story_box.text_frame
story_tf.word_wrap = True

story_content = """
Yesterday, while I __________ (walk) to the local market, it suddenly __________ (start) to rain heavily. I __________ (run) into a nearby cafe to take shelter. While I __________ (wait) in line, I __________ (see) an old friend of mine, Sarah. She __________ (order) a cup of coffee. We __________ (sit) down and __________ (chat) about our weekend plans. Just then, the rain __________ (stop), and the sun __________ (come) out. We __________ (decide) to continue our walk together.
"""

p = story_tf.paragraphs[0]
p.text = story_content.strip()
p.font.size = Pt(16)

# Save the presentation
prs.save('output.pptx')