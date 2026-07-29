from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "Lessons that we can learn from Elon Musk's life"
title_run = title_para.runs[0]
title_run.font.size = Pt(34)
title_run.font.bold = True

# 2. Add Bullet Points
bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(8.0), Inches(3.5))
bullets_tf = bullets_box.text_frame
bullets_tf.word_wrap = True

lessons = [
    "First Principles Thinking: Boil things down to fundamental truths.",
    "Embrace Failure: Treat failure as necessary feedback.",
    "Long-Term Vision: Focus on the future over immediate gains.",
    "Relentless Work Ethic: Dedicate significant time and energy."
]

for i, lesson in enumerate(lessons):
    if i == 0:
        para = bullets_tf.paragraphs[0]
    else:
        para = bullets_tf.add_paragraph()
    para.text = "• " + lesson
    para.space_after = Pt(12)
    for run in para.runs:
        run.font.size = Pt(18)

# 3. Add Quote
quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(8.0), Inches(1.2))
quote_tf = quote_box.text_frame
quote_tf.word_wrap = True
quote_para = quote_tf.paragraphs[0]
quote_para.text = "\"Failure is an option here. If things are not failing, you are not innovating enough.\" - Elon Musk"
quote_run = quote_para.runs[0]
quote_run.font.size = Pt(20)
quote_run.font.italic = True

# 4. Add Image
# Positioning the image on the right side of the slide
slide.shapes.add_picture('image.png', Inches(9.0), Inches(1.5), Inches(4.0), Inches(4.2))

# Save the presentation
prs.save('output.pptx')