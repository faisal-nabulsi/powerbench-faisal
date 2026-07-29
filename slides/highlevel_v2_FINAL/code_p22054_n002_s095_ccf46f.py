from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set Background Color (Dark Slate Blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x48, 0x3D, 0x8B)

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Innovations (Future Products or Services)"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Description
txBoxDesc = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(2.0))
tfDesc = txBoxDesc.text_frame
tfDesc.word_wrap = True
pDesc = tfDesc.paragraphs[0]
pDesc.text = "In 2021, Airbnb launched 'Work Different' to support remote travelers and expanded 'Experiences' to deepen local connections. The company also utilized AI to enhance search relevance and introduced new safety features to foster trust within the community."
pDesc.font.size = Pt(20)
pDesc.font.color.rgb = RGBColor(255, 255, 255)
pDesc.alignment = PP_ALIGN.CENTER

# Quote
txBoxQuote = slide.shapes.add_textbox(Inches(2.5), Inches(5), Inches(8.333), Inches(1.5))
tfQuote = txBoxQuote.text_frame
tfQuote.word_wrap = True
pQuote = tfQuote.paragraphs[0]
pQuote.text = "\"We're not just a company that helps you book a place to stay. We're a company that builds products and services to help people belong anywhere.\""
pQuote.font.size = Pt(24)
pQuote.font.italic = True
pQuote.font.color.rgb = RGBColor(255, 165, 0) # Orange
pQuote.alignment = PP_ALIGN.CENTER
pAuthor = tfQuote.add_paragraph()
pAuthor.text = "- Brian Chesky"
pAuthor.font.size = Pt(18)
pAuthor.font.color.rgb = RGBColor(255, 255, 255)
pAuthor.alignment = PP_ALIGN.CENTER

prs.save('output.pptx')