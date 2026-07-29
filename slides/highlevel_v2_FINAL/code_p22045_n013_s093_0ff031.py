from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title "VIEWING ACTIVITY" at the top
# Positioned at the top center
title_left = Inches(0)
title_top = Inches(0.5)
title_width = Inches(13.333)
title_height = Inches(1.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_box.text_frame
title_tf.word_wrap = True

title_paragraph = title_tf.paragraphs[0]
title_paragraph.text = "VIEWING ACTIVITY"
title_paragraph.alignment = PP_ALIGN.CENTER

title_run = title_paragraph.runs[0]
title_run.bold = True
title_run.font.size = Pt(36)
title_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 2. Add Video Camera Icon (using the provided placeholder image)
# Positioned centrally below the title
icon_width = Inches(3.0)
icon_height = Inches(3.0)
# Center X calculation: (13.333 - 3.0) / 2 = 5.1665
icon_left = Inches(5.1665)
icon_top = Inches(2.2)

slide.shapes.add_picture('image.png', icon_left, icon_top, icon_width, icon_height)

# 3. Add Subtitle "Friends | Joey Doesn't Share Food!" below the icon
# Positioned below the icon
sub_left = Inches(0)
sub_top = Inches(5.5)
sub_width = Inches(13.333)
sub_height = Inches(1.5)

sub_box = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True

sub_paragraph = sub_tf.paragraphs[0]
sub_paragraph.text = "Friends | Joey Doesn't Share Food!"
sub_paragraph.alignment = PP_ALIGN.CENTER

sub_run = sub_paragraph.runs[0]
sub_run.font.size = Pt(24)
sub_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Save the presentation
prs.save('output.pptx')