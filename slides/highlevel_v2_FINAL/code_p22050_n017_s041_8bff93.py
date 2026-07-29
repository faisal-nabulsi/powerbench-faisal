from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set the background color to dark
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x2C, 0x2C, 0x2C)

# Add circular images (simulated with rectangular placeholders as python-pptx doesn't support easy circular cropping)
# We will place 3 images on the left side
image_path = 'image.png'

# Image 1
left1 = Inches(0.5)
top1 = Inches(1.0)
width1 = Inches(2.5)
height1 = Inches(2.5)
slide.shapes.add_picture(image_path, left1, top1, width1, height1)

# Image 2
left2 = Inches(0.5)
top2 = Inches(3.8)
width2 = Inches(2.5)
height2 = Inches(2.5)
slide.shapes.add_picture(image_path, left2, top2, width2, height2)

# Image 3
left3 = Inches(0.5)
top3 = Inches(6.6)
width3 = Inches(2.5)
height3 = Inches(2.5)
slide.shapes.add_picture(image_path, left3, top3, width3, height3)

# Add the title "Types of Markets" at the bottom right
title_text = "Types of Markets"
title_left = Inches(8.0)
title_top = Inches(6.0)
title_width = Inches(5.0)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_box.text_frame
title_para = title_tf.paragraphs[0]
title_run = title_para.add_run()
title_run.text = title_text
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text
title_para.alignment = PP_ALIGN.RIGHT

# Save the presentation
prs.save('output.pptx')