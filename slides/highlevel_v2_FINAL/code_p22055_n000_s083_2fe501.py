from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

# Create presentation with 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Gradient Background ---
bg = slide.background
bg_xml = bg._element

# Remove any existing fill elements
for child in list(bg_xml):
    if child.tag.endswith('solidFill') or child.tag.endswith('gradFill') or child.tag.endswith('pictFill'):
        bg_xml.remove(child)

# Define DrawingML namespace
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Create gradient fill element
grad_fill = etree.SubElement(bg_xml, f'{{{ns_a}}}gradFill')
gs = etree.SubElement(grad_fill, f'{{{ns_a}}}gsLst')

# Gradient stop 1 (dark blue)
stop1 = etree.SubElement(gs, f'{{{ns_a}}}gs')
stop1.set('pos', '0')
srgb1 = etree.SubElement(stop1, f'{{{ns_a}}}srgbClr')
srgb1.set('val', '1E3A5F')

# Gradient stop 2 (lighter blue)
stop2 = etree.SubElement(gs, f'{{{ns_a}}}gs')
stop2.set('pos', '100000')
srgb2 = etree.SubElement(stop2, f'{{{ns_a}}}srgbClr')
srgb2.set('val', '4A90D9')

# Linear gradient direction
lin = etree.SubElement(grad_fill, f'{{{ns_a}}}lin')
lin.set('ang', '54000')

# --- Title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.add_paragraph()
title_run = title_para.add_run()
title_run.text = "Home Fun:"
title_run.font.size = Pt(44)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(255, 255, 255)

# --- Bullet Points ---
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# First bullet point
p1 = content_frame.add_paragraph()
p1.text = "Market opportunities are expanding rapidly in emerging economies, offering significant growth potential for innovative businesses."
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.level = 0

# Second bullet point
p2 = content_frame.add_paragraph()
p2.text = "Addressing inequality remains a critical challenge, requiring targeted policies and inclusive strategies to ensure equitable development."
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.level = 0

# Save the presentation
prs.save('output.pptx')