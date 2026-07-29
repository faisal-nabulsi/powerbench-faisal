from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Set the slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- 1. Title ---
# Position: Top center
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "Various Media drive Various forms of Global Integration"
title_run.font.size = Pt(36)
title_run.font.bold = True

# --- 2. Section Label ---
# Position: Below title, acting as a header
section_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
section_tf = section_box.text_frame
section_para = section_tf.paragraphs[0]
section_run = section_para.add_run()
section_run.text = "GLOBAL INTEGRATION"
section_run.font.size = Pt(24)
section_run.font.bold = True

# --- 3. Bullet Points ---
# Position: Indented below the section label
content_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(10.333), Inches(3))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Bullet Point 1: Definition
p1 = content_tf.paragraphs[0]
r1 = p1.add_run()
# Adding a manual bullet character for visual clarity
r1.text = "• Definition: The increasing interconnectedness of countries and regions through the exchange of goods, services, capital, technology, and information."
r1.font.size = Pt(20)
p1.space_after = Pt(14)

# Bullet Point 2: Processes
p2 = content_tf.add_paragraph()
r2 = p2.add_run()
r2.text = "• Processes: Driven by media and technology to reduce communication barriers, facilitating the flow of ideas, migration, and cross-border collaboration."
r2.font.size = Pt(20)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')