from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import xml.etree.ElementTree as ET

# Initialize presentation with 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Add the background image (placeholder for forklift/shipping containers)
# This covers the entire slide
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Create a semi-transparent overlay rectangle for text clarity
# Dimensions: Centered horizontally, positioned near the top
overlay_width = Inches(10)
overlay_height = Inches(2.5)
overlay_left = (prs.slide_width - overlay_width) / 2
overlay_top = Inches(1.5)

# Add the shape
overlay_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    overlay_left, 
    overlay_top, 
    overlay_width, 
    overlay_height
)

# Style the overlay: Black fill with no border
overlay_shape.fill.solid()
overlay_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
overlay_shape.line.fill.background()

# Apply transparency via XML manipulation (50% opacity)
# Note: python-pptx does not support shape opacity directly in the API
overlay_element = overlay_shape._element
solid_fill = overlay_element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')

if solid_fill is not None:
    color_node = solid_fill[0]
    alpha_elem = ET.SubElement(color_node, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
    alpha_elem.set('val', '50000') # 50% opacity (values range 0 to 100000)

# 3. Add the Title Text Box
# Positioned directly over the overlay area
txBox = slide.shapes.add_textbox(overlay_left, overlay_top, overlay_width, overlay_height)
tf = txBox.text_frame
tf.word_wrap = True

# Style the text
p = tf.paragraphs[0]
p.text = "Transportation and Storage"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255) # White text
p.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('output.pptx')