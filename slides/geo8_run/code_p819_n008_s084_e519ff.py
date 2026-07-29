from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Initialize Presentation
    prs = Presentation()
    
    # Set standard 16:9 dimensions
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Background: Colorful abstract image
    # Using 'image.png' as the placeholder for the abstract background image
    left = Inches(0)
    top = Inches(0)
    width = slide_width
    height = slide_height
    bg_image = slide.shapes.add_picture('image.png', left, top, width, height)
    
    # 2. White Layer: To ensure text visibility
    # Create a white rectangle that acts as a backdrop for text and images
    # Dimensions: 85% width, 80% height, centered with top bias
    layer_width = Inches(11.33)
    layer_height = Inches(6.0)
    layer_left = Inches(1.0)
    layer_top = Inches(0.75)
    
    white_layer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, layer_left, layer_top, layer_width, layer_height)
    white_layer.fill.solid()
    white_layer.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
    white_layer.line.fill.background() # No border
    
    # 3. Subtitle: "The Subject and Content of Art"
    # Top-left of the content area (white layer)
    sub_text = "The Subject and Content of Art"
    sub_left = layer_left + Inches(0.5)
    sub_top = layer_top + Inches(0.5)
    sub_width = Inches(10)
    sub_height = Inches(0.8)
    
    txBox_sub = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
    tf_sub = txBox_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = sub_text
    run_sub.font.size = Pt(24) # Smaller than main title
    run_sub.font.bold = True
    # Contrasting color (e.g., Deep Red) to stand out from black title
    run_sub.font.color.rgb = RGBColor(180, 0, 0) 
    
    # 4. Main Title: "D. Mythology and religion, dreams and fantasies."
    # Centered, Large, Bold
    title_text = "D. Mythology and religion, dreams and fantasies."
    title_left = layer_left + Inches(0.5)
    title_top = layer_top + Inches(1.5)
    title_width = Inches(10.33)
    title_height = Inches(1.5)
    
    txBox_title = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.add_run()
    run_title.text = title_text
    run_title.font.size = Pt(44) # Large
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0, 0, 0) # Black for high contrast
    
    # 5. Images: Two images side by side
    # Positioned below the title within the white layer
    img_top = layer_top + Inches(3.2)
    img_width = Inches(4.5)
    img_height = Inches(2.5)
    gap = Inches(0.5)
    
    # Left Image: Iconic portrayal of beauty and mythology (Botticelli)
    # Using 'image.png' as placeholder
    img1_left = layer_left + Inches(0.5)
    slide.shapes.add_picture('image.png', img1_left, img_top, img_width, img_height)
    
    # Right Image: Dramatic baroque portrayal of heroism (Baroque)
    # Using 'image.png' as placeholder
    img2_left = img1_left + img_width + gap
    slide.shapes.add_picture('image.png', img2_left, img_top, img_width, img_height)
    
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()