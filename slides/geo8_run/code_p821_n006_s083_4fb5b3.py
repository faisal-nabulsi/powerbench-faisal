from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_slide():
    # Initialize the presentation
    prs = Presentation()

    # Add a blank slide (Layout 6 is typically 'Blank')
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Left Side: Blue Background and Title ---
    
    # Create a background shape for the left side (Blue)
    # Dimensions covering the full height of a standard slide (7.5 inches)
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.0), Inches(7.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Blue
    left_bg.line.color.rgb = RGBColor(0, 51, 102)       # No border

    # Add the title text box on top of the blue background
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.0), Inches(3.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Principles of graphic design"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p_title.alignment = PP_ALIGN.LEFT

    # --- Right Side: Subtitle and List Items ---
    
    right_start_x = Inches(5.2)
    item_width = Inches(4.5)

    # 1. Subtitle Text Box
    # "The principles of graphic design related to the areas are,"
    subtitle_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_start_x, Inches(1.0), item_width, Inches(1.2))
    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light Blue
    subtitle_shape.line.color.rgb = RGBColor(173, 216, 230)
    
    tf_sub = subtitle_shape.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "The principles of graphic design related to the areas are,"
    p_sub.font.size = Pt(14)  # Smaller font
    p_sub.font.color.rgb = RGBColor(0, 0, 0)  # Black text for legibility on light background
    p_sub.alignment = PP_ALIGN.CENTER

    # 2. List Items
    items = ["Arrangement", "Proximity", "Repetition", "Contrast", "Balance"]
    
    # Varying shades of blue and green for the backgrounds
    colors = [
        RGBColor(52, 152, 219),  # Blue
        RGBColor(39, 174, 96),   # Green
        RGBColor(22, 160, 133),  # Teal
        RGBColor(44, 62, 80),    # Dark Blue
        RGBColor(46, 204, 113)   # Bright Green
    ]
    
    item_height = Inches(0.7)
    gap = Inches(0.15)
    start_y_list = Inches(2.5)

    for i, item_text in enumerate(items):
        # Calculate Y position for each item
        y_pos = start_y_list + i * (item_height + gap)
        
        # Create a shape to act as the text box with background
        item_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_start_x, y_pos, item_width, item_height)
        item_shape.fill.solid()
        item_shape.fill.fore_color.rgb = colors[i]
        item_shape.line.color.rgb = colors[i]
        
        # Add text to the shape
        tf_item = item_shape.text_frame
        tf_item.word_wrap = True
        p_item = tf_item.paragraphs[0]
        p_item.text = item_text
        p_item.font.size = Pt(20)  # Consistent, clear font size
        p_item.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p_item.alignment = PP_ALIGN.CENTER

    # Save the presentation to 'output.pptx'
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()