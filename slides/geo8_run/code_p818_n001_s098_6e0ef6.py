from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Define slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add a blank slide
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set dark blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 128)  # Dark blue

    # Add main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = "Three perspectives on global cultural flows:"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Add perspective title
    perspective_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.8))
    perspective_title_frame = perspective_title_box.text_frame
    perspective_title_frame.word_wrap = True
    perspective_title_para = perspective_title_frame.paragraphs[0]
    perspective_title_para.text = "1. CULTURAL DIFFERENTIALISM"
    perspective_title_para.font.size = Pt(24)
    perspective_title_para.font.bold = True
    perspective_title_para.font.color.rgb = RGBColor(255, 255, 255)
    perspective_title_para.alignment = PP_ALIGN.LEFT

    # Add content with bullet points
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(8.5), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # First bullet point
    bullet1 = content_frame.add_paragraph()
    bullet1.text = "emphasizes the fact that cultures are essentially different and are only superficially affected by global flows."
    bullet1.font.size = Pt(18)
    bullet1.font.color.rgb = RGBColor(255, 255, 255)
    bullet1.level = 0
    bullet1.space_after = Pt(12)

    # Second bullet point
    bullet2 = content_frame.add_paragraph()
    bullet2.text = "It also involves barriers that prevent flows that serve to make cultures more a line; cultures tend to remain stubbornly different from one another."
    bullet2.font.size = Pt(18)
    bullet2.font.color.rgb = RGBColor(255, 255, 255)
    bullet2.level = 0
    bullet2.space_after = Pt(12)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()