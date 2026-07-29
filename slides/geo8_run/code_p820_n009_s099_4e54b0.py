from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Define colors
BG_COLOR = RGBColor(30, 30, 30)       # Dark background
TITLE_COLOR = RGBColor(255, 255, 255) # White
SUBHEADING_COLOR = RGBColor(255, 215, 0) # Gold/Yellow
BODY_COLOR = RGBColor(220, 220, 220)  # Light Gray

# Create presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BG_COLOR

# Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
tf_title = title_box.text_frame
tf_title.word_wrap = True

p_title = tf_title.paragraphs[0]
p_title.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = TITLE_COLOR
p_title.alignment = PP_ALIGN.LEFT

# Content Data
content_data = [
    {
        "heading": "Financial Literacy",
        "points": [
            "Skills include budgeting, saving, investing, and making informed financial decisions."
        ]
    },
    {
        "heading": "Leadership and Teamwork",
        "points": [
            "Leadership skills involve motivating and guiding others toward common goals.",
            "Teamwork skills enable collaboration and the ability to work effectively in diverse groups."
        ]
    },
    {
        "heading": "Time Management",
        "points": [
            "Effective time management and organizational skills improve productivity and reduce stress."
        ]
    },
    {
        "heading": "Creativity and Innovation",
        "points": [
            "These skills lead to the development of new ideas, products, and solutions."
        ]
    },
    {
        "heading": "Cultural Competence and Diversity",
        "points": [
            "In an increasingly globalized world, understanding and respecting cultural differences is essential.",
            "Cultural competence promotes inclusivity and effective interactions in multicultural environments."
        ]
    }
]

# Add Content Text Box
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)
body_box = slide.shapes.add_textbox(left, top, width, height)
tf_body = body_box.text_frame
tf_body.word_wrap = True

first_paragraph = True

for item in content_data:
    # Add Heading Paragraph
    if first_paragraph:
        p_heading = tf_body.paragraphs[0]
        first_paragraph = False
    else:
        p_heading = tf_body.add_paragraph()
    
    p_heading.text = item["heading"]
    p_heading.level = 0
    p_heading.font.size = Pt(20)
    p_heading.font.bold = True
    p_heading.font.color.rgb = SUBHEADING_COLOR
    p_heading.space_after = Pt(6)
    
    # Add Bullet Points
    for point_text in item["points"]:
        p_point = tf_body.add_paragraph()
        p_point.text = point_text
        p_point.level = 1
        p_point.font.size = Pt(16)
        p_point.font.bold = False
        p_point.font.color.rgb = BODY_COLOR
        p_point.space_after = Pt(4)

# Save presentation
prs.save('output.pptx')