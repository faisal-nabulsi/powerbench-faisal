from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize the presentation
prs = Presentation()

# 2. Add a slide using a blank layout
# Index 6 is typically the "Blank" layout in standard templates
try:
    slide_layout = prs.slide_layouts[6]
except IndexError:
    # Fallback to first layout if blank is not found
    slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(slide_layout)

# 3. Set the background image
# Use the placeholder image provided covering the whole slide
slide_width = prs.slide_width
slide_height = prs.slide_height
try:
    # Adding picture first places it at the bottom of the Z-order (background)
    slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
except Exception:
    pass # Continue if image file is not found

# 4. Create the smaller title (Top-Left)
# Text: "The Subject and Content of Art"
# Background: Yellow, Text: Black
left_st = Inches(0.3)
top_st = Inches(0.3)
width_st = Inches(3.5)
height_st = Inches(0.5)

shape_st = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_st, top_st, width_st, height_st)
shape_st.fill.solid()
shape_st.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow
shape_st.line.fill.background()  # Remove border

tf_st = shape_st.text_frame
tf_st.word_wrap = True
p_st = tf_st.paragraphs[0]
p_st.text = "The Subject and Content of Art"
p_st.alignment = PP_ALIGN.CENTER
run_st = p_st.runs[0]
run_st.font.color.rgb = RGBColor(0, 0, 0)  # Black
run_st.font.size = Pt(18)

# 5. Create the main title
# Text: "THE CONTENT OF ART"
# Background: Yellow, Text: Black
left_mt = Inches(0.5)
top_mt = Inches(1.0)
width_mt = Inches(9.0)
height_mt = Inches(0.8)

shape_mt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_mt, top_mt, width_mt, height_mt)
shape_mt.fill.solid()
shape_mt.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow
shape_mt.line.fill.background()  # Remove border

tf_mt = shape_mt.text_frame
tf_mt.word_wrap = True
p_mt = tf_mt.paragraphs[0]
p_mt.text = "THE CONTENT OF ART"
p_mt.alignment = PP_ALIGN.CENTER
run_mt = p_mt.runs[0]
run_mt.font.color.rgb = RGBColor(0, 0, 0)  # Black
run_mt.font.size = Pt(36)
run_mt.font.bold = True

# 6. Create the main content box
# Background: White
left_c = Inches(0.5)
top_c = Inches(2.2)
width_c = Inches(9.0)
height_c = Inches(4.5)

shape_c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_c, top_c, width_c, height_c)
shape_c.fill.solid()
shape_c.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
shape_c.line.fill.background()  # Remove border

tf_c = shape_c.text_frame
tf_c.word_wrap = True

# Introductory Text
p_intro = tf_c.paragraphs[0]
p_intro.text = "It is the mass of ideas associated with each artwork and communicated through the following:"
run_intro = p_intro.runs[0]
run_intro.font.size = Pt(16)
run_intro.font.color.rgb = RGBColor(0, 0, 0)
p_intro.space_after = Pt(12)

# Bullet Points
bullets = [
    "1. The art’s imagery",
    "2. The symbolic meaning",
    "3. Its surroundings where it is used or displayed",
    "4. The customs, beliefs and values of the culture that uses it"
]

# Unicode checkmark
checkmark = "\u2713"

for bullet_text in bullets:
    p_new = tf_c.add_paragraph()
    p_new.text = f"{checkmark} {bullet_text}"
    run_new = p_new.runs[0]
    run_new.font.size = Pt(16)
    run_new.font.color.rgb = RGBColor(0, 0, 0)
    p_new.space_after = Pt(6)

# 7. Save the presentation
prs.save('output.pptx')