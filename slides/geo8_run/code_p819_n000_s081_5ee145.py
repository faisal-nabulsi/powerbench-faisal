from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# Create a new presentation
prs = Presentation()

# Add a blank slide (Index 6 is typically the 'Blank' layout in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a clean, professional background color (Light Gray)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

# --- Add Title ---
# Position: Top, Left aligned
left = Inches(0.5)
top = Inches(0.5)
width = Inches(12.33)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "The Technology used in Blockchain"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.LEFT

# --- Add Bulleted List ---
# Position: Left side of the slide
left = Inches(0.5)
top = Inches(2)
width = Inches(6)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

items = [
    "Cryptographic Keys",
    "Network Protocol",
    "Distributed Ledger Technology",
    "Hashing"
]

for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    p.text = item
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Manually add bullet character to the paragraph
    pPr = p._p.get_or_add_pPr()
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
    pPr.append(buChar)

# --- Add Image ---
# Position: Right side of the slide
left = Inches(7)
top = Inches(2)
width = Inches(5.5)
height = Inches(4.5)
slide.shapes.add_picture('image.png', left, top, width, height)

# Save the presentation
prs.save('output.pptx')