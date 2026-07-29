from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_still_life_slide():
    # 1. Setup Presentation
    prs = Presentation()
    # Set standard 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use a blank layout to ensure no default placeholders interfere
    # Index 6 is typically 'Blank' in standard templates
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        slide_layout = prs.slide_layouts[0] # Fallback
    
    slide = prs.slides.add_slide(slide_layout)

    # 2. Background: Colorful abstract image
    # Add the image first so it renders behind other shapes.
    # We use the provided 'image.png' as the placeholder.
    # Adding it to the slide shapes list first handles the z-order naturally in most cases.
    background_shape = slide.shapes.add_picture(
        'image.png', 
        0, 0, 
        prs.slide_width, 
        prs.slide_height
    )

    # 3. Top-Left Label: "The Subject and Content of Art" with yellow background
    # Positioned at top-left with a small margin for aesthetics
    label_left = Inches(0.5)
    label_top = Inches(0.2)
    label_width = Inches(5.5)
    label_height = Inches(0.6)
    
    label_box = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
    
    # Set yellow background for the label
    label_fill = label_box.fill
    label_fill.solid()
    label_fill.fore_color.rgb = RGBColor(255, 255, 0)
    
    # Set text content
    label_tf = label_box.text_frame
    label_tf.word_wrap = False
    label_p = label_tf.paragraphs[0]
    label_p.alignment = PP_ALIGN.LEFT
    label_run = label_p.add_run()
    label_run.text = "The Subject and Content of Art"
    label_run.font.size = Pt(16)
    label_run.font.color.rgb = RGBColor(0, 0, 0) # Black text for contrast on yellow

    # 4. Title: "Still Life" centered at the top
    # Positioned below the label to prevent overlap, but still in the top area
    title_top = Inches(0.9)
    title_height = Inches(1.0)
    # Width spans the slide to allow text alignment to center it
    title_box = slide.shapes.add_textbox(Inches(0), title_top, prs.slide_width, title_height)
    
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.add_run()
    run_title.text = "Still Life"
    run_title.font.size = Pt(44)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(255, 255, 255) # White text for visibility on abstract bg

    # 5. Insert two images side by side
    # Dimensions for the two images
    img_width = Inches(4.5)
    img_height = Inches(3.5)
    gap = Inches(1.0)
    
    # Vertical position (below title and label)
    img_y = Inches(2.0)
    
    # Calculate start X to center the pair of images on the slide
    total_images_width = (img_width * 2) + gap
    start_x = (prs.slide_width - total_images_width) / 2
    
    # Left Image (Flowers in a vase) - Using placeholder
    left_img = slide.shapes.add_picture(
        'image.png', 
        start_x, 
        img_y, 
        img_width, 
        img_height
    )
    
    # Right Image (Fruits and a vase) - Using placeholder
    right_img = slide.shapes.add_picture(
        'image.png', 
        start_x + img_width + gap, 
        img_y, 
        img_width, 
        img_height
    )

    # 6. White border around images for separation
    white_border_color = RGBColor(255, 255, 255)
    border_thickness = Inches(0.15) # Approx 10pt border
    
    for img in [left_img, right_img]:
        img.line.color.rgb = white_border_color
        img.line.width = border_thickness

    # 7. Save the presentation
    prs.save('output.pptx')

if __name__ == '__main__':
    create_still_life_slide()