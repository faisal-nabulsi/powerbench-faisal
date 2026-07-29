from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_shape_transparency(shape, alpha):
    """
    Sets the transparency for a solid-filled shape.
    alpha: 0 (opaque) to 100000 (fully transparent)
    """
    try:
        color_el = shape.fill.fore_color.element
        if color_el is not None:
            color_el.set('alpha', str(alpha))
    except Exception:
        pass

# Initialize presentation with standard 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use a blank layout to maintain full creative control
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 1. Background Image
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Left Side Image (Warehouse)
left_img = slide.shapes.add_picture('image.png', 0, 0, Inches(6.0), prs.slide_height)
left_img.line.fill.background()

# 3. Right Side Image (Ship)
right_img = slide.shapes.add_picture('image.png', Inches(6.0), 0, Inches(7.333), prs.slide_height)
right_img.line.fill.background()

# 4. Curved Transition Shape
# Placed at the boundary to create a smooth visual blend between sections
transition_shape = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, 
    Inches(5.0), 0, Inches(2.0), prs.slide_height
)
transition_shape.fill.solid()
transition_shape.fill.fore_color.rgb = RGBColor(50, 50, 50)
transition_shape.line.fill.background()
set_shape_transparency(transition_shape, 60000) # 60% transparent

# 5. Semi-Transparent Overlay for the Right Side
# Darkens the background to ensure the title stands out
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(6.0), 0, Inches(7.333), prs.slide_height
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
overlay.line.fill.background()
set_shape_transparency(overlay, 50000) # 50% transparent

# 6. Title Text Box
# Positioned prominently on the right side
txBox = slide.shapes.add_textbox(Inches(8.0), Inches(3.0), Inches(5.5), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Transportation and Storage"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Save the final presentation
prs.save('output.pptx')