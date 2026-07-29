from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create the presentation
prs = Presentation()

# Use a blank slide layout (index 6 is typically blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Set Dark Background
background = slide.background
fill = background.fill
fill.solid()
# Dark gray/black background
fill.fore_color.rgb = RGBColor(30, 30, 30)

# 2. Add Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(255, 255, 255)  # White
p_title.alignment = PP_ALIGN.CENTER

# 3. Add Content Box
content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5.5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Define Color Variables
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 255, 50)  # Bright yellow for contrast

# Data Structure
bullet_data = [
    {
        "header": "Effective Communication",
        "sub_bullets": [
            "Strong communication skills encompass the ability to convey ideas clearly, listen actively, and adapt communication styles to different audiences."
        ]
    },
    {
        "header": "Problem Solving and Critical Thinking",
        "sub_bullets": [
            "Critical thinking skills enable individuals to analyze information, evaluate options, and make informed decisions.",
            "Problem-solving skills are essential for addressing challenges and finding innovative solutions in a wide range of situations."
        ]
    },
    {
        "header": "Emotional Intelligence",
        "sub_bullets": [
            "Emotional intelligence enhances empathy, interpersonal relationships, and the ability to work effectively in teams."
        ]
    },
    {
        "header": "Adaptability and Resilience",
        "sub_bullets": [
            "These skills help individuals navigate unexpected setbacks, cope with stress, and embrace change as an opportunity for growth."
        ]
    }
]

# Populate Content
# The text box creates one empty paragraph by default.
# We will reuse it for the first item.
first_paragraph = True

for item in bullet_data:
    # --- Main Bullet Point ---
    if first_paragraph:
        p_main = content_tf.paragraphs[0]
        first_paragraph = False
    else:
        p_main = content_tf.add_paragraph()
    
    p_main.level = 0
    run_main = p_main.add_run()
    # Using a bullet symbol for consistency
    run_main.text = "\u2022 " + item["header"]
    run_main.font.size = Pt(24)
    run_main.font.bold = True
    run_main.font.color.rgb = WHITE
    p_main.space_before = Pt(14)
    p_main.space_after = Pt(6)

    # --- Sub-bullets ---
    for sub_text in item["sub_bullets"]:
        p_sub = content_tf.add_paragraph()
        p_sub.level = 1
        run_sub = p_sub.add_run()
        # Using indentation spaces and bullet symbol
        run_sub.text = "      \u2022 " + sub_text
        run_sub.font.size = Pt(18)
        run_sub.font.bold = False
        run_sub.font.color.rgb = YELLOW
        p_sub.space_before = Pt(4)
        p_sub.space_after = Pt(4)

# Save the presentation
prs.save('output.pptx')