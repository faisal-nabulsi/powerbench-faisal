from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize a presentation with 16:9 slide size
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide (Index 6 is typically 'Blank' in default templates)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Set a clean light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 2. Add Title "PART 3"
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(5)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PART 3"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # 3. Add Instruction Text
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(10)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Choose true (T), false (F), or not given (N) according to the information in the video."
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    # 4. Create Table
    # 4 Rows (1 Header + 3 Statements), 4 Columns (Statement, T, F, N)
    rows, cols = 4, 4
    left = Inches(0.5)
    top = Inches(2.1)
    width = Inches(8.5)
    height = Inches(3.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Adjust column widths to prioritize statement text
    table.columns[0].width = Inches(6.0)
    table.columns[1].width = Inches(0.833)
    table.columns[2].width = Inches(0.833)
    table.columns[3].width = Inches(0.833)
    
    # Define Headers and Data
    headers = ["Statement", "T", "F", "N"]
    statements = [
        "1. Phoebe sarcastically says that Sarah is a monster for her etiquette.",
        "2. Joey prefers eating French fries with his fingers.",
        "3. Rachel is surprised to hear about Joey’s food sharing rule."
    ]
    
    # Unicode checkbox character for empty checkboxes
    checkbox_char = "☐"
    
    # Style and Populate Header Row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(70, 130, 180) # Steel Blue
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
            
    # Style and Populate Data Rows
    for i, stmt in enumerate(statements):
        row_idx = i + 1
        
        # Statement Cell
        cell_stmt = table.cell(row_idx, 0)
        cell_stmt.text = stmt
        cell_stmt.fill.solid()
        cell_stmt.fill.fore_color.rgb = RGBColor(255, 255, 255)
        for paragraph in cell_stmt.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
        # Checkboxes Cells (T, F, N)
        for col in range(1, 4):
            cell = table.cell(row_idx, col)
            cell.text = checkbox_char
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(22) # Larger font for visible box
                paragraph.alignment = PP_ALIGN.CENTER
                
    # 5. Add Image of Joey on the right side
    # Positioned to the right of the table
    try:
        slide.shapes.add_picture('image.png', Inches(9.5), Inches(0.5), Inches(3.5), Inches(5.5))
    except FileNotFoundError:
        # Fallback placeholder if image.png is missing
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0.5), Inches(3.5), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        shape.line.fill.background()
        shape.text_frame.paragraphs[0].text = "Joey Image"
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    # 6. Add "VIEWING ACTIVITY" Button at the bottom
    left = Inches(5.5)
    top = Inches(6.5)
    width = Inches(2.5)
    height = Inches(0.6)
    
    # Create Rectangle Shape for Button
    btn_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    btn_tf = btn_shape.text_frame
    btn_tf.paragraphs[0].text = "VIEWING ACTIVITY"
    btn_tf.paragraphs[0].font.size = Pt(14)
    btn_tf.paragraphs[0].font.bold = True
    btn_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    btn_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Style Button
    btn_shape.fill.solid()
    btn_shape.fill.fore_color.rgb = RGBColor(46, 139, 87) # Sea Green
    btn_shape.line.fill.background() # Remove border
    
    # 7. Save Presentation
    prs.save('output.pptx')

if __name__ == '__main__':
    create_presentation()