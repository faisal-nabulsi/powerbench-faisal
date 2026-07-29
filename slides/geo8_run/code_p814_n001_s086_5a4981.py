from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Set slide dimensions to 16:9 (13.33" x 7.5") for a modern look
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Use a blank layout. Index 6 is typically blank in standard templates.
    # We use a try-except block to handle cases where layout indices might vary.
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        # Fallback to the first layout usually available
        slide_layout = prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(slide_layout)

    # Set white background to ensure clean look
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Add Yellow Vertical Bar on the far left for emphasis
    # Coordinates: x=0.3, top=0, width=0.25, height=7.5 (full height)
    bar_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.3), 
        Inches(0), 
        Inches(0.25), 
        Inches(7.5)
    )
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    bar_shape.line.fill.background() # Remove border/outline

    # Add Title on the left side
    # Text: "Obstacles that Elon faced"
    # Style: Bold, Large font (44pt), Left aligned
    # Coordinates: Left=1.5, Top=1.0 (aligned with content), Width=4, Height=1.5
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(4), Inches(1.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_title = tf_title.paragraphs[0]
    p_title.text = "Obstacles that Elon faced"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 0, 0)
    p_title.alignment = PP_ALIGN.LEFT

    # Add Bullet Points on the right side
    # Coordinates: Left=6.5, Top=1.0, Width=6, Height=5.5
    content_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.0), Inches(6), Inches(5.5))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True

    bullets = [
        "Rejection from Netscape and Ousted at ZIP2: In 1996, Elon Musk was ousted as the CEO of the company.",
        "PayPal was voted 'worst business concept of the year’: When PayPal was issued as security software for the PalmPilot and other portable devices in 1999, it was voted the worst business idea of the year.",
        "Tesla crisis: The company also had serious financial problems, and it was on the verge of going out of business.",
        "SpaceX failure: The failure of the company's first three launches generated skepticism among many investors, but the fourth one got a huge success.",
        "Running on empty: 2008 was the worst year of his life. After investing his whole money in Tesla and SpaceX, Musk was reliant on personal loans from friends."
    ]

    # Add each bullet point to the text box
    for i, text in enumerate(bullets):
        # Get existing paragraph or add new one
        if i == 0:
            p = tf_content.paragraphs[0]
        else:
            p = tf_content.add_paragraph()
        
        # Add bullet character "• " manually to ensure consistent rendering across viewers
        p.text = "• " + text
        
        # Format text for legibility
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0) # Black text
        p.space_after = Pt(10) # Spacing between bullet points

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()