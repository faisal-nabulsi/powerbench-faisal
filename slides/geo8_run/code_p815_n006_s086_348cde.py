from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # Select a blank slide layout to ensure no default placeholders interfere
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. Slide Background: Light Gray
    # Add a rectangle covering the entire slide to set the background color
    bg_shape = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        0, 0, slide_width, slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0xD3, 0xD3, 0xD3) # Light Gray
    bg_shape.line.fill.background() # Remove outline

    # 2. Title "PART 4"
    # Top left, blue background, bold font
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(2.5)
    title_height = Inches(0.7)

    title_shape = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        title_left, title_top, title_width, title_height
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0) # Blue
    title_shape.line.fill.background()

    # Set text for title
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 4"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text
    p.alignment = PP_ALIGN.LEFT

    # 3. Icon of human speaking
    # Below the title (using Unicode emoji as icon)
    icon_top = title_top + title_height + Inches(0.2)
    icon_left = title_left
    
    speaking_icon_shape = slide.shapes.add_textbox(
        icon_left, icon_top, Inches(1), Inches(0.6)
    )
    tf_icon = speaking_icon_shape.text_frame
    p_icon = tf_icon.paragraphs[0]
    p_icon.text = "🗣️" # Speaking head emoji
    p_icon.font.size = Pt(28)
    p_icon.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # 4. Instruction text
    # Below the title/icon area
    instr_top = icon_top + Inches(0.6)
    instr_width = slide_width - Inches(1)
    
    instr_box = slide.shapes.add_textbox(
        title_left, instr_top, instr_width, Inches(1.0)
    )
    tf_instr = instr_box.text_frame
    tf_instr.word_wrap = True
    p_instr = tf_instr.paragraphs[0]
    p_instr.text = "Tell the short story below out loud using the pronunciation rule you studied on the previous slide."
    p_instr.font.size = Pt(16)
    p_instr.font.name = "Arial"
    p_instr.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # 5. Story Text
    # Center-left area
    story_top = instr_top + Inches(0.8)
    story_width = slide_width * 0.6
    story_height = Inches(4.0)

    story_box = slide.shapes.add_textbox(
        title_left, story_top, story_width, story_height
    )
    tf_story = story_box.text_frame
    tf_story.word_wrap = True
    p_story = tf_story.paragraphs[0]
    p_story.text = "Last week, I went on a date. We were skating in the park, we were laughing, and the guy was making some really funny jokes...and then, out of nowhere, I fell down and I broke my leg!"
    p_story.font.size = Pt(20)
    p_story.font.name = "Calibri" 
    p_story.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 6. Image of people skating
    # Right side
    img_left = slide_width * 0.65
    img_top = story_top
    img_width = slide_width * 0.3
    img_height = Inches(4.0)

    try:
        slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
    except Exception:
        pass # Ignore if image not found

    # 7. Bottom Section
    # Audio Icon and Instruction
    bottom_top = slide_height - Inches(1.5)
    
    # Audio Icon (Unicode)
    audio_icon_left = title_left
    audio_icon_top = bottom_top + Inches(0.2)
    
    audio_icon_shape = slide.shapes.add_textbox(
        audio_icon_left, audio_icon_top, Inches(0.5), Inches(0.5)
    )
    tf_audio = audio_icon_shape.text_frame
    p_audio = tf_audio.paragraphs[0]
    p_audio.text = "🔊" # Speaker emoji
    p_audio.font.size = Pt(24)
    p_audio.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Bottom Instruction Text
    text_bottom_left = title_left + Inches(0.6)
    text_bottom_top = bottom_top + Inches(0.2)
    bottom_width = slide_width - Inches(1)
    
    bottom_text_box = slide.shapes.add_textbox(
        text_bottom_left, text_bottom_top, bottom_width, Inches(0.8)
    )
    tf_bottom = bottom_text_box.text_frame
    tf_bottom.word_wrap = True
    p_bottom = tf_bottom.paragraphs[0]
    p_bottom.text = "Play and listen to the audio clip to check your answers. Repeat if necessary."
    p_bottom.font.size = Pt(16)
    p_bottom.font.name = "Arial"
    p_bottom.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()