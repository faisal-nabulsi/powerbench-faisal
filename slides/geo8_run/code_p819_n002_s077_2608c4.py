from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# Get slide dimensions (Standard Widescreen is 13.333 x 7.5 inches)
slide_width = prs.slide_width
slide_height = prs.slide_height

# Use blank layout
slide_layout = prs.slide_layouts[6] 
slide = prs.slides.add_slide(slide_layout)

# 1. Set Background Image
# The instruction requires a "colorful abstract image" for the background.
# We use the provided placeholder 'image.png' for this.
bg_img = slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)

# 2. Text Area Overlay
# To ensure text is clear and legible against the potentially busy background,
# we add a semi-transparent (simulated dark grey) rectangle behind the titles.
overlay_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(0.4), Inches(0.4), Inches(7.5), Inches(2.5)
)
overlay_shape.fill.solid()
overlay_shape.fill.fore_color.rgb = RGBColor(40, 40, 40) # Dark Grey
overlay_shape.line.color.rgb = RGBColor(40, 40, 40) # Match fill to hide border

# 3. Small Title (Top-Left)
# "At the top-left corner, add a smaller title 'The Subject and Content of Art.'"
txBox_small = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(7), Inches(0.8))
tf = txBox_small.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The Subject and Content of Art"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255) # White text for contrast

# 4. Main Title (Below Small Title)
# "Create a slide titled 'Non-representational or Non-objective Art' in a large font."
txBox_main = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(7), Inches(1))
tf = txBox_main.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Non-representational or Non-objective Art"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255) # White text
p.font.bold = True

# 5. Left Image
# "Below the title, on the left side, place an image showing an abstract expression of vibrant splashes."
# Position: Left aligned (0.5), Y starts below titles (3.2)
left_img = slide.shapes.add_picture('image.png', Inches(0.5), Inches(3.2), Inches(6), Inches(4))
left_img.line.color.rgb = RGBColor(255, 255, 255) # White border
left_img.line.width = Pt(5)

# 6. Right Top Image
# "On the right side, first put an image showing a symphony of colors in motion on the top"
# Position: Right aligned (7.0), Top aligned (0.5)
rt_img = slide.shapes.add_picture('image.png', Inches(7.0), Inches(0.5), Inches(5.5), Inches(3.5))
rt_img.line.color.rgb = RGBColor(255, 255, 255) # White border
rt_img.line.width = Pt(5)

# 7. Right Bottom Image
# "then place an image of a vibrant abstract symphony on the bottom."
# Position: Right aligned (7.0), Below top image + gap (4.5)
rb_img = slide.shapes.add_picture('image.png', Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5))
rb_img.line.color.rgb = RGBColor(255, 255, 255) # White border
rb_img.line.width = Pt(5)

prs.save('output.pptx')