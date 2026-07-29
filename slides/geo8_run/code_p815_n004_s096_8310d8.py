from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create a presentation
prs = Presentation()

# Add a blank slide (index 6 is typically the blank layout in standard themes)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background to white
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# 1. Add a vertical yellow bar for emphasis on the left side
# Positioned at x=0.5 inches, starting below the title area
bar_left = Inches(0.5)
bar_top = Inches(1.2)
bar_width = Inches(0.15)
bar_height = Inches(5.0)

bar_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    bar_left, bar_top, bar_width, bar_height
)
bar_shape.fill.solid()
bar_shape.fill.fore_color.rgb = RGBColor(255, 204, 0) # Yellow
bar_shape.line.fill.background() # Remove border

# 2. Add the Title "Content"
# Positioned to the right of the yellow bar
title_left = Inches(0.9)
title_top = Inches(0.5)
title_width = Inches(8.5)
title_height = Inches(0.8)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_tf = title_box.text_frame
title_tf.word_wrap = True

p = title_tf.paragraphs[0]
p.text = "Content"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 0, 0) # Black
p.alignment = PP_ALIGN.LEFT

# 3. Add the Bullet Points
# Positioned below the title, aligned with the bar and title
bullets_left = Inches(0.9)
bullets_top = Inches(1.2)
bullets_width = Inches(8.5)
bullets_height = Inches(5.5)

bullets_box = slide.shapes.add_textbox(bullets_left, bullets_top, bullets_width, bullets_height)
bullets_tf = bullets_box.text_frame
bullets_tf.word_wrap = True

bullet_list = [
    "Brief highlights of Elon Musk",
    "Biography",
    "Early life of Elon Musk",
    "How he came up with his ideas",
    "Elon’s current stage",
    "Obstacles that Elon faced",
    "Lessons that we can learn from Elon Musk's life",
    "Elon’s Future Plans"
]

# Add bullet points
# The text box comes with one empty paragraph initially
p = bullets_tf.paragraphs[0]
p.text = bullet_list[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0, 0, 0) # Black
p.alignment = PP_ALIGN.LEFT
p.level = 0

for text in bullet_list[1:]:
    p = bullets_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black
    p.alignment = PP_ALIGN.LEFT
    p.level = 0

# Save the presentation
prs.save('output.pptx')