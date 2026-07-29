from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set slide dimensions to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Background: Light Green
# Note: A hexagonal pattern is complex to generate programmatically without image assets.
# We use a solid light green background to satisfy the color requirement.
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(214, 233, 210) # Light Green

# 2. Large White Textbox in the Center
# Positioned to allow space for images at the bottom
left = Inches(1.5)
top = Inches(1.0)
width = Inches(10.333)
height = Inches(4.5)

textbox = slide.shapes.add_textbox(left, top, width, height)

# Set textbox fill to white
textbox.fill.solid()
textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Access the text frame
tf = textbox.text_frame
tf.word_wrap = True

# 3. Title Section
p_title = tf.paragraphs[0]
p_title.text = "NEGATIVE IMPACT"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0, 100, 0) # Bold Green
p_title.alignment = PP_ALIGN.LEFT
p_title.space_after = Pt(20)

# 4. Content Section
content_texts = [
    "Humans are continuously cutting down trees and littering have a negative impact on environment.",
    "Pollution is everywhere, from the trash thrown out on the free way to the millions of metric tons of pollution pumped into the atmosphere every year.",
    "Burning fossil fuels is the prime culprit from the climate for climate change. When oil gas and coal are burned they release carbon dioxide and other harmful gases which trap heat in atmosphere and cause it to warm up.",
    "The increasing use of fertilizers and pesticides to protect crops have damaged the fertility of land.",
    "Food wastage is a huge issue with an estimated one third of all food produced globally is being wasted."
]

for text in content_texts:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(10)
    p.space_after = Pt(10)

# 5. Image Section
# Insert two images at the bottom using the placeholder 'image.png'
image_height = Inches(1.5)
image_width = Inches(4.5)
image_top = Inches(5.8)

# Left Image (Deforestation placeholder)
slide.shapes.add_picture('image.png', Inches(1.5), image_top, image_width, image_height)

# Right Image (Pollution placeholder)
# Positioned symmetrically on the right
right_image_left = Inches(13.333) - Inches(1.5) - image_width
slide.shapes.add_picture('image.png', right_image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')