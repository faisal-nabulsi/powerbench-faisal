import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide():
    # Create a blank presentation
    prs = Presentation()

    # Get slide dimensions to ensure responsive layout
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add a blank slide layout (Index 6 is typically 'Blank')
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        # Fallback to first layout if blank is not at index 6
        slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(slide_layout)

    # 1. Background Image
    # Place 'image.png' to cover the entire slide background
    try:
        slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
    except FileNotFoundError:
        print("Warning: 'image.png' not found. Background will be default.")

    # 2. Title Section with Legibility Enhancement
    # Create a semi-transparent white shape for the title to ensure it contrasts with the background image
    title_card_width = slide_width * 0.80
    title_card_height = Inches(1.2)
    title_card_left = (slide_width - title_card_width) / 2
    title_card_top = Inches(0.8)

    title_bg = slide.shapes.add_shape(1, title_card_left, title_card_top, title_card_width, title_card_height)
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Attempt to set transparency (Alpha channel) for a modern look
    try:
        title_bg.fill.fore_color.set_alpha(0.85)
    except AttributeError:
        pass
    
    title_bg.line.color.rgb = RGBColor(0, 0, 0)
    title_bg.line.width = Pt(0)

    # Add Title Text
    title_box = slide.shapes.add_textbox(
        title_card_left + Inches(0.3), 
        title_card_top + Inches(0.2), 
        title_card_width - Inches(0.6), 
        title_card_height - Inches(0.4)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_para = title_tf.paragraphs[0]
    title_para.text = "Components of visual communication."
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(50, 50, 50) # Dark text for contrast
    title_para.alignment = PP_ALIGN.LEFT

    # 3. Content Section with Dark Overlay
    # Create a semi-transparent black overlay for the bottom portion of the slide
    content_card_height = slide_height * 0.60
    content_card_top = slide_height - content_card_height

    content_bg = slide.shapes.add_shape(1, 0, content_card_top, slide_width, content_card_height)
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(30, 30, 30) # Dark Grey
    try:
        content_bg.fill.fore_color.set_alpha(0.80)
    except AttributeError:
        pass
    content_bg.line.color.rgb = RGBColor(0, 0, 0)
    content_bg.line.width = Pt(0)

    # 4. Bullet Points
    bullet_points = [
        "Components of craftsmanship are the essential units of any visual plan that structure the plan's design and pass on its visual messages.",
        "The components of visual computerization are line, shape, size, color, Typography, Space and Texture."
    ]

    content_box = slide.shapes.add_textbox(
        Inches(1.5), 
        content_card_top + Inches(0.8), 
        slide_width - Inches(3), 
        content_card_height - Inches(1.5)
    )
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    for i, text in enumerate(bullet_points):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        
        # Prepend bullet character manually as default style might not render on blank layout
        p.text = "• " + text
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255) # White text for dark overlay
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # Save the presentation
    prs.save('output.pptx')
    print("Presentation saved to 'output.pptx'")

if __name__ == "__main__":
    create_slide()