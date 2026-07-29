from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Create a presentation object
    prs = Presentation()

    # Set slide width and height (16:9 is default, but good to be explicit if needed)
    # Default is 13.333 x 7.5 inches
    
    # Use a blank layout to have full control
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Background: Colorful abstract image
    # Using the provided placeholder 'image.png'
    # We add it as a shape covering the whole slide
    bg_shape = slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)
    # Send to back
    sp = bg_shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp) # Insert after background elements

    # 2. Small Title: "The Subject and Content of Art"
    # Top-left corner
    small_title_left = Inches(0.5)
    small_title_top = Inches(0.5)
    small_title_width = Inches(6)
    small_title_height = Inches(0.6)

    txBox_small = slide.shapes.add_textbox(small_title_left, small_title_top, small_title_width, small_title_height)
    tf_small = txBox_small.text_frame
    tf_small.text = "The Subject and Content of Art"
    
    # Formatting small title
    p_small = tf_small.paragraphs[0]
    p_small.font.size = Pt(16)
    p_small.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p_small.font.bold = True
    
    # Yellow background for small title
    txBox_small.fill.solid()
    txBox_small.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    txBox_small.line.fill.background() # No border

    # 3. Main Title: "Portrait"
    # Centered below the small title
    # Let's center it horizontally on the slide for better aesthetics, 
    # but ensure it is vertically below the small title.
    # Or strictly centered below the small title's center? 
    # "Centered below the small title" -> The text "Portrait" is centered, and it is below the small title.
    
    main_title_width = Inches(4)
    main_title_height = Inches(0.8)
    main_title_top = Inches(1.3) # Below small title (0.5 + 0.6 + gap)
    
    # Center horizontally on slide
    slide_width_inches = prs.slide_width / 914400 # Convert EMU to inches approx
    # Actually easier to use Inches() and math
    # Slide width is 13.333 inches.
    # Center X = 6.666
    # Box Left = 6.666 - (4/2) = 4.666
    main_title_left = Inches(4.66) 

    txBox_main = slide.shapes.add_textbox(main_title_left, main_title_top, main_title_width, main_title_height)
    tf_main = txBox_main.text_frame
    tf_main.text = "Portrait"
    
    # Formatting main title
    p_main = tf_main.paragraphs[0]
    p_main.alignment = PP_ALIGN.CENTER
    p_main.font.size = Pt(40)
    p_main.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p_main.font.bold = True
    
    # Yellow background for main title
    txBox_main.fill.solid()
    txBox_main.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    txBox_main.line.fill.background()

    # 4. Images: Mona Lisa (Left) and Child Portrait (Right)
    # Using 'image.png' for both
    # White background for images
    
    img_width = Inches(4)
    img_height = Inches(5)
    img_top = Inches(2.5) # Below titles
    
    # Left Image (Mona Lisa)
    # Position: Left side, centered vertically in remaining space or just below titles
    # Let's place it at x=1.5
    left_img_x = Inches(1.5)
    
    # White background shape
    rect_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_img_x, img_top, img_width, img_height)
    rect_left.fill.solid()
    rect_left.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_left.line.fill.background()
    
    # Image
    slide.shapes.add_picture('image.png', left_img_x, img_top, img_width, img_height)
    
    # Right Image (Child Portrait)
    # Position: Right side, aligned with left image
    # Gap of 1 inch between images
    right_img_x = Inches(1.5 + 4 + 1) # 6.5
    
    # White background shape
    rect_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_img_x, img_top, img_width, img_height)
    rect_right.fill.solid()
    rect_right.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect_right.line.fill.background()
    
    # Image
    slide.shapes.add_picture('image.png', right_img_x, img_top, img_width, img_height)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()