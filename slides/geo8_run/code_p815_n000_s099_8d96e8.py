from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()

# Set slide dimensions (16:9 Aspect Ratio)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set slide background to white (Main content background)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# --- 1. Title Textbox ---
# Positioned at the top
title_left = Inches(0.5)
title_top = Inches(0.5)
title_width = Inches(12.333)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf = title_box.text_frame
tf.word_wrap = True

# Apply blue background to the title textbox
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 51, 153) # Blue
title_box.line.fill.background() # Remove border

# Apply white text to the title
p = tf.paragraphs[0]
p.text = "Global and Local Cultural Products"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255) # White
p.alignment = PP_ALIGN.CENTER

# --- 2. Subtitle ---
# Positioned below the title
sub_left = Inches(0.5)
sub_top = Inches(1.9)
sub_width = Inches(12.333)
sub_height = Inches(0.6)

sub_box = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
tf_sub = sub_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "Cultural Products"
p_sub.font.size = Pt(28)
p_sub.font.color.rgb = RGBColor(0, 0, 0) # Black
p_sub.alignment = PP_ALIGN.LEFT

# --- 3. Bullet Point ---
# Content text below subtitle
bullet_left = Inches(0.5)
bullet_top = Inches(2.7)
bullet_width = Inches(12.333)
bullet_height = Inches(1.5)

bullet_box = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
tf_bullet = bullet_box.text_frame
tf_bullet.word_wrap = True
p_bullet = tf_bullet.paragraphs[0]
p_bullet.text = "✓ are goods and services such as arts, architectures, museums etc. that showcase the history and information of certain which belong to the country’s cultural heritage."
p_bullet.font.size = Pt(20)
p_bullet.font.color.rgb = RGBColor(0, 0, 0) # Black
p_bullet.alignment = PP_ALIGN.LEFT

# --- 4. Images ---
# Three images aligned horizontally at the bottom
# Using 'image.png' as the placeholder for all images as instructed
img_width = Inches(3.5)
img_height = Inches(2.5)
gap = Inches(0.5)

# Position at the bottom with a margin
bottom_margin = Inches(0.5)
img_top = prs.slide_height - img_height - bottom_margin

# Calculate start X to center the three images
total_img_width = (3 * img_width) + (2 * gap)
start_x = (prs.slide_width - total_img_width) / 2

# Add Image 1
slide.shapes.add_picture('image.png', start_x, img_top, img_width, img_height)
# Add Image 2
slide.shapes.add_picture('image.png', start_x + img_width + gap, img_top, img_width, img_height)
# Add Image 3
slide.shapes.add_picture('image.png', start_x + 2 * (img_width + gap), img_top, img_width, img_height)

# Save the presentation
prs.save('output.pptx')