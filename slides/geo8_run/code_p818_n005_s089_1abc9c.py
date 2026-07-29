from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Create Presentation
prs = Presentation()

# Set slide size to 16:9 (Standard for modern presentations)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (index 6 is usually the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 2. Title Section with Dark Blue Background
# Create a rectangle for the title background to satisfy "dark blue background for only the title"
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    0, 0, 
    Inches(13.333), Inches(1.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_bg.line.color.rgb = RGBColor(0, 51, 102) # Remove border to blend in

# Add the main title text
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9)
)
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "5 Stages of Development of Media"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(255, 255, 255) # White Text
p_title.alignment = PP_ALIGN.LEFT

# 3. Subtitle
subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.8), Inches(12), Inches(0.5)
)
tf_sub = subtitle_box.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "4. ELECTRONIC MEDIA"
p_sub.font.size = Pt(24)
p_sub.font.bold = True
p_sub.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
p_sub.alignment = PP_ALIGN.LEFT

# 4. Bullet Points
bullet_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.6), Inches(6), Inches(4)
)
tf_bullets = bullet_box.text_frame
tf_bullets.word_wrap = True

# First Bullet
p1 = tf_bullets.paragraphs[0]
# Using Unicode bullet character \u2022 to ensure it displays correctly on blank layouts
p1.text = "\u2022 It includes the telegraphs, telephone, radio, film and television."
p1.font.size = Pt(18)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.alignment = PP_ALIGN.LEFT

# Second Bullet
p2 = tf_bullets.add_paragraph()
p2.text = "\u2022 The wide range of these media continue to open up new perspectives on economic, political and cultural processes of globalization."
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0, 0, 0)
p2.alignment = PP_ALIGN.LEFT

# 5. Image
# Place image on the right side, roughly centered vertically with the text
img_left = Inches(7.8)
img_top = Inches(2.2)
img_width = Inches(5.2)
img_height = Inches(4.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# 6. Save Presentation
prs.save('output.pptx')