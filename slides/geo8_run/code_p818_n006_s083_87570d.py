from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize Presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 25, 60) # A professional dark blue

# Get slide dimensions for dynamic positioning
slide_width = prs.slide_width
slide_height = prs.slide_height

# --- Top Left Section ---

# 1. Movie Icon (placed to the left of the title)
# Using 'image.png' as placeholder
icon_left = Inches(0.5)
icon_top = Inches(0.5)
icon_width = Inches(1.5)
icon_height = Inches(1.5)
slide.shapes.add_picture('image.png', icon_left, icon_top, icon_width, icon_height)

# 2. Title: "VIEWING ACTIVITY"
# Positioned to the right of the icon
title_left = icon_left + icon_width + Inches(0.2)
title_top = Inches(0.6)
title_width = Inches(6.0)
title_height = Inches(0.8)

txBox_title = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf_title = txBox_title.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "VIEWING ACTIVITY"
p_title.font.bold = True
p_title.font.size = Pt(36)
p_title.font.color.rgb = RGBColor(255, 255, 255) # White text for contrast

# 3. Instruction: "Watch the first 2 minutes of the video."
# Below the title
instruction_top = title_top + Inches(0.8)
txBox_instr = slide.shapes.add_textbox(title_left, instruction_top, title_width, Inches(0.6))
tf_instr = txBox_instr.text_frame
tf_instr.word_wrap = True
p_instr = tf_instr.paragraphs[0]
p_instr.text = "Watch the first 2 minutes of the video."
p_instr.font.size = Pt(20)
p_instr.font.color.rgb = RGBColor(220, 220, 220) # Light grey

# --- Right Side Section ---

# Determine center position for right-aligned content (approx 75% of slide width)
right_center_x = int(slide_width * 0.75)
content_width = Inches(4.0)
# Calculate left position to center the content at right_center_x
content_left = right_center_x - (content_width / 2)

# 4. Text Box: "Friends | Joey Doesn’t Share Food!"
text_top = Inches(0.5)
text_height = Inches(1.2)
txBox_right = slide.shapes.add_textbox(content_left, text_top, content_width, text_height)
tf_right = txBox_right.text_frame
p_right = tf_right.paragraphs[0]
p_right.text = "Friends | Joey Doesn’t Share Food!"
p_right.alignment = PP_ALIGN.CENTER
p_right.font.bold = True
p_right.font.size = Pt(34) # Prominent font
p_right.font.color.rgb = RGBColor(255, 255, 255) # White

# 5. Image: Joey (below text, center-aligned)
# Reusing 'image.png' for the image placeholder
image_top = text_top + text_height + Inches(0.2)
image_width = Inches(4.0) # Same width as text for alignment
image_height = Inches(3.0)
slide.shapes.add_picture('image.png', content_left, image_top, image_width, image_height)

# --- Bottom Left Logo ---

# 6. Logo representing activity
logo_width = Inches(1.0)
logo_height = Inches(1.0)
logo_left = Inches(0.5)
logo_top = slide_height - logo_height - Inches(0.5) # Bottom alignment with padding

slide.shapes.add_picture('image.png', logo_left, logo_top, logo_width, logo_height)

# Save the presentation
prs.save('output.pptx')