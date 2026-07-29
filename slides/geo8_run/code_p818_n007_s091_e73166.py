from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize the presentation
prs = Presentation()

# Use a blank layout for full control
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a light background color for readability
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

# 1. Title "PART 3"
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.3), Inches(8.0), Inches(1.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PART 3"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
p.alignment = PP_ALIGN.CENTER

# 2. Instruction
instr_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.8))
tf = instr_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Choose true (T), false (F), or not given (N) according to the information in the video."
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = RGBColor(50, 50, 50)
p.alignment = PP_ALIGN.LEFT

# 3. Image of Joey (Placeholder image.png)
# Position on the right side of the slide
try:
    slide.shapes.add_picture('image.png', Inches(6.2), Inches(2.2), Inches(3.2), Inches(3.5))
except FileNotFoundError:
    pass

# 4. Table
# 4 Rows (Header + 3 statements), 5 Columns (No, Statement, T, F, N)
rows = 4
cols = 5
left = Inches(0.8)
top = Inches(2.2)
width = Inches(5.0)
height = Inches(2.8)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set Column Widths
table.columns[0].width = Inches(0.4) # No
table.columns[1].width = Inches(3.6) # Statement
table.columns[2].width = Inches(0.4) # T
table.columns[3].width = Inches(0.4) # F
table.columns[4].width = Inches(0.4) # N

# Headers
headers = ["No.", "Statement", "T", "F", "N"]
for i, text in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = text
    # Style Header Cell
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(70, 130, 180) # Steel Blue
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 255, 255)

# Data Statements
statements = [
    "1. Phoebe sarcastically says that Sarah is a monster for her etiquette.",
    "2. Joey prefers eating French fries with his fingers.",
    "3. Rachel is surprised to hear about Joey’s food sharing rule."
]
checkbox = "☐" # Unicode Checkbox

for r, stmt in enumerate(statements):
    row_idx = r + 1
    
    # Split number and text
    # Format expected: "N. Text"
    if ". " in stmt:
        parts = stmt.split(". ", 1)
        num = parts[0]
        text = parts[1]
    else:
        num = str(r + 1)
        text = stmt
    
    # Column 0: No
    cell = table.cell(row_idx, 0)
    cell.text = num
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    
    # Column 1: Statement
    cell = table.cell(row_idx, 1)
    cell.text = text
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(13)
    
    # Columns 2, 3, 4: Checkboxes
    for c in range(2, 5):
        cell = table.cell(row_idx, c)
        cell.text = checkbox
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(16)

# 5. Button "VIEWING ACTIVITY"
# Position at the bottom center
btn_left = Inches(3.5)
btn_top = Inches(6.5)
btn_width = Inches(3.0)
btn_height = Inches(0.6)

btn_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, btn_left, btn_top, btn_width, btn_height)
btn_shape.text = "VIEWING ACTIVITY"

# Style Button
bf = btn_shape.fill
bf.solid()
bf.fore_color.rgb = RGBColor(128, 128, 128) # Grey
tf = btn_shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for run in p.runs:
    run.font.color.rgb = RGBColor(255, 255, 255) # White Text
    run.font.bold = True
    run.font.size = Pt(16)

# Save the presentation
prs.save('output.pptx')