from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Set slide dimensions to standard widescreen (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide layout
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Slide Background
    # "Use an colorful abstract image describing a vivid dance of colors and emotions as the slide background."
    # We use 'image.png' as the placeholder for any image requested.
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    background_image = slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
    
    # Move background image to the very back
    bg_sp = background_image._element
    slide.shapes._spTree.remove(bg_sp)
    slide.shapes._spTree.insert(2, bg_sp)
    
    # 2. Top-Left Title Box
    # "Add a smaller title on the top-left corner: 'The Subject and Content of Art.'"
    # "Color the background of the title textbox with yellow and the text color with black."
    left_title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(4.8), Inches(0.7)
    )
    
    # Yellow Background
    fill = left_title_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    
    # No Border
    left_title_box.line.fill.background()
    
    # Black Text
    tf = left_title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Subject and Content of Art."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # 3. Main Title
    # "Create a slide titled 'Subject of Art'"
    main_title_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(0.1), Inches(6), Inches(0.8)
    )
    
    tf = main_title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Subject of Art"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) # White for visibility on abstract bg
    p.alignment = PP_ALIGN.CENTER
    
    # 4. Main Content Text
    # "A cat catching a bird is the subject of both works above."
    # "For the main content, use a white background."
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(1.2), Inches(6.5), Inches(0.8)
    )
    
    # White Background for Content
    fill = content_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box.line.fill.background()
    
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A cat catching a bird is the subject of both works above."
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # 5. Images Side by Side
    # "On the left, put an image showing ferocious cat capturing a bird..."
    # "On the right, put an image showing artistic expressions..."
    # Using 'image.png' for both placeholders.
    
    # Define dimensions for the two images
    img_width = Inches(3.5)
    img_height = Inches(2.5)
    img_top = Inches(3.0)
    
    # Left Image
    left_image = slide.shapes.add_picture(
        'image.png', 
        Inches(2.2), img_top, 
        img_width, img_height
    )
    
    # Right Image
    right_image = slide.shapes.add_picture(
        'image.png', 
        Inches(7.8), img_top, 
        img_width, img_height
    )
    
    # Save the presentation
    prs.save('output.pptx')
    print("Presentation created successfully as 'output.pptx'.")

if __name__ == "__main__":
    create_presentation()