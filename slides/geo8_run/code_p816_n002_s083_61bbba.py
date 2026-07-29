from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide():
    # Create presentation
    prs = Presentation()
    
    # Add a blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    # Add magnifying glass icon on the left
    icon_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(1.0), Inches(1.0))
    icon_tf = icon_box.text_frame
    icon_tf.word_wrap = True
    icon_para = icon_tf.paragraphs[0]
    icon_para.alignment = PP_ALIGN.CENTER
    icon_run = icon_para.add_run()
    icon_run.text = "🔍"  # Magnifying glass emoji
    icon_run.font.size = Pt(52)
    icon_run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add main title
    title_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(6.0), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_para = title_tf.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.add_run()
    title_run.text = "QUIZ & REVIEW ACTIVITY"
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add instruction text
    inst_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.7), Inches(6.0), Inches(0.6))
    inst_tf = inst_box.text_frame
    inst_tf.word_wrap = True
    inst_para = inst_tf.paragraphs[0]
    inst_para.alignment = PP_ALIGN.CENTER
    inst_run = inst_para.add_run()
    inst_run.text = "Work with a partner."
    inst_run.font.size = Pt(28)
    inst_run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add FLUENT logo placeholder in top right corner
    slide.shapes.add_picture('image.png', Inches(9.0), Inches(0.3), Inches(1.5), Inches(0.8))
    
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()