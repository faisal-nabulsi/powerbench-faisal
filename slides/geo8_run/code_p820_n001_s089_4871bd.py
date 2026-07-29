from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

# Initialize Presentation with standard 16:9 aspect ratio
prs = Presentation()

# Create a blank slide layout (index 6)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Get slide dimensions for centering calculations
slide_width, slide_height = prs.slide_width, prs.slide_height

# 1. Background: Light Green
# Note: Complex patterns like hexagons require external image assets. 
# A solid light green color is used to fulfill the color requirement.
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(200, 230, 200) # Light Green

# 2. White Textbox Container
# Box Dimensions
box_width = Inches(9.0)
box_height = Inches(4.0)
# Position (Centered horizontally)
box_left = (slide_width - box_width) / 2
box_top = Inches(1.5)

white_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height)
white_box.fill.solid()
white_box.fill.fore_color.rgb = RGBColor(255, 255, 255) # White Fill
white_box.line.color.rgb = RGBColor(255, 255, 255) # Hide border

# 3. Text Content
# Access Text Frame
tf = white_box.text_frame
tf.clear()
tf.word_wrap = True
tf.auto_size = None

# Title Section
pTitle = tf.paragraphs[0]
pTitle.text = "NATURAL ENVIRONMENT"
pTitle.font.size = Pt(28)
pTitle.font.bold = True
pTitle.font.color.rgb = RGBColor(0, 90, 0) # Bold Green
pTitle.alignment = PP_ALIGN.LEFT

# Content Section: Item 1
p1 = tf.add_paragraph()
p1.text = "The natural environment means the non human made surroundings conditions in which all biotic and abiotic things exist on earth."
p1.font.size = Pt(18)
p1.font.color.rgb = RGBColor(0, 0, 0) # Black Font
p1.alignment = PP_ALIGN.LEFT
p1.space_before = Pt(20) # Spacing from title

# Add Bullet to Item 1
pPr1 = p1._p.get_or_add_pPr()
buChar1 = etree.SubElement(pPr1, qn('a:buChar'))
buChar1.set('char', '\u2022') # Standard bullet character
pPr1.set('marL', '114300') # Paragraph left margin
pPr1.set('indent', '304800') # Text indent

# Content Section: Item 2
p2 = tf.add_paragraph()
p2.text = "Natural environment can further be divided into 4 parts."
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0, 0, 0) # Black Font
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(16) # Consistent spacing

# Add Bullet to Item 2
pPr2 = p2._p.get_or_add_pPr()
buChar2 = etree.SubElement(pPr2, qn('a:buChar'))
buChar2.set('char', '\u2022')
pPr2.set('marL', '114300')
pPr2.set('indent', '304800')

# 4. Image Section: Tree inside bubble below texts
# Bubble Dimensions
bubble_size = Inches(2.5)
bubble_left = (slide_width - bubble_size) / 2
bubble_top = box_top + box_height + Inches(0.6) # Positioned below the textbox

# Create Bubble (Circle)
bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, bubble_left, bubble_top, bubble_size, bubble_size)
bubble.fill.solid()
bubble.fill.fore_color.rgb = RGBColor(255, 255, 255) # White Bubble
bubble.line.color.rgb = RGBColor(150, 200, 150) # Subtle border
bubble.line.width = Pt(2)

# Add Image inside Bubble
img_padding = Inches(0.2)
img_shape = slide.shapes.add_picture(
    'image.png', 
    bubble_left + img_padding, 
    bubble_top + img_padding, 
    bubble_size - 2 * img_padding, 
    bubble_size - 2 * img_padding
)

# Save the presentation
prs.save('output.pptx')