from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_slide():
    # Initialize Presentation
    prs = Presentation()
    
    # Get slide dimensions (default to 4:3 usually, but we plan relatively)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Background: Gradient of black and gray
    # Simulating gradient with a dark gray base rectangle
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(40, 40, 40) # Dark Grey/Black
    bg_shape.line.fill.background() # No border

    # 2. Textured image on the right side
    image_path = 'image.png'
    if os.path.exists(image_path):
        # Position on the right 45% of the slide
        img_left = slide_width * 0.55
        img_width = slide_width * 0.45
        try:
            slide.shapes.add_picture(image_path, img_left, 0, img_width, slide_height)
        except Exception:
            pass # Ignore if image cannot be loaded

    # 3. Title: "Now: I, CAN" at the top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Now: I, CAN"
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255) # White text
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 4. White background for bullet points with thin orange border
    # Positioned on the left side to avoid image overlap
    box_left = Inches(1.0)
    box_top = Inches(1.5)
    box_width = Inches(6.5)
    box_height = Inches(4.5)
    
    bullet_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, box_left, box_top, box_width, box_height
    )
    
    # Styling: White fill, Orange border
    bullet_box.fill.solid()
    bullet_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bullet_box.line.color.rgb = RGBColor(255, 165, 0) # Orange
    bullet_box.line.width = Pt(1)

    # Content Coordinates
    text_left = box_left + Inches(0.5)
    text_width = Inches(5.5)

    # Bullet Point 1
    txt_box_1 = slide.shapes.add_textbox(text_left, box_top + Inches(0.8), text_width, Inches(0.6))
    tf1 = txt_box_1.text_frame
    tf1.word_wrap = True
    p1 = tf1.add_paragraph()
    p1.text = "• Define the term Marketing"
    p1.font.size = Pt(24)
    p1.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p1.font.bold = True

    # Arrows: Orange and Yellow separating the bullet points
    arrow_size = Inches(0.4)
    
    # Orange Arrow
    arrow1_top = box_top + Inches(1.9)
    arrow1_left = text_left
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow1_left, arrow1_top, arrow_size, arrow_size)
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(255, 165, 0) # Orange
    arrow1.line.fill.background()

    # Yellow Arrow (stacked slightly or next to it)
    arrow2_top = box_top + Inches(2.3)
    arrow2_left = text_left + Inches(0.2) 
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow2_left, arrow2_top, arrow_size, arrow_size)
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    arrow2.line.fill.background()

    # Bullet Point 2
    txt_box_2 = slide.shapes.add_textbox(text_left, box_top + Inches(2.9), text_width, Inches(1.2))
    tf2 = txt_box_2.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = "• Explain various things involved how products reach to market"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.font.bold = True

    # Save to output.pptx
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()