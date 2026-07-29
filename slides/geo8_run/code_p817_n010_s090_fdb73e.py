from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Helper to get width/height in inches safely
# Although prs.slide_width returns EMU (int), we convert to float inches for layout logic
EMU = 914400 
slide_w_inch = float(prs.slide_width) / EMU
slide_h_inch = float(prs.slide_height) / EMU

# Left side width
left_w = Inches(slide_w_inch / 2)
full_h = Inches(slide_h_inch)
full_w = Inches(slide_w_inch)

# Blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Colors
COLOR_BLUE_BG = RGBColor(30, 80, 160)
COLOR_LIGHT_BLUE = RGBColor(210, 235, 250)
COLOR_LIST = [
    RGBColor(70, 130, 180),   # Steel Blue
    RGBColor(60, 139, 113),   # Medium Turquoise
    RGBColor(30, 144, 255),   # Dodger Blue
    RGBColor(50, 130, 80),    # Medium Sea Green
    RGBColor(0, 100, 148)     # Cadet Blue
]

# 1. Left Background
# Shape covering left half
left_bg = slide.shapes.add_shape(
    1, # Rectangle
    0, 0, left_w, full_h
)
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = COLOR_BLUE_BG
left_bg.line.fill.background()

# 2. Left Title
# Centered in left half
# Width of text box: 80% of left_w? Or just fixed inches?
# Let's use relative inches for stability, but bound by left_w.
title_tb_w = Inches(slide_w_inch * 0.4) 
title_tb_h = Inches(2.0)
title_tb_l = Inches(slide_w_inch * 0.25) # Centered: (0.5 - 0.4)/2 = 0.1 ? 
# Left half is 0 to 0.5. Width 0.4. Center is 0.1.
title_tb_t = Inches(slide_h_inch * 0.35)

title_box = slide.shapes.add_textbox(title_tb_l, title_tb_t, title_tb_w, title_tb_h)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Principles of graphic design"
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# 3. Right Header
header_h = Inches(2.0)
# Header Textbox (acting as background)
# Spans right half
header_l = Inches(slide_w_inch / 2)
header_tb = slide.shapes.add_textbox(header_l, 0, left_w, header_h) # width same as left_w (half)
# Actually right half width is left_w.

tf = header_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The principles of graphic design related to the areas are,"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.LEFT

# Style header textbox background
header_tb.fill.solid()
header_tb.fill.fore_color.rgb = COLOR_LIGHT_BLUE
header_tb.line.color.rgb = RGBColor(150, 200, 255)
header_tb.line.width = Pt(0.5)

# Internal margins for text
header_tb.left = header_l + Inches(0.2)
header_tb.top = Inches(0.2)
header_tb.width = left_w - Inches(0.4)
header_tb.height = header_h - Inches(0.4)

# 4. List Items
list_items = ["Arrangement", "Proximity", "Repetition", "Contrast", "Balance"]
list_top = header_h
list_height = full_h - header_h
item_count = len(list_items)
item_h = list_height / item_count
margin_x = Inches(0.2)
item_w = left_w - (2 * margin_x) # width of list item area

for i, text in enumerate(list_items):
    color = COLOR_LIST[i]
    top_pos = list_top + (i * item_h) + Inches(0.1)
    
    # Create Text Box
    # x = header_l (start of right side)
    txBox = slide.shapes.add_textbox(
        header_l + margin_x,
        top_pos,
        item_w,
        item_h - Inches(0.2)
    )
    
    # Fill
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = color
    txBox.line.fill.background()
    
    # Text
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

prs.save('output.pptx')