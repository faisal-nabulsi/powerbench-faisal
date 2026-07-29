from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Get a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background to dark gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33) # Dark Gray

# 1. Add three circular images (using placeholders as rectangles due to tool limits)
# Note: System instruction requires using slide.shapes.add_picture('image.png', ...)
# Image 1: Fruits and vegetables - Top Left
left = Inches(0.5)
top = Inches(0.5)
width = Inches(3)
height = Inches(3)
slide.shapes.add_picture('image.png', left, top, width, height)

# Image 2: Sugar - Bottom Left
left = Inches(0.5)
top = Inches(4.0)
slide.shapes.add_picture('image.png', left, top, width, height)

# Image 3: Hanging clothes - To the right (Top Right)
left = Inches(9.0)
top = Inches(0.5)
slide.shapes.add_picture('image.png', left, top, width, height)

# 2. Center the text on the right side
# Positioning text box in the middle-right area, below the clothes image
text_left = Inches(6.0)
text_top = Inches(3.5)
text_width = Inches(6.5)
text_height = Inches(1.5)
textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
tf = textbox.text_frame
tf.text = "From where you would like to buy the following products:"

# Font styling
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
p.font.name = 'Arial'
p.alignment = PP_ALIGN.CENTER

# 3. Layout check (balanced)
# Images on left and top-right. Text centered right below image.
# Adjust text position slightly to ensure center alignment on the right side visually
# If text box starts at 6.0 and is 6.5 wide, it ends at 12.5.
# Slide width is 13.33. Margin right is ~0.8.
# Left margin of text box is 6.0.
# Center of text box is 9.25.
# Slide center is 6.66.
# It is on the right side.

# Save presentation
prs.save('output.pptx')