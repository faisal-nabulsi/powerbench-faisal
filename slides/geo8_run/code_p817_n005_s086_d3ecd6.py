from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()

# Select a blank layout to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define Colors
DARK_BLUE = RGBColor(0, 51, 102)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# 1. Set Slide Background to White (for the rest of the page)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# 2. Create Dark Blue Header Background for the Title
# We add a rectangle covering the top of the slide.
# Using prs.slide_width ensures it covers the full width regardless of aspect ratio.
header_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(0), Inches(0), 
    prs.slide_width, Inches(1.5)
)
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = DARK_BLUE
# Remove the border line
header_shape.line.fill.background()

# 3. Add Title Text (White, Large, Bold)
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.3), 
    Inches(9), Inches(1.0)
)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "5 Stages of Development of Media"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

# 4. Add Subtitle "3. PRINTING PRESS" (Dark Blue, Bold)
# Positioned below the header
subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.8), 
    Inches(5), Inches(0.8)
)
tf_sub = subtitle_box.text_frame
tf_sub.text = "3. PRINTING PRESS"
p_sub = tf_sub.paragraphs[0]
p_sub.font.size = Pt(28)
p_sub.font.bold = True
p_sub.font.color.rgb = DARK_BLUE

# 5. Add Main Content with Bullet Point (Black)
content_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.8), 
    Inches(4.5), Inches(3.0)
)
tf_con = content_box.text_frame
tf_con.word_wrap = True
p_con = tf_con.paragraphs[0]
# Using Unicode bullet for compatibility across viewers
p_con.text = "\u2022 " + "It allowed the continuous production, reproduction and circulation of print materials."
p_con.font.size = Pt(18)
p_con.font.color.rgb = BLACK
p_con.space_after = Pt(6)

# 6. Add Illustration (Image) on the Right
# Positioning on the right side of the slide
img_left = Inches(5.5)
img_top = Inches(2.0)
img_width = Inches(4.5)
img_height = Inches(3.5)

slide.shapes.add_picture(
    'image.png', 
    left=img_left, 
    top=img_top, 
    width=img_width, 
    height=img_height
)

# Save the presentation
prs.save('output.pptx')