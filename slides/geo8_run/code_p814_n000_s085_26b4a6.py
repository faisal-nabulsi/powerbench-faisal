from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_elon_musk_timeline():
    # Create a presentation object
    prs = Presentation()
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Brief highlights of Elon Musk"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    # Timeline Data
    events = [
        ("1971", "Born in South Africa"),
        ("1983", "Sold his first video game"),
        ("1999", "Sold Zip2, first company to Compaq"),
        ("2002", "eBay acquired PayPal"),
        ("2002", "Founded SpaceX"),
        ("2004", "Join with Tesla"),
        ("2006", "Cofounded Solarcity"),
        ("2013", "Develop the Hyperloop concept"),
        ("2016", "Cofounded Neuralink"),
        ("2021", "Become world’s richest man")
    ]
    
    # Layout parameters
    slide_width = prs.slide_width
    margin = Inches(1)
    timeline_y = Inches(3.5)
    circle_size = Inches(0.7)
    line_thickness = Inches(0.04)
    
    # Draw horizontal line
    line_left = margin
    line_top = timeline_y - (line_thickness / 2)
    line_width = slide_width - (margin * 2)
    
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        line_left, line_top, line_width, line_thickness
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(180, 180, 180)
    line_shape.line.fill.background()
    
    # Calculate spacing
    num_events = len(events)
    step = (slide_width - (margin * 2)) / num_events
    
    for i, (year, desc) in enumerate(events):
        # Calculate center X for the current event
        center_x = margin + (step * i) + (step / 2)
        
        # Add Yellow Circle
        circle_left = center_x - (circle_size / 2)
        circle_top = timeline_y - (circle_size / 2)
        
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left, circle_top, circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(255, 215, 0) # Gold/Yellow
        circle.line.fill.background()
        
        # Add Year Text inside Circle
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        # Add Description Text below
        text_width = step - Inches(0.2)
        text_left = center_x - (text_width / 2)
        text_top = timeline_y + (circle_size / 2) + Inches(0.2)
        
        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, Inches(1.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.alignment = PP_ALIGN.CENTER
        
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_elon_musk_timeline()