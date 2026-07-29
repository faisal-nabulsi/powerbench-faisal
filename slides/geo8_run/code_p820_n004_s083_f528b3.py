from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Add a blank slide (slide_layouts[6] is usually the blank layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background to solid red
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# Define text box properties
# Slide dimensions are typically 10x7.5 inches.
# The vertical center is at 3.75 inches.
# We want the text centered vertically.
# Font size 44pt is approx 0.61 inches.
# To center the text line, the top of the box should be at Center - (Height/2).
# 3.75 - 0.305 = 3.445 inches.

left = Inches(1)
top = Inches(3.45)
width = Inches(8)
height = Inches(1)

# Add the text box
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

# Set the text content
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "Where do you buy your products from?"

# Format the text
run = p.runs[0]
run.font.size = Pt(44)
run.font.color.rgb = RGBColor(255, 255, 255) # White text
run.font.name = 'Calibri' # Clean, modern font
run.font.bold = True # Make it bold for readability

# Save the presentation to 'output.pptx'
prs.save('output.pptx')