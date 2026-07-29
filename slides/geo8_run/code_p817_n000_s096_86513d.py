from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

def create_slide():
    # Create a presentation object
    prs = Presentation()

    # Use a blank slide layout (index 6) to have full control
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 20) # Dark Gray/Black

    # Define Colors
    COLOR_TITLE = RGBColor(255, 255, 255)       # White
    COLOR_HEADING = RGBColor(255, 215, 0)       # Gold/Yellow for contrast
    COLOR_TEXT = RGBColor(240, 240, 240)        # Light Gray/White for body text

    # Add Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT

    # Add Content Text Box
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Helper function to set bullet character via XML
    def set_bullet(paragraph, char):
        pPr = paragraph._p.get_or_add_pPr()
        buChar = etree.Element(qn('a:buChar'))
        buChar.set('char', char)
        pPr.append(buChar)

    # Helper function to add formatted paragraph
    def add_paragraph(text_frame, text, level, color, bold=False, size=Pt(18)):
        # Use the first paragraph if it's empty, otherwise add a new one
        if len(text_frame.paragraphs) == 1 and text_frame.paragraphs[0].text == '':
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.level = level
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        
        # Set bullet character based on level
        if level == 0:
            set_bullet(p, '•')
        else:
            set_bullet(p, '–')

    # Content Data
    content = [
        {
            "heading": "Financial Literacy",
            "points": ["Skills include budgeting, saving, investing, and making informed financial decisions."]
        },
        {
            "heading": "Leadership and Teamwork",
            "points": [
                "Leadership skills involve motivating and guiding others toward common goals.",
                "Teamwork skills enable collaboration and the ability to work effectively in diverse groups."
            ]
        },
        {
            "heading": "Time Management",
            "points": ["Effective time management and organizational skills improve productivity and reduce stress."]
        },
        {
            "heading": "Creativity and Innovation",
            "points": ["These skills lead to the development of new ideas, products, and solutions."]
        },
        {
            "heading": "Cultural Competence and Diversity",
            "points": [
                "In an increasingly globalized world, understanding and respecting cultural differences is essential.",
                "Cultural competence promotes inclusivity and effective interactions in multicultural environments."
            ]
        }
    ]

    # Populate the slide
    for item in content:
        # Add Heading (Level 0) - Yellow and Bold
        add_paragraph(
            tf, 
            item["heading"], 
            level=0, 
            color=COLOR_HEADING, 
            bold=True, 
            size=Pt(22)
        )
        
        # Add Sub-points (Level 1) - White/Light Gray
        for point in item["points"]:
            add_paragraph(
                tf, 
                point, 
                level=1, 
                color=COLOR_TEXT, 
                bold=False, 
                size=Pt(16)
            )

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()