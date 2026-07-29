from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a Presentation object
prs = Presentation()

# Add a blank slide (Index 6 is typically the Blank layout in default Office templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Set Dark Background ---
# We create a full-slide rectangle to serve as the dark background
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, 0, 
    prs.slide_width, 
    prs.slide_height
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = RGBColor(30, 30, 30) # Dark Gray/Black background
bg_shape.line.fill.background() # No border line

# --- 2. Add Title ---
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1.2))
tf_title = title_box.text_frame
tf_title.word_wrap = True

p_title = tf_title.paragraphs[0]
p_title.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(255, 255, 255) # White text for title
p_title.alignment = PP_ALIGN.LEFT

# --- 3. Define Content Data ---
# Sections with headers (Yellow) and bullet points (Light Gray)
content_data = [
    {
        "header": "Health and Wellness",
        "points": [
            "Skills related to physical and mental health, such as stress management, nutrition, and self-care, contribute to overall well-being and a higher quality of life."
        ]
    },
    {
        "header": "Digital Literacy",
        "points": [
            "With the digitalization of many aspects of life, digital literacy skills are essential for navigating technology, staying safe online, and using digital tools for various purposes."
        ]
    },
    {
        "header": "Networking and Relationship Building",
        "points": [
            "Building a strong professional and social network is vital for career advancement and personal growth.",
            "Effective networking skills can open doors to opportunities and resources."
        ]
    },
    {
        "header": "Conflict Resolution",
        "points": [
            "Conflict is a natural part of life, and the ability to resolve conflicts constructively is valuable in both personal and professional relationships."
        ]
    }
]

# Define Colors
COLOR_YELLOW = RGBColor(255, 193, 7)   # Warm Yellow
COLOR_GRAY = RGBColor(210, 210, 210)    # Light Gray for body text

# --- 4. Add Content Text Box ---
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(5.5))
tf_content = content_box.text_frame
tf_content.word_wrap = True

# Loop through data to create formatted paragraphs
for i, section in enumerate(content_data):
    # Add Section Header (Subtitle)
    p_header = tf_content.add_paragraph()
    # Add spacing before the header, except for the first item
    p_header.space_before = Pt(16) if i > 0 else 0
    p_header.space_after = Pt(4) # Small space after header before bullets
    
    run_header = p_header.add_run()
    run_header.text = section["header"]
    run_header.font.size = Pt(18)
    run_header.font.bold = True
    run_header.font.color.rgb = COLOR_YELLOW

    # Add Bullet Points
    for point in section["points"]:
        p_bullet = tf_content.add_paragraph()
        p_bullet.level = 1 # Indent to simulate bullet level
        
        run_bullet = p_bullet.add_run()
        # We prepend a bullet character "• " manually to ensure visibility 
        # as blank layouts sometimes lack default bullet styles.
        run_bullet.text = "  • " + point
        run_bullet.font.size = Pt(14)
        run_bullet.font.color.rgb = COLOR_GRAY
        run_bullet.font.bold = False

# Save the presentation
prs.save('output.pptx')