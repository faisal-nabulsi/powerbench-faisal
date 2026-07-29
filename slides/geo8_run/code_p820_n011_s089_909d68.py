from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create the presentation
prs = Presentation()

# Select a blank layout to allow custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Set background to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Slide dimensions for calculations
slide_width = prs.slide_width
slide_height = prs.slide_height

# 2. Add the title: "TOP GRADES ARE NOT THE SOLE DETERMINANT OF SUCCESS"
# Positioning at the top
title_left = Inches(1)
title_top = Inches(0.5)
title_width = Inches(8)
title_height = Inches(1.5)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_run = title_para.add_run()
title_run.text = "TOP GRADES ARE NOT THE SOLE DETERMINANT OF SUCCESS"
title_run.font.size = Pt(36) # Large font
title_run.font.color.rgb = RGBColor(211, 211, 211) # Light Gray
title_run.font.bold = True
title_run.font.name = 'Arial'

# 3. Add the main content area
# Positioning below the title
content_left = Inches(1)
content_top = Inches(2.5)
content_width = Inches(8)
content_height = Inches(3)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

content_para = content_frame.paragraphs[0]
content_para.alignment = PP_ALIGN.CENTER
content_para.font.size = Pt(24)
content_para.font.name = 'Arial'
# Default color for standard text (White)
content_para.font.color.rgb = RGBColor(255, 255, 255)

# Run 1: "Success" in bold, yellow font
run1 = content_para.add_run()
run1.text = "Success"
run1.font.bold = True
run1.font.color.rgb = RGBColor(255, 255, 0) # Yellow

# Run 2: " is multifaceted and depends on a variety of factors, including " in standard font
run2 = content_para.add_run()
run2.text = " is multifaceted and depends on a variety of factors, including "
run2.font.color.rgb = RGBColor(255, 255, 255) # White

# Run 3: "Skills," in bold, green font
run3 = content_para.add_run()
run3.text = "Skills,"
run3.font.bold = True
run3.font.color.rgb = RGBColor(0, 200, 0) # Green

# Run 4: " Experiences," in bold, blue font (Space added at start for separation)
run4 = content_para.add_run()
run4.text = " Experiences,"
run4.font.bold = True
run4.font.color.rgb = RGBColor(0, 100, 255) # Blue

# Run 5: " and Personal Attributes" in bold, red font (Space added at start for separation)
run5 = content_para.add_run()
run5.text = " and Personal Attributes"
run5.font.bold = True
run5.font.color.rgb = RGBColor(255, 50, 50) # Red

# 4. Include an image of a trophy and a dartboard on the bottom-right corner
# Using the placeholder image.png provided in the working directory
image_width = Inches(2.5)
image_height = Inches(2.5)

# Calculate position for bottom-right corner (with some padding)
image_left = slide_width - image_width - Inches(0.5)
image_top = slide_height - image_height - Inches(0.5)

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# Save the presentation
prs.save('output.pptx')
print("Presentation saved successfully as output.pptx")