from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add the image as a background
# The prompt asks for an image as main content but also mentions text standing out against a background image.
# Placing the image to fill the slide ensures the best visual balance and readability.
image_path = 'image.png'
slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)

# Helper function to add centered white text
def add_centered_text(slide, left, top, width, height, text, font_size, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255) # White color
    run.font.name = 'Arial'
    
    return txBox

# 1. Main Title
# "Airbnb is all about people and not about the places at all."
title_text = "Airbnb is all about people and not about the places at all."
# Position: Top center
add_centered_text(slide, Inches(1), Inches(0.5), Inches(11.33), Inches(1), title_text, 36, bold=True)

# 2. Subtitle
# "The passion of its founders..."
subtitle_text = "The passion of its founders and the love and generosity of people around the world was changing the culture of travel with places and personal experiences unlike any other. Belonging the world over was Airbnb’s truth and its differentiator."
# Position: Below title
add_centered_text(slide, Inches(1), Inches(1.5), Inches(11.33), Inches(1.5), subtitle_text, 18, bold=False)

# 3. Quote
# "IT'S ABOUT THE PEOPLE, NOT THE PLACES"
quote_text = "IT'S ABOUT THE PEOPLE, NOT THE PLACES"
# Position: Center of the slide
quote_width = Inches(10)
quote_height = Inches(1.5)
quote_left = (slide_width - quote_width) / 2
quote_top = (slide_height - quote_height) / 2

add_centered_text(slide, quote_left, quote_top, quote_width, quote_height, quote_text, 48, bold=True)

# Save the presentation
prs.save('output.pptx')