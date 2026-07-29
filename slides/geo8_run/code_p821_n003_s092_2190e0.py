from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Create a Presentation object
prs = Presentation()
# Add a blank slide layout (Index 6 is typically the Blank layout in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Set Background Color to Dark (Dark Grey)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(40, 40, 40)

# 2. Add Main Title at the Top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
tf = title_box.text_frame
tf.paragraphs[0].text = "KEY reasons why academic success is important in society"
tf.paragraphs[0].font.size = Pt(36)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # White text
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Constants for styling and layout
YELLOW_COLOR = RGBColor(255, 215, 0)
WHITE_COLOR = RGBColor(240, 240, 240)
SECTION_WIDTH = Inches(2.7)
IMAGE_HEIGHT = Inches(1.5)
Y_START = Inches(1.3) # Start vertical position for sections

# Helper function to add section titles
def add_section_title(slide, left, top, title_text):
    txBox = slide.shapes.add_textbox(left, top, SECTION_WIDTH, Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.color.rgb = YELLOW_COLOR # Yellow color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return txBox

# Helper function to add bullet points
def add_bullets(slide, left, top, bullets):
    txBox = slide.shapes.add_textbox(left, top, SECTION_WIDTH, Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.LEFT
        # Add bullet character
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')
    return txBox

# Helper function to add images
def add_section_image(slide, left, top):
    # Using the placeholder 'image.png' as requested
    slide.shapes.add_picture('image.png', left, top, SECTION_WIDTH, IMAGE_HEIGHT)

# ---------------------------------------------------------
# Section 1: Left
# ---------------------------------------------------------
x1 = Inches(0.5)
add_section_title(slide, x1, Y_START, "Personal Growth and Development")
add_section_image(slide, x1, Y_START + Inches(0.6))
add_bullets(slide, x1, Y_START + Inches(2.4), [
    "Academic success promotes personal growth through critical thinking, problem-solving, and intellectual curiosity.",
    "It leads to fulfilling careers and a sense of purpose."
])

# ---------------------------------------------------------
# Section 2: Middle
# ---------------------------------------------------------
x2 = x1 + SECTION_WIDTH + Inches(0.2) # Offset for spacing
add_section_title(slide, x2, Y_START, "Employability and Career Opportunities")
add_section_image(slide, x2, Y_START + Inches(0.6))
add_bullets(slide, x2, Y_START + Inches(2.4), [
    "Academic foundation is key for many careers.",
    "Achievements can lead to better jobs, earnings, and advancement."
])

# ---------------------------------------------------------
# Section 3: Right
# ---------------------------------------------------------
x3 = x2 + SECTION_WIDTH + Inches(0.2) # Offset for spacing
add_section_title(slide, x3, Y_START, "Economic Impact")
add_section_image(slide, x3, Y_START + Inches(0.6))

# The third section requested text (not explicitly bullets, but text block)
txBox = slide.shapes.add_textbox(x3, Y_START + Inches(2.4), SECTION_WIDTH, Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Well-educated societies lead to stronger economies by creating skilled workers and reducing unemployment."
p.font.size = Pt(16)
p.font.color.rgb = WHITE_COLOR
p.alignment = PP_ALIGN.LEFT

# Save the presentation to 'output.pptx'
prs.save('output.pptx')