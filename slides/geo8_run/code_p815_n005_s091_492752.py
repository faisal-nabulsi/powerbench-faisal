from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Create a presentation object
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use a blank layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set light background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 1. Title "PART 2" at top left
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(2)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PART 2"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # 2. Chat Icon (Image) below title
    img_left = Inches(0.5)
    img_top = Inches(1.2)
    img_width = Inches(1.0)
    img_height = Inches(1.0)
    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
    
    # 3. Instruction "Discuss the questions." right to icon
    inst_left = Inches(1.8)
    inst_top = Inches(1.3)
    inst_width = Inches(5)
    inst_height = Inches(0.8)
    txBox = slide.shapes.add_textbox(inst_left, inst_top, inst_width, inst_height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Discuss the questions."
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # 4. Dialogue Text
    dial_left = Inches(1.0)
    dial_top = Inches(2.5)
    dial_width = Inches(11.0)
    dial_height = Inches(3.0)
    txBox = slide.shapes.add_textbox(dial_left, dial_top, dial_width, dial_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Helper function to add dialogue lines with bold names
    def add_dialogue(tf, name, text):
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        
        # Opening quote
        run1 = p.add_run()
        run1.text = '"'
        run1.font.bold = False
        
        # Name (Bold)
        run2 = p.add_run()
        run2.text = f'{name}: '
        run2.font.bold = True
        
        # Dialogue text
        run3 = p.add_run()
        run3.text = text
        run3.font.bold = False
        
        # Closing quote
        run4 = p.add_run()
        run4.text = '"'
        run4.font.bold = False

    # Joey
    add_dialogue(tf, "JOEY", "We were out to dinner. We were getting along, having a really nice time, I was thinking she was really cool and then, out of nowhere, (she reached over and took some of my fries from my plate!)")
    
    # Phoebe
    add_dialogue(tf, "PHOEBE", "So she took some fries, big deal!")
    
    # Rachel
    add_dialogue(tf, "RACHEL", "Oh yeah, Joey doesn't share food. I mean, just last week, we were having breakfast, and...and he had a couple of grapes on his plate...")

    # 5. Two numbered questions in textboxes at the bottom
    q1_left = Inches(1.0)
    q1_top = Inches(5.8)
    q1_width = Inches(11.0)
    q1_height = Inches(0.5)
    
    txBox1 = slide.shapes.add_textbox(q1_left, q1_top, q1_width, q1_height)
    tf1 = txBox1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "1. When is the past continuous used in comparison with the past simple?"
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(0, 0, 0)
    
    q2_left = Inches(1.0)
    q2_top = Inches(6.3)
    q2_width = Inches(11.0)
    q2_height = Inches(0.5)
    
    txBox2 = slide.shapes.add_textbox(q2_left, q2_top, q2_width, q2_height)
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2. How is the structure for the past simple different than the past continuous?"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0, 0, 0)

    # 6. Button "VIEWING FOLLOW-UP" at bottom right
    btn_left = Inches(10.5)
    btn_top = Inches(6.5)
    btn_width = Inches(2.5)
    btn_height = Inches(0.6)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, btn_left, btn_top, btn_width, btn_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192) # Blue
    shape.line.color.rgb = RGBColor(0, 0, 0)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "VIEWING FOLLOW-UP"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()