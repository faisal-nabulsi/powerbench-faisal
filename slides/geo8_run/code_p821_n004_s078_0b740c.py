import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Initialize the presentation
    prs = Presentation()

    # Get slide dimensions (Standard 13.33 x 7.5 inches)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Create a blank slide layout (Index 6 is typically Blank in standard theme)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Insert the colorful abstract background image
    # The prompt specifies using 'image.png' as the placeholder for images
    slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)

    # 2. Add a white layer behind texts and images to ensure visibility
    # We define a large white rectangle with margins
    margin = Inches(0.5)
    box_width = slide_width - (margin * 2)
    box_height = slide_height - (margin * 2)

    white_layer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, margin, box_width, box_height)
    white_layer.fill.solid()
    white_layer.fill.fore_color.rgb = RGBColor(255, 255, 255)
    white_layer.line.fill.background() # Remove border

    # 3. Sub-title: "The Subject and Content of Art"
    # Location: Top-left corner (inside the white layer bounds)
    sub_left = Inches(0.7)
    sub_top = Inches(0.7)
    sub_width = Inches(6.0)
    sub_height = Inches(0.5)

    sub_box = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Subject and Content of Art"
    p.font.size = Pt(18)
    p.font.bold = True
    # Contrasting color: Dark Blue against the white background
    p.font.color.rgb = RGBColor(0, 0, 139) 
    p.alignment = PP_ALIGN.LEFT

    # 4. Main Title: "D. Mythology and religion, dreams and fantasies."
    # Location: Below subtitle, clearly visible
    main_left = Inches(1.0)
    main_top = Inches(1.4)
    main_width = Inches(11.0)
    main_height = Inches(0.8)

    main_box = slide.shapes.add_textbox(main_left, main_top, main_width, main_height)
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "D. Mythology and religion, dreams and fantasies."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    # 5. Insert two images side by side
    # Using 'image.png' for both placeholders (Botticelli and Baroque)
    img_width = Inches(4.0)
    img_height = Inches(3.0)
    img_top = Inches(3.0)

    # Calculate positioning to ensure even spacing and alignment
    # Left image X position (approx 2.25 inches from left)
    left_img_x = Inches(2.25)
    # Right image X position (Left X + Width + Spacing)
    right_img_x = left_img_x + img_width + Inches(1.25)

    # Left Image: Iconic portrayal of beauty/mythology (Botticelli)
    slide.shapes.add_picture('image.png', left_img_x, img_top, img_width, img_height)

    # Right Image: Dramatic baroque portrayal (Heroism)
    slide.shapes.add_picture('image.png', right_img_x, img_top, img_width, img_height)

    # 6. Save the presentation
    prs.save('output.pptx')

if __name__ == '__main__':
    create_slide()