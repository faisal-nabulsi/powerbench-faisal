from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Use blank layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- 1. Background Image ---
# Add image.png as background
# Using 'image.png' as the colorful abstract background placeholder
try:
    bg_picture = slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)
    # Ensure it is at the very back (index 0)
    # Although adding first usually puts it at back, explicit move is safer.
    slide.shapes._spTree.remove(bg_picture._element)
    slide.shapes._spTree.insert(0, bg_picture._element)
except Exception:
    pass

# --- 2. White Overlay ---
# Add a white rectangle covering the slide to ensure text readability
# Note: Standard shapes are opaque, so this covers the background image.
overlay_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
)
overlay_shape.fill.solid()
overlay_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
overlay_shape.line.fill.background() # Remove border

# --- 3. Text Content ---
# Main Title: "The Subject and Content of Art"
title_box = slide.shapes.add_textbox(
    Inches(1), # Left margin
    Inches(0.8), # Top margin
    slide_width - Inches(2), # Width
    Inches(1.5) # Height
)
tf = title_box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "The Subject and Content of Art"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0) # Black text
p.alignment = PP_ALIGN.CENTER

# Subtitle: "Animals"
# Centered below the main title
subtitle_box = slide.shapes.add_textbox(
    Inches(1), 
    Inches(2.5), # Spacing below title
    slide_width - Inches(2), 
    Inches(1)
)
tf2 = subtitle_box.text_frame
tf2.clear()
p2 = tf2.paragraphs[0]
p2.text = "Animals"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0, 0, 0) # Black text
p2.alignment = PP_ALIGN.CENTER

# --- 4. Images ---
# Two images side by side: Left (Rabbit), Right (Animals with person)
# Using 'image.png' for both as per placeholder constraint.
# Formatting: Equal size, white border.

# Dimensions
img_height = Inches(3.5)
img_width = Inches(3.5)
gap = Inches(0.5)

# Calculate starting X position to center the pair
total_imgs_width = (img_width * 2) + gap
start_x = (slide_width - total_imgs_width) / 2
start_y = Inches(4.0) # Position below subtitle

# Left Image
pic_left = slide.shapes.add_picture('image.png', start_x, start_y, img_width, img_height)
# Add white border
pic_left.line.color.rgb = RGBColor(255, 255, 255)
pic_left.line.width = Inches(0.15)
pic_left.line.fill.solid()

# Right Image
pic_right = slide.shapes.add_picture('image.png', 
                                     start_x + img_width + gap, 
                                     start_y, 
                                     img_width, 
                                     img_height)
# Add white border
pic_right.line.color.rgb = RGBColor(255, 255, 255)
pic_right.line.width = Inches(0.15)
pic_right.line.fill.solid()

# Save Presentation
prs.save('output.pptx')