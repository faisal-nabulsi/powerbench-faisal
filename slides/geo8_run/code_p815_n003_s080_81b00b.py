from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    
    # Add a blank slide (layout index 6 is typically blank in standard templates)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set the background to a dark color (simulating dark texture for contrast)
    # A solid dark grey provides high contrast for the white/yellow text.
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark Gray

    # --- Main Title ---
    # Add the main title at the top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_title = tf_title.paragraphs[0]
    p_title.text = "ACADEMIC WORLD AND PROFESSIONAL WORLD: DEMANDS AND CHARACTERISTICS"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)  # White
    p_title.alignment = PP_ALIGN.CENTER

    # --- Left Section: Academic World ---
    # Add textbox for the left section
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
    ac_tf = left_box.text_frame
    ac_tf.clear()
    ac_tf.word_wrap = True

    # Section Title: ACADEMIC WORLD (Blue)
    p_ac_title = ac_tf.add_paragraph()
    p_ac_title.text = "ACADEMIC WORLD"
    p_ac_title.font.size = Pt(22)
    p_ac_title.font.bold = True
    p_ac_title.font.color.rgb = RGBColor(0, 112, 192) # Standard PowerPoint Blue

    # Subheading: Outcome (Bold Yellow)
    p_ac_subhead = ac_tf.add_paragraph()
    p_ac_subhead.text = "Outcome"
    p_ac_subhead.font.size = Pt(18)
    p_ac_subhead.font.bold = True
    p_ac_subhead.font.color.rgb = RGBColor(255, 255, 0) # Yellow
    p_ac_subhead.space_before = Pt(6)

    # Subsection: Degrees and Certifications
    p_ac_sub = ac_tf.add_paragraph()
    p_ac_sub.text = "Degrees and Certifications"
    p_ac_sub.font.size = Pt(14)
    p_ac_sub.font.bold = True
    p_ac_sub.font.color.rgb = RGBColor(255, 255, 255) # White
    p_ac_sub.space_before = Pt(10)

    # Bullet Point
    p_ac_bullet = ac_tf.add_paragraph()
    p_ac_bullet.text = "\u2022 Academic success leads to degrees and certifications (e.g., diplomas, degrees, Ph.Ds.), which serve as qualifications for future opportunities."
    p_ac_bullet.font.size = Pt(12)
    p_ac_bullet.font.color.rgb = RGBColor(220, 220, 220) # Light Gray
    p_ac_bullet.level = 1
    p_ac_bullet.space_before = Pt(4)

    # --- Right Section: Professional World ---
    # Add textbox for the right section
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(5))
    prof_tf = right_box.text_frame
    prof_tf.clear()
    prof_tf.word_wrap = True

    # Section Title: PROFESSIONAL WORLD (Green)
    p_prof_title = prof_tf.add_paragraph()
    p_prof_title.text = "PROFESSIONAL WORLD"
    p_prof_title.font.size = Pt(22)
    p_prof_title.font.bold = True
    p_prof_title.font.color.rgb = RGBColor(0, 176, 80) # Standard PowerPoint Green

    # Subheading: Outcome (Bold Yellow)
    p_prof_subhead = prof_tf.add_paragraph()
    p_prof_subhead.text = "Outcome"
    p_prof_subhead.font.size = Pt(18)
    p_prof_subhead.font.bold = True
    p_prof_subhead.font.color.rgb = RGBColor(255, 255, 0) # Yellow
    p_prof_subhead.space_before = Pt(6)

    # Subsection: Achievement and Success
    p_prof_sub1 = prof_tf.add_paragraph()
    p_prof_sub1.text = "Achievement and Success"
    p_prof_sub1.font.size = Pt(14)
    p_prof_sub1.font.bold = True
    p_prof_sub1.font.color.rgb = RGBColor(255, 255, 255) # White
    p_prof_sub1.space_before = Pt(10)

    # Bullet 1
    p_prof_bullet1 = prof_tf.add_paragraph()
    p_prof_bullet1.text = "\u2022 Professional success is measured by job performance, career advancement, and the impact of one's work on organizations and clients."
    p_prof_bullet1.font.size = Pt(12)
    p_prof_bullet1.font.color.rgb = RGBColor(220, 220, 220) # Light Gray
    p_prof_bullet1.level = 1
    p_prof_bullet1.space_before = Pt(4)

    # Subsection: Compensation
    p_prof_sub2 = prof_tf.add_paragraph()
    p_prof_sub2.text = "Compensation"
    p_prof_sub2.font.size = Pt(14)
    p_prof_sub2.font.bold = True
    p_prof_sub2.font.color.rgb = RGBColor(255, 255, 255) # White
    p_prof_sub2.space_before = Pt(10)

    # Bullet 2
    p_prof_bullet2 = prof_tf.add_paragraph()
    p_prof_bullet2.text = "\u2022 Compensation and benefits are typically tied to professional success and experience."
    p_prof_bullet2.font.size = Pt(12)
    p_prof_bullet2.font.color.rgb = RGBColor(220, 220, 220) # Light Gray
    p_prof_bullet2.level = 1
    p_prof_bullet2.space_before = Pt(4)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == '__main__':
    create_presentation()