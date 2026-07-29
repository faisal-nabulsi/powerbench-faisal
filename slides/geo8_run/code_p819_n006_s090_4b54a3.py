from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Initialize presentation
    prs = Presentation()

    # Determine slide dimensions (adapts to default presentation size)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Try to find a Blank layout, fallback to index 6 (standard blank)
    slide_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            slide_layout = layout
            break
    if slide_layout is None:
        slide_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(slide_layout)

    # 1. Background: White
    # Create a full-size rectangle to serve as a white background.
    # Added first so it sits at the bottom of the object stack.
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg_shape.line.fill.background() # Remove border

    # 2. Vertical Yellow Bar on the Left
    bar_width = Inches(0.75)
    yellow_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, bar_width, slide_height)
    yellow_bar.fill.solid()
    yellow_bar.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
    yellow_bar.line.fill.background() # Remove border

    # 3. Title: "Content" with bold font
    # Positioned at top, centered horizontally
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_width - Inches(1.0), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.text = "Content"
    title_para = title_tf.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.bold = True
    title_run.font.size = Pt(36)
    title_run.font.color.rgb = RGBColor(0, 0, 0) # Black

    # 4. Bullet Points
    bullet_items = [
        "Brief highlights of Elon Musk",
        "Biography",
        "Early life of Elon Musk",
        "How he came up with his ideas",
        "Elon’s current stage",
        "Obstacles that Elon faced",
        "Lessons that we can learn from Elon Musk's life",
        "Elon’s Future Plans"
    ]

    # Positioned to the right of the yellow bar to avoid overlap
    content_left = Inches(1.5)
    content_top = Inches(2.0)
    content_width = slide_width - content_left - Inches(0.5)
    content_height = Inches(5.0)

    text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = text_box.text_frame
    tf.word_wrap = True

    # Add bullet points
    for i, item in enumerate(bullet_items):
        if i == 0:
            # Use the first paragraph created with the textbox
            p = tf.paragraphs[0]
        else:
            # Add new paragraphs for subsequent items
            p = tf.add_paragraph()
        
        # Prepend bullet character for visual bullet points
        p.text = "• " + item
        p.alignment = PP_ALIGN.LEFT
        
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0, 0, 0) # Black text

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()