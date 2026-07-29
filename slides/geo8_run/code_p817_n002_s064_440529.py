from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation with a standard slide size
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# --- Dark blue background ---
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 30, 70)  # Dark blue

# --- Light bulb icon on the left side ---
slide.shapes.add_picture(
    'image.png',
    left=Inches(1.0),
    top=Inches(2.5),
    width=Inches(1.8),
    height=Inches(1.8)
)

# --- Centered title text ---
title_box = slide.shapes.add_textbox(
    left=Inches(2.5),
    top=Inches(2.8),
    width=Inches(6.0),
    height=Inches(1.2)
)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "'DATE GONE WRONG' STORY"
p.alignment = PP_ALIGN.LEFT

run = p.runs[0]
run.font.size = Pt(42)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)

# --- FLUENT logo in the top right corner ---
slide.shapes.add_picture(
    'image.png',
    left=Inches(6.5),
    top=Inches(0.4),
    width=Inches(2.2),
    height=Inches(0.9)
)

# Save the presentation
prs.save('output.pptx')