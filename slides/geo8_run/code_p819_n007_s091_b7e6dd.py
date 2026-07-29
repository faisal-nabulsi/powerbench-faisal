from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Use a blank layout to allow for custom background and positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add the background image
# The instruction asks for a colorful abstract image. We use the provided placeholder.
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Add the smaller title on the top-left corner
left_small = Inches(0.5)
top_small = Inches(0.5)
width_small = Inches(6)
height_small = Inches(0.5)

txBox_small = slide.shapes.add_textbox(left_small, top_small, width_small, height_small)
tf_small = txBox_small.text_frame
p_small = tf_small.paragraphs[0]
p_small.text = "The Subject and Content of Art"
p_small.font.size = Pt(14)
p_small.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Set background of the small title textbox to yellow
txBox_small.fill.solid()
txBox_small.fill.fore_color.rgb = RGBColor(255, 255, 0)

# 3. Add the main title "Three levels of meaning"
left_main = Inches(1)
top_main = Inches(1.2)
width_main = Inches(8)
height_main = Inches(0.8)

txBox_main = slide.shapes.add_textbox(left_main, top_main, width_main, height_main)
tf_main = txBox_main.text_frame
p_main = tf_main.paragraphs[0]
p_main.text = "Three levels of meaning"
p_main.font.size = Pt(28)
p_main.font.bold = True
p_main.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Set background of the main title textbox to yellow
txBox_main.fill.solid()
txBox_main.fill.fore_color.rgb = RGBColor(255, 255, 0)

# 4. Add the main content area with a white background
left_content = Inches(1)
top_content = Inches(2.2)
width_content = Inches(8)
height_content = Inches(4.5)

txBox_content = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)

# Set background of the content textbox to white
txBox_content.fill.solid()
txBox_content.fill.fore_color.rgb = RGBColor(255, 255, 255)

tf_content = txBox_content.text_frame
tf_content.word_wrap = True

# Add the first numbered point
p1 = tf_content.paragraphs[0]
p1.text = "1. Factual meaning - the literal statement or narrative content in the work that can be directly apprehended because the objects presented are easily recognized."
p1.font.size = Pt(18)
p1.font.color.rgb = RGBColor(0, 0, 0)

# Add the second numbered point
p2 = tf_content.add_paragraph()
p2.text = "2. Conventional meaning - refers to the special meaning that the certain object or color has for a particular culture or group of people when it is shown in an artwork."
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0, 0, 0)
p2.space_before = Pt(12)  # Add spacing between points

# Save the presentation
prs.save('output.pptx')