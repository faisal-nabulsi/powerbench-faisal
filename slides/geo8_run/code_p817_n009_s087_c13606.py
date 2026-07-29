import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from lxml import etree

def create_slide():
    # 1. Initialize Presentation
    prs = Presentation()
    
    # Set slide dimensions to Widescreen (13.33 x 7.5 inches is standard)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide (Index 6 is typically the 'Blank' layout in default templates)
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        # Fallback to first layout if index out of range
        slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(slide_layout)

    # 2. Background Image
    # Use the placeholder image file 'image.png'
    image_path = 'image.png'
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
    else:
        # If image is missing, the background remains default white
        print(f"Warning: {image_path} not found. Background image skipped.")

    # 3. Dark Overlay
    # Create a black rectangle covering the entire slide to act as a dark overlay
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0) # Black
    overlay.line.fill.background() # No border around the overlay
    
    # Set transparency/opacity using lxml manipulation (50% opacity)
    # python-pptx does not have a direct method for shape fill opacity
    try:
        fill_elem = overlay.fill._xFill
        srgb_clr = fill_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb_clr is not None:
            alpha = etree.SubElement(srgb_clr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '50000') # 50000 = 50% opacity (50% transparent)
    except Exception:
        pass

    # 4. Text Box (Red Cubic Shape)
    # Dimensions for the text box
    box_width = Inches(8)
    box_height = Inches(2.5)
    
    # Center the box on the slide
    left = (prs.slide_width - box_width) / 2
    top = (prs.slide_height - box_height) / 2
    
    # Add the shape (Rectangle for "cubic" box feel)
    text_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
    
    # Style the box
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = RGBColor(255, 0, 0) # Red fill
    
    # Add a border for definition
    text_box.line.color.rgb = RGBColor(255, 255, 255)
    text_box.line.width = Pt(2)

    # 5. Text Content & Formatting
    tf = text_box.text_frame
    tf.word_wrap = True
    
    # Access the paragraph
    p = tf.paragraphs[0]
    p.text = "Which type of market is beneficial for the customers? Why?"
    
    # Apply formatting: Large, Bold, White
    p.font.size = Pt(36) # Large size
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) # White text
    
    # Center text horizontally
    p.alignment = PP_ALIGN.CENTER
    
    # Center text vertically using XML anchor attribute
    try:
        body_pr = tf._xTextBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if body_pr is not None:
            body_pr.set('anchor', 'ctr') # 'ctr' centers text vertically
    except Exception:
        pass

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()