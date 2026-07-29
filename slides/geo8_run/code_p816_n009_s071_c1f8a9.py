from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Setup
prs = Presentation()
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Dimensions
s_width = prs.slide_width
s_height = prs.slide_height
emu_per_inch = 914400

# Helper to get Inches object from float
def inches(val):
    return Inches(val)

# 2. Background (Red)
# Using shape to force background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, s_width, s_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(204, 0, 0) # Red
bg.line.fill.background()

# 3. Black Stripe
stripe_h = 3.0 # inches
stripe_top_val = (s_height - Inches(stripe_h)) / 2
# Arithmetic check: (int - Length) might fail. 
# Let's do it in Emu to be safe or just use Inches logic if supported.
# Actually, pptx.util.Inches returns a Length. 
# Length supports __rsub__ with int? 
# Let's assume yes. If not, fallback.
# Fallback calculation:
stripe_top_emu = (s_height - int(Inches(stripe_h))) // 2
stripe_top = inches(stripe_top_emu / emu_per_inch)

stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, stripe_top, s_width, Inches(stripe_h))
stripe.fill.solid()
stripe.fill.fore_color.rgb = RGBColor(0, 0, 0) # Black
stripe.line.fill.background()

# 4. Title "BIG QUESTION"
# Blue, top of stripe.
# Stripe top is at ~2.25 inches (on 7.5 height slide).
# Place title 0.5 inches above stripe.
title_y = inches( (stripe_top_emu / emu_per_inch) - 0.8 ) # 0.8 inches above? 
# Let's refine. 
# Stripe top approx 2.25. 
# Title box height 1.0.
# If top is 1.5, bottom is 2.5 (overlap).
# If top is 1.2, bottom is 2.2 (just touching).
title_top = inches( (stripe_top_emu / emu_per_inch) - 1.2 )

title_w = 8.0
title_x = inches( (s_width / emu_per_inch - title_w) / 2 )
title_box = slide.shapes.add_textbox(title_x, title_top, inches(title_w), inches(1.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "BIG QUESTION"
p.font.size = Pt(42)
p.font.color.rgb = RGBColor(0, 0, 255) # Blue
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# 5. Main Question
# Centered in stripe.
# Stripe center Y = stripe_top + 1.5 inches.
stripe_center_y_in = (stripe_top_emu / emu_per_inch) + (stripe_h / 2)
q_box_h = 2.0
q_top = inches( stripe_center_y_in - (q_box_h / 2) )
q_w = 9.0
q_x = inches( (s_width / emu_per_inch - q_w) / 2 )

q_box = slide.shapes.add_textbox(q_x, q_top, inches(q_w), inches(q_box_h))
tf_q = q_box.text_frame
tf_q.word_wrap = True
p_q = tf_q.paragraphs[0]
p_q.text = "HOW DO MARKETS HELP US?"
p_q.font.size = Pt(50)
p_q.font.bold = True
p_q.font.color.rgb = RGBColor(255, 255, 255) # White
p_q.alignment = PP_ALIGN.CENTER

# 6. Ribbon
# Top left, Orange (Image placeholder)
r_x = inches(0.5)
r_y = inches(0.5)
r_w = inches(1.2)
r_h = inches(1.2)
try:
    slide.shapes.add_picture('image.png', r_x, r_y, r_w, r_h)
except:
    pass

# Save
prs.save('output.pptx')