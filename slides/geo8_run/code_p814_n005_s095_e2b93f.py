from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation with standard 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Use a blank slide layout to have full control over positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the slide background to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# --- Separator Line ---
# A thin vertical yellow line separating the text and image
# Positioned roughly in the center (x=6.6 inches)
separator = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    left=Inches(6.6),
    top=Inches(1.2),
    width=Inches(0.04),
    height=Inches(5.0)
)
separator.fill.solid()
separator.fill.fore_color.rgb = RGBColor(255, 203, 0) # Yellow accent color
separator.line.fill.background() # No border on the line

# --- Title ---
# Left aligned, Yellow accent text, Sans-serif font
title_box = slide.shapes.add_textbox(
    left=Inches(0.8),
    top=Inches(0.5),
    width=Inches(5.5),
    height=Inches(0.8)
)
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "Elon Musk Current Stage"
p_title.font.size = Pt(36)
p_title.font.bold = True
p_title.font.name = "Arial"
# Using a darker yellow/gold for visibility on white background
p_title.font.color.rgb = RGBColor(255, 180, 0) 
p_title.alignment = PP_ALIGN.LEFT

# --- Bullet Points ---
# Positioned on the left side of the slide
text_box = slide.shapes.add_textbox(
    left=Inches(0.8),
    top=Inches(1.8),
    width=Inches(5.5),
    height=Inches(4.5)
)
tf_text = text_box.text_frame
tf_text.word_wrap = True

bullet_items = [
    "CEO and Chief Engineer at SpaceX",
    "CEO and Product Architect of Tesla",
    "Founder of The Boring Company",
    "Cofounder of Neuralink",
    "Cofounder of OpenAI"
]

for i, item in enumerate(bullet_items):
    if i == 0:
        p = tf_text.paragraphs[0]
    else:
        p = tf_text.add_paragraph()
    
    # Prepend bullet character for visual clarity in blank layout
    p.text = "• " + item
    p.font.size = Pt(22)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p.space_after = Pt(12)

# --- Image ---
# Positioned on the right side of the slide
# Using the placeholder image 'image.png'
image_shape = slide.shapes.add_picture(
    'image.png',
    left=Inches(7.1),
    top=Inches(1.8),
    width=Inches(5.5),
    height=Inches(4.5)
)
# Ensure no border around the image
image_shape.line.fill.background()

# Save the presentation to the current directory
prs.save('output.pptx')