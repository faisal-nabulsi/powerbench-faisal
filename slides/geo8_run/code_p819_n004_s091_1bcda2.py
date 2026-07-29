from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Use a blank slide layout to allow full customization
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Background Image
# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Add the image covering the entire slide
# The image 'image.png' is assumed to be in the current directory
slide.shapes.add_picture('image.png', 0, 0, slide_width, slide_height)

# 2. Top-Left Title
# "The Subject and Content of Art"
top_left_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.5))
tf = top_left_title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The Subject and Content of Art"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 0, 0) # Black text

# 3. Main Title
# "Non-representational or Non-objective Art"
# Yellow background, Black text
main_title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = main_title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Non-representational or Non-objective Art"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

# Apply Yellow Background to Main Title
main_title_box.fill.solid()
main_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)

# 4. Subtitle
# "Sources of the Subject" (Bold)
# Yellow background, Black text
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.8))
tf = subtitle_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Sources of the Subject"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

# Apply Yellow Background to Subtitle
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 0)

# 5. Main Content (Bullet Points)
# White background
content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(3))
tf = content_box.text_frame
tf.word_wrap = True

# Apply White Background to Content Box
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)

# List items
items = [
    "Nature",
    "History",
    "Greek and Roman Mythology",
    "Religion",
    "Sacred Oriental Text"
]

for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    # Add bullet character
    p.text = "• " + item
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.level = 0

# Save the presentation
prs.save('output.pptx')