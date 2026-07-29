from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation with standard 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set dark blue background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 25, 70)

# Add title text: large, bold, white, centered
title_box = slide.shapes.add_textbox(
    left=Inches(0),
    top=Inches(1.8),
    width=Inches(13.333),
    height=Inches(2.2)
)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "VIEWING FOLLOW-UP"
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(64)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)

# Add light bulb icon with brain (left side)
slide.shapes.add_picture(
    'image.png',
    left=Inches(1.5),
    top=Inches(3.8),
    width=Inches(3.5),
    height=Inches(3.5)
)

# Add FLUENT logo (top right corner)
slide.shapes.add_picture(
    'image.png',
    left=Inches(10.8),
    top=Inches(0.6),
    width=Inches(2.2),
    height=Inches(0.9)
)

# Save the presentation
prs.save('output.pptx')