import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

def create_presentation():
    # 1. Create the Presentation object
    prs = Presentation()
    
    # Use a blank slide layout (index 6 is typically blank in default templates)
    try:
        slide_layout = prs.slide_layouts[6]
    except IndexError:
        slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(slide_layout)

    # 2. Set Background: Gradient Teal
    # We manipulate the XML directly to create a linear gradient background.
    # Colors: Teal (#008080) fading to a Darker Teal (#004d40).
    bg = slide.background
    try:
        # Access the underlying XML element
        if bg._element is None:
            # Create background element if missing
            ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            bg_elem = etree.SubElement(slide._element, '{%s}background' % ns_p)
        else:
            bg_elem = bg._element
            
        # Clear existing fill content
        for child in list(bg_elem):
            bg_elem.remove(child)
            
        # Create fillToFill container
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        fill_to_fill = etree.SubElement(bg_elem, '{%s}fillToFill' % ns_p)
        
        # Create Gradient Fill
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        grad_fill = etree.SubElement(fill_to_fill, '{%s}gradFill' % ns_a)
        
        # Set Linear Gradient direction (Top to Bottom, 90 degrees -> ang='5400000')
        lin = etree.SubElement(grad_fill, '{%s}lin' % ns_a, ang='5400000', scaled="1")
        
        # Gradient Stops List
        gs_lst = etree.SubElement(grad_fill, '{%s}gsLst' % ns_a)
        
        # Start Color (Teal)
        gs1 = etree.SubElement(gs_lst, '{%s}gs' % ns_a, pos='0')
        clr1 = etree.SubElement(gs1, '{%s}srgbClr' % ns_a, val='008080')
        
        # End Color (Dark Teal)
        gs2 = etree.SubElement(gs_lst, '{%s}gs' % ns_a, pos='100000')
        clr2 = etree.SubElement(gs2, '{%s}srgbClr' % ns_a, val='004d40')
        
    except Exception:
        # Fallback to solid color if XML manipulation fails
        try:
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0, 128, 128)
        except Exception:
            pass

    # 3. Dimensions and Layout Calculations
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Margins
    left_margin = Inches(0.5)
    right_margin = Inches(0.5)
    top_margin = Inches(0.5)
    bottom_padding = Inches(0.5)
    
    # Vertical Positioning
    title_top = top_margin
    title_height = Inches(0.8)
    
    label_top = title_top + title_height + Inches(0.2)
    label_height = Inches(0.5)
    
    img_top = label_top + label_height + Inches(0.2)
    img_height = slide_height - img_top - bottom_padding
    
    # Horizontal Division: 3 Equal Sections
    usable_width = slide_width - left_margin - right_margin
    col_width = usable_width / 3
    
    # 4. Add Title: "Types of Fixed Retailers"
    title_box = slide.shapes.add_textbox(
        Inches(0), 
        title_top, 
        slide_width, 
        title_height
    )
    tf = title_box.text_frame
    tf.text = "Types of Fixed Retailers"
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255) # White
    run.font.name = 'Calibri'

    # 5. Add Content Sections (Labels + Images)
    section_data = [
        {"label": "General Stores", "col_idx": 0},
        {"label": "Department Store", "col_idx": 1},
        {"label": "Chain Store", "col_idx": 2}
    ]
    
    for section in section_data:
        label_text = section["label"]
        col_idx = section["col_idx"]
        
        # Calculate X position based on column index
        x_start = left_margin + (col_idx * col_width)
        
        # Helper to center content within the column
        # Center of the column
        center_x = x_start + (col_width / 2)
        
        # --- Label ---
        label_width = col_width * 0.8
        label_left = center_x - (label_width / 2)
        
        label_box = slide.shapes.add_textbox(
            label_left,
            label_top,
            label_width,
            label_height
        )
        tf_label = label_box.text_frame
        tf_label.text = label_text
        tf_label.word_wrap = True
        
        p_label = tf_label.paragraphs[0]
        p_label.alignment = PP_ALIGN.CENTER
        run_label = p_label.add_run()
        run_label.font.size = Pt(20) # Slightly smaller than title
        run_label.font.color.rgb = RGBColor(255, 255, 255) # White
        run_label.font.name = 'Calibri'
        # Optional: Make labels bold for visibility, or regular. 
        # Keeping bold for clear contrast on gradient.
        run_label.font.bold = True 
        
        # --- Image ---
        # Using 'image.png' as the placeholder for all store types
        img_width = col_width * 0.9
        img_left = center_x - (img_width / 2)
        
        try:
            slide.shapes.add_picture(
                'image.png', 
                img_left, 
                img_top, 
                img_width, 
                img_height
            )
        except Exception:
            # If image is missing, the script continues
            pass

    # 6. Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()