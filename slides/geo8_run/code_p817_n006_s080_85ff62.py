from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# Initialize the presentation
prs = Presentation()

# Use a blank layout for full control over background elements
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Create a dark blue background for the title area
# Dimensions: Full width of slide, height of 2.3 inches
left = Inches(0)
top = Inches(0)
width = prs.slide_width
height = Inches(2.3)

# Add a rectangle shape
title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x80) # Dark Blue
title_bg.line.fill.background() # Remove border

# 2. Add the Title Text
# Positioned inside the blue bar
title_left = Inches(1.0)
title_top = Inches(0.3)
title_width = Inches(11.3)
title_height = Inches(0.8)

txBoxTitle = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tfTitle = txBoxTitle.text_frame
tfTitle.clear()
pTitle = tfTitle.paragraphs[0]
pTitle.text = "Global and Local Cultural Products"
pTitle.font.size = Pt(34)
pTitle.font.bold = True
pTitle.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text
pTitle.alignment = PP_ALIGN.CENTER

# 3. Add the Subtitle Text
# Positioned below the title, still within the blue area
sub_left = Inches(1.0)
sub_top = Inches(1.2)
sub_width = Inches(11.3)
sub_height = Inches(0.5)

txBoxSub = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
tfSub = txBoxSub.text_frame
tfSub.clear()
pSub = tfSub.paragraphs[0]
pSub.text = "Global Product"
pSub.font.size = Pt(24)
pSub.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text
pSub.alignment = PP_ALIGN.CENTER

# 4. Add Main Content with Bullet Point
# Positioned below the title area (white background by default)
content_left = Inches(1.0)
content_top = Inches(2.8)
content_width = Inches(11.3)
content_height = Inches(1.5)

txBoxContent = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
tfContent = txBoxContent.text_frame
tfContent.clear()
pContent = tfContent.paragraphs[0]
content_text = "Those products that are marketed internationally under the same brand name, features and specifications across countries."
pContent.text = content_text
pContent.font.size = Pt(18)
pContent.font.color.rgb = RGBColor(0x00, 0x00, 0x00) # Black text

# Apply bullet formatting programmatically
# Access XML properties of the paragraph to define a bullet
pPr = pContent._p.get_or_add_pPr()
# Create a bullet definition element
defBullet = pPr.makeelement(qn('a:defBullet'), {'char': '\u2022'})
pPr.append(defBullet)
# Set indentation for the bullet
pPr.set('lvl', '0')
pPr.set('indent', '-36000') # EMUs (approx -0.25 inches)
pPr.set('marL', '72000')   # EMUs (approx 0.5 inches left margin)

# 5. Add Logos (Placeholders)
# Arrange 4 images horizontally below the text
image_path = 'image.png'
img_w = Inches(2.2)
img_h = Inches(2.2)
img_top_pos = Inches(4.8)

# Calculate horizontal positions to center the images
# 4 images + 3 gaps. Gap = 0.3 inches.
# Total content width = 4 * 2.2 + 3 * 0.3 = 9.7 inches
# Slide width = 13.333 inches
# Start X = (13.333 - 9.7) / 2 = 1.8165 inches
start_x = Inches(1.8)
gap = Inches(0.3)
current_x = start_x

for _ in range(4):
    # Add placeholder image
    slide.shapes.add_picture(image_path, current_x, img_top_pos, img_w, img_h)
    current_x += img_w + gap

# Save the presentation
prs.save('output.pptx')