from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Initialize Presentation
prs = Presentation()

# Add a slide using the blank layout (usually index 6) for full control
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set a light background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(246, 248, 250)

# Define Colors
DARK_GRAY = RGBColor(33, 33, 33)
ACCENT_BLUE = RGBColor(0, 56, 112)
WHITE = RGBColor(255, 255, 255)

# 1. Add Title "PART 2"
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(5), Inches(1))
title_tf = title_box.text_frame
title_tf.clear()
p_title = title_tf.add_paragraph()
p_title.text = "PART 2"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = DARK_GRAY
p_title.font.name = "Arial"

# 2. Add Instruction Text
instr_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(5), Inches(0.5))
instr_tf = instr_box.text_frame
instr_tf.clear()
p_instr = instr_tf.add_paragraph()
p_instr.text = "Discuss the questions below about dating."
p_instr.font.size = Pt(18)
p_instr.font.color.rgb = RGBColor(60, 60, 60)
p_instr.font.name = "Arial"

# 3. List Questions with Numbers in Circles
questions = [
    "What do you think are some other good (or bad) date ideas not listed on the previous slide?",
    "What are some things that could go wrong on a first date?",
    "What are some common etiquette mistakes to avoid on a first date?"
]

y_start = Inches(2.5)
y_step = Inches(1.2)
circle_size = Inches(0.4)

for i, question_text in enumerate(questions):
    y_pos = y_start + (i * y_step)
    
    # Create a Circle Shape for the Number
    # MSO_SHAPE.OVAL with equal width and height creates a circle
    circle_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y_pos + Inches(0.05), circle_size, circle_size)
    
    # Style the Circle
    circle_shape.fill.solid()
    circle_shape.fill.fore_color.rgb = ACCENT_BLUE
    circle_shape.line.fill.background() # No border
    
    # Add Number Text inside the Circle
    circle_tf = circle_shape.text_frame
    circle_tf.clear()
    p_num = circle_tf.add_paragraph()
    p_num.text = str(i + 1)
    p_num.font.size = Pt(14)
    p_num.font.color.rgb = WHITE
    p_num.font.bold = True
    p_num.alignment = PP_ALIGN.CENTER
    
    # Add Question Text next to the circle
    q_box = slide.shapes.add_textbox(Inches(2.1), y_pos, Inches(5.5), Inches(1.0))
    q_tf = q_box.text_frame
    q_tf.clear()
    p_q = q_tf.add_paragraph()
    p_q.text = question_text
    p_q.font.size = Pt(16)
    p_q.font.color.rgb = DARK_GRAY
    p_q.font.name = "Arial"

# 4. Add Image (Graphic) on the Right Side
if os.path.exists('image.png'):
    slide.shapes.add_picture('image.png', Inches(7.5), Inches(1.5), Inches(3.5), Inches(4.0))

# 5. Add "PREVIEW ACTIVITY" Button at Bottom Right
# Create a Rectangle Shape
btn_left = Inches(8.5)
btn_top = Inches(6.5)
btn_width = Inches(3.0)
btn_height = Inches(0.7)

btn_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, btn_left, btn_top, btn_width, btn_height)

# Style Button Background
btn_shape.fill.solid()
btn_shape.fill.fore_color.rgb = ACCENT_BLUE
btn_shape.line.fill.background() # No border

# Add Button Text
btn_tf = btn_shape.text_frame
btn_tf.clear()
p_btn = btn_tf.add_paragraph()
p_btn.text = "PREVIEW ACTIVITY"
p_btn.font.size = Pt(14)
p_btn.font.bold = True
p_btn.font.color.rgb = WHITE
p_btn.alignment = PP_ALIGN.CENTER

# 6. Save the Presentation
prs.save('output.pptx')