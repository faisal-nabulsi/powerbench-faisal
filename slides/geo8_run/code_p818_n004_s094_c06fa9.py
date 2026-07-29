from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

def set_bullet(paragraph):
    """Adds a standard bullet point to the paragraph using XML manipulation."""
    # Get paragraph properties element
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    
    # Create XML element for bullet character
    buChar = etree.fromstring('<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>')
    pPr.append(buChar)

# Create a Presentation object
prs = Presentation()

# Add a slide using the Blank layout (index 6 is typically Blank in default theme)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Title: "Blockchain in Cryptocurrency" at the top center
title_left = Inches(1.5)
title_top = Inches(0.5)
title_width = Inches(7)
title_height = Inches(1)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]

title_para.text = "Blockchain in Cryptocurrency"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 0, 0)
title_para.alignment = PP_ALIGN.CENTER

# 2. Add Bullet Points in the main content area (Left side)
content_left = Inches(1)
content_top = Inches(1.5)
content_width = Inches(6)
content_height = Inches(5)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Clear the default empty paragraph
content_frame.paragraphs[0].text = ""

bullet_texts = [
    "The term blockchain is often used to refer to cryptocurrency. Cryptocurrency is a medium of exchange such as US dollars.",
    "It is just an application in the form of e-currency using blockchain.",
    "It is not governed by any financial institution.",
    "The main difference between blockchain and cryptocurrency is that cryptocurrency is created and held electronically in forms such as a virtual wallet.",
    "It is decentralized and it is not governed by anyone whereas blockchain is an advanced record and it has all information related to cryptocurrency exchanges over a shared system."
]

for text in bullet_texts:
    p = content_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    set_bullet(p)

# 3. Place an image of a Bitcoin on the right side
# 'image.png' is the placeholder file available in the working directory
img_left = Inches(7.5)
img_top = Inches(2.5)
img_width = Inches(2.5)
img_height = Inches(2.5)

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# 4. Save the presentation
prs.save('output.pptx')