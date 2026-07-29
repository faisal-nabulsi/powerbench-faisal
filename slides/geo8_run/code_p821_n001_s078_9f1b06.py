from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

def main():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    DARK_BLUE = RGBColor(0, 51, 102)
    WHITE = RGBColor(255, 255, 255)
    
    # Get slide dimensions (Standard 16:9)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # --- 1. Background & Title Bar ---
    # The prompt requests a dark blue background for the title, 
    # but a light color for the rest of the slide for readability.
    # We achieve this by adding a dark blue rectangle at the top as a header.
    
    # Create Header Bar (Dark Blue)
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        slide_width,
        Inches(1.5)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.fill.background() # Remove border
    
    # The rest of the slide background remains the default light color (usually white/off-white)
    
    # --- 2. Title ---
    # "5 Stages of Development of Media"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5 Stages of Development of Media"
    p.font.size = Pt(42)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # --- 3. Subtitle ---
    # "1. ORAL COMMUNICATION"
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7), Inches(1))
    tf_sub = subtitle_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "1. ORAL COMMUNICATION"
    p_sub.font.size = Pt(30)
    p_sub.font.color.rgb = DARK_BLUE
    p_sub.font.bold = True
    
    # --- 4. Bullet Points ---
    bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(7), Inches(3))
    tf_bullets = bullets_box.text_frame
    tf_bullets.word_wrap = True
    
    # Bullet 1
    p1 = tf_bullets.paragraphs[0]
    p1.text = "Language allowed humans to communicate and share information."
    p1.font.size = Pt(20)
    p1.font.color.rgb = DARK_BLUE
    p1.level = 0
    
    # Bullet 2
    p2 = tf_bullets.add_paragraph()
    p2.text = "Language became the most important tool for exploring the world and different cultures."
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK_BLUE
    p2.level = 0
    
    # Apply bullet formatting via XML (robust method for blank text boxes)
    for p in tf_bullets.paragraphs:
        pPr = p._p.get_or_add_pPr()
        # Clear existing bullet elements to ensure clean rendering
        for elem in list(pPr):
            if 'bu' in elem.tag:
                pPr.remove(elem)
        
        # Bullet Character
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)
        
        # Bullet Font
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': '+Minor'})
        pPr.append(buFont)
        
        # Bullet Color (Dark Blue)
        buClr = pPr.makeelement(qn('a:solidFill'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '003366'})
        buClr.append(srgbClr)
        pPr.append(buClr)

    # --- 5. Image ---
    # Right side, placeholder 'image.png'
    # Positioning: Start below title (y=2.0), aligned right (x=8.0)
    left_img = Inches(8.0)
    top_img = Inches(2.0)
    width_img = Inches(5.0)
    height_img = Inches(4.5)
    
    if os.path.exists('image.png'):
        slide.shapes.add_picture('image.png', left_img, top_img, width_img, height_img)
    else:
        # Fallback in case image is missing during execution in restricted env
        # We add a placeholder shape instead
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_img, top_img, width_img, height_img)
        shape.text = "Image Placeholder"
        
    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    main()