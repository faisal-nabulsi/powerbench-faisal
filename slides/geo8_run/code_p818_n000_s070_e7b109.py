from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen layout
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Get blank layout (index 6 usually)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set light background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 240)

# Function to add text box
def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                font_color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

# Title: "PART 2"
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(1.2),
            "PART 2", font_size=44, bold=True,
            font_color=RGBColor(30, 30, 30), font_name='Calibri')

# Subtitle
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.6),
            "Discuss the questions below about dating.", font_size=20,
            bold=False, font_color=RGBColor(80, 80, 80), font_name='Calibri')

# Decorative line under title
line_left = Inches(0.8)
line_top = Inches(2.0)
line_width = Inches(5.5)
line_height = Inches(0.04)
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height)
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = RGBColor(200, 120, 80)
line_shape.line.fill.background()

# Questions with circled numbers
questions = [
    "What do you think are some other good (or bad) date ideas not listed on the previous slide?",
    "What are some things that could go wrong on a first date?",
    "What are some common etiquette mistakes to avoid on a first date?"
]

circle_numbers = ['①', '②', '③']

question_top = Inches(2.35)
question_spacing = Inches(1.35)

for i, (q, circ_num) in enumerate(zip(questions, circle_numbers)):
    y_pos = question_top + Inches(i * question_spacing)

    # Create circle shape for number
    circle_left = Inches(0.8)
    circle_top = y_pos + Inches(0.08)
    circle_size = Inches(0.55)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(200, 120, 80)
    circle.line.fill.background()

    # Add number inside circle
    circle_tf = circle.text_frame
    circle_tf.word_wrap = False
    circle_p = circle_tf.paragraphs[0]
    circle_p.text = str(i + 1)
    circle_p.font.size = Pt(22)
    circle_p.font.bold = True
    circle_p.font.color.rgb = RGBColor(255, 255, 255)
    circle_p.font.name = 'Calibri'
    circle_p.alignment = PP_ALIGN.CENTER
    circle_tf.paragraphs[0].space_before = Pt(2)

    # Question text
    add_textbox(slide, Inches(1.55), y_pos, Inches(7.5), Inches(1.1),
                q, font_size=17, bold=False,
                font_color=RGBColor(50, 50, 50), font_name='Calibri')

    # Small orange dot separator
    if i < 2:
        dot_y = y_pos + Inches(1.05)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), dot_y, Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(200, 120, 80)
        dot.line.fill.background()

# Communication-related image on the right side
img_left = Inches(8.8)
img_top = Inches(1.5)
img_width = Inches(4.0)
img_height = Inches(4.0)

# Add a subtle semi-transparent rectangle behind the image
img_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 img_left - Inches(0.15), img_top - Inches(0.15),
                                 img_width + Inches(0.3), img_height + Inches(0.3))
img_bg.fill.solid()
img_bg.fill.fore_color.rgb = RGBColor(230, 228, 225)
img_bg.line.fill.background()

# Add the image
slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# "PREVIEW ACTIVITY" button at bottom right
btn_left = Inches(9.0)
btn_top = Inches(6.3)
btn_width = Inches(3.8)
btn_height = Inches(0.7)

btn_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, btn_left, btn_top, btn_width, btn_height)
btn_shape.fill.solid()
btn_shape.fill.fore_color.rgb = RGBColor(30, 60, 120)
btn_shape.line.fill.background()

btn_tf = btn_shape.text_frame
btn_tf.word_wrap = True
btn_p = btn_tf.paragraphs[0]
btn_p.text = "PREVIEW ACTIVITY"
btn_p.font.size = Pt(16)
btn_p.font.bold = True
btn_p.font.color.rgb = RGBColor(255, 255, 255)
btn_p.font.name = 'Calibri'
btn_p.alignment = PP_ALIGN.CENTER

btn_tf.paragraphs[0].space_before = Pt(0)
btn_tf.paragraphs[0].space_after = Pt(0)

# Access icon/label along the image
icon_left = Inches(9.2)
icon_top = Inches(5.7)
add_textbox(slide, icon_left, icon_top, Inches(3.5), Inches(0.5),
            "💬 Communication Tips", font_size=14, bold=True,
            font_color=RGBColor(200, 120, 80), alignment=PP_ALIGN.LEFT, font_name='Calibri')

# Save the presentation
prs.save('output.pptx')
print("Presentation saved to 'output.pptx'")