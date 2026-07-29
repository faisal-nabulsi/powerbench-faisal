from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Create a Presentation object
prs = Presentation()

# Get the default slide dimensions to assist with centering calculations
slide_width = prs.slide_width
slide_height = prs.slide_height

# 2. Add a blank slide (Layout index 6 is standard for blank slides)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 3. Set the background to solid red
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# 4. Define text content and box dimensions
text_content = "Where do you buy your products from?"
box_width = Inches(9)   # Width of the text box
box_height = Inches(2)  # Height of the text box

# 5. Calculate position to center the text box on the slide
# Position = (Slide_Dimension - Box_Dimension) / 2
left = (slide_width - box_width) / 2
top = (slide_height - box_height) / 2

# 6. Add the text box to the slide
textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
text_frame = textbox.text_frame

# 7. Format the text
# Allow text wrapping
text_frame.word_wrap = True

# Access the paragraph and center alignment
paragraph = text_frame.paragraphs[0]
paragraph.alignment = PP_ALIGN.CENTER

# Add the text run
run = paragraph.add_run()
run.text = text_content

# Apply font styles: Clean modern font (Calibri), Large, White, Bold
run.font.name = "Calibri"
run.font.size = Pt(48)           # Large size for readability
run.font.color.rgb = RGBColor(255, 255, 255) # White color
run.font.bold = True

# 8. Save the presentation to 'output.pptx'
prs.save('output.pptx')