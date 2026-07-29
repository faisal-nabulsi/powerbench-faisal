from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Create Presentation
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height

# 2. Add a blank slide
slide_layout = prs.slide_layouts[6] # 6 is usually Blank
slide = prs.slides.add_slide(slide_layout)

# 3. Set Background Color to Coral (RGB: 255, 127, 80)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 127, 80)

# 4. Add Image (Airbnb Icon) on the Left Side
# "Entire left side" implies full height and 50% width
# Image placeholder name is 'image.png'
image_left = Inches(0)
image_top = Inches(0)
image_width = slide_width / 2
image_height = slide_height

# Check if image.png exists? Prompt implies it is available.
try:
    slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)
except FileNotFoundError:
    # Fallback if image not found, though prompt guarantees it.
    # We could skip or add a placeholder shape, but script should assume availability.
    pass

# 5. Define Right Side Text Area
# Start right after the image, with a small margin
text_left_start = slide_width / 2 + Inches(0.5)
text_width = slide_width / 2 - Inches(1.0) # Leave right margin

# Define Vertical Positions
# Slide height is 7.5 inches. Center is 3.75 inches.

# Element 2: Text "airbnb" logo
# "Place the text 'airbnb' logo on the right side of the slide, centered vertically."
# We interpret this as the text box being centered vertically on the slide.
airbnb_font_size = Pt(28)
airbnb_box_height = Inches(1.0) # Estimated height for the text
airbnb_box_top = (slide_height / 2) - (airbnb_box_height / 2)

txBox_airbnb = slide.shapes.add_textbox(text_left_start, airbnb_box_top, text_width, airbnb_box_height)
tf_airbnb = txBox_airbnb.text_frame
tf_airbnb.word_wrap = True
p_airbnb = tf_airbnb.paragraphs[0]
p_airbnb.text = "airbnb"
p_airbnb.font.size = airbnb_font_size
p_airbnb.font.color.rgb = RGBColor(255, 255, 255)
p_airbnb.font.name = "Arial" # Standard sans-serif
p_airbnb.alignment = PP_ALIGN.RIGHT # Aligning right to match names/style of "right side"

# Element 3: Title "Business Case"
# "Add the text 'Business Case' in a prominent font, centered vertically and below the text logo."
# Since it must be below the logo, we position it below the airbnb box.
# Note: Strict vertical centering conflicts with "below logo". 
# We prioritize the logical flow (Below Logo) while keeping it generally centered/prominent.
title_gap = Inches(0.3)
title_box_top = airbnb_box_top + airbnb_box_height + title_gap
title_font_size = Pt(44) # Prominent
title_box_height = Inches(1.5)

txBox_title = slide.shapes.add_textbox(text_left_start, title_box_top, text_width, title_box_height)
tf_title = txBox_title.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
p_title.text = "Business Case"
p_title.font.size = title_font_size
p_title.font.color.rgb = RGBColor(255, 255, 255)
p_title.font.name = "Arial"
p_title.font.bold = True
p_title.alignment = PP_ALIGN.RIGHT

# Element 4: Names
# "Below the title... smaller font size... aligned to the right"
names_gap = Inches(0.2)
names_box_top = title_box_top + title_box_height + names_gap
names_font_size = Pt(18) # Smaller
names_box_height = Inches(1.5) # Space for 3 names

txBox_names = slide.shapes.add_textbox(text_left_start, names_box_top, text_width, names_box_height)
tf_names = txBox_names.text_frame
tf_names.word_wrap = True
p_names = tf_names.paragraphs[0]
p_names.text = "Daniel Consuegra\nAlejandra Del Chiaro\nMaria Camila Echeverri"
p_names.font.size = names_font_size
p_names.font.color.rgb = RGBColor(255, 255, 255)
p_names.font.name = "Arial"
p_names.alignment = PP_ALIGN.RIGHT

# 6. Save Presentation
prs.save('output.pptx')