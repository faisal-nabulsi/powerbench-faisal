from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Add a blank slide (Index 6 is typically the 'Blank' layout)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set background to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# --- Left Section ---
# Insert the first image (blockchain graphic) on the left side
# Using 'image.png' as the placeholder
left_image = slide.shapes.add_picture(
    'image.png', 
    left=Inches(0.5), 
    top=Inches(1.5), 
    width=Inches(4), 
    height=Inches(4)
)

# --- Right Section ---
# 1. Title "Blockchain Technology"
# Position: Top right
title_box = slide.shapes.add_textbox(
    left=Inches(5.5), 
    top=Inches(1.0), 
    width=Inches(4), 
    height=Inches(1)
)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Blockchain Technology"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0) # Black text
p.alignment = PP_ALIGN.LEFT

# 2. Second image (keyboard and letters spelling "BLOCKCHAIN")
# Position: Below the title on the right side
# Using 'image.png' as the placeholder
right_image = slide.shapes.add_picture(
    'image.png', 
    left=Inches(5.5), 
    top=Inches(3.0), 
    width=Inches(4), 
    height=Inches(4)
)

# Save the presentation
prs.save('output.pptx')