from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a new presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Define Colors
DARK_BLUE = RGBColor(0, 51, 102)    # Dark Blue for Title Background and Subtitle Text
WHITE = RGBColor(255, 255, 255)     # White for Title Text
DARK_GRAY = RGBColor(50, 50, 50)    # Dark Gray for Body Text

# 1. Set Slide Background to White
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# 2. Create Title Bar with Dark Blue Background
# Dimensions: Full width, 1.5 inches height
title_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(0), Inches(0), 
    prs.slide_width, Inches(1.5)
)
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = DARK_BLUE
title_bar.line.fill.background()  # Remove border line

# Add Title Text
tf_title = title_bar.text_frame
p_title = tf_title.paragraphs[0]
p_title.text = "5 Stages of Development of Media"
p_title.font.size = Pt(42)
p_title.font.bold = True
p_title.font.color.rgb = WHITE
p_title.alignment = PP_ALIGN.CENTER

# 3. Add Subtitle ("4. ELECTRONIC MEDIA")
# Positioned below the title bar
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(1))
tf_sub = subtitle_box.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.text = "4. ELECTRONIC MEDIA"
p_sub.font.size = Pt(32)
p_sub.font.bold = True
p_sub.font.color.rgb = DARK_BLUE
p_sub.alignment = PP_ALIGN.LEFT

# 4. Add Content Text (Bullet Points)
# Positioned on the left side
txt_width = Inches(6.5)
txt_height = Inches(4.5)
txt_left = Inches(0.5)
txt_top = Inches(2.8)

body_box = slide.shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
tf_body = body_box.text_frame
tf_body.word_wrap = True

# First Bullet Point
p1 = tf_body.paragraphs[0]
p1.text = "It includes the telegraphs, telephone, radio, film and television."
p1.font.size = Pt(20)
p1.font.color.rgb = DARK_GRAY
p1.bullet = True
p1.space_after = Pt(12)

# Second Bullet Point
p2 = tf_body.add_paragraph()
p2.text = "The wide range of these media continue to open up new perspectives on economic, political and cultural processes of globalization."
p2.font.size = Pt(20)
p2.font.color.rgb = DARK_GRAY
p2.bullet = True

# 5. Add Image on the Right
img_width = Inches(6.0)
img_height = Inches(5.0)
img_left = prs.slide_width - Inches(0.5) - img_width  # 0.5 inch right margin
img_top = Inches(2.5)

try:
    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
except Exception:
    pass

# Save the presentation
prs.save('output.pptx')