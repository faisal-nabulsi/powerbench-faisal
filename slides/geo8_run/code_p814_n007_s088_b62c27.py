from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Use a blank layout (index 6)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Get slide dimensions to ensure elements fit properly
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. Left Side: Main Title "Principles of graphic design"
    # Create a shape for the background on the left side (approx 35% width)
    left_width = Inches(slide_width.inches * 0.35)
    left_bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        left_width, slide_height
    )
    left_bg_shape.fill.solid()
    left_bg_shape.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
    left_bg_shape.line.fill.background() # Remove border

    # Add text to the left background shape
    tf_left = left_bg_shape.text_frame
    tf_left.word_wrap = True
    p_left = tf_left.paragraphs[0]
    p_left.text = "Principles of graphic design"
    p_left.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_left.font.size = Pt(40)
    p_left.font.bold = True
    p_left.font.name = "Arial"
    p_left.alignment = PP_ALIGN.CENTER
    
    # Center text vertically roughly
    tf_left.margin_top = Inches(2.5)

    # 2. Right Side: Secondary Title
    # Position: Start after the left section
    right_start_x = left_width + Inches(0.5)
    right_width = slide_width - right_start_x
    title_y = Inches(0.5)
    title_height = Inches(1.8)

    right_title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, right_start_x, title_y, right_width, title_height
    )
    right_title_shape.fill.solid()
    right_title_shape.fill.fore_color.rgb = RGBColor(173, 216, 230) # Light Blue
    right_title_shape.line.fill.background()

    tf_title = right_title_shape.text_frame
    tf_title.text = "The principles of graphic design related to the areas are,"
    p_title = tf_title.paragraphs[0]
    p_title.font.color.rgb = RGBColor(0, 0, 0) # Dark text for contrast
    p_title.font.size = Pt(18)
    p_title.font.name = "Arial"
    p_title.alignment = PP_ALIGN.LEFT
    
    # Margins for the title text box
    tf_title.margin_left = Inches(0.2)
    tf_title.margin_top = Inches(0.3)

    # 3. Right Side: List Items
    # Varying shades of blue and green
    items = ["Arrangement", "Proximity", "Repetition", "Contrast", "Balance"]
    item_colors = [
        RGBColor(70, 130, 180),  # Steel Blue
        RGBColor(46, 139, 87),   # Sea Green
        RGBColor(30, 144, 255),  # Dodger Blue
        RGBColor(0, 128, 128),   # Teal
        RGBColor(0, 100, 0)      # Dark Green
    ]

    item_start_y = title_y + title_height + Inches(0.5)
    item_height = Inches(1.2)
    item_spacing = Inches(0.2)

    for i, item_text in enumerate(items):
        # Create shape for each list item
        item_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            right_start_x, 
            item_start_y + (i * (item_height + item_spacing)), 
            right_width, 
            item_height
        )
        
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = item_colors[i]
        item_box.line.fill.background()
        
        # Add text to the item box
        tf_item = item_box.text_frame
        tf_item.text = item_text
        p_item = tf_item.paragraphs[0]
        p_item.font.color.rgb = RGBColor(255, 255, 255) # White text
        p_item.font.size = Pt(24) # Consistent font size
        p_item.font.name = "Arial"
        p_item.font.bold = True
        p_item.alignment = PP_ALIGN.LEFT
        
        # Margins for list item text
        tf_item.margin_left = Inches(0.3)
        tf_item.margin_top = Inches(0.2)

    # Save the presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()