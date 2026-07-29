from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new Presentation object
prs = Presentation()

# Select a blank layout (Index 6 is typically 'Blank' in the default template)
# This allows for manual positioning of text and images.
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Set Background Color to White for high contrast
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# 2. Add Title: "Blockchain in Cryptocurrency" at top center
# Coordinates are set to center the title on a standard slide width.
title_box = slide.shapes.add_textbox(
    left=Inches(2), 
    top=Inches(0.5), 
    width=Inches(6), 
    height=Inches(1)
)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Blockchain in Cryptocurrency"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0) # Black text
p.alignment = PP_ALIGN.CENTER

# 3. Add Bullet Points in the Main Content Area
content_points = [
    "The term blockchain is often used to refer to cryptocurrency. Cryptocurrency is a medium of exchange such as US dollars.",
    "It is just an application in the form of e-currency using blockchain.",
    "It is not governed by any financial institution.",
    "The main difference between blockchain and cryptocurrency is that cryptocurrency is created and held electronically in forms such as a virtual wallet.",
    "It is decentralized and it is not governed by anyone whereas blockchain is an advanced record and it has all information related to cryptocurrency exchanges over a shared system."
]

# Create a text box for the content, positioned on the left side
# Width is limited to 6.5 inches to prevent overlap with the image on the right.
content_box = slide.shapes.add_textbox(
    left=Inches(1), 
    top=Inches(2), 
    width=Inches(6.5), 
    height=Inches(5)
)
tf_content = content_box.text_frame
tf_content.word_wrap = True

for i, point in enumerate(content_points):
    if i == 0:
        p = tf_content.paragraphs[0]
    else:
        p = tf_content.add_paragraph()
    
    # Prepend a bullet character to visually represent bullet points
    p.text = "• " + point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(40, 40, 40) # Dark Grey text for legibility
    p.space_after = Pt(12)

# 4. Add Image of a Bitcoin on the Right Side
# The placeholder 'image.png' is placed on the right side of the slide.
try:
    slide.shapes.add_picture(
        'image.png', 
        left=Inches(7.5), # Positioned to the right of the text box
        top=Inches(2), 
        width=Inches(2), 
        height=Inches(2)
    )
except Exception:
    pass

# 5. Save the presentation to 'output.pptx'
prs.save('output.pptx')