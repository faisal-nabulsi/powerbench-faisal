from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize the presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 20, 60)

# 1. Title "Viewing Activity" with a film icon in the top left corner
# Using a unicode character for the film icon as no specific icon file is provided
title_text = "🎬 Viewing Activity"
left = Inches(0.5)
top = Inches(0.5)
width = Inches(5)
height = Inches(0.6)
title_box = slide.shapes.add_textbox(left, top, width, height)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = title_text
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.bold = True
p.font.name = "Arial"

# 2. Insert an image from the "Friends" episode on the right side
# Using the provided placeholder image 'image.png'
image_left = Inches(6.0)
image_top = Inches(1.5)
image_width = Inches(3.5)
image_height = Inches(2.5)
slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

# 3. Add the title "Friends | Joey Doesn't Share Food!" above the image
img_title_left = Inches(6.0)
img_title_top = Inches(1.0)
img_title_width = Inches(3.5)
img_title_height = Inches(0.5)
img_title_box = slide.shapes.add_textbox(img_title_left, img_title_top, img_title_width, img_title_height)
tf_img = img_title_box.text_frame
tf_img.word_wrap = True
p_img = tf_img.paragraphs[0]
p_img.text = "Friends | Joey Doesn't Share Food!"
p_img.font.size = Pt(20)
p_img.font.color.rgb = RGBColor(255, 255, 255)
p_img.font.bold = True
p_img.font.name = "Arial"

# 4. Include the text "Finish watching the video." on the left side
instruction_left = Inches(0.5)
instruction_top = Inches(2.5)
instruction_width = Inches(4)
instruction_height = Inches(1)
instruction_box = slide.shapes.add_textbox(instruction_left, instruction_top, instruction_width, instruction_height)
tf_instr = instruction_box.text_frame
tf_instr.word_wrap = True
p_instr = tf_instr.paragraphs[0]
p_instr.text = "Finish watching the video."
p_instr.font.size = Pt(28)
p_instr.font.color.rgb = RGBColor(255, 255, 255)
p_instr.font.name = "Arial"

# 5. Add the logo or text "FLUENT" in the bottom left corner
logo_left = Inches(0.5)
logo_top = Inches(6.5)
logo_width = Inches(2)
logo_height = Inches(0.5)
logo_box = slide.shapes.add_textbox(logo_left, logo_top, logo_width, logo_height)
tf_logo = logo_box.text_frame
tf_logo.word_wrap = True
p_logo = tf_logo.paragraphs[0]
p_logo.text = "FLUENT"
p_logo.font.size = Pt(24)
p_logo.font.color.rgb = RGBColor(255, 215, 0) # Gold color for contrast
p_logo.font.bold = True
p_logo.font.name = "Arial"

# Save the presentation
prs.save('output.pptx')