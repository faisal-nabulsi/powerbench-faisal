from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 3. Set Light Gray Background for Slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 240)

# 4. Title "PART 4" with Blue Background (Top Left)
# Create a rectangle shape for the blue background
title_bg = slide.shapes.add_shape(
    1, # MSO_SHAPE.RECTANGLE
    Inches(0.5), Inches(0.5), Inches(3.5), Inches(1.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = RGBColor(0, 112, 192) # Blue color
title_bg.line.fill.background()

# Add Text to Title Shape
tf = title_bg.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "PART 4"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255) # White text for contrast
p.alignment = PP_ALIGN.CENTER

# 5. Icon of Human Speaking (Below Title)
# Represented by a simple shape (Circle) as no specific icon image file is available
icon_circle = slide.shapes.add_shape(
    5, # MSO_SHAPE.OVAL
    Inches(1.75), Inches(2.2), Inches(1), Inches(1)
)
icon_circle.fill.solid()
icon_circle.fill.fore_color.rgb = RGBColor(0, 112, 192)
icon_circle.line.fill.background()

# 6. Instruction 1
instruction1_text = "Tell the short story below out loud using the pronunciation rule you studied on the previous slide."
txBox1 = slide.shapes.add_textbox(
    Inches(0.5), Inches(3.4), Inches(8.5), Inches(0.8)
)
tf1 = txBox1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = instruction1_text
p1.font.size = Pt(16)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.alignment = PP_ALIGN.LEFT

# 7. Story Text
story_text = "Last week, I went on a date. We were skating in the park, we were laughing, and the guy was making some really funny jokes...and then, out of nowhere, I fell down and I broke my leg!"
story_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(4.4), Inches(8), Inches(2)
)
tf_story = story_box.text_frame
tf_story.word_wrap = True
p_story = tf_story.paragraphs[0]
p_story.text = story_text
p_story.font.size = Pt(20)
p_story.font.color.rgb = RGBColor(50, 50, 50) # Dark Gray for legibility
p_story.alignment = PP_ALIGN.LEFT

# 8. Image: People Skating (Right Side)
# Using the mandatory 'image.png' placeholder
slide.shapes.add_picture(
    'image.png',
    Inches(9), Inches(2.5), Inches(4), Inches(4)
)

# 9. Bottom Section: Audio Icon + Instruction
# Audio Icon (Rectangle shape representing an icon next to text)
audio_icon = slide.shapes.add_shape(
    1, # MSO_SHAPE.RECTANGLE
    Inches(1), Inches(6.2), Inches(0.8), Inches(0.8)
)
audio_icon.fill.solid()
audio_icon.fill.fore_color.rgb = RGBColor(0, 112, 192)
audio_icon.line.fill.background()

# Instruction Text next to the icon
bottom_instruction_text = "Play and listen to the audio clip to check your answers. Repeat if necessary."
txBoxBottom = slide.shapes.add_textbox(
    Inches(2.2), Inches(6.2), Inches(10), Inches(1)
)
tf_bot = txBoxBottom.text_frame
tf_bot.word_wrap = True
p_bot = tf_bot.paragraphs[0]
p_bot.text = bottom_instruction_text
p_bot.font.size = Pt(16)
p_bot.font.bold = True
p_bot.font.color.rgb = RGBColor(0, 0, 0)
p_bot.alignment = PP_ALIGN.LEFT

# 10. Save the presentation
prs.save('output.pptx')