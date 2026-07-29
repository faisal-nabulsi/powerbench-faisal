from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_slide():
    # Initialize Presentation
    prs = Presentation()

    # Add a blank slide (Layout 6 is typically blank)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Set Slide Background to Black
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # 2. Incorporate relevant image in the background
    # Placed on the right side to complement the text area without overwhelming it.
    # Dimensions adjusted to fit within the right portion of the slide.
    if True: # Placeholder check allowed by prompt context, assuming file exists
        try:
            slide.shapes.add_picture('image.png', Inches(7), Inches(0.5), Inches(6), Inches(6.5))
        except FileNotFoundError:
            pass # Silently pass if image missing, though prompt guarantees it.

    # 3. Add Title "Now: I, CAN"
    # Positioned at the top, centered text for a title look.
    title_left = Inches(1.5)
    title_top = Inches(0.5)
    title_width = Inches(10)
    title_height = Inches(1.2)

    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf_title = title_shape.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = "Now: I, CAN"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_title.alignment = PP_ALIGN.CENTER

    # 4. Add Bullet Points as Rounded Rectangles
    # Common dimensions for the text boxes
    box_left = Inches(1.5)
    box_width = Inches(10)
    box_height = Inches(1.3)

    # Exact text content requested
    bullet1 = "• Differentiate between different types of market."
    bullet2 = "• Explain why society need different types of market."

    # --- Top Box (Orange) ---
    box1_top = Inches(2.2)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box1_top, box_width, box_height)

    # Set Orange Background
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 165, 0)
    shape1.line.fill.background() # Remove border for cleaner look

    # Set Text inside Box 1
    tf1 = shape1.text_frame
    tf1.clear()
    p1 = tf1.paragraphs[0]
    p1.text = bullet1
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p1.alignment = PP_ALIGN.CENTER

    # --- Bottom Box (Yellow) ---
    box2_top = Inches(3.8)
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box2_top, box_width, box_height)

    # Set Yellow Background
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 255, 0)
    shape2.line.fill.background() # Remove border

    # Set Text inside Box 2
    tf2 = shape2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = bullet2
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p2.alignment = PP_ALIGN.CENTER

    # Save Presentation
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()