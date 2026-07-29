from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 slide size
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Set light background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(248, 249, 250)

# Title: "PART 2" in bold, dark font
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(5), Inches(1.0))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_run = title_para.add_run()
title_run.text = "PART 2"
title_run.font.bold = True
title_run.font.size = Pt(42)
title_run.font.color.rgb = RGBColor(26, 26, 26)

# Add decorative line under title
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.8),
    Inches(1.35),
    Inches(1.8),
    Inches(0.06)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0, 102, 178)
line.line.fill.background()

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(7), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
sub_para = subtitle_frame.paragraphs[0]
sub_run = sub_para.add_run()
sub_run.text = "Discuss the questions below about dating."
sub_run.font.size = Pt(20)
sub_run.font.color.rgb = RGBColor(85, 85, 85)

# Questions list with numbered circles
questions = [
    "What do you think are some other good (or bad) date ideas not listed on the previous slide?",
    "What are some things that could go wrong on a first date?",
    "What are some common etiquette mistakes to avoid on a first date?"
]

start_y = Inches(2.3)
y_step = Inches(1.35)
circle_size = Inches(0.55)

for idx, q_text in enumerate(questions):
    # Create circle shape for number
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.9),
        start_y + Inches(idx * 1.35) + Inches(0.08),
        circle_size,
        circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 102, 178)
    circle.line.fill.background()
    
    # Set text inside circle
    circle.text_frame.word_wrap = False
    circle.text_frame.auto_size = None
    num_para = circle.text_frame.paragraphs[0]
    num_para.alignment = PP_ALIGN.CENTER
    num_run = num_para.add_run()
    num_run.text = str(idx + 1)
    num_run.font.size = Pt(20)
    num_run.font.bold = True
    num_run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add question text
    q_box = slide.shapes.add_textbox(
        Inches(1.7),
        start_y + Inches(idx * 1.35),
        Inches(6.2),
        Inches(1.0)
    )
    q_frame = q_box.text_frame
    q_frame.word_wrap = True
    q_para = q_frame.paragraphs[0]
    q_run = q_para.add_run()
    q_run.text = q_text
    q_run.font.size = Pt(17)
    q_run.font.color.rgb = RGBColor(45, 45, 45)

# Add communication-related graphic on the right side using placeholder image
comm_graphic = slide.shapes.add_picture(
    'image.png',
    Inches(8.5),
    Inches(1.8),
    Inches(3.8),
    Inches(3.2)
)

# Add subtle decorative elements related to communication
# Small circle accents
accent1 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(12.5),
    Inches(1.0),
    Inches(0.4),
    Inches(0.4)
)
accent1.fill.solid()
accent1.fill.fore_color.rgb = RGBColor(0, 102, 178)
accent1.line.fill.background()
accent1.fill.fore_color.brightness = -0.2

accent2 = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(8.0),
    Inches(5.0),
    Inches(0.35),
    Inches(0.35)
)
accent2.fill.solid()
accent2.fill.fore_color.rgb = RGBColor(70, 130, 180)
accent2.line.fill.background()

# "PREVIEW ACTIVITY" button at bottom right in dark blue rectangle
button_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(10.0),
    Inches(6.2),
    Inches(3.0),
    Inches(0.75)
)
button_shape.fill.solid()
button_shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
button_shape.line.fill.background()

# Add rounded corners effect via smaller inner shape simulation (using text alignment instead)
btn_frame = button_shape.text_frame
btn_frame.word_wrap = False
btn_para = btn_frame.paragraphs[0]
btn_para.alignment = PP_ALIGN.CENTER
btn_run = btn_para.add_run()
btn_run.text = "PREVIEW ACTIVITY"
btn_run.font.size = Pt(15)
btn_run.font.bold = True
btn_run.font.color.rgb = RGBColor(255, 255, 255)
btn_para.space_before = Pt(2)

# Save the presentation
prs.save('output.pptx')