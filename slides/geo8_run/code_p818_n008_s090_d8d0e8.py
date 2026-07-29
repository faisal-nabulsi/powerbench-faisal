from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Access slide dimensions to set background image correctly
slide_width = prs.slide_width
slide_height = prs.slide_height

# Add a blank slide (Layout 6 is usually blank in default templates)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# --- 1. Background Image ---
# Add the placeholder image 'image.png' covering the entire slide
picture = slide.shapes.add_picture('image.png', Inches(0), Inches(0), slide_width, slide_height)

# Move the picture shape to the back (z-order 0) so text appears on top
# We manipulate the XML directly to reorder the shape
pic_sp = picture._element
pic_sp_tree = pic_sp.getparent()
pic_sp_tree.remove(pic_sp)
pic_sp_tree.insert(0, pic_sp)

# --- Helper Function for Title Styling ---
def style_title(shape, text, font_size=18, alignment=PP_ALIGN.LEFT):
    """Sets yellow background and black bold text for a text shape."""
    # Set background to Yellow
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 0)
    
    # Configure Text
    tf = shape.text_frame
    tf.clear() # Clear default content
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p.font.bold = True
    p.alignment = alignment
    return shape

# --- 2. Top-Left Title ---
# "The Subject and Content of Art"
top_left_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3.5), Inches(0.5))
style_title(top_left_tb, "The Subject and Content of Art", font_size=14)

# --- 3. Main Title ---
# "THE CONTENT OF ART"
main_title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.8))
style_title(main_title_tb, "THE CONTENT OF ART", font_size=32, alignment=PP_ALIGN.CENTER)

# --- 4. Main Content ---
# White background content box
content_tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(3))
content_tb.fill.solid()
content_tb.fill.fore_color.rgb = RGBColor(255, 255, 255) # White background

tf = content_tb.text_frame
tf.clear()

# Intro text
p_intro = tf.paragraphs[0]
p_intro.text = "It is the mass of ideas associated with each artwork and communicated through the following:"
p_intro.font.size = Pt(18)
p_intro.font.color.rgb = RGBColor(0, 0, 0)
p_intro.space_after = Pt(10)

# List items with checkmarks
# The prompt asks for checkmarks (✓) before each bullet point.
items = [
    "✓ 1. The art’s imagery",
    "✓ 2. The symbolic meaning",
    "✓ 3. Its surroundings where it is used or displayed",
    "✓ 4. The customs, beliefs and values of the culture that uses it"
]

for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_before = Pt(6)

# Save the presentation to 'output.pptx'
prs.save('output.pptx')