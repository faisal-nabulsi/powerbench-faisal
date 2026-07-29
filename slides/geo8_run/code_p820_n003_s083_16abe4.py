from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation
prs = Presentation()
# Set slide dimensions to 16:9 (13.33 x 7.5 inches)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Add a blank slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 1. Add Background Image
# Using 'image.png' as specified in the prompt
slide.shapes.add_picture('image.png', 0, 0, prs.slide_width, prs.slide_height)

# 2. Add White Overlay for Content Area
# A white rectangle to ensure text readability against the background image
overlay_left = Inches(1.5)
overlay_top = Inches(2.5)
overlay_width = Inches(10.33)
overlay_height = Inches(4.5)

overlay_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    overlay_left, overlay_top, overlay_width, overlay_height
)
overlay_shape.fill.solid()
overlay_shape.fill.fore_color.rgb = RGBColor(255, 255, 255) # White
overlay_shape.line.fill.background() # No border

# 3. Add Top-Left Title
txBox_top = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
tf_top = txBox_top.text_frame
tf_top.text = "The Subject and Content of Art"
tf_top.paragraphs[0].font.size = Pt(18)
tf_top.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0) # Black
tf_top.paragraphs[0].font.bold = True

# 4. Add Main Title Box (Yellow Background)
title_left = Inches(1.5)
title_top = Inches(1.0)
title_width = Inches(10.33)
title_height = Inches(1.0)

title_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    title_left, title_top, title_width, title_height
)
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
title_shape.line.fill.background()

title_tf = title_shape.text_frame
title_tf.word_wrap = True
p_title = title_tf.paragraphs[0]
p_title.text = "Two kinds of Art as to Subject"
p_title.font.size = Pt(28)
p_title.font.color.rgb = RGBColor(0, 0, 0) # Black
p_title.font.bold = True
p_title.alignment = PP_ALIGN.CENTER

# 5. Add Content Text
# Text box placed over the white overlay
content_left = Inches(2.0)
content_top = Inches(3.0)
content_width = Inches(9.33)
content_height = Inches(3.5)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
content_tf = content_box.text_frame
content_tf.word_wrap = True

# Main Point Header
p_main = content_tf.paragraphs[0]
p_main.text = "2. Non-representational or Non-objective Art"
p_main.font.size = Pt(22)
p_main.font.bold = True
p_main.font.color.rgb = RGBColor(0, 0, 0)
p_main.space_after = Pt(10)

# Bullet Points with Checkmarks
# Using the checkmark character directly in the text string for reliability
bullet_1 = "✓ Also known as non-objective art"
bullet_2 = "✓ The subject is not obviously or directly represented"
bullet_3 = "✓ Example: Total Abstractions- non-representational or non-objective because they stray away from the reality. Not all abstracts are non-representational."

p_b1 = content_tf.add_paragraph()
p_b1.text = bullet_1
p_b1.font.size = Pt(18)
p_b1.font.color.rgb = RGBColor(0, 0, 0)
p_b1.space_before = Pt(5)

p_b2 = content_tf.add_paragraph()
p_b2.text = bullet_2
p_b2.font.size = Pt(18)
p_b2.font.color.rgb = RGBColor(0, 0, 0)
p_b2.space_before = Pt(5)

p_b3 = content_tf.add_paragraph()
p_b3.text = bullet_3
p_b3.font.size = Pt(18)
p_b3.font.color.rgb = RGBColor(0, 0, 0)
p_b3.space_before = Pt(5)

# Save the presentation
prs.save('output.pptx')