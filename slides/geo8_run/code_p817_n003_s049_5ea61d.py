from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_slide():
    # Create a new presentation
    prs = Presentation()
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Background: Clean, White
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 2. Title: "ELON MUSK"
    # Large, bold, left side
    title_box = slide.shapes.add_textbox(
        Inches(1), 
        Inches(2), 
        Inches(7), 
        Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ELON MUSK"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)
    p.alignment = PP_ALIGN.LEFT

    # 3. Subtitle: "Lahiru Herath"
    # Smaller, below title, left side
    sub_box = slide.shapes.add_textbox(
        Inches(1), 
        Inches(3.8), 
        Inches(7), 
        Inches(1)
    )
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Lahiru Herath"
    p_sub.font.size = Pt(24)
    p_sub.font.color.rgb = RGBColor(60, 60, 60)
    p_sub.alignment = PP_ALIGN.LEFT

    # 4. Image: Portrait on Right
    # Positioning on right side
    # Image placeholder 'image.png'
    pic = slide.shapes.add_picture('image.png', Inches(6.5), Inches(1), Inches(5.5), Inches(5.5))
    
    # White border
    pic.line.color.rgb = RGBColor(255, 255, 255)
    pic.line.width = Inches(0.1)

    # 5. Accent: Yellow corner border on top right
    # Interpreted as yellow bars at top and right edges
    accent_color = RGBColor(255, 215, 0) # Yellow

    # Top bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, 
        0, 
        prs.slide_width, 
        Inches(0.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()

    # Right bar
    # Left position = Full Width - Bar Width
    right_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        prs.slide_width - Inches(0.2), 
        0, 
        Inches(0.2), 
        prs.slide_height
    )
    right_bar.fill.solid()
    right_bar.fill.fore_color.rgb = accent_color
    right_bar.line.fill.background()

    # Save
    prs.save('output.pptx')

if __name__ == "__main__":
    create_slide()