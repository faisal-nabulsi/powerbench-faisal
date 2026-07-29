from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Add a blank slide to allow for custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 128)

# Get slide dimensions for positioning calculations
slide_width = prs.slide_width
slide_height = prs.slide_height

# 1. Add the title "PRONUNCIATION ACTIVITY"
# Centered vertically and horizontally
# We define a text box size and calculate the center position
title_box_width = Inches(8)
title_box_height = Inches(1.5)
title_left = (slide_width - title_box_width) / 2
title_top = (slide_height - title_box_height) / 2

title_box = slide.shapes.add_textbox(title_left, title_top, title_box_width, title_box_height)
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "PRONUNCIATION ACTIVITY"
p.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255) # White text

# 2. Add a graphic of a speaking head on the left side of the title
# Using 'image.png' as the placeholder image as per instructions
img_width = Inches(1.5)
img_height = Inches(1.5)
# Position to the left of the title with a small gap
img_left = title_left - img_width - Inches(0.5)
# Vertically center the image to align with the title
img_top = (slide_height - img_height) / 2

slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)

# 3. Include the "FLUENT" logo in the top right corner
# Using 'image.png' as the placeholder image as per instructions
logo_width = Inches(2)
logo_height = Inches(0.8)
# Position in the top right corner with padding
logo_left = slide_width - logo_width - Inches(0.5)
logo_top = Inches(0.5)

slide.shapes.add_picture('image.png', logo_left, logo_top, logo_width, logo_height)

# Save the presentation
prs.save('output.pptx')