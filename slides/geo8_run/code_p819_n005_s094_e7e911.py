from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation object
prs = Presentation()

# Add a blank slide (Layout index 6 is typically the Blank layout in default templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Get slide width for dynamic right-aligned positioning
slide_width = prs.slide_width

# --- 1. Title Section ---
# Title: "Blockchain in Cryptocurrency" at top center
# Dimensions: Left 1", Top 0.5", Width 8" (centered on standard 10" slide), Height 1"
title_box = slide.shapes.add_textbox(
    left=Inches(1),
    top=Inches(0.5),
    width=Inches(8),
    height=Inches(1)
)

title_frame = title_box.text_frame
title_frame.word_wrap = True
title_para = title_frame.paragraphs[0]
title_para.text = "Blockchain in Cryptocurrency"
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 0, 0) # Black text for high contrast
title_para.alignment = PP_ALIGN.CENTER

# --- 2. Content Section ---
# Main content area on the left side for bullet points
# Dimensions: Left 0.5", Top 1.8", Width 6", Height 4.5"
content_box = slide.shapes.add_textbox(
    left=Inches(0.5),
    top=Inches(1.8),
    width=Inches(6),
    height=Inches(4.5)
)

content_frame = content_box.text_frame
content_frame.word_wrap = True

# List of bullet points as requested
bullet_points = [
    "The term blockchain is often used to refer to cryptocurrency. Cryptocurrency is a medium of exchange such as US dollars.",
    "It is just an application in the form of e-currency using blockchain.",
    "It is not governed by any financial institution.",
    "The main difference between blockchain and cryptocurrency is that cryptocurrency is created and held electronically in forms such as a virtual wallet.",
    "It is decentralized and it is not governed by anyone whereas blockchain is an advanced record and it has all information related to cryptocurrency exchanges over a shared system."
]

# Add paragraphs to the text box
first = True
for text in bullet_points:
    if first:
        p = content_frame.paragraphs[0]
        first = False
    else:
        p = content_frame.add_paragraph()
    
    # Prepend a bullet character for visual formatting
    p.text = "• " + text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black text for contrast
    p.space_after = Pt(10)

# --- 3. Image Section ---
# Place the image 'image.png' on the right side
# Dimensions: 3.5" x 3.5"
img_width = Inches(3.5)
img_height = Inches(3.5)
# Position: Align to the right edge with a 0.5" margin
img_left = slide_width - img_width - Inches(0.5)
img_top = Inches(2.0)

try:
    slide.shapes.add_picture('image.png', img_left, img_top, img_width, img_height)
except FileNotFoundError:
    print("Note: 'image.png' was not found in the directory. Skipping image insertion.")

# Save the presentation to 'output.pptx'
prs.save('output.pptx')