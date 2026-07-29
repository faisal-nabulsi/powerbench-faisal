from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_presentation():
    # 1. Setup Presentation
    prs = Presentation()
    # Use Blank layout (index 6) to have full control
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 2. Background: Simple, light
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250) # Light professional background

    # 3. Title: Large, bold, top
    # Centered on top of slide
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Advantages of Blockchain"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
    p.alignment = PP_ALIGN.CENTER

    # 4. Checklist (Left Side)
    left_x = Inches(0.5)
    list_y = Inches(1.5)
    list_width = Inches(4.2)
    list_height = Inches(6.5)

    text_box_left = slide.shapes.add_textbox(left_x, list_y, list_width, list_height)
    tf_list = text_box_left.text_frame
    tf_list.word_wrap = True

    items = [
        "Greater Transparency",
        "Highly secure",
        "Easily traceable",
        "High efficiency and speed",
        "Low cost",
        "Zero percentage of fraud",
        "Extremely volatile"
    ]

    for i, item in enumerate(items):
        if i == 0:
            p = tf_list.paragraphs[0]
        else:
            p = tf_list.add_paragraph()
        # Format as checklist using checkmark character
        p.text = f"✓  {item}"
        p.font.size = Pt(18) # Smaller font than title
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)
        p.line_spacing = Pt(24) # Consistent spacing

    # 5. Diagram (Right Side)
    # Container bounds for the diagram
    right_x_start = Inches(5.2)
    diagram_top = Inches(1.8)
    diagram_width = Inches(4.3)

    # Node 1: Root (Top Center of Right Side)
    root_w = Inches(2.4)
    root_h = Inches(0.7)
    root_x = right_x_start + (diagram_width - root_w) / 2
    
    # Create Root Shape
    shape_root = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        root_x, diagram_top, root_w, root_h
    )
    shape_root.fill.solid()
    shape_root.fill.fore_color.rgb = RGBColor(27, 94, 32) # Dark Green
    shape_root.line.color.rgb = RGBColor(27, 94, 32)
    
    p_root = shape_root.text_frame.paragraphs[0]
    p_root.text = "Blockchain Benefits"
    p_root.font.color.rgb = RGBColor(255, 255, 255)
    p_root.font.bold = True
    p_root.font.size = Pt(14)
    p_root.alignment = PP_ALIGN.CENTER

    # Nodes 2, 3, 4: Leaves (Bottom Row)
    # Grouping benefits into categories for the diagram nodes
    leaves_data = [
        "Security\n& Trust",
        "Transparency\n& Traceability",
        "Efficiency\n& Cost"
    ]
    
    leaf_w = Inches(1.2)
    leaf_h = Inches(0.9)
    # Calculate spacing to center 3 boxes within diagram_width
    gap = (diagram_width - (3 * leaf_w)) / 4
    
    for i, text in enumerate(leaves_data):
        leaf_x = right_x_start + gap + i * (leaf_w + gap)
        leaf_y = diagram_top + root_h + Inches(0.6) # Vertical distance from root

        # Add Connector (Line) first so it is under the box
        # From Root bottom center to Leaf top center
        x_start = root_x + (root_w / 2)
        y_start = diagram_top + root_h
        x_end = leaf_x + (leaf_w / 2)
        y_end = leaf_y
        
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x_start, y_start, x_end, y_end
        )
        connector.line.color.rgb = RGBColor(150, 150, 150)
        connector.line.width = Pt(2)

        # Add Leaf Shape (Green Box)
        shape_leaf = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            leaf_x, leaf_y, leaf_w, leaf_h
        )
        shape_leaf.fill.solid()
        shape_leaf.fill.fore_color.rgb = RGBColor(40, 167, 69) # Consistent Green
        shape_leaf.line.color.rgb = RGBColor(40, 167, 69)

        p_leaf = shape_leaf.text_frame.paragraphs[0]
        p_leaf.text = text
        p_leaf.font.color.rgb = RGBColor(255, 255, 255)
        p_leaf.font.size = Pt(12)
        p_leaf.alignment = PP_ALIGN.CENTER

    prs.save('output.pptx')

if __name__ == "__main__":
    create_presentation()