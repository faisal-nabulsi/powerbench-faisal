from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize the presentation
prs = Presentation()

# Add a slide using the Blank layout (Index 6 is typically Blank)
# Fallback to the first layout if Blank is not available
try:
    slide_layout = prs.slide_layouts[6]
except IndexError:
    slide_layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(slide_layout)

# 1. Set Dark Background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 20, 20)  # Dark Grey Background

# 2. Add Title Box
# Coordinates: Left, Top, Width, Height
title_left = Inches(0.5)
title_top = Inches(0.4)
title_width = Inches(11)
title_height = Inches(1.2)

title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
tf_title = title_box.text_frame
tf_title.word_wrap = True

# Format Title
p_title = tf_title.paragraphs[0]
p_title.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
p_title.font.bold = True
p_title.font.size = Pt(36)  # Large font
p_title.font.color.rgb = RGBColor(255, 255, 255)  # White text
p_title.alignment = PP_ALIGN.LEFT

# 3. Add Content Body Box
content_left = Inches(0.5)
content_top = Inches(1.8)
content_width = Inches(11)
content_height = Inches(5.5)

content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
tf_content = content_box.text_frame
tf_content.word_wrap = True

# Clear the default empty paragraph created by add_textbox
tf_content.paragraphs[0].text = ""

# Define Colors
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_YELLOW = RGBColor(255, 255, 0)  # Bright Yellow for contrast

# Define Content Data
content_data = [
    {
        "header": "1. Effective Communication",
        "bullets": [
            "Strong communication skills encompass the ability to convey ideas clearly, listen actively, and adapt communication styles to different audiences."
        ]
    },
    {
        "header": "2. Problem Solving and Critical Thinking",
        "bullets": [
            "Critical thinking skills enable individuals to analyze information, evaluate options, and make informed decisions.",
            "Problem-solving skills are essential for addressing challenges and finding innovative solutions in a wide range of situations."
        ]
    },
    {
        "header": "3. Emotional Intelligence",
        "bullets": [
            "Emotional intelligence enhances empathy, interpersonal relationships, and the ability to work effectively in teams."
        ]
    },
    {
        "header": "4. Adaptability and Resilience",
        "bullets": [
            "These skills help individuals navigate unexpected setbacks, cope with stress, and embrace change as an opportunity for growth."
        ]
    }
]

# Helper function to add formatted paragraphs
def add_paragraph(text_frame, text, color, size, bold=False, use_bullet=False):
    p = text_frame.add_paragraph()
    run = p.add_run()
    
    # Set text content
    if use_bullet:
        run.text = "\u2022  " + text  # Using unicode bullet for consistency
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    else:
        run.text = text
        p.space_before = Pt(14)
        p.space_after = Pt(6)
    
    # Set font properties
    run.font.color.rgb = color
    run.font.size = size
    run.font.bold = bold

# Populate the Content Body
for item in content_data:
    # Add Header (Yellow, Bold, No Bullet)
    # This acts as the "sub-title" in main content
    add_paragraph(
        tf_content, 
        item["header"], 
        COLOR_YELLOW, 
        Pt(24), 
        bold=True, 
        use_bullet=False
    )
    
    # Add Sub-points (White, Regular, With Bullet)
    for bullet_text in item["bullets"]:
        add_paragraph(
            tf_content, 
            bullet_text, 
            COLOR_WHITE, 
            Pt(18), 
            bold=False, 
            use_bullet=True
        )

# Save the presentation
prs.save('output.pptx')