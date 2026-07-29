from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height

# Set light green background
background = slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xC8, 0xE6, 0xC9)  # light green

# Create subtle hexagonal pattern
def add_hexagon(slide, left, top, size, fill_color):
    hexagon = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, left, top, size, size
    )
    hexagon.fill.solid()
    hexagon.fill.fore_color.rgb = fill_color
    hexagon.line.fill.background()
    hexagon.shadow.inherit = False
    return hexagon

# Create a subtle hex pattern with slightly lighter/darker green hexagons
hex_size = Inches(0.8)
row_spacing = Inches(0.6)
col_spacing = Inches(1.38)
hex_colors = [
    RGBColor(0xBF, 0xDB, 0xBE),
    RGBColor(0xA5, 0xD6, 0xA7),
    RGBColor(0xCD, 0xE1, 0xCD),
    RGBColor(0xB0, 0xE0, 0xB0),
]

for row in range(10):
    for col in range(12):
        left = Emu(col * int(col_spacing)) + (Emu(int(row_spacing / 2)) if row % 2 == 1 else Emu(0))
        top = Emu(row * int(row_spacing))
        color = hex_colors[(row + col) % len(hex_colors)]
        add_hexagon(slide, left, top, hex_size, color)

# Create large white textbox in the middle
textbox_left = Inches(1.5)
textbox_top = Inches(1.0)
textbox_width = Inches(10.3)
textbox_height = Inches(5.5)

textbox = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
txBox = textbox.text_frame
txBox.word_wrap = True

# Style the textbox with white background
text_shape = textbox
text_shape.fill.solid()
text_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
text_shape.line.color.rgb = RGBColor(0x4A, 0x7C, 0x4A)
text_shape.line.width = Pt(2)

# Add rounded corners via XML manipulation
try:
    sp = textbox._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is not None:
                avLst.set('{http://schemas.openxmlformats.org/drawingml/2006/main}round', '160000')
except:
    pass

# Title Section
title_para = txBox.add_paragraph()
title_para.alignment = PP_ALIGN.LEFT
title_para.space_before = Pt(15)
title_para.space_after = Pt(15)
title_run = title_para.add_run()
title_run.text = "NATURAL ENVIRONMENT"
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)  # bold green
title_run.font.size = Pt(32)
title_run.font.name = "Calibri"

# Add a separator line effect with a paragraph
sep_para = txBox.add_paragraph()
sep_para.space_before = Pt(2)
sep_para.space_after = Pt(15)
sep_run = sep_para.add_run()
sep_run.text = "━" * 14
sep_run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
sep_run.font.size = Pt(14)

# Content Section
content_list = [
    "The natural environment means the non human made surroundings conditions in which all biotic and abiotic things exist on earth.",
    "Natural environment can further be divided into 4 parts."
]

for item in content_list:
    para = txBox.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    para.space_before = Pt(8)
    para.space_after = Pt(8)
    para.level = 0
    run = para.add_run()
    run.text = item
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    run.font.name = "Calibri"
    
    # Add custom bullet using character
    bullet_para = txBox.add_paragraph()
    bullet_para.alignment = PP_ALIGN.LEFT
    bullet_para.space_before = Pt(2)
    bullet_para.space_after = Pt(0)
    bullet_run = bullet_para.add_run()
    bullet_run.text = "● " + item
    bullet_run.font.size = Pt(16)
    bullet_run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    bullet_run.font.name = "Calibri"
    # Remove the duplicate paragraph approach - let me redo this

# Clear and redo the textbox content properly
# Reset textbox
txBox.text = ""
txBox.word_wrap = True

# Title
title_para = txBox.add_paragraph()
title_para.alignment = PP_ALIGN.LEFT
title_para.space_before = Pt(15)
title_para.space_after = Pt(10)
title_run = title_para.add_run()
title_run.text = "NATURAL ENVIRONMENT"
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
title_run.font.size = Pt(32)
title_run.font.name = "Calibri"

# Separator
sep_para = txBox.add_paragraph()
sep_para.alignment = PP_ALIGN.LEFT
sep_para.space_before = Pt(2)
sep_para.space_after = Pt(12)
sep_run = sep_para.add_run()
sep_run.text = "━━━━━━━━━━━━━━━━━━━━"
sep_run.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
sep_run.font.size = Pt(14)

# Bullet points
bullet_texts = [
    "The natural environment means the non human made surroundings conditions in which all biotic and abiotic things exist on earth.",
    "Natural environment can further be divided into 4 parts."
]

for i, text in enumerate(bullet_texts):
    para = txBox.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    para.space_before = Pt(6)
    para.space_after = Pt(6)
    run = para.add_run()
    run.text = "  ●  " + text
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    run.font.name = "Calibri"

# Image Section - Bubble with tree image
bubble_left = Inches(5.0)
bubble_top = Inches(6.2)
bubble_size = Inches(2.5)

# Create circular/oval bubble
bubble_shape = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, bubble_left, bubble_top, bubble_size, bubble_size
)
bubble_shape.fill.solid()
bubble_shape.fill.fore_color.rgb = RGBColor(0x4C, 0xAF, 0x50)  # green bubble
bubble_shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
bubble_shape.line.width = Pt(4)

# Add image inside the bubble
image_width = Inches(1.8)
image_height = Inches(1.8)
image_left = bubble_left + (bubble_size - image_width) / 2
image_top = bubble_top + (bubble_size - image_height) / 2

slide.shapes.add_picture('image.png', image_left, image_top, image_width, image_height)

prs.save('output.pptx')