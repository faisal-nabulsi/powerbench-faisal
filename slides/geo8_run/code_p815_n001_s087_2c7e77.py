from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set clean light background for readability
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(247, 248, 250)

# Define dark blue color
dark_blue = RGBColor(11, 32, 75)

# Add dark blue banner at the top
banner = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    left=Inches(0),
    top=Inches(0),
    width=Inches(13.333),
    height=Inches(1.5)
)
banner.fill.solid()
banner.fill.fore_color.rgb = dark_blue
banner.line.fill.background()

# Add title in large white font
title_box = slide.shapes.add_textbox(
    left=Inches(1),
    top=Inches(0.35),
    width=Inches(11),
    height=Inches(0.8)
)
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_para = title_tf.paragraphs[0]
title_para.text = "5 Stages of Development of Media"
title_para.font.size = Pt(40)
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.LEFT

# Add subtitle in slightly smaller dark blue font
sub_box = slide.shapes.add_textbox(
    left=Inches(1.2),
    top=Inches(1.8),
    width=Inches(5),
    height=Inches(0.6)
)
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_para = sub_tf.paragraphs[0]
sub_para.text = "1. ORAL COMMUNICATION"
sub_para.font.size = Pt(28)
sub_para.font.color.rgb = dark_blue
sub_para.font.bold = True

# Add bullet points
bullets_box = slide.shapes.add_textbox(
    left=Inches(1.2),
    top=Inches(2.6),
    width=Inches(5.5),
    height=Inches(2.5)
)
bullets_tf = bullets_box.text_frame
bullets_tf.word_wrap = True

p1 = bullets_tf.paragraphs[0]
p1.text = "\u2022 Language allowed humans to communicate and share information."
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(45, 45, 45)
p1.space_after = Pt(16)

p2 = bullets_tf.add_paragraph()
p2.text = "\u2022 Language became the most important tool for exploring the world and different cultures."
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(45, 45, 45)

# Insert image placeholder on the right side
slide.shapes.add_picture(
    'image.png',
    left=Inches(7.3),
    top=Inches(2),
    width=Inches(5),
    height=Inches(4.6)
)

# Save the presentation
prs.save('output.pptx')