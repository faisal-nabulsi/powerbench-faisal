from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
prs = Presentation()

# Use a blank slide layout to allow custom positioning
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# --- Title Section ---
# 1. Dark Blue Background for Title
title_height = Inches(1.5)
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), slide_width, title_height
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
title_bg.line.fill.background() # Remove border

# 2. Title Text
# Position text box in the center of the title background
title_text_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(0.2), slide_width - Inches(1), Inches(1.1)
)
tf = title_text_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Dynamics of Local and Global Culture"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255) # White text
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# --- Content Section ---
# 1. White Background for Content
content_top = title_height
content_height = slide_height - title_height
content_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), content_top, slide_width, content_height
)
content_bg.fill.solid()
content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255) # White background
content_bg.line.fill.background()

# 2. Content Text
# Add a text box with some padding from the edges
padding = Inches(1)
text_box_left = padding
text_box_top = content_top + padding
text_box_width = slide_width - (padding * 2)
text_box_height = content_height - (padding * 2)

content_text_box = slide.shapes.add_textbox(
    text_box_left, text_box_top, text_box_width, text_box_height
)
tf = content_text_box.text_frame
tf.word_wrap = True

# Bullet 1
p1 = tf.paragraphs[0]
p1.text = "Global flows of culture tend to move more easily around the globe than ever before, especially through non-material digital forms."
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(0, 0, 0)
p1.space_after = Pt(14)

# Bullet 2
p2 = tf.add_paragraph()
p2.text = "There are three perspectives on global cultural flows:"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0, 0, 0)
p2.space_after = Pt(14)

# Bullet 3
p3 = tf.add_paragraph()
p3.text = "These are cultural differentialism, hybridization, and convergence."
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
prs.save('output.pptx')