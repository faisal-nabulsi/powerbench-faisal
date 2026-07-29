from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Use the 'Title and Content' layout (index 1)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# --- Title Configuration ---
title_shape = slide.shapes.title
title_shape.text = "Origin of Blockchain Technology"

# Apply formatting: Large, Bold, Left-Aligned, Professional Font
if title_shape.has_text_frame:
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(40)  # Large font
            run.font.bold = True
            run.font.name = 'Arial' # Clean, professional font

# --- Content Configuration (Bullet Points) ---
# Access the content placeholder (index 1)
body_shape = slide.placeholders[1]
text_frame = body_shape.text_frame

bullet_points = [
    "The blockchain was created by a person (or group of people) using the name (or pseudonym) Satoshi Nakamoto in 2008 to serve as the public distributed ledger for bitcoin cryptocurrency transactions.",
    "The implementation of the blockchain within bitcoin made it the first digital currency to solve the double-spending problem without the need of a trusted authority or central server.",
    "Private blockchains have been proposed for business use."
]

# The content placeholder typically contains a single empty paragraph initially.
# We will reuse this first paragraph and append new ones for subsequent bullets.
paragraphs = list(text_frame.paragraphs)

if paragraphs:
    # Update the first existing paragraph with the first bullet point
    p = paragraphs[0]
    p.text = bullet_points[0]
    p.alignment = PP_ALIGN.LEFT
    # Apply font styling to the run created by setting text
    if p.runs:
        p.runs[0].font.size = Pt(18)  # Smaller font
        p.runs[0].font.name = 'Arial'
else:
    # Safety fallback if the placeholder is unexpectedly empty
    p = text_frame.add_paragraph()
    p.text = bullet_points[0]
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.name = 'Arial'

# Add the remaining bullet points
for text in bullet_points[1:]:
    new_p = text_frame.add_paragraph()
    new_p.text = text
    new_p.alignment = PP_ALIGN.LEFT
    # Apply font styling
    if new_p.runs:
        new_p.runs[0].font.size = Pt(18)
        new_p.runs[0].font.name = 'Arial'

# Save the presentation
prs.save('output.pptx')