from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation with standard 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Add a blank slide (index 6 is typically blank in most templates)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set dark background for the entire slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(25, 25, 25)

# Add main title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
tf = title_shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KEY REASONS WHY ACADEMIC SUCCESS IS IMPORTANT IN SOCIETY"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Helper function to create each vertical section
def add_section(left, title_text, content_text):
    # Add placeholder image on top
    slide.shapes.add_picture('image.png', left, Inches(1.4), Inches(4.0), Inches(3.5))
    
    # Add yellow section title
    title_box = slide.shapes.add_textbox(left, Inches(5.1), Inches(4.0), Inches(0.6))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = RGBColor(255, 215, 0)  # Yellow/Gold for contrast
    p_t.alignment = PP_ALIGN.CENTER
    
    # Add white section content
    content_box = slide.shapes.add_textbox(left, Inches(5.8), Inches(4.0), Inches(1.3))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = content_text
    p_c.font.size = Pt(13)
    p_c.font.color.rgb = RGBColor(255, 255, 255)
    p_c.alignment = PP_ALIGN.CENTER

# Define horizontal positions for three equal columns
col1_left = Inches(0.3)
col2_left = Inches(4.6)
col3_left = Inches(8.9)

# Add the three sections
add_section(col1_left, 
            "Innovation and Technological Advancement", 
            "Academic success fuels innovation, driving progress in various fields, including healthcare, communication, and transportation.")

add_section(col2_left, 
            "Social Mobility and Equality", 
            "Education can promote social mobility and reduce disparities by providing opportunities for academic success.")

add_section(col3_left, 
            "Civic Engagement and Informed Citizenship", 
            "Education promotes critical thinking, enabling individuals to engage in democracy, contribute to community development, and solve social issues.")

# Save the presentation
prs.save('output.pptx')