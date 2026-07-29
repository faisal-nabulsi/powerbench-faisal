from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a new presentation
prs = Presentation()

# Add a blank slide layout
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Set the background to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# 1. Title: "Content" with bold font
# Positioning: Top of the slide
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)

tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Content"
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 0, 0) # Black text

# 2. Content Area: Vertical Yellow Bar + Bullet Points
# Define the start position for the content block
content_top = Inches(2.0)
content_left = Inches(1.0)
bar_height = Inches(4.2) # Height sufficient for the list

# Add Vertical Yellow Bar
bar_width = Inches(0.15)
bar_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    content_left,
    content_top,
    bar_width,
    bar_height
)
bar_shape.fill.solid()
bar_shape.fill.fore_color.rgb = RGBColor(255, 255, 0) # Yellow
bar_shape.line.fill.background() # Remove border

# Add Text Box for Bullets
text_left = Inches(1.3) # Position to the right of the bar
text_width = Inches(7.5)
text_box = slide.shapes.add_textbox(text_left, content_top, text_width, bar_height)

tf = text_box.text_frame
tf.word_wrap = True

# List of bullet points
items = [
    "Brief highlights of Elon Musk",
    "Biography",
    "Early life of Elon Musk",
    "How he came up with his ideas",
    "Elon’s current stage",
    "Obstacles that Elon faced",
    "Lessons that we can learn from Elon Musk's life",
    "Elon’s Future Plans"
]

# Add each item as a paragraph in the text box
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0) # Black text
    p.font.name = 'Calibri' # Standard font

# Save the presentation
prs.save('output.pptx')