#!/usr/bin/env python3
"""SlideBench Track B (Malleability Index) — four binary production-readiness checks.

From the SlideBench methodology: Track B scores the exported PPTX 0-4 on
  1. native modularity      — real editable objects, not rasterised images
  2. template compatibility — uses layouts/theme so a deck can be re-themed
  3. asset utility          — pictures/charts are genuine, replaceable assets
  4. export stability       — the file opens and round-trips cleanly

Final SlideBench score = TrackA x (1 + 0.03 x B), so 4/4 is a 12% multiplicative uplift.
This script estimates B for decks we already produce, to see whether that uplift is
"free" for a python-pptx generator.
"""
import glob
import os
import sys

from pptx import Presentation
from pptx.util import Emu  # noqa: F401


def check_native_modularity(prs):
    """Real, individually-editable objects rather than one flattened image."""
    native = pictures = total = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            total += 1
            st = str(sh.shape_type or "")
            if "PICTURE" in st:
                pictures += 1
            elif sh.has_text_frame or "TABLE" in st or "CHART" in st or "AUTO_SHAPE" in st:
                native += 1
    if total == 0:
        return False, "no shapes"
    # fail if the deck is essentially just images
    frac = native / float(total)
    return frac >= 0.5, "%d/%d native objects (%.0f%%)" % (native, total, 100 * frac)


def check_template_compatibility(prs):
    """Slides reference real layouts/masters, so a theme can be swapped in."""
    ok = 0
    for slide in prs.slides:
        try:
            if slide.slide_layout is not None and slide.slide_layout.slide_master is not None:
                ok += 1
        except Exception:
            pass
    n = len(prs.slides)
    return (n > 0 and ok == n), "%d/%d slides bound to a layout+master" % (ok, n)


def check_asset_utility(prs):
    """Any images are genuine embedded assets that can be swapped, not decoration-only."""
    imgs = bad = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if "PICTURE" not in str(sh.shape_type or ""):
                continue
            imgs += 1
            try:
                blob = sh.image.blob
                if not blob or len(blob) < 100:
                    bad += 1
            except Exception:
                bad += 1
    if imgs == 0:
        return True, "no images (vacuously ok)"
    return bad == 0, "%d images, %d unusable" % (imgs, bad)


def check_export_stability(path):
    """Opens cleanly, has slides, and survives a save/reload round-trip."""
    try:
        prs = Presentation(path)
        if len(prs.slides) == 0:
            return False, "0 slides"
        tmp = path + ".roundtrip.pptx"
        prs.save(tmp)
        prs2 = Presentation(tmp)
        n = len(prs2.slides)
        os.remove(tmp)
        return n == len(prs.slides), "round-trip ok (%d slides)" % n
    except Exception as e:
        return False, "unstable: %s" % str(e)[:60]


def score_deck(path):
    try:
        prs = Presentation(path)
    except Exception as e:
        return None, {"open": (False, str(e)[:60])}
    res = {}
    res["native_modularity"] = check_native_modularity(prs)
    res["template_compatibility"] = check_template_compatibility(prs)
    res["asset_utility"] = check_asset_utility(prs)
    res["export_stability"] = check_export_stability(path)
    b = sum(1 for v in res.values() if v[0])
    return b, res


def main():
    d = sys.argv[1] if len(sys.argv) > 1 else "."
    decks = sorted(glob.glob(os.path.join(d, "deck_*.pptx")))
    if not decks:
        decks = sorted(glob.glob(os.path.join(d, "*.pptx")))
    if not decks:
        print("no decks found in", d)
        return
    dist = {0: 0, 1: 0, 2: 0, 3: 0, 4: 0}
    fails = {}
    for p in decks:
        b, res = score_deck(p)
        if b is None:
            dist[0] += 1
            continue
        dist[b] += 1
        for k, (ok, why) in res.items():
            if not ok:
                fails.setdefault(k, []).append(why)
    n = len(decks)
    print("SlideBench Track B (Malleability) on %d decks" % n)
    print("-" * 58)
    for b in (4, 3, 2, 1, 0):
        if dist[b]:
            print("  B=%d : %3d decks (%.0f%%)   uplift x%.2f" % (b, dist[b], 100 * dist[b] / n, 1 + 0.03 * b))
    mean_b = sum(b * c for b, c in dist.items()) / float(n)
    print("-" * 58)
    print("  mean B = %.2f  ->  score uplift x%.3f" % (mean_b, 1 + 0.03 * mean_b))
    if fails:
        print("\n  failing checks:")
        for k, v in sorted(fails.items(), key=lambda kv: -len(kv[1])):
            print("    %-24s %d decks  e.g. %s" % (k, len(v), v[0]))
    else:
        print("\n  no failures — every check passed on every deck")


if __name__ == "__main__":
    main()
