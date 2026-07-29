"""Build a verl GRPO dataset from SlidesBench instructions.

The public AutoPresent repo ships instructions but NOT reference decks (only food.pptx),
so we can't run the exact reference-based metric on most tasks. Instead we derive a
deterministic, program-scored target from each instruction: the set of literal text
strings it requires on the slide (quoted strings + bulleted list items). The reward
(slidesbench_reward.compute_score) executes the model's generated code and scores
text coverage + execution + layout sanity. Faithful-in-spirit to SlidesBench's
program-scored philosophy; the exact reference metric returns once we have reference decks.
"""
import os, re, json, glob, argparse
import pyarrow as pa
import pyarrow.parquet as pq

SYSTEM = (
    "You are an expert presentation designer who builds PowerPoint slides with python-pptx. "
    "Read the instruction and write a single Python script using python-pptx that creates the slide. "
    "The script MUST save the presentation to the file 'output.pptx' in the current directory. "
    "If the instruction asks for an image, photo, logo, or picture, a placeholder image "
    "file named 'image.png' is available in your working directory -- use it via "
    "slide.shapes.add_picture('image.png', left, top, width, height). Do not reference any "
    "other image filename. "
    "The slide canvas is 16:9 widescreen: set prs.slide_width = Inches(13.333) and "
    "prs.slide_height = Inches(7.5) before adding slides, and keep all content inside it. "
    "python-pptx API notes (this environment runs python-pptx 1.0.2 -- these are the ONLY "
    "styling calls that exist; inventing others makes the script crash and the slide score zero): "
    "imports: from pptx.util import Inches, Pt, Emu; from pptx.dml.color import RGBColor; "
    "from pptx.enum.text import PP_ALIGN; from pptx.enum.dml import MSO_THEME_COLOR. "
    "Solid fill: shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x1F,0x3B,0x73). "
    "Gradient (exactly 2 stops, accessed BY INDEX -- there is no .add()): shape.fill.gradient(); "
    "shape.fill.gradient_stops[0].color.rgb = RGBColor(...); shape.fill.gradient_stops[0].position = 0.0; "
    "shape.fill.gradient_stops[1].color.rgb = RGBColor(...); shape.fill.gradient_stops[1].position = 1.0; "
    "shape.fill.gradient_angle = 45.0. "
    "There is NO pptx.dml.gradient module, NO MSO_GRADIENT_DIRECTION, NO GradientDirection, "
    "and FillFormat has no .picture attribute. Text: tf = shape.text_frame; tf.word_wrap = True; "
    "p = tf.paragraphs[0] or tf.add_paragraph(); p.text = '...'; p.font.size = Pt(28); p.font.bold = True; "
    "p.alignment = PP_ALIGN.CENTER. Prefer solid fills; only use a gradient if you follow the exact form above. "
    "Put the final script in one ```python code block."
)

def extract_required_texts(instr: str) -> list:
    req = []
    # 1) double-quoted literal strings
    req += re.findall(r'"([^"]{1,60})"', instr)
    # 2) bulleted list items that look like literal content (proper-noun-ish, short)
    for line in instr.splitlines():
        m = re.match(r'^\s*[-*]\s+(.+)$', line)
        if m:
            item = m.group(1).strip().rstrip('.')
            # keep short, non-imperative items (likely literal names/labels)
            if 2 <= len(item) <= 40 and not re.match(r'(?i)^(use|add|place|ensure|make|set|include|align|create|keep)\b', item):
                req.append(item)
    # dedup, drop empties/very short
    seen, out = set(), []
    for t in req:
        t = t.strip()
        if len(t) >= 2 and t.lower() not in seen:
            seen.add(t.lower()); out.append(t)
    return out

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--examples_dir", required=True)
    ap.add_argument("--out_dir", required=True)
    ap.add_argument("--n_eval", type=int, default=45)
    args = ap.parse_args()

    rows = []
    for instr_path in sorted(glob.glob(os.path.join(args.examples_dir, "*", "slide_*", "instruction_high_level.txt"))):
        instr = open(instr_path).read().strip()
        req = extract_required_texts(instr)
        if len(req) < 1:
            continue  # need at least one checkable target for signal
        topic = instr_path.split(os.sep)[-3]
        slide = instr_path.split(os.sep)[-2]
        rows.append({
            "data_source": "slidesbench",
            "prompt": [{"role": "system", "content": SYSTEM},
                       {"role": "user", "content": instr + "\n\nSave the presentation to 'output.pptx'."}],
            "ability": "slide_generation",
            "reward_model": {"style": "rule", "ground_truth": json.dumps({"required_texts": req})},
            "extra_info": {"task": f"{topic}/{slide}", "n_required": len(req)},
        })

    # deterministic shuffle (fixed seed) then split
    import random
    random.Random(0).shuffle(rows)
    for i, r in enumerate(rows):
        r["extra_info"]["index"] = i
    eval_rows = rows[:args.n_eval]
    train_rows = rows[args.n_eval:]
    os.makedirs(args.out_dir, exist_ok=True)
    for split, data in [("train", train_rows), ("test", eval_rows)]:
        pq.write_table(pa.Table.from_pylist(data), os.path.join(args.out_dir, f"{split}.parquet"))
    avg_req = sum(len(json.loads(r["reward_model"]["ground_truth"])["required_texts"]) for r in rows)/max(1,len(rows))
    print(f"built {len(rows)} tasks: train {len(train_rows)}, eval {len(eval_rows)}; avg required_texts/task={avg_req:.1f}")
    print("sample required_texts:", json.loads(rows[0]["reward_model"]["ground_truth"])["required_texts"][:6])

if __name__ == "__main__":
    main()
