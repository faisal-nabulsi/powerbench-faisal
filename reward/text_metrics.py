"""
text_metrics.py -- accurate, render-free text measurement for the PPTX geometric reward.

WHY THIS EXISTS
---------------
The incumbent estimator assumes every glyph advances 0.5 * font_size. That is a
*character-count* approximation: it cannot tell "WWWWWWWWWW" from "iiiiillllll",
and it cannot know that words do not split at line ends. This module measures the
real advance width with Pillow's FreeType metrics instead, wraps at real word
boundaries, and reads font size from BOTH run level and paragraph level.

MEASURED CALIBRATION (Liberation Sans, which is metric-compatible with Arial --
the family real model decks declare most often; see self-test output):

    string                        pt    real     naive(0.5)   error
    "iiiiillllll"                 18    44.00    99.00        -55.6%
    "Our strategy focuses on..."  18   289.22   333.00        -13.1%
    "Quarterly Business Review"   24   288.09   300.00         -4.0%
    "Implementation Roadmap..."   20   371.38   370.00         +0.4%
    "Key Takeaways"               28   194.06   182.00         +6.6%
    "REVENUE GROWTH ACCELERATED"  32   550.59   416.00        +32.4%
    "WWWWWWWWWW"                  18   169.84    90.00        +88.7%

So the naive ratio is not biased one way -- it is *unreliable per string*, from
-56% to +89%. It happens to average near 0.5 for mixed-case prose, which is
exactly why it survived. It fails hardest on ALL-CAPS titles (+32%), which is a
very common slide element, and that is where text silently runs off the box.

NOTE ON THE 0.5-vs-"360.2px" BRIEF FIGURE: the task brief cites "Quarterly
Business Review" at 24pt as "real 360.2px vs estimated 300px (~20% under)".
That 360.2 reproduces exactly with **LiberationMono-Regular.ttf** (360.16pt --
a monospace face is 0.6 em/char, so 0.6*24*25 = 360.0). Measured with a
*proportional* sans (Liberation Sans / Arial metrics) the same string is
288.09pt, i.e. the naive estimate OVER-states it by 4%. Defaulting to a
monospace face would inflate every width by ~25% and over-penalise the whole
corpus, so this module deliberately prefers a proportional sans and only falls
back to mono if nothing else exists. See PREFERRED_FACES below.

SCORE CONVENTION
----------------
  textfit_shape()   -> [0,1] SCORE, 1.0 = GOOD (text fits).
  text_offcanvas()  -> [0,1] PENALTY, 0.0 = GOOD (nothing off-slide).
                       This inversion is per the assigned spec and matches the
                       incumbent `_text_offcanvas_penalty`. Callers that need a
                       1.0=GOOD score should use text_offcanvas_score().
"""

from __future__ import annotations

import math
import os
from functools import lru_cache
from typing import List, Optional, Sequence, Tuple

# --------------------------------------------------------------------------- #
# constants
# --------------------------------------------------------------------------- #
EMU_PER_PT = 12700.0
DEFAULT_FONT_PT = 18.0          # python-pptx default when nothing sets a size
DEFAULT_LINE_SPACING = 1.2      # line height as a multiple of point size
FALLBACK_CHAR_W_RATIO = 0.50    # only used when NO font can be loaded at all

# PowerPoint's default text-frame internal margins, in points.
# (0.1 in left/right, 0.05 in top/bottom)
DEFAULT_MARGIN_LR_PT = 7.2
DEFAULT_MARGIN_TB_PT = 3.6

# Fonts are loaded once at this pixel size and widths are scaled linearly.
# Advance width is exactly linear in em size, so this is both more accurate than
# re-loading at every integer size (it handles fractional point sizes such as
# 10.5pt, which truncate to 10 if you load the face at int(size)) and far
# cheaper -- one ImageFont object per face instead of one per (face, size).
_REF_SIZE = 1000


# --------------------------------------------------------------------------- #
# font discovery
# --------------------------------------------------------------------------- #
# Ordered preference. Liberation Sans first: it is metric-compatible with Arial,
# and Arial is what real model decks declare. Mono faces are deliberately absent
# -- they are ~25% wider per character and would skew every measurement.
PREFERRED_FACES: Sequence[Tuple[str, str]] = (
    ("LiberationSans-Regular.ttf", "LiberationSans-Bold.ttf"),
    ("DejaVuSans.ttf", "DejaVuSans-Bold.ttf"),
    ("FreeSans.ttf", "FreeSansBold.ttf"),
    ("NotoSans-Regular.ttf", "NotoSans-Bold.ttf"),
    ("Lato-Regular.ttf", "Lato-Bold.ttf"),
    ("Arial.ttf", "Arial_Bold.ttf"),
    ("Helvetica.ttf", "Helvetica-Bold.ttf"),
)

_FONT_DIRS: Sequence[str] = (
    "/usr/share/fonts",
    "/usr/local/share/fonts",
    "/usr/share/matplotlib/mpl-data/fonts/ttf",
    "/Library/Fonts",
    "/System/Library/Fonts",
    os.path.expanduser("~/.fonts"),
    os.path.expanduser("~/.local/share/fonts"),
)

# Map a declared PPTX font name onto a face stem we may have locally.
_FAMILY_ALIASES = {
    "arial": "sans", "helvetica": "sans", "helvetica neue": "sans",
    "liberation sans": "sans", "sans-serif": "sans", "segoe ui": "sans",
    "calibri": "calibri", "carlito": "calibri",
    "verdana": "dejavu", "tahoma": "dejavu", "dejavu sans": "dejavu",
    "times new roman": "serif", "times": "serif", "georgia": "serif",
    "liberation serif": "serif", "garamond": "serif", "serif": "serif",
    "courier new": "mono", "courier": "mono", "consolas": "mono",
    "monaco": "mono", "menlo": "mono", "liberation mono": "mono",
    "monospace": "mono",
}

# stem -> (regular basename candidates, bold basename candidates)
_STEM_FILES = {
    "sans": (("LiberationSans-Regular.ttf", "FreeSans.ttf", "DejaVuSans.ttf"),
             ("LiberationSans-Bold.ttf", "FreeSansBold.ttf", "DejaVuSans-Bold.ttf")),
    "calibri": (("Carlito-Regular.ttf", "LiberationSans-Regular.ttf"),
                ("Carlito-Bold.ttf", "LiberationSans-Bold.ttf")),
    "dejavu": (("DejaVuSans.ttf", "LiberationSans-Regular.ttf"),
               ("DejaVuSans-Bold.ttf", "LiberationSans-Bold.ttf")),
    "serif": (("LiberationSerif-Regular.ttf", "DejaVuSerif.ttf", "FreeSerif.ttf"),
              ("LiberationSerif-Bold.ttf", "DejaVuSerif-Bold.ttf", "FreeSerifBold.ttf")),
    "mono": (("LiberationMono-Regular.ttf", "DejaVuSansMono.ttf", "FreeMono.ttf"),
             ("LiberationMono-Bold.ttf", "DejaVuSansMono-Bold.ttf", "FreeMonoBold.ttf")),
}


def _index_fonts() -> dict:
    """basename -> absolute path, for every TTF/OTF we can see. Walked once."""
    idx: dict = {}
    for d in _FONT_DIRS:
        try:
            if not os.path.isdir(d):
                continue
            for root, _dirs, files in os.walk(d):
                for fn in files:
                    if fn.lower().endswith((".ttf", ".otf")):
                        idx.setdefault(fn, os.path.join(root, fn))
        except Exception:
            continue
    return idx


_FONT_INDEX = _index_fonts()


def _first_present(names: Sequence[str]) -> Optional[str]:
    for n in names:
        p = _FONT_INDEX.get(n)
        if p:
            return p
    return None


def _generic_sans() -> Optional[str]:
    """Last-ditch: any proportional face on the box. Mono/narrow faces are
    ranked last because their metrics are unrepresentative of slide text."""
    best = None
    best_rank = 99
    for fn, path in _FONT_INDEX.items():
        low = fn.lower()
        if "mono" in low or "narrow" in low:
            rank = 3
        elif "sans" in low:
            rank = 0
        elif "serif" in low:
            rank = 2
        else:
            rank = 1
        if rank < best_rank:
            best_rank, best = rank, path
    return best


def _resolve_default_faces() -> Tuple[Optional[str], Optional[str]]:
    for reg, bold in PREFERRED_FACES:
        p = _FONT_INDEX.get(reg)
        if p:
            return p, (_FONT_INDEX.get(bold) or p)
    g = _generic_sans()
    return (g, g) if g else (None, None)


DEFAULT_FONT_PATH, DEFAULT_BOLD_FONT_PATH = _resolve_default_faces()


@lru_cache(maxsize=64)
def _load_font(path: str, size: int):
    """Cached ImageFont per (path, size). Small cache by design: with reference
    scaling we normally need exactly one entry per face."""
    from PIL import ImageFont  # imported lazily so import never hard-fails
    return ImageFont.truetype(path, size)


def _probe_real_metrics() -> bool:
    if not DEFAULT_FONT_PATH:
        return False
    try:
        f = _load_font(DEFAULT_FONT_PATH, _REF_SIZE)
        return float(f.getlength("Mg")) > 0.0
    except Exception:
        return False


#: True when real FreeType metrics are live; False when we fell back to the
#: 0.5 * font_pt * len(text) approximation. Log this so a silent regression to
#: the approximation is visible in training logs.
USING_REAL_METRICS = _probe_real_metrics()

#: Human-readable description of the measuring face, for logging.
FONT_DESCRIPTION = (
    os.path.basename(DEFAULT_FONT_PATH) if (USING_REAL_METRICS and DEFAULT_FONT_PATH)
    else "NONE (0.5*pt char-count fallback)"
)


@lru_cache(maxsize=256)
def _face_path(family: Optional[str], bold: bool) -> Optional[str]:
    """Resolve a declared PPTX font name to a local TTF, or the default face."""
    if not USING_REAL_METRICS:
        return None
    if family:
        stem = _FAMILY_ALIASES.get(str(family).strip().lower())
        if stem:
            reg, bld = _STEM_FILES.get(stem, ((), ()))
            p = _first_present(bld if bold else reg)
            if p:
                return p
    return DEFAULT_BOLD_FONT_PATH if bold else DEFAULT_FONT_PATH


@lru_cache(maxsize=200_000)
def _ref_width(text: str, path: str) -> float:
    """Advance width of `text` at the reference em size, in reference units."""
    return float(_load_font(path, _REF_SIZE).getlength(text))


# --------------------------------------------------------------------------- #
# public: width
# --------------------------------------------------------------------------- #
def measure_text_width(text: str, font_pt: float, bold: bool = False,
                       family: Optional[str] = None) -> float:
    """Rendered advance width of `text` in POINTS.

    Uses PIL ImageFont.truetype(...).getlength(), which returns the true sum of
    glyph advances (including kerning) rather than a character count.

    `family` is optional and additive to the assigned signature: pass a PPTX
    font name (e.g. "Arial", "Courier New") to measure with a matching local
    face. Omitted -> the default proportional sans.
    """
    if not text:
        return 0.0
    try:
        font_pt = float(font_pt)
    except (TypeError, ValueError):
        return 0.0
    if not math.isfinite(font_pt) or font_pt <= 0.0:
        return 0.0

    path = _face_path(family, bool(bold))
    if path is None:
        return FALLBACK_CHAR_W_RATIO * font_pt * len(text)
    try:
        return _ref_width(text, path) * font_pt / _REF_SIZE
    except Exception:
        return FALLBACK_CHAR_W_RATIO * font_pt * len(text)


# --------------------------------------------------------------------------- #
# wrapping
# --------------------------------------------------------------------------- #
def _break_long_word(word: str, font_pt: float, box_w_pt: float,
                     bold: bool, family: Optional[str]) -> List[Tuple[str, float]]:
    """A single word wider than the line. PowerPoint breaks it mid-word, so we do
    too -- otherwise one long URL would report a single impossibly wide line."""
    out: List[Tuple[str, float]] = []
    cur = ""
    cur_w = 0.0
    for ch in word:
        cand = cur + ch
        w = measure_text_width(cand, font_pt, bold, family)
        if cur and w > box_w_pt:
            out.append((cur, cur_w))
            cur, cur_w = ch, measure_text_width(ch, font_pt, bold, family)
        else:
            cur, cur_w = cand, w
    if cur:
        out.append((cur, cur_w))
    return out or [(word, measure_text_width(word, font_pt, bold, family))]


def wrap_lines(text: str, font_pt: float, box_w_pt: float, bold: bool = False,
               family: Optional[str] = None) -> List[Tuple[str, float]]:
    """Greedy word wrap at REAL measured widths. Returns [(line_text, width_pt)].

    Words are the atomic unit -- this is the part a character-count estimator
    cannot model. A ragged right edge means real text needs more lines than
    `ceil(len(text)/chars_per_line)` predicts.
    """
    words = str(text).split()
    if not words:
        return [("", 0.0)]
    if not (box_w_pt and box_w_pt > 0 and math.isfinite(box_w_pt)):
        w = measure_text_width(" ".join(words), font_pt, bold, family)
        return [(" ".join(words), w)]

    space_w = measure_text_width(" ", font_pt, bold, family)
    lines: List[Tuple[str, float]] = []
    cur: List[str] = []
    cur_w = 0.0

    for word in words:
        ww = measure_text_width(word, font_pt, bold, family)
        if ww > box_w_pt:
            if cur:
                lines.append((" ".join(cur), cur_w))
                cur, cur_w = [], 0.0
            chunks = _break_long_word(word, font_pt, box_w_pt, bold, family)
            lines.extend(chunks[:-1])
            cur, cur_w = [chunks[-1][0]], chunks[-1][1]
            continue
        cand_w = ww if not cur else cur_w + space_w + ww
        if cur and cand_w > box_w_pt:
            lines.append((" ".join(cur), cur_w))
            cur, cur_w = [word], ww
        else:
            cur.append(word)
            cur_w = cand_w
    if cur:
        lines.append((" ".join(cur), cur_w))
    return lines or [("", 0.0)]


# --------------------------------------------------------------------------- #
# public: height
# --------------------------------------------------------------------------- #
def estimate_text_height(text: str, font_pt: float, box_w_pt: float,
                         word_wrap: bool = True,
                         line_spacing: float = DEFAULT_LINE_SPACING,
                         bold: bool = False,
                         family: Optional[str] = None) -> float:
    """Height in POINTS required to lay `text` out in a box `box_w_pt` wide.

    Explicit newlines always start a new line; when `word_wrap` is True the text
    is additionally wrapped at word boundaries using measured widths.
    """
    if text is None:
        return 0.0
    try:
        font_pt = float(font_pt)
        line_spacing = float(line_spacing)
    except (TypeError, ValueError):
        return 0.0
    if not math.isfinite(font_pt) or font_pt <= 0:
        return 0.0
    if not math.isfinite(line_spacing) or line_spacing <= 0:
        line_spacing = DEFAULT_LINE_SPACING

    n_lines = 0
    for para in str(text).split("\n"):
        if not para.strip():
            n_lines += 1
            continue
        if word_wrap:
            n_lines += len(wrap_lines(para, font_pt, box_w_pt, bold, family))
        else:
            n_lines += 1
    return n_lines * font_pt * line_spacing


# --------------------------------------------------------------------------- #
# pptx introspection helpers
# --------------------------------------------------------------------------- #
def _pt(length, default: float) -> float:
    """python-pptx Length (EMU) -> points, tolerating None/garbage."""
    try:
        if length is None:
            return default
        return float(length) / EMU_PER_PT
    except (TypeError, ValueError):
        return default


def _shape_max_font_pt(tf) -> float:
    """Largest font size anywhere in the frame, checking BOTH paragraph-level
    (`para.font.size`) and run-level (`run.font.size`).

    The paragraph-level branch is the real bug this replaces: a deck that sets
    64pt via `para.font.size` and never touches a run was read as the 18pt
    default and declared to fit.
    """
    size = 0.0
    try:
        for para in tf.paragraphs:
            try:
                if para.font.size is not None:
                    size = max(size, float(para.font.size.pt))
            except Exception:
                pass
            try:
                for run in para.runs:
                    try:
                        if run.font.size is not None:
                            size = max(size, float(run.font.size.pt))
                    except Exception:
                        pass
            except Exception:
                pass
    except Exception:
        pass
    return size


def _para_spec(para, shape_size: float) -> dict:
    """Per-paragraph size / bold / family / spacing, with inheritance."""
    size = 0.0
    bold = False
    family = None
    try:
        if para.font.size is not None:
            size = max(size, float(para.font.size.pt))
    except Exception:
        pass
    try:
        bold = bool(para.font.bold)
    except Exception:
        pass
    try:
        family = para.font.name or None
    except Exception:
        pass
    try:
        for run in para.runs:
            try:
                if run.font.size is not None:
                    size = max(size, float(run.font.size.pt))
            except Exception:
                pass
            try:
                if run.font.bold:
                    bold = True
            except Exception:
                pass
            try:
                if family is None and run.font.name:
                    family = run.font.name
            except Exception:
                pass
    except Exception:
        pass
    if size <= 0:
        size = shape_size if shape_size > 0 else DEFAULT_FONT_PT

    spacing = DEFAULT_LINE_SPACING
    absolute = None
    try:
        ls = para.line_spacing
        if ls is not None:
            if isinstance(ls, (int, float)):
                spacing = float(ls)
            else:                      # a Length -> absolute points per line
                absolute = float(ls) / EMU_PER_PT
    except Exception:
        pass
    if not math.isfinite(spacing) or spacing <= 0:
        spacing = DEFAULT_LINE_SPACING

    before = after = 0.0
    try:
        if para.space_before is not None:
            before = float(para.space_before) / EMU_PER_PT
    except Exception:
        pass
    try:
        if para.space_after is not None:
            after = float(para.space_after) / EMU_PER_PT
    except Exception:
        pass

    text = ""
    try:
        text = "".join(r.text or "" for r in para.runs)
        if not text:
            text = para.text or ""
    except Exception:
        try:
            text = para.text or ""
        except Exception:
            text = ""

    align = None
    try:
        align = str(para.alignment) if para.alignment is not None else None
    except Exception:
        pass

    return {"text": text, "size": size, "bold": bold, "family": family,
            "spacing": spacing, "absolute": absolute,
            "before": max(0.0, before), "after": max(0.0, after),
            "align": align}


def _frame_geometry(shape):
    """(box_w_pt, box_h_pt, margin_l, margin_t, word_wrap) or None."""
    try:
        tf = shape.text_frame
        w, h = shape.width, shape.height
        if not w or not h:
            return None
        w = float(w)
        h = float(h)
        if w <= 0 or h <= 0 or not math.isfinite(w) or not math.isfinite(h):
            return None
    except Exception:
        return None

    ml = _pt(getattr(tf, "margin_left", None), DEFAULT_MARGIN_LR_PT)
    mr = _pt(getattr(tf, "margin_right", None), DEFAULT_MARGIN_LR_PT)
    mt = _pt(getattr(tf, "margin_top", None), DEFAULT_MARGIN_TB_PT)
    mb = _pt(getattr(tf, "margin_bottom", None), DEFAULT_MARGIN_TB_PT)

    box_w = w / EMU_PER_PT - (ml + mr)
    box_h = h / EMU_PER_PT - (mt + mb)
    # Degenerate/negative after margins: fall back to the raw extent so a very
    # small box does not produce an infinite ratio.
    if box_w <= 1.0:
        box_w = max(1.0, w / EMU_PER_PT)
        ml = 0.0
    if box_h <= 1.0:
        box_h = max(1.0, h / EMU_PER_PT)
        mt = 0.0

    wrap = True
    try:
        if tf.word_wrap is False:      # None means "inherit" -> wrapping
            wrap = False
    except Exception:
        pass
    return box_w, box_h, ml, mt, wrap


def _layout(shape):
    """Lay a shape's text out. Returns (need_h_pt, widest_pt, lines, geom) or None."""
    try:
        if not shape.has_text_frame:
            return None
        tf = shape.text_frame
        if not (tf.text or "").strip():
            return None
    except Exception:
        return None

    geom = _frame_geometry(shape)
    if geom is None:
        return None
    box_w, _box_h, _ml, _mt, wrap = geom

    shape_size = _shape_max_font_pt(tf)
    need_h = 0.0
    widest = 0.0
    lines: List[Tuple[str, float, dict]] = []
    try:
        paragraphs = list(tf.paragraphs)
    except Exception:
        return None

    for para in paragraphs:
        spec = _para_spec(para, shape_size)
        text = spec["text"]
        if not text.strip():
            lh = spec["absolute"] or spec["size"] * spec["spacing"]
            need_h += lh + spec["before"] + spec["after"]
            lines.append(("", 0.0, spec))
            continue
        if wrap:
            wl = wrap_lines(text, spec["size"], box_w, spec["bold"], spec["family"])
        else:
            wl = [(text, measure_text_width(text, spec["size"], spec["bold"],
                                            spec["family"]))]
        lh = spec["absolute"] or spec["size"] * spec["spacing"]
        need_h += len(wl) * lh + spec["before"] + spec["after"]
        for t, w in wl:
            widest = max(widest, w)
            lines.append((t, w, spec))
    return need_h, widest, lines, geom


def _ratio_score(ratio: float) -> float:
    """1.0 while it fits; linear decay so 2x the available space scores 0.0."""
    if not math.isfinite(ratio) or ratio <= 1.0:
        return 1.0
    return max(0.0, min(1.0, 1.0 - (ratio - 1.0)))


# --------------------------------------------------------------------------- #
# public: scores
# --------------------------------------------------------------------------- #
def textfit_shape(shape) -> float:
    """[0,1] SCORE -- 1.0 = the shape's text fits inside its own box.

    Vertical: measured wrapped height vs box height.
    Horizontal: when `word_wrap` is False the text does NOT wrap and runs off
    sideways, so we penalise by how far the widest line exceeds the box width.
    Both are evaluated and the worse of the two wins.

    Returns 1.0 for shapes with nothing to measure (no text frame, empty text,
    unusable geometry) so it is safe to average. Use textfit_shape_opt() if you
    want None for "not applicable" -- that is the drop-in shape for the
    incumbent `_textfit_shape` call site, which filters None before averaging.
    """
    v = textfit_shape_opt(shape)
    return 1.0 if v is None else v


def textfit_shape_opt(shape) -> Optional[float]:
    """As textfit_shape(), but None when there is nothing to measure."""
    try:
        laid = _layout(shape)
        if laid is None:
            return None
        need_h, widest, _lines, (box_w, box_h, _ml, _mt, wrap) = laid

        score = _ratio_score(need_h / max(1e-6, box_h))
        if not wrap and widest > 0:
            score = min(score, _ratio_score(widest / max(1e-6, box_w)))
        return max(0.0, min(1.0, score))
    except Exception:
        return None


def text_offcanvas(shape, slide_w_emu: float, slide_h_emu: float) -> float:
    """[0,1] PENALTY -- 0.0 = fine, 1.0 = severe. Text that overflows its box and
    then runs off the SLIDE.

    The box itself can sit legally on-canvas while the rendered text spills past
    the bottom or right edge, so a box-based overflow check reports a clean 1.00
    on exactly the decks a reader sees cut off. We lay the text out for real and
    ask where the glyphs actually land, honouring vertical anchor and paragraph
    alignment.
    """
    try:
        sw = float(slide_w_emu) / EMU_PER_PT
        sh = float(slide_h_emu) / EMU_PER_PT
        if not (math.isfinite(sw) and math.isfinite(sh)) or sw <= 0 or sh <= 0:
            return 0.0

        laid = _layout(shape)
        if laid is None:
            return 0.0
        need_h, _widest, lines, (box_w, box_h, ml, mt, _wrap) = laid

        left = float(shape.left) / EMU_PER_PT
        top = float(shape.top) / EMU_PER_PT
        if not (math.isfinite(left) and math.isfinite(top)):
            return 0.0

        # --- vertical: where does the text block start? ---
        anchor = ""
        try:
            anchor = str(shape.text_frame.vertical_anchor or "")
        except Exception:
            pass
        if "MIDDLE" in anchor:
            y0 = top + mt + (box_h - need_h) / 2.0
        elif "BOTTOM" in anchor:
            y0 = top + mt + (box_h - need_h)
        else:
            y0 = top + mt
        v_spill = max(0.0, (y0 + need_h) - sh) + max(0.0, -y0)
        v_pen = v_spill / max(1e-6, need_h)

        # --- horizontal: per line, honouring alignment ---
        h_pen = 0.0
        for _t, lw, spec in lines:
            if lw <= 0:
                continue
            align = spec.get("align") or ""
            if "CENTER" in align:
                x0 = left + ml + (box_w - lw) / 2.0
            elif "RIGHT" in align or "END" in align:
                x0 = left + ml + (box_w - lw)
            else:
                x0 = left + ml
            spill = max(0.0, (x0 + lw) - sw) + max(0.0, -x0)
            h_pen = max(h_pen, spill / max(1e-6, lw))

        return max(0.0, min(1.0, max(v_pen, h_pen)))
    except Exception:
        return 0.0


def text_offcanvas_score(shape, slide_w_emu: float, slide_h_emu: float) -> float:
    """Convenience: 1.0 = GOOD form of text_offcanvas(), for callers that keep
    every metric on the '1.0 is good' convention."""
    return 1.0 - text_offcanvas(shape, slide_w_emu, slide_h_emu)


# --------------------------------------------------------------------------- #
# self-test
# --------------------------------------------------------------------------- #
if __name__ == "__main__":  # pragma: no cover
    import glob
    import statistics
    import sys

    def rule(t):
        print("\n" + "=" * 78)
        print(t)
        print("=" * 78)

    rule("0. FONT RESOLUTION")
    print(f"USING_REAL_METRICS = {USING_REAL_METRICS}")
    print(f"regular face       = {DEFAULT_FONT_PATH}")
    print(f"bold face          = {DEFAULT_BOLD_FONT_PATH}")
    print(f"fonts indexed      = {len(_FONT_INDEX)}")

    rule("1. measure_text_width: real metrics vs 0.5*size*len")
    cases = [
        ("Quarterly Business Review", 24), ("Quarterly Business Review", 44),
        ("REVENUE GROWTH ACCELERATED", 32), ("Key Takeaways", 28),
        ("Our strategy focuses on three pillars", 18),
        ("Implementation Roadmap and Next Steps", 20),
        ("Q1 2026 Results", 40), ("iiiiillllll", 18), ("WWWWWWWWWW", 18),
    ]
    print(f"{'string':38s} {'pt':>4s} {'real':>9s} {'naive':>9s} {'diff':>8s} "
          f"{'bold':>9s}")
    print("-" * 78)
    errs = []
    for s, pt in cases:
        real = measure_text_width(s, pt)
        naive = 0.5 * pt * len(s)
        bold = measure_text_width(s, pt, bold=True)
        d = 100.0 * (real - naive) / naive
        errs.append(abs(d))
        print(f"{s[:38]:38s} {pt:4d} {real:9.2f} {naive:9.2f} {d:+7.1f}% {bold:9.2f}")
    print("-" * 78)
    print(f"mean |error| of the 0.5 approximation: {statistics.mean(errs):.1f}%  "
          f"(range {min(errs):.1f}% .. {max(errs):.1f}%)")

    print("\nThe brief's '360.2' benchmark, resolved:")
    for nm in ("LiberationMono-Regular.ttf", "LiberationSans-Regular.ttf"):
        p = _FONT_INDEX.get(nm)
        if p:
            w = _ref_width("Quarterly Business Review", p) * 24.0 / _REF_SIZE
            print(f"  {nm:32s} 24pt -> {w:7.2f} pt")

    rule("2. SYNTHETIC DECKS (a) fits  (b) crammed paragraph  (c) 64pt at PARA level")
    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches, Pt

        def deck(path, build):
            prs = Presentation()
            prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            build(sl)
            prs.save(path)
            return Presentation(path)

        def a(sl):
            tb = sl.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.add_run().text = "Key Takeaways"
            p.runs[0].font.size = Pt(18)

        def b(sl):
            tb = sl.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.4))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.add_run().text = (
                "Our operating strategy for the coming fiscal year focuses on three "
                "pillars: disciplined cost management, sustained investment in the "
                "core platform, and a deliberate expansion into adjacent enterprise "
                "segments where we already hold a credible right to win.")
            p.runs[0].font.size = Pt(18)

        def c(sl):
            tb = sl.shapes.add_textbox(Inches(1), Inches(1), Inches(2.5), Inches(0.8))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]
            p.font.size = Pt(64)            # PARAGRAPH level only -- no run size
            p.add_run().text = "Annual Revenue Growth"

        specs = [("a) short text, big box", a, "~1.0"),
                 ("b) long paragraph, one-line box", b, "low"),
                 ("c) 64pt set at PARA level, small box", c, "low")]
        print(f"{'case':40s} {'expect':>7s} {'score':>7s}   detail")
        print("-" * 78)
        for i, (name, fn, exp) in enumerate(specs):
            path = f"/tmp/tm_selftest_{i}.pptx"
            prs = deck(path, fn)
            for shp in prs.slides[0].shapes:
                if not shp.has_text_frame:
                    continue
                sc = textfit_shape(shp)
                laid = _layout(shp)
                need, wid, _ln, (bw, bh, _, _, wr) = laid
                oc = text_offcanvas(shp, prs.slide_width, prs.slide_height)
                print(f"{name:40s} {exp:>7s} {sc:7.3f}   "
                      f"need_h={need:6.1f}pt box_h={bh:6.1f}pt "
                      f"widest={wid:6.1f}pt box_w={bw:6.1f}pt offcanvas={oc:.2f}")

        # paragraph-level regression check, stated explicitly
        prs = Presentation("/tmp/tm_selftest_2.pptx")
        shp = [s for s in prs.slides[0].shapes if s.has_text_frame][0]
        print(f"\n  regression check -- font size seen by _shape_max_font_pt: "
              f"{_shape_max_font_pt(shp.text_frame):.0f}pt "
              f"(run-level-only reader would have said {DEFAULT_FONT_PT:.0f}pt)")
    except Exception as e:
        print("SYNTHETIC DECK TEST FAILED:", type(e).__name__, e)
        raise

    rule("3. DISTRIBUTION OVER REAL GALLERY DECKS")
    try:
        from pptx import Presentation
        files = sorted(glob.glob(
            "/home/ubuntu/powerbench/agentic/gallery/deck_*.pptx"))[:20]
        fit, ocp, per_deck = [], [], []
        n_shapes = 0
        for fp in files:
            try:
                prs = Presentation(fp)
            except Exception:
                continue
            dvals = []
            for sl in prs.slides:
                for shp in sl.shapes:
                    v = textfit_shape_opt(shp)
                    if v is None:
                        continue
                    n_shapes += 1
                    fit.append(v)
                    dvals.append(v)
                    ocp.append(text_offcanvas(shp, prs.slide_width, prs.slide_height))
            if dvals:
                per_deck.append(sum(dvals) / len(dvals))

        def stats(v, label):
            if not v:
                print(f"{label}: no data")
                return
            sd = statistics.pstdev(v) if len(v) > 1 else 0.0
            print(f"{label}: n={len(v):4d} mean={statistics.mean(v):.3f} "
                  f"sd={sd:.3f} min={min(v):.3f} max={max(v):.3f}")

        print(f"decks read: {len(files)}, measurable text shapes: {n_shapes}\n")
        stats(fit, "textfit_shape  (per shape)")
        stats(per_deck, "textfit_shape  (per deck) ")
        stats(ocp, "text_offcanvas (per shape)")
        if fit:
            lo = sum(1 for v in fit if v < 0.9)
            zero = sum(1 for v in fit if v == 0.0)
            one = sum(1 for v in fit if v >= 0.999)
            print(f"\nshapes < 0.9 : {lo}/{len(fit)} ({100.0*lo/len(fit):.1f}%)")
            print(f"shapes == 0.0: {zero}/{len(fit)} ({100.0*zero/len(fit):.1f}%)")
            print(f"shapes == 1.0: {one}/{len(fit)} ({100.0*one/len(fit):.1f}%)")
        if ocp:
            bad = sum(1 for v in ocp if v > 0.0)
            print(f"shapes with text off-slide: {bad}/{len(ocp)} "
                  f"({100.0*bad/len(ocp):.1f}%)")
        if per_deck and len(per_deck) > 1:
            print(f"\nper-deck spread (the GRPO-relevant number -- a metric that is "
                  f"constant\nacross rollouts contributes no gradient): "
                  f"sd={statistics.pstdev(per_deck):.3f}, "
                  f"{min(per_deck):.3f}..{max(per_deck):.3f}")
    except Exception as e:
        print("GALLERY TEST FAILED:", type(e).__name__, e)
        raise

    rule("4. FALLBACK PATH (no font available)")
    _saved = (DEFAULT_FONT_PATH, DEFAULT_BOLD_FONT_PATH, USING_REAL_METRICS)
    try:
        USING_REAL_METRICS = False
        _face_path.cache_clear()
        s, pt = "Quarterly Business Review", 24
        print(f"measure_text_width({s!r}, {pt}) with no font = "
              f"{measure_text_width(s, pt):.2f} pt (expect {0.5*pt*len(s):.2f})")
    finally:
        USING_REAL_METRICS = _saved[2]
        _face_path.cache_clear()
    print("\nSELF-TEST COMPLETE")
