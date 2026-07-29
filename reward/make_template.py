"""Build the corporate template deck the enterprise task makes the model conform to.

Why this file exists
--------------------
Template-conformance rewards (use the layouts, fill the placeholders, inherit the theme,
leave the master alone) can only produce a gradient if the task actually HANDS the model a
template. Without one, every rollout scores the same zero and the metric teaches nothing.
This script builds that template once, deterministically, into assets/template.pptx.

Approach
--------
python-pptx cannot author slide layouts from nothing, so we start from the default template
that ships with python-pptx (``Presentation()`` already carries one master and 11 layouts)
and customise it:

  * slide size -> 16:9 (13.333in x 7.5in), and every explicitly-positioned placeholder on
    the master and on the layouts is rescaled horizontally by 4/3 so content still spans the
    wider canvas instead of hugging the left 75% of it;
  * the theme part (theme1.xml) gets a brand colour scheme and a brand font scheme, so a
    model that inherits the theme automatically renders in brand colours/fonts and a model
    that hard-codes RGBColor/font.name visibly deviates;
  * the slide master background is filled with a theme colour (BACKGROUND_2 -> lt2);
  * three named brand shapes are added to the master: an accent bar and a two-part logo
    (coloured rounded rectangle + wordmark). They are named, so a tamper-check can notice
    them being deleted or recoloured:  BrandAccentBar / BrandLogoMark / BrandLogoText.

Everything a downstream reward or dataset builder needs is exported as a module constant
(BRAND, THEME_FONT_*, MASTER_BRAND_SHAPES, REQUIRED_LAYOUTS) or via ``layout_inventory()``.

Run it:
    /home/ubuntu/powerbench/.venv/bin/python make_template.py            # build + report
    /home/ubuntu/powerbench/.venv/bin/python make_template.py --inspect  # report only
"""

from __future__ import annotations

import argparse
import os
import re
import zipfile

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as _PresentationCls
from pptx.shapes.autoshape import AutoShapeType
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------------------
# Constants (imported by build_enterprise_data.py and available to reward code)
# --------------------------------------------------------------------------------------

DEFAULT_TEMPLATE_PATH = "/home/ubuntu/powerbench/agentic/assets/template.pptx"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

BRAND_NAME = "GOLDILOCKS LABS"
LOGO_INITIALS = "GL"

# Theme colour scheme written into theme1.xml. Keys are the DrawingML slot names.
BRAND = {
    "dk1": "111827",       # ink / body text
    "lt1": "FFFFFF",       # paper
    "dk2": "0B2545",       # deep navy - headings, wordmark
    "lt2": "F2F5F9",       # page tint - master background
    "accent1": "0B5FFF",   # brand blue - accent bar, logo mark
    "accent2": "00A3A3",
    "accent3": "F2A900",
    "accent4": "D64545",
    "accent5": "6B4EFF",
    "accent6": "2E7D32",
    "hlink": "0B5FFF",
    "folHlink": "6B4EFF",
}
BRAND_THEME_NAME = "Goldilocks Corporate"
THEME_FONT_MAJOR = "Segoe UI Semibold"   # headings
THEME_FONT_MINOR = "Segoe UI"            # body

# Named shapes added to the slide master. A tamper-check looks for exactly these.
ACCENT_BAR_NAME = "BrandAccentBar"
LOGO_MARK_NAME = "BrandLogoMark"
LOGO_TEXT_NAME = "BrandLogoText"
MASTER_BRAND_SHAPES = (ACCENT_BAR_NAME, LOGO_MARK_NAME, LOGO_TEXT_NAME)
MASTER_LOGO_SHAPES = (LOGO_MARK_NAME, LOGO_TEXT_NAME)

# Layouts the enterprise task promises the model. All exist in the python-pptx default deck.
REQUIRED_LAYOUTS = (
    "Title Slide",
    "Title and Content",
    "Two Content",
    "Section Header",
    "Blank",
)

# Placeholder types that live on the master/layout but are NOT cloned onto a new slide by
# python-pptx, so the model can never fill them and the prompt must not advertise them.
LATENT_PH_TYPES = (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER)

# Bottom band geometry (inches) for the footer row + logo lockup on the master.
_FOOTER_TOP_IN = 6.86
_FOOTER_H_IN = 0.32
_CONTENT_RIGHT_IN = 12.667   # right edge of the content column after the 4/3 rescale


# --------------------------------------------------------------------------------------
# Theme rewriting (zip level: no dependence on python-pptx internals)
# --------------------------------------------------------------------------------------

def _clr_scheme_xml() -> str:
    def slot(tag: str) -> str:
        return f'<a:{tag}><a:srgbClr val="{BRAND[tag]}"/></a:{tag}>'

    return (
        f'<a:clrScheme name="{BRAND_THEME_NAME}">'
        + "".join(
            slot(k)
            for k in ("dk1", "lt1", "dk2", "lt2",
                      "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                      "hlink", "folHlink")
        )
        + "</a:clrScheme>"
    )


def _font_scheme_xml() -> str:
    return (
        f'<a:fontScheme name="{BRAND_THEME_NAME}">'
        f'<a:majorFont><a:latin typeface="{THEME_FONT_MAJOR}"/>'
        f'<a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
        f'<a:minorFont><a:latin typeface="{THEME_FONT_MINOR}"/>'
        f'<a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
        f"</a:fontScheme>"
    )


def _theme_partname(prs) -> str:
    """Zip entry name of the theme the slide master actually uses (e.g. ppt/theme/theme1.xml)."""
    theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    return str(theme_part.partname).lstrip("/")


def _rewrite_theme(pptx_path: str, theme_name: str) -> None:
    """Swap the colour scheme and font scheme inside the saved .pptx."""
    with zipfile.ZipFile(pptx_path) as zin:
        names = zin.namelist()
        blobs = {n: zin.read(n) for n in names}
    if theme_name not in blobs:
        raise RuntimeError(f"{theme_name} not found in {pptx_path}")

    xml = blobs[theme_name].decode("utf-8")
    new = re.sub(r"<a:clrScheme .*?</a:clrScheme>", lambda m: _clr_scheme_xml(),
                 xml, count=1, flags=re.S)
    if new == xml:
        raise RuntimeError("colour scheme not found in theme XML")
    xml2 = re.sub(r"<a:fontScheme .*?</a:fontScheme>", lambda m: _font_scheme_xml(),
                  new, count=1, flags=re.S)
    if xml2 == new:
        raise RuntimeError("font scheme not found in theme XML")
    blobs[theme_name] = xml2.encode("utf-8")

    tmp = pptx_path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, blobs[n])
    os.replace(tmp, pptx_path)


# --------------------------------------------------------------------------------------
# Geometry helpers
# --------------------------------------------------------------------------------------

def _explicit_xfrm(shape):
    """Return the shape's own <a:xfrm>, or None when it inherits its position.

    Reading ``shape.left`` is not enough: layout placeholders report the value they
    INHERIT from the master, and writing it back would freeze an explicit position onto
    every layout. We only rescale shapes that carry their own geometry.
    """
    el = shape._element
    spPr = el.find(qn("p:spPr"))
    if spPr is None:
        return el.find(qn("a:xfrm"))  # graphicFrame / pic style geometry
    return spPr.find(qn("a:xfrm"))


def _scale_x(shapes, ratio: float) -> int:
    """Scale left/width of every explicitly-positioned shape. Returns how many changed."""
    n = 0
    for shape in shapes:
        xfrm = _explicit_xfrm(shape)
        if xfrm is None:
            continue
        off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
        if off is None or ext is None:
            continue
        off.set("x", str(int(round(int(off.get("x")) * ratio))))
        ext.set("cx", str(int(round(int(ext.get("cx")) * ratio))))
        n += 1
    return n


def _place_master_footers(master) -> None:
    """Lay the date/footer/slide-number row out across the bottom, clear of the logo."""
    band = {
        PP_PLACEHOLDER.DATE: (0.667, 2.40),
        PP_PLACEHOLDER.FOOTER: (3.40, 4.50),
        PP_PLACEHOLDER.SLIDE_NUMBER: (8.10, 1.25),
    }
    for ph in master.placeholders:
        spec = band.get(ph.placeholder_format.type)
        if spec is None:
            continue
        left_in, width_in = spec
        ph.left, ph.width = Inches(left_in), Inches(width_in)
        ph.top, ph.height = Inches(_FOOTER_TOP_IN), Inches(_FOOTER_H_IN)


def _master_add_autoshape(master, mso_shape, name, left, top, width, height):
    """MasterShapes has no add_shape()/add_textbox(); go through the spTree factory.

    This is the same path SlideShapes.add_shape() takes, minus the slide-only extent
    recalculation, so the resulting <p:sp> is an ordinary autoshape on the master.
    """
    shapes = master.shapes
    sp = shapes._spTree.add_autoshape(
        shapes._next_shape_id, name, AutoShapeType(mso_shape).prst, left, top, width, height
    )
    return shapes._shape_factory(sp)


def _master_add_textbox(master, name, left, top, width, height):
    shapes = master.shapes
    sp = shapes._spTree.add_textbox(shapes._next_shape_id, name, left, top, width, height)
    return shapes._shape_factory(sp)


def _no_outline(shape) -> None:
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:  # pragma: no cover - older python-pptx
        pass


def _add_brand_shapes(master) -> None:
    """Accent bar + logo lockup, drawn on the master so they appear on every slide."""
    bar = _master_add_autoshape(
        master, MSO_SHAPE.RECTANGLE, ACCENT_BAR_NAME,
        Emu(0), Emu(0), Inches(SLIDE_W_IN), Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    _no_outline(bar)

    mark_size = 0.34
    mark_left = _CONTENT_RIGHT_IN - 2.72
    mark = _master_add_autoshape(
        master, MSO_SHAPE.ROUNDED_RECTANGLE, LOGO_MARK_NAME,
        Inches(mark_left), Inches(_FOOTER_TOP_IN), Inches(mark_size), Inches(mark_size),
    )
    mark.fill.solid()
    mark.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    _no_outline(mark)
    tf = mark.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = LOGO_INITIALS
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1

    word_left = mark_left + mark_size + 0.10
    word = _master_add_textbox(
        master, LOGO_TEXT_NAME,
        Inches(word_left), Inches(_FOOTER_TOP_IN),
        Inches(_CONTENT_RIGHT_IN - word_left), Inches(mark_size),
    )
    tf = word.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = BRAND_NAME
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.theme_color = MSO_THEME_COLOR.TEXT_2


# --------------------------------------------------------------------------------------
# Public API
# --------------------------------------------------------------------------------------

def build_template(path: str = DEFAULT_TEMPLATE_PATH) -> str:
    """Build the corporate template at ``path`` and return that path."""
    prs = Presentation()  # python-pptx default deck: 1 master, 11 layouts
    old_w = prs.slide_width
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    ratio = prs.slide_width / float(old_w)

    master = prs.slide_masters[0]
    _scale_x(master.shapes, ratio)
    for layout in prs.slide_layouts:
        _scale_x(layout.shapes, ratio)
    _place_master_footers(master)

    master.background.fill.solid()
    master.background.fill.fore_color.theme_color = MSO_THEME_COLOR.BACKGROUND_2

    _add_brand_shapes(master)

    theme_name = _theme_partname(prs)
    parent = os.path.dirname(os.path.abspath(path))
    if parent:
        os.makedirs(parent, exist_ok=True)
    prs.save(path)
    _rewrite_theme(path, theme_name)
    return path


def layout_inventory(source=DEFAULT_TEMPLATE_PATH) -> list:
    """[{index, name, placeholders:[{idx, name, type}]}] for every layout in the template.

    Only placeholders python-pptx actually clones onto a new slide are listed -- the
    date/footer/slide-number ones stay latent on the layout and the model can never reach
    them, so advertising them in a prompt would make the task unsatisfiable.
    """
    prs = source if isinstance(source, _PresentationCls) else Presentation(source)
    out = []
    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            pf = ph.placeholder_format
            if pf.type in LATENT_PH_TYPES:
                continue
            ph_type = getattr(pf.type, "name", str(pf.type))
            phs.append({"idx": pf.idx, "name": ph.name, "type": ph_type})
        phs.sort(key=lambda d: d["idx"])
        out.append({"index": i, "name": layout.name, "placeholders": phs})
    return out


def describe(path: str = DEFAULT_TEMPLATE_PATH) -> None:
    """Reopen the saved file and print everything a reviewer needs to trust it."""
    prs = Presentation(path)
    w, h = prs.slide_width, prs.slide_height
    print(f"file            : {path} ({os.path.getsize(path)} bytes)")
    print(f"slide size      : {w} x {h} EMU = {w / 914400:.3f}in x {h / 914400:.3f}in "
          f"(aspect {w / float(h):.4f})")
    print(f"masters/layouts : {len(prs.slide_masters)} master, {len(prs.slide_layouts)} layouts")

    master = prs.slide_masters[0]
    bg_fill = master.background.fill
    print(f"master bg       : fill={bg_fill.type}, theme_color={bg_fill.fore_color.theme_color}")

    names = [s.name for s in master.shapes]
    print("master shapes   :")
    for s in master.shapes:
        tag = "  <-- BRAND" if s.name in MASTER_BRAND_SHAPES else ""
        print(f"   {s.name:<28} {s.shape_type} "
              f"[{s.left / 914400:.2f},{s.top / 914400:.2f} "
              f"{s.width / 914400:.2f}x{s.height / 914400:.2f}in]{tag}")
    missing = [n for n in MASTER_BRAND_SHAPES if n not in names]
    print(f"brand shapes    : {'ALL PRESENT ' + str(list(MASTER_BRAND_SHAPES)) if not missing else 'MISSING ' + str(missing)}")

    theme_name = _theme_partname(prs)
    with zipfile.ZipFile(path) as z:
        theme_xml = z.read(theme_name).decode("utf-8")
    hexes = re.findall(r'<a:(dk1|lt1|dk2|lt2|accent[1-6])><a:srgbClr val="([0-9A-Fa-f]{6})"', theme_xml)
    print(f"theme part      : {theme_name}")
    print(f"theme colours   : {', '.join(f'{k}=#{v}' for k, v in hexes)}")
    fonts = re.findall(r'<a:(major|minor)Font><a:latin typeface="([^"]*)"', theme_xml)
    print(f"theme fonts     : {', '.join(f'{k}={v}' for k, v in fonts)}")

    print("layouts         :")
    for lay in layout_inventory(prs):
        phs = ", ".join(f"{p['idx']}:{p['name']} ({p['type']})" for p in lay["placeholders"]) or "(none)"
        print(f"   [{lay['index']:>2}] {lay['name']:<26} -> {phs}")
    have = {lay["name"] for lay in layout_inventory(prs)}
    lacking = [n for n in REQUIRED_LAYOUTS if n not in have]
    print(f"required layouts: {'ALL PRESENT' if not lacking else 'MISSING ' + str(lacking)}")


def main() -> None:
    ap = argparse.ArgumentParser(description="Build the corporate template deck.")
    ap.add_argument("--path", default=DEFAULT_TEMPLATE_PATH)
    ap.add_argument("--inspect", action="store_true", help="report an existing file, do not rebuild")
    args = ap.parse_args()
    if not args.inspect:
        build_template(args.path)
        print(f"built {args.path}\n")
    describe(args.path)


if __name__ == "__main__":
    main()
