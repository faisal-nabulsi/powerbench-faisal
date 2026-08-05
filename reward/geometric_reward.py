"""
Geometric PPTX reward — the ungameable replacement for the text-coverage proxy.
RENDER-FREE version: everything is computed from python-pptx, no soffice/pixels.

WHY RENDER-FREE
---------------
The pixel-based whitespace metric required rendering every rollout with soffice. At 64
rollouts/step that spawned 64 concurrent soffice processes, which saturated the box CPU:
~83 min/step and the machine became unresponsive. The content-density signal we actually
need (is the slide appropriately full — not empty, not crammed) can be estimated directly
from the deck's text volume and shape coverage, with no rendering. ~50x faster and it does
not fork a subprocess storm.

WHY THIS REWARD RESISTS GAMING
------------------------------
Every metric is computed from slide GEOMETRY / text VOLUME, never from text presence:
  collision  overlapping shapes                -> padding overlaps -> worse
  overflow   shapes spilling off the canvas    -> padding spills   -> worse
  imbalance  visual centroid off-centre        -> padding unbalances -> worse
  density    text volume in a healthy band     -> padding = too much text -> worse
                                                -> empty  = too little     -> worse
Padding a slide moves EVERY metric the wrong way. That is what "hard to game" means.

DESIGN RULES (AeSlides arXiv 2604.22840 + our failure analysis)
  1. Invalid output scores 0.0 — NEVER negative. A negative floor makes every rollout in a
     group identical -> zero relative advantage -> zero gradient (arXiv 2601.03525). That is
     how the previous run died.
  2. Every metric is on a common [0,1] scale (1.0 = good).
  3. density is weighted 2x the shape metrics: the shape metrics only measure ABSENCE OF
     DEFECTS, which a near-empty slide satisfies for free; density is the one metric that
     requires real, appropriately-sized content, so it separates a real slide from a lazy one.
  4. Weighted SUM, not product: a product sends most decks to ~0, collapsing group variance
     (the gradient GRPO needs). Smooth partial credit keeps the spread.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# NOTE on "aspect": its variance is ~0 *because the model currently complies* (the system
# prompt states the 16:9 canvas), NOT because it is useless. It is kept at a small weight as a
# REGRESSION GUARD: training from a fresh checkpoint, a different base model, or dropping the
# prompt hint all make it vary again, and it is the only metric that catches a drift back to
# 4:3. A constant-because-satisfied metric is not the same as a constant-because-uninformative
# one -- only the latter should be dropped.
WEIGHTS = {
    # --- v2 (2026-07-29) --------------------------------------------------------
    # Re-weighted against the 45 blind-graded audit decks. rho = Spearman vs a
    # designer's grade. The v1 grader scored rho -0.139 overall: WORSE than chance.
    #
    # "content" is the single strongest signal found (rho +0.544, AUC 0.826) and is the
    # term that fixes the structural failure: every other metric is a DEFECT detector,
    # and a near-empty slide has no defects, so v1 paid the model to delete content.
    # Measured consequence of that: over 30 GRPO steps text/slide fell 41% and font size
    # 54% while the score rose 95%.
    #
    # collision (rho -0.534) and textfit (rho -0.431) are the two most ANTI-correlated
    # metrics with designer judgement, so both are demoted.
    "adherence": 0.10,   # NEW (2026-07-31): does the slide say what the instruction asked?
                         # The only content-CORRECTNESS signal; every other metric is form.
                         # Title-level anchor on high-level prompts; a guard against the
                         # generic-template attack a form-only reward cannot see. Weights now
                         # sum to 1.10 and are renormalized by wsum below (auto).
    "content":   0.20,   # NEW, render-based: real ink coverage, band-limited both sides
    "collision": 0.12,   # was 0.19 -- anti-correlated, demoted
    "overflow":  0.10,   # was 0.15
    "density":   0.10,   # was 0.13 -- partly superseded by content
    "imbalance": 0.09,   # was 0.11
    "picfit":    0.09,   # was 0.10
    "textfit":   0.08,   # was 0.13 -- anti-correlated, demoted
    "contrast":  0.08,   # was 0.09
    "clipping":  0.07,   # NEW, render-based: ink cut off at the canvas edge.
                         # Deliberately SMALL -- it is another defect detector, and defect
                         # detectors are what caused the sparseness bias. It exists to catch
                         # what geometry provably cannot see (a text box legally on-canvas
                         # whose rendered text is not), not to carry the calibration.
    "alignment": 0.04,   # was 0.05
    "aspect":    0.03,   # was 0.05 -- constant-because-satisfied regression guard
}

# --- aspect-ratio target (16:9 widescreen) -----------------------------------
# The model calls Presentation() and inherits python-pptx's 4:3 default (10x7.5in)
# instead of setting 16:9 (13.333x7.5in). We reward closeness to 16:9 on a log-ratio
# curve: 16:9 -> 1.0, 16:10 -> ~0.80, 4:3 -> ~0.19, square -> ~0.
ASPECT_TARGET = 16.0 / 9.0
ASPECT_K = 20.0

INVALID_REWARD = 0.0

# Worst-critical-metric multiplier floor: min_critical=0 -> score x0.55; =1 -> unchanged.
SOFT_FLOOR = 0.55

# Emptiness gate: a near-empty deck gets every geometric metric for free, so density alone
# cannot punish it enough. Continuous at the knee (dens=KNEE -> multiplier 1.0, no cliff).
EMPTY_KNEE = 0.25
EMPTY_FLOOR = 0.30

# --- content-density band (characters of text per slide) --------------------
# Tuned against the fixtures: good deck ~232 chars/slide (healthy), padded ~4900
# chars/slide (crammed), empty 0, a one-word "Q3" deck ~2. We want good high, the rest low.
DENS_EMPTY = 0.0      # <= this -> score 0 (blank)
DENS_LO = 60.0        # ramp 0->1 between EMPTY and LO
DENS_HI = 900.0       # full score through the [LO, HI] band
DENS_CRAM = 2600.0    # >= this -> score 0 (crammed / padded)


def _smoothstep(edge0, edge1, x):
    if edge1 == edge0:
        return 1.0 if x >= edge1 else 0.0
    t = (x - edge0) / (edge1 - edge0)
    t = max(0.0, min(1.0, t))
    return t * t * (3.0 - 2.0 * t)


def _slide_chars(slide):
    n = 0
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    n += len(run.text or "")
    return n


def score_density(pptx_path):
    """Render-free content density: is each slide appropriately full?

    Returns a float in [0,1], 1.0 = good. Both extremes (empty and crammed) score low, so
    it punishes BOTH the padding exploit and the empty-slide exploit.
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    if not slides:
        return 0.0
    per = []
    for slide in slides:
        c = _slide_chars(slide)
        low = _smoothstep(DENS_EMPTY, DENS_LO, c)       # 0 when empty, 1 once past LO
        high = 1.0 - _smoothstep(DENS_HI, DENS_CRAM, c)  # 1 in-band, 0 once crammed
        per.append(max(0.0, min(1.0, low * high)))
    return sum(per) / len(per)


def score_chart_content(pptx_path):
    """Render-free: does each chart carry real data, or is it empty chrome?

    WHY THIS EXISTS. The render-based `content` metric is `detail_cov` -- the fraction of
    pixels with local structure (edges/strokes). A chart renders a plot area, both axes,
    gridlines, tick labels and a legend: a dense mesh of thin lines that scores high
    `detail_cov` even when the chart has NO data. The blind study confirmed the hole:
    item_30 ("the chart is empty") scored grader 0.97 / content 0.98 while a human gave it
    6/10. Geometry cannot see this (the chart shape is a single well-placed box) and pixels
    cannot see it (chrome and bars are both edges). Only the DATA can, and python-pptx
    exposes it directly.

    A chart is "empty" if no series holds a value that is not None and not 0. Returns
    {n_charts, n_empty_charts, chart_ok} where chart_ok in [0,1] is the fraction of charts
    carrying real data (1.0 when there are no charts -- text slides are untouched).
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    n_charts = n_empty = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_chart", False):
                continue
            n_charts += 1
            has_data = False
            try:
                for plot in sh.chart.plots:
                    for ser in plot.series:
                        if any(v not in (None, 0) for v in ser.values):
                            has_data = True
                            break
                    if has_data:
                        break
            except Exception:
                has_data = True  # unreadable chart: don't punish, we can't prove it empty
            if not has_data:
                n_empty += 1
    chart_ok = 1.0 if n_charts == 0 else 1.0 - (n_empty / float(n_charts))
    return {"n_charts": n_charts, "n_empty_charts": n_empty, "chart_ok": chart_ok}


# An empty chart's chrome inflates render `content` (detail_cov) to ~0.98. Cap content to
# this floor when ALL of a slide's charts are empty, so chrome-without-data cannot farm the
# one metric that is supposed to certify real content. Small non-zero floor because a title
# / caption may still be present -- we punish the empty chart, not the whole slide to zero.
EMPTY_CHART_CONTENT_FLOOR = 0.15


import re as _re_adh


def _norm_text(s):
    return " ".join(_re_adh.sub(r"[^a-z0-9 ]+", " ", (s or "").lower()).split())


# Minimum on-canvas area (fraction of the slide) and font size for text to count toward
# adherence. Without these, the required strings could be stashed in a 1pt / white / off-canvas
# "ghost" textbox the collector would read but a VIEWER never sees -- reinstating the exact
# generic-template attack adherence exists to block (red-team, 2026-07-31). We can only credit
# text a viewer would actually read.
ADH_MIN_AREA_FRAC = 0.0004   # shape's on-canvas area below this = invisible speck
ADH_MIN_FONT_EMU = 76200     # 6pt in EMU (Pt(6)=76200); runs smaller than this don't count


def _shape_on_canvas_area_ok(sh, slide_w, slide_h):
    """Does the shape occupy a non-trivial area INSIDE the canvas? Rejects off-canvas and
    zero/near-zero-size shapes. Unknown geometry -> True (never punish on missing data)."""
    try:
        L, T, W, H = int(sh.left), int(sh.top), int(sh.width), int(sh.height)
    except Exception:
        return True
    if W <= 0 or H <= 0:
        return False
    ox = max(0, min(L + W, slide_w) - max(L, 0))
    oy = max(0, min(T + H, slide_h) - max(T, 0))
    if ox <= 0 or oy <= 0:
        return False   # bbox does not intersect the canvas at all
    return (ox * oy) >= ADH_MIN_AREA_FRAC * slide_w * slide_h


def _visible_shape_text(sh):
    """Text from a shape's runs, dropping runs whose font is explicitly tiny (<6pt). A run
    with size None inherits a normal size -> counted."""
    out = []
    for para in sh.text_frame.paragraphs:
        psz = para.font.size
        runs = list(para.runs)
        if runs:
            for run in runs:
                sz = run.font.size if run.font.size is not None else psz
                if sz is not None and int(sz) < ADH_MIN_FONT_EMU:
                    continue
                out.append(run.text or "")
        elif para.text:  # text set without explicit runs
            if psz is None or int(psz) >= ADH_MIN_FONT_EMU:
                out.append(para.text)
    return " ".join(out)


def _deck_text_norm(pptx_path):
    """Normalized text a VIEWER would see: only shapes on-canvas at non-trivial size, and only
    runs whose font is not explicitly tiny. This is the adherence haystack -- hidden text must
    not earn adherence credit. (White-on-white is a residual a render-free check can't see; it
    is a narrower, harder-to-discover exploit and would also tank the `contrast` metric.)"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    parts = []
    for slide in prs.slides:
        sw, sh_h = prs.slide_width, prs.slide_height
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            if not _shape_on_canvas_area_ok(sh, sw, sh_h):
                continue
            parts.append(_visible_shape_text(sh))
    return _norm_text(" ".join(parts))


# A required phrase counts as present if it appears verbatim OR >= this fraction of its
# significant (>=3-char) tokens appear -- robust to punctuation/wording drift.
ADHERENCE_TOKEN_FRAC = 0.8


def score_adherence(pptx_path, required_texts):
    """Fraction of the instruction's required strings present on the slide.

    WHY THIS EXISTS. Every other metric scores the FORM of the slide (geometry, ink, text
    volume) and none scores whether it says WHAT IT WAS ASKED TO SAY. That gap is the biggest
    reward-hack surface: a generic, well-formatted, instruction-ignoring template maxes every
    form metric. This ties the slide to the task's own required content -- the dataset already
    carries it in ground_truth as `required_texts`, the reward just never used it.

    PADDING-PROOF (unlike the old text-coverage reward it revives): it rewards presence of the
    SPECIFIC required strings, which generic filler does not contain, and the existing density
    band + emptiness gate already punish the padding that broke text-coverage. No required_texts
    -> 1.0 (nothing to check; never punish for missing ground truth)."""
    if not required_texts:
        return 1.0
    hay = _deck_text_norm(pptx_path)
    hay_words = set(hay.split())   # word-boundary set: 'cost' must not match 'forecast'
    hits = 0
    for t in required_texts:
        nt = _norm_text(t)
        if not nt:
            hits += 1
            continue
        # keep >=3-char tokens AND any token containing a digit (numbers/units are the
        # quantitative specifics an off-instruction deck would drop -- don't let it).
        toks = [w for w in nt.split() if len(w) >= 3 or any(c.isdigit() for c in w)]
        if nt in hay:                          # exact normalized phrase present verbatim
            hits += 1
        elif toks and sum(1 for w in toks if w in hay_words) / len(toks) >= ADHERENCE_TOKEN_FRAC:
            hits += 1
    return hits / len(required_texts)


def _parse_required(ground_truth):
    """Pull required_texts out of the verl ground_truth (JSON str or dict). None if absent."""
    import json
    try:
        gt = json.loads(ground_truth) if isinstance(ground_truth, str) else (ground_truth or {})
        req = gt.get("required_texts")
        return req or None
    except Exception:
        return None


def score_aspect(pptx_path):
    """Render-free 16:9 compliance from the deck's slide dimensions. [0,1], 1.0 = 16:9."""
    import math
    from pptx import Presentation
    prs = Presentation(pptx_path)
    w, h = prs.slide_width, prs.slide_height
    if not w or not h or h <= 0:
        return 0.0
    ratio = float(w) / float(h)
    d = abs(math.log(ratio / ASPECT_TARGET))
    return max(0.0, min(1.0, math.exp(-ASPECT_K * d * d)))


def score_deck(pptx_path, png_paths=None, required_texts=None):
    """Score a single deck. png_paths is accepted but ignored (render-free geometry; the
    render for content/clipping happens internally). required_texts is the task's ground-truth
    content the slide should contain (None/[] -> adherence not checked).

    Returns {'score', 'metrics', 'valid', 'reason'}. An unopenable deck scores 0.0.
    """
    if not pptx_path or not os.path.isfile(pptx_path):
        return {"score": INVALID_REWARD, "metrics": {}, "valid": False,
                "reason": "no pptx produced"}

    try:
        from geom_shapes import score_shapes
        shape_scores = score_shapes(pptx_path)
    except Exception as e:
        return {"score": INVALID_REWARD, "metrics": {}, "valid": False,
                "reason": "pptx unopenable: %s" % e}

    metrics = {}
    for k in ("collision", "overflow", "imbalance", "textfit", "contrast", "picfit", "alignment"):
        metrics[k] = float(shape_scores.get(k, 0.0))
    try:
        metrics["density"] = float(score_density(pptx_path))
    except Exception:
        metrics["density"] = 0.0
    try:
        metrics["aspect"] = float(score_aspect(pptx_path))
    except Exception:
        metrics["aspect"] = 0.0
    try:
        metrics["adherence"] = float(score_adherence(pptx_path, required_texts))
    except Exception:
        metrics["adherence"] = 1.0   # never punish the slide for a scorer error

    # ---- render-based metrics (v2) --------------------------------------------
    # Geometry reads the shape BOX; a viewer sees the rendered INK. That gap is not a
    # detail -- a deck whose third bullet is visibly cut off at the slide edge scored
    # collision=1.000, overflow=1.000, textfit=1.000, total 0.9204, because its text box
    # was legally inside the canvas. Only pixels can see that.
    # Cost: 0.34 s/deck batched = ~22 s on a 64-rollout step = 0.8% of a 44-min step.
    # The "rendering is too slow for RL" assumption was never measured and was wrong.
    # If rendering fails we fall back to neutral 1.0 rather than crashing: a reward
    # exception kills a training run.
    try:
        import render_metrics as _rm
        _r = _rm.score_render(pptx_path)
        if _r is None:
            _r = _rm.score_render(pptx_path, use_cache=False)  # retry once: transient soffice flake
    except Exception:
        _r = None
    if _r:
        metrics["content"] = float(_r["content"])
        metrics["clipping"] = float(_r["clipping"])
        render_ok = True
    else:
        # A non-rendering deck shows the VIEWER nothing, so it must be CONDEMNED, not given a
        # neutral score. The old 0.5 fallback was a reward-hack surface (red-team 2026-07-31):
        # forcing a render failure (soffice timeout / malformed-but-openable OOXML) banked
        # content=clipping=0.5 -- BETTER than an honestly-blank rendered slide (content~0) --
        # and lifted the soft-min floor from 0.55 to 0.775, and dodged the emptiness gate too.
        # content=0.0 is a critical-gate member, so a failed render now floors the score (x0.55).
        # With the per-render profile fix + one retry above, a persistent failure is almost
        # always the deck's own fault, so scoring it 0 is correct, not unfair.
        metrics["content"] = 0.0
        metrics["clipping"] = 0.0
        render_ok = False

    # ---- empty-chart correction (render-free) ---------------------------------
    # detail_cov cannot tell a chart's data from its chrome. Read the series data and,
    # when all of a slide's charts are empty, cap the render `content` metric so the
    # chrome cannot certify content that is not there. `chart_ok` also enters the critical
    # set below, so an empty chart trips the soft-min gate the way any catastrophe does.
    try:
        _ch = score_chart_content(pptx_path)
    except Exception:
        _ch = {"n_charts": 0, "n_empty_charts": 0, "chart_ok": 1.0}
    metrics["chart_ok"] = float(_ch["chart_ok"])
    if _ch["n_charts"] and _ch["chart_ok"] < 1.0:
        # blend toward the empty-chart floor in proportion to how many charts are empty
        cap = EMPTY_CHART_CONTENT_FLOOR + (1.0 - EMPTY_CHART_CONTENT_FLOOR) * _ch["chart_ok"]
        metrics["content"] = min(metrics["content"], cap)

    total = wsum = 0.0
    for k, w in WEIGHTS.items():
        if k in metrics:
            total += w * max(0.0, min(1.0, metrics[k]))
            wsum += w
    score = (total / wsum) if wsum else INVALID_REWARD

    # SOFT-MIN GATE (from hand-audit): a weighted mean lets ONE catastrophic defect be
    # averaged away by five good metrics -- content fully off-screen (overflow=0.00) still
    # scored 0.65. We multiply by a factor driven by the WORST critical metric, so a
    # catastrophe roughly halves the score. Kept multiplicative-but-smooth (never a cliff,
    # never zero) so group variance -- the GRPO gradient -- survives.
    # `clipping` joins the critical set in v2: rendered content cut off at the canvas edge
    # is information the viewer never sees, which is exactly the catastrophe this gate is
    # for. Without it the falsifying deck only fell 0.9204 -> 0.8649, because clipping
    # carries just 0.07 of the weighted sum.
    # v2 critical set. Two changes, both measured:
    #  + content  -- a slide that says nothing is the catastrophe this reward missed for 30
    #                steps. On food/slide_22 the model replaced real quiz answers with
    #                "Option A/B/C"; content correctly reads 0.000 there vs 0.480 before.
    #  - textfit  -- REMOVED. It is anti-correlated with designer judgement (rho -0.431) and
    #                is the metric the policy games: shrinking type makes it perfect. It was
    #                gating the GOOD slide (textfit 0.240 -> factor 0.658) while the gutted
    #                one sailed through at 1.000. A gate driven by the exploited metric
    #                rewards the exploit.
    #  + chart_ok -- an all-empty chart is a catastrophe the viewer sees as a blank plot;
    #                render `content` alone was fooled by chrome (item_30: content 0.98).
    #  + adherence -- a slide that ignores the instruction entirely is a catastrophe no form
    #                 metric sees; it belongs in the soft-min gate so an off-instruction slide
    #                 is multiplicatively penalized, not just docked its 0.10 additive weight.
    crit = [metrics.get(k, 1.0) for k in
            ("collision", "overflow", "clipping", "content", "chart_ok", "adherence")
            if k in metrics]
    if crit:
        score *= SOFT_FLOOR + (1.0 - SOFT_FLOOR) * min(crit)

    # ---- emptiness gate -------------------------------------------------------
    # An almost-empty slide trivially satisfies EVERY geometric constraint: nothing can
    # collide, overflow, misfit, or clash if nothing is on the canvas. Before this gate a
    # 1-shape deck scored 0.85 (density 0.07 x 0.13 weight is only a 0.12 penalty, and the
    # critical gate above sees collision/overflow/textfit all at a free 1.00).
    # That is a degenerate optimum, and GRPO finds those.
    # Found by hand-audit: technology/slide_12, "lowkey empty", scored 0.85.
    # Only bites below the knee -- the corpus has a clean gap at 0.15-0.20 and 86% of real
    # decks sit above 0.70, so normal slides are untouched.
    dens = metrics.get("density", 1.0)
    if dens < EMPTY_KNEE:
        score *= EMPTY_FLOOR + (1.0 - EMPTY_FLOOR) * (dens / EMPTY_KNEE)

    return {"score": round(score, 4),
            "render_ok": render_ok,
            "metrics": {k: round(v, 4) for k, v in metrics.items()},
            "valid": True, "reason": "ok"}


def compute_score(data_source, solution_str, ground_truth, extra_info=None, **kwargs):
    """Direct verl hook (unused by the single-turn path, which extracts+runs code first)."""
    extra_info = extra_info or {}
    result = score_deck(extra_info.get("pptx_path"), required_texts=_parse_required(ground_truth))
    out = {"score": result["score"], "valid": 1.0 if result["valid"] else 0.0}
    for k in ("collision", "overflow", "imbalance", "density"):
        out["m_" + k] = float(result["metrics"].get(k, 0.0))
    return out
