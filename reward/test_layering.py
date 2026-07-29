#!/usr/bin/env python3
"""Regression tests for the three grader false positives found by hand-audit + Fable decks.

Written BEFORE the fix (systematic-debugging Phase 4.1). Every FALSE POSITIVE case below
currently fails; every TRUE POSITIVE case currently passes and MUST CONTINUE to pass -- the
whole risk of this fix is that loosening collision/overflow also stops catching real defects.

Run:  ~/powerbench/.venv/bin/python test_layering.py
Exit 0 = both halves hold.
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

FAILS = []


def check(name, ok, detail=""):
    print("  [%s] %-46s %s" % ("PASS" if ok else "FAIL", name, detail))
    if not ok:
        FAILS.append(name)


def _deck(path, build):
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    s = p.slides.add_slide(p.slide_layouts[6])
    build(s)
    p.save(path)
    return path


def _text(s, l, t, w, h, txt, pt=24, rgb=(0, 0, 0)):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True   # add_textbox defaults to wrap="none" (one long line)
    tb.text_frame.text = txt
    for para in tb.text_frame.paragraphs:
        for r in para.runs:
            r.font.size = Pt(pt)
            r.font.color.rgb = RGBColor(*rgb)
    return tb


def _card(s, l, t, w, h, rgb=(0x1E, 0x2A, 0x44)):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(*rgb)
    return c


# ---------------------------------------------------------------- FALSE POSITIVES
def fp_text_on_card(s):
    """A label on a filled card. The most common business-deck idiom there is."""
    _card(s, 1.0, 2.0, 5.5, 3.5)
    _text(s, 1.4, 2.4, 4.7, 1.0, "Quarterly Highlights", rgb=(255, 255, 255))


def fp_decorative_bleed(s):
    """Accent circles deliberately running off the canvas corners."""
    _text(s, 1.5, 3.0, 8.0, 1.2, "Design with edge bleed", pt=28)
    for (l, t, d) in ((-0.9, -0.9, 2.2), (12.5, 5.9, 1.5), (-0.5, 5.9, 1.7)):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFF, 0x45, 0x45)


def fp_rich_two_column(s):
    """Shape-rich but well-built: dark bg, accent bar, 3 bulleted cards, framed image.
    Structurally this is social_media/slide_3, which a human rated GOOD and we scored 0.28."""
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2E)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFF, 0x5A, 0x5A)
    _text(s, 0.6, 0.45, 7.5, 0.9, "GLOBAL MEDIA CULTURES", pt=32, rgb=(255, 255, 255))
    for i, top in enumerate((1.9, 3.65, 5.4)):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), Inches(top + 0.13), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(0xFF, 0x5A, 0x5A)
        _text(s, 1.10, top, 6.5, 1.45, "Section %d\nSupporting detail line." % (i + 1),
              pt=15, rgb=(255, 255, 255))
    frame = _card(s, 7.98, 0.48, 4.84, 6.54, (0xFF, 0x5A, 0x5A))  # noqa: F841


# ---------------------------------------------------------------- TRUE POSITIVES
def tp_image_over_text(s):
    """A picture dumped on top of a text block. A REAL collision -- must stay punished."""
    _text(s, 1.0, 2.0, 6.0, 3.0, "Important body copy that is being covered up " * 6, pt=18)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.4), Inches(5.0), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)


def tp_text_clipped_off_bottom(s):
    """Text running off the bottom edge. Content is LOST -- must stay punished."""
    # must genuinely exceed the remaining canvas: starts at 6.3in of a 7.5in slide, so only
    # ~1.2in remains; this text needs ~3in of wrapped height. Earlier version wrapped to
    # ~1.2in and did NOT clip -- the grader was right and the test was wrong.
    _text(s, 0.6, 6.3, 12.0, 3.4, "This paragraph starts near the bottom of the slide and "
                                  "continues on past the edge where nobody can read it. " * 12,
          pt=20)


def tp_two_textboxes_stacked(s):
    """Text on text -- containment must NOT excuse this, both shapes carry content."""
    _text(s, 2.0, 2.0, 6.0, 2.0, "First block of body text sitting here", pt=20)
    _text(s, 2.2, 2.1, 5.6, 1.8, "Second block written right on top of it", pt=20)


def main():
    from geometric_reward import score_deck
    from geom_shapes import score_shapes

    print("FALSE POSITIVES — good design our grader wrongly punishes (must RISE)\n")
    f1 = _deck("/tmp/fp_card.pptx", fp_text_on_card)
    m = score_shapes(f1)
    check("text on a filled card: collision >= 0.90", m["collision"] >= 0.90,
          "collision=%.3f" % m["collision"])

    f2 = _deck("/tmp/fp_bleed.pptx", fp_decorative_bleed)
    m = score_shapes(f2)
    check("decorative edge bleed: overflow >= 0.90", m["overflow"] >= 0.90,
          "overflow=%.3f" % m["overflow"])

    f3 = _deck("/tmp/fp_rich.pptx", fp_rich_two_column)
    r = score_deck(f3)
    check("rich two-column deck: score >= 0.75", r["score"] >= 0.75,
          "score=%.3f coll=%.2f over=%.2f" % (r["score"], r["metrics"]["collision"],
                                              r["metrics"]["overflow"]))

    print("\nTRUE POSITIVES — real defects that must STAY caught (must NOT rise)\n")
    t1 = _deck("/tmp/tp_cover.pptx", tp_image_over_text)
    m = score_shapes(t1)
    check("shape covering text: collision <= 0.60", m["collision"] <= 0.60,
          "collision=%.3f" % m["collision"])

    t2 = _deck("/tmp/tp_clip.pptx", tp_text_clipped_off_bottom)
    m = score_shapes(t2)
    check("text clipped off bottom: overflow <= 0.60", m["overflow"] <= 0.60,
          "overflow=%.3f" % m["overflow"])

    t3 = _deck("/tmp/tp_stack.pptx", tp_two_textboxes_stacked)
    m = score_shapes(t3)
    check("text stacked on text: collision <= 0.60", m["collision"] <= 0.60,
          "collision=%.3f" % m["collision"])

    print("\nREAL DECKS — human-labelled anchors\n")
    anchors = [
        ("fable_agentic/deck_social_media_slide_3.pptx", 0.75, 1.01, "human: GOOD"),
        ("fable_agentic/deck_business_slide_4.pptx", 0.00, 0.50, "human: empty frame"),
        ("gallery/deck_p88354_n001_s068_02f479.pptx", 0.00, 0.60, "render-confirmed clipping"),
    ]
    for path, lo, hi, why in anchors:
        if not os.path.isfile(path):
            check("anchor %s" % os.path.basename(path)[:34], False, "MISSING")
            continue
        sc = score_deck(path)["score"]
        check("%s in [%.2f,%.2f]" % (os.path.basename(path)[5:34], lo, hi),
              lo <= sc <= hi, "%.3f — %s" % (sc, why))

    print("\n" + "=" * 68)
    if FAILS:
        print("FAILING (%d): %s" % (len(FAILS), ", ".join(FAILS)))
    else:
        print("ALL LAYERING INVARIANTS HOLD")
    print("=" * 68)
    return 1 if FAILS else 0


if __name__ == "__main__":
    sys.exit(main())
