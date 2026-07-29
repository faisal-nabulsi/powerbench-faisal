"""
Enterprise template-conformance metric — ONE blended score, render-free.

WHY THIS EXISTS
---------------
For an enterprise deck the CONTENT is given; the skill being trained is filling the
CORPORATE TEMPLATE correctly.  Measured on our model's real output: 37/40 slides use the
"Blank" layout, ~96% of shapes are free-floating text boxes rather than template
placeholders, and 100% of the colours that are set are hardcoded RGB with zero theme
colours.  The model produces tidy slides that ignore the template completely.  In a
corporate setting that IS the failure: re-theme the deck, change the brand palette, or
hand it to the template's owner and everything the model made falls apart, because none
of it inherits from the master.

WHY ONE METRIC AND NOT FIVE
---------------------------
Layout usage, placeholder usage, theme colours and theme fonts are not four independent
skills.  They are four observable consequences of ONE behaviour: "did you open the
corporate template and work inside it, or did you start from a blank canvas".  A model
that learns the behaviour moves all four together.  Given them four separate weights in
the reward we would be counting a single signal four times, which mechanically dilutes
every genuinely independent metric (collision, density, overflow ...).  So we blend
internally and expose ONE number.  The sub-scores are returned for DIAGNOSIS ONLY — log
them, do not weight them.

BLEND
    conformance = 0.35*placeholder_usage + 0.30*theme_color
                + 0.20*layout_usage      + 0.15*theme_font
    conformance *= master_integrity          # anti-tamper multiplier

  placeholder_usage carries the most weight because it is the hardest to fake and the
  most consequential: a real placeholder inherits position, size, and font from the
  layout and survives re-theming.  A text box at hardcoded coordinates does not.

ANTI-TAMPER (master_integrity)
------------------------------
Every other metric in this reward gets EASIER if you vandalise the template.  Delete the
master's logo and you free up canvas area (better overlap, better balance).  Delete the
layouts and no slide can be judged against them.  master_integrity compares the produced
deck against the supplied template and multiplies the score down if the master's shapes
or the layout set have been stripped.  It is a PRODUCT of three retention ratios, not a
weighted sum, because any single form of stripping is by itself enough to destroy the
template — partial credit for "I only deleted the logo" would be the wrong lesson.
With no template_path supplied it is 1.0 (nothing to compare against, no penalty).

EMPTY-DECK GUARD
----------------
Two sub-metrics score inheritance, so "set nothing at all" is their perfect answer
(theme_color=1.0 when no colour is set anywhere; theme_font=1.0 when no font is set).
That is correct for a real deck and a reward hole for an empty one: a deck with no
content would collect 0.45 for producing nothing.  So a deck with no content-bearing
shapes scores 0.0 across the board.  Inheritance only counts as a virtue once there is
something to inherit.

API
    score_conformance(pptx_path, template_path=None)
        -> {'conformance': float, 'sub': {...}}      all values in [0,1], 1.0 = good

Raises whatever python-pptx raises if pptx_path cannot be opened — the caller's validity
gate owns that.  Never raises on a file that opens.
"""

import os
import sys

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

# ---------------------------------------------------------------------------
# Blend weights.  These sum to 1.0 and are then scaled by master_integrity.
# ---------------------------------------------------------------------------
SUB_WEIGHTS = {
    "placeholder_usage": 0.35,
    "theme_color": 0.30,
    "layout_usage": 0.20,
    "theme_font": 0.15,
}

SUB_KEYS = ("placeholder_usage", "theme_color", "layout_usage", "theme_font",
            "master_integrity")

# Layout names that mean "no template structure at all".  Matched case-insensitively
# after stripping.  A layout is ALSO treated as blank if it structurally contains no
# placeholders, whatever it is called — otherwise renaming "Blank" to "Content" in a
# hand-built template would buy free points.
_BLANK_LAYOUT_NAMES = {"blank", "blank slide", "empty", ""}

# Date / footer / slide-number placeholders are template FURNITURE, not slide content.
# Every stock layout carries all three, including "Blank" — so they must be excluded
# from both the layout check (or a renamed Blank layout would read as a content layout)
# and from placeholder_usage (or a policy could farm placeholder credit by stuffing its
# content into the footer).  Filling a footer is neither rewarded nor punished.
_CHROME_PH_TYPES = frozenset((
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
))


def _is_chrome_placeholder(shape):
    """True for a date / footer / slide-number placeholder."""
    try:
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.type in _CHROME_PH_TYPES
    except Exception:
        return False


# ---------------------------------------------------------------------------
# shape walking
# ---------------------------------------------------------------------------
def _is_group(shape):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _iter_leaf_shapes(container):
    """Yield every non-group shape, recursing into groups.

    Group members are yielded individually: a shape inside a group is never a template
    placeholder, so flattening keeps the placeholder fraction honest.
    """
    try:
        shapes = list(container)
    except Exception:
        return
    for shape in shapes:
        if _is_group(shape):
            try:
                for inner in _iter_leaf_shapes(shape.shapes):
                    yield inner
            except Exception:
                continue
        else:
            yield shape


def _shape_text_runs(shape):
    """All text runs belonging to a shape, including table cell runs."""
    runs = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                runs.extend(list(para.runs))
    except Exception:
        pass
    try:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        runs.extend(list(para.runs))
    except Exception:
        pass
    return runs


def _is_content_bearing(shape):
    """True if the shape carries actual content the audience reads.

    Content = non-empty text, a picture, a table, or a chart.  Empty boxes, connectors
    and decorative rules are excluded: they are not what a placeholder would have held,
    so counting them would just add noise to placeholder_usage.
    """
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return True
    except Exception:
        pass
    for attr in ("has_table", "has_chart"):
        try:
            if getattr(shape, attr, False):
                return True
        except Exception:
            pass
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
            return True
    except Exception:
        pass
    return False


# ---------------------------------------------------------------------------
# colour classification
# ---------------------------------------------------------------------------
def _color_kind(color):
    """Classify one ColorFormat as 'none' (inherited), 'theme', or 'explicit' (RGB).

    NOTE on the enum: the brief describes a theme colour as
    ``color.type is MSO_THEME_COLOR``.  In python-pptx (checked on 1.0.2) that is not the
    value you get — ``ColorFormat.type`` returns a member of MSO_COLOR_TYPE, and a theme
    colour reads as ``MSO_COLOR_TYPE.SCHEME`` while ``color.theme_color`` holds the
    MSO_THEME_COLOR member.  We accept either spelling so this keeps working if the
    library changes.
    """
    try:
        ctype = color.type
    except Exception:
        return "none"
    if ctype is None:
        return "none"
    try:
        if ctype == MSO_COLOR_TYPE.SCHEME:
            return "theme"
    except Exception:
        pass
    name = str(getattr(ctype, "name", ctype)).upper()
    if "SCHEME" in name or "THEME" in name:
        return "theme"
    # Belt and braces: some paths expose the theme slot even when type reads oddly.
    try:
        tc = color.theme_color
        if tc is not None and str(getattr(tc, "name", tc)).upper() != "NOT_THEME_COLOR":
            return "theme"
    except Exception:
        pass
    return "explicit"


def _tally_fill(fill, tally):
    """Count an explicitly-set solid fill as theme or explicit."""
    try:
        ftype = fill.type
    except Exception:
        return
    if ftype is None:  # inherited from layout/master — not an explicit choice
        return
    try:
        kind = _color_kind(fill.fore_color)
    except Exception:
        return
    if kind == "theme":
        tally[0] += 1
        tally[1] += 1
    elif kind == "explicit":
        tally[1] += 1


# ---------------------------------------------------------------------------
# layout classification
# ---------------------------------------------------------------------------
def _layout_is_real(layout):
    """True if the slide's layout is a real content layout rather than 'Blank'."""
    if layout is None:
        return False
    try:
        name = (layout.name or "").strip().lower()
    except Exception:
        name = ""
    if name in _BLANK_LAYOUT_NAMES or name.startswith("blank"):
        return False
    # Structural check, so renaming "Blank" to "Executive Content" buys nothing: a real
    # content layout must offer at least one placeholder that is not date/footer/slide
    # number.  The stock "Blank" layout carries all three of those and nothing else.
    try:
        for ph in layout.placeholders:
            try:
                if ph.placeholder_format.type not in _CHROME_PH_TYPES:
                    return True
            except Exception:
                return True  # unreadable type: assume it is a content placeholder
        return False
    except Exception:
        return True


# ---------------------------------------------------------------------------
# anti-tamper
# ---------------------------------------------------------------------------
def _master_fingerprint(prs):
    """(total master shapes, {layout name: total shapes}) across all masters."""
    master_shapes = 0
    layouts = {}
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    for master in masters:
        try:
            master_shapes += len(list(master.shapes))
        except Exception:
            pass
        try:
            for layout in master.slide_layouts:
                try:
                    name = (layout.name or "").strip().lower()
                except Exception:
                    name = ""
                try:
                    count = len(list(layout.shapes))
                except Exception:
                    count = 0
                layouts[name] = layouts.get(name, 0) + count
        except Exception:
            pass
    return master_shapes, layouts


def _master_integrity(prs, template_path):
    """0..1 anti-tamper score: has the produced deck stripped the template?

    Product of three retention ratios (master shapes kept, layouts kept, per-layout
    shapes kept), each capped at 1.0 so ADDING shapes never scores above intact.
    Returns 1.0 when no template is supplied or the template cannot be read — the
    metric can only accuse when it has something to compare against.
    """
    if not template_path:
        return 1.0
    try:
        tpl = Presentation(template_path)
    except Exception:
        return 1.0  # unreadable template: no evidence of tampering, do not punish

    try:
        got_shapes, got_layouts = _master_fingerprint(prs)
        want_shapes, want_layouts = _master_fingerprint(tpl)
    except Exception:
        return 1.0

    shape_ratio = 1.0
    if want_shapes > 0:
        shape_ratio = min(1.0, got_shapes / float(want_shapes))

    layout_count_ratio = 1.0
    if want_layouts:
        kept = sum(1 for name in want_layouts if name in got_layouts)
        layout_count_ratio = min(1.0, kept / float(len(want_layouts)))

    per_layout = []
    for name, want_n in want_layouts.items():
        if name not in got_layouts:
            continue  # already penalised by layout_count_ratio
        if want_n <= 0:
            per_layout.append(1.0)
        else:
            per_layout.append(min(1.0, got_layouts[name] / float(want_n)))
    layout_shape_ratio = sum(per_layout) / len(per_layout) if per_layout else 1.0

    return max(0.0, min(1.0, shape_ratio * layout_count_ratio * layout_shape_ratio))


# ---------------------------------------------------------------------------
# public API
# ---------------------------------------------------------------------------
def score_conformance(pptx_path: str, template_path: str = None) -> dict:
    """Score how well a deck fills its corporate template.

    Returns {'conformance': float, 'sub': {placeholder_usage, theme_color,
    layout_usage, theme_font, master_integrity}} — every value in [0,1], 1.0 = good.
    'conformance' is the only number to reward on; 'sub' is for diagnosis.

    Raises whatever python-pptx raises if pptx_path cannot be opened (the caller's
    validity gate handles that).  Never raises on a file that opens.
    """
    prs = Presentation(pptx_path)  # let a bad file raise: caller owns the gate

    integrity = 1.0
    try:
        integrity = _master_integrity(prs, template_path)
    except Exception:
        integrity = 1.0

    try:
        slides = list(prs.slides)
    except Exception:
        slides = []

    zero = {k: 0.0 for k in SUB_KEYS}
    zero["master_integrity"] = integrity
    if not slides:
        return {"conformance": 0.0, "sub": zero}

    n_content = 0            # content-bearing shapes
    n_content_ph = 0         # ... that are real placeholders
    color_theme = 0          # explicit colours that are theme colours
    color_explicit_total = 0  # all explicitly-set colours
    n_runs = 0
    n_runs_inherit_font = 0
    n_real_layout = 0

    for slide in slides:
        try:
            if _layout_is_real(slide.slide_layout):
                n_real_layout += 1
        except Exception:
            pass

        for shape in _iter_leaf_shapes(getattr(slide, "shapes", [])):
            if _is_chrome_placeholder(shape):
                continue  # template furniture: neither rewarded nor punished
            try:
                content = _is_content_bearing(shape)
            except Exception:
                content = False
            if content:
                n_content += 1
                try:
                    if bool(shape.is_placeholder):
                        n_content_ph += 1
                except Exception:
                    pass

            # explicit shape fill
            try:
                tally = [0, 0]
                _tally_fill(shape.fill, tally)
                color_theme += tally[0]
                color_explicit_total += tally[1]
            except Exception:
                pass

            # table cell fills
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            tally = [0, 0]
                            _tally_fill(cell.fill, tally)
                            color_theme += tally[0]
                            color_explicit_total += tally[1]
            except Exception:
                pass

            # runs: font colour + font name
            for run in _shape_text_runs(shape):
                n_runs += 1
                try:
                    if run.font.name is None:
                        n_runs_inherit_font += 1
                except Exception:
                    n_runs_inherit_font += 1  # unreadable == not explicitly set
                try:
                    kind = _color_kind(run.font.color)
                except Exception:
                    kind = "none"
                if kind == "theme":
                    color_theme += 1
                    color_explicit_total += 1
                elif kind == "explicit":
                    color_explicit_total += 1

    # Empty-deck guard: inheritance is only a virtue when there is content to inherit.
    if n_content == 0:
        return {"conformance": 0.0, "sub": zero}

    sub = {
        "placeholder_usage": n_content_ph / float(n_content),
        # nothing coloured at all == pure inheritance == ideal
        "theme_color": (color_theme / float(color_explicit_total)
                        if color_explicit_total else 1.0),
        "layout_usage": n_real_layout / float(len(slides)),
        # no runs at all (picture-only deck) sets no fonts == pure inheritance
        "theme_font": (n_runs_inherit_font / float(n_runs) if n_runs else 1.0),
        "master_integrity": integrity,
    }
    for k in SUB_KEYS:
        sub[k] = max(0.0, min(1.0, float(sub[k])))

    blended = sum(SUB_WEIGHTS[k] * sub[k] for k in SUB_WEIGHTS)
    conformance = max(0.0, min(1.0, blended * sub["master_integrity"]))
    return {"conformance": conformance, "sub": sub}


# ===========================================================================
# self-test / CLI
# ===========================================================================
def _mean(xs):
    return sum(xs) / float(len(xs)) if xs else 0.0


def _sd(xs):
    if len(xs) < 2:
        return 0.0
    m = _mean(xs)
    return (sum((x - m) ** 2 for x in xs) / float(len(xs) - 1)) ** 0.5


def _build_conformant(path, template_path):
    """A deck that does it right: real layout, filled placeholders, no explicit
    colours, no explicit fonts."""
    prs = Presentation(template_path)
    facts = [
        ("Q3 Revenue Review", ["Revenue up 12% QoQ", "EMEA led growth",
                               "Churn flat at 2.1%"]),
        ("Cost Structure", ["COGS down 3 pts", "Headcount steady",
                            "Cloud spend consolidated"]),
        ("Outlook", ["Q4 guidance unchanged", "Two launches in flight"]),
    ]
    layout = prs.slide_layouts[1]  # "Title and Content"
    for title, bullets in facts:
        slide = prs.slides.add_slide(layout)
        slide.placeholders[0].text = title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            body.add_paragraph().text = b
    prs.save(path)
    return path


def _build_tampered(path, template_path, wipe_master=True, drop_layouts=0):
    """Conformant content, but the template has been vandalised."""
    _build_conformant(path, template_path)
    prs = Presentation(path)
    master = prs.slide_master
    if wipe_master:
        for shape in list(master.shapes):
            shape._element.getparent().remove(shape._element)
    else:  # delete a single master shape (e.g. the logo)
        shapes = list(master.shapes)
        if shapes:
            shapes[0]._element.getparent().remove(shapes[0]._element)
    if drop_layouts:
        # SlideLayout is not hashable, so compare the underlying XML elements.
        used = [s.slide_layout._element for s in prs.slides]
        removed = 0
        for layout in list(master.slide_layouts):
            if removed >= drop_layouts:
                break
            if any(layout._element is el for el in used):
                continue
            try:
                master.slide_layouts.remove(layout)
                removed += 1
            except Exception:
                pass
    prs.save(path)
    return path


def _self_test():
    import glob
    import tempfile

    here = os.path.dirname(os.path.abspath(__file__))
    gallery = sorted(glob.glob(os.path.join(here, "gallery", "deck_*.pptx")))[:20]
    import pptx as _pptx_pkg
    # Stand-in "corporate template": the stock python-pptx template, which is exactly
    # what the model's decks were built from (plain Presentation()), so any integrity
    # loss on a gallery deck would be real tampering rather than a mismatched baseline.
    tpl = os.path.join(os.path.dirname(os.path.abspath(_pptx_pkg.__file__)),
                       "templates", "default.pptx")

    print("=" * 74)
    print("1. REAL MODEL DECKS (n=%d)  template=%s" % (len(gallery),
                                                       os.path.basename(tpl)))
    print("=" * 74)
    vals, subs = [], {k: [] for k in SUB_KEYS}
    for p in gallery:
        r = score_conformance(p, template_path=tpl)
        vals.append(r["conformance"])
        for k in SUB_KEYS:
            subs[k].append(r["sub"][k])
        print("  %-44s %.4f  ph=%.2f col=%.2f lay=%.2f fnt=%.2f mi=%.2f"
              % (os.path.basename(p)[:44], r["conformance"],
                 r["sub"]["placeholder_usage"], r["sub"]["theme_color"],
                 r["sub"]["layout_usage"], r["sub"]["theme_font"],
                 r["sub"]["master_integrity"]))
    print("  ---")
    print("  conformance: mean=%.4f  sd=%.4f  min=%.4f  max=%.4f  distinct=%d"
          % (_mean(vals), _sd(vals), min(vals), max(vals), len(set(round(v, 6) for v in vals))))
    for k in SUB_KEYS:
        print("  sub %-18s mean=%.4f sd=%.4f min=%.4f max=%.4f"
              % (k, _mean(subs[k]), _sd(subs[k]), min(subs[k]), max(subs[k])))

    tmp = tempfile.mkdtemp(prefix="conformance_selftest_")
    print()
    print("=" * 74)
    print("2. HAND-BUILT CONFORMANT DECK (real layout + placeholders, no colours/fonts)")
    print("=" * 74)
    good = _build_conformant(os.path.join(tmp, "conformant.pptx"), tpl)
    rg = score_conformance(good, template_path=tpl)
    print("  %s" % good)
    print("  conformance=%.4f" % rg["conformance"])
    for k in SUB_KEYS:
        print("    %-18s %.4f" % (k, rg["sub"][k]))

    print()
    print("=" * 74)
    print("3. TAMPERED DECKS (same content, template vandalised)")
    print("=" * 74)
    t1 = _build_tampered(os.path.join(tmp, "tampered_wipe.pptx"), tpl, wipe_master=True)
    r1 = score_conformance(t1, template_path=tpl)
    print("  a) master shapes WIPED       conformance=%.4f  master_integrity=%.4f"
          % (r1["conformance"], r1["sub"]["master_integrity"]))
    t2 = _build_tampered(os.path.join(tmp, "tampered_logo.pptx"), tpl,
                         wipe_master=False, drop_layouts=3)
    r2 = score_conformance(t2, template_path=tpl)
    print("  b) 1 master shape + 3 layouts deleted"
          "  conformance=%.4f  master_integrity=%.4f"
          % (r2["conformance"], r2["sub"]["master_integrity"]))
    print("  (undamaged same deck scored %.4f — tamper cannot pay)" % rg["conformance"])
    for label, rr in (("wipe", r1), ("partial", r2)):
        print("    %-8s subs: " % label
              + " ".join("%s=%.3f" % (k[:9], rr["sub"][k]) for k in SUB_KEYS))

    print()
    print("=" * 74)
    print("4. ROBUSTNESS")
    print("=" * 74)
    print("  no template arg on a real deck: %.4f"
          % score_conformance(gallery[0])["conformance"])
    empty_path = os.path.join(tmp, "empty.pptx")
    Presentation(tpl).save(empty_path)
    print("  zero-slide deck: %.4f (empty-deck guard)"
          % score_conformance(empty_path, tpl)["conformance"])
    blank_path = os.path.join(tmp, "blank_slides.pptx")
    pb = Presentation(tpl)
    pb.slides.add_slide(pb.slide_layouts[6])  # Blank, no content
    pb.save(blank_path)
    print("  blank-layout slide with no shapes: %.4f"
          % score_conformance(blank_path, tpl)["conformance"])
    bad = os.path.join(tmp, "not_a_deck.pptx")
    with open(bad, "wb") as fh:
        fh.write(b"this is not a pptx")
    try:
        score_conformance(bad, tpl)
        print("  corrupt file: NO RAISE  <-- BUG")
    except Exception as exc:
        print("  corrupt file raises %s (validity gate handles it) OK"
              % type(exc).__name__)

    print()
    print("=" * 74)
    print("5. GRADIENT VERDICT")
    print("=" * 74)
    sd = _sd(vals)
    print("  sd across the %d decks above = %.4f, range [%.4f, %.4f], %d distinct values"
          % (len(vals), sd, min(vals), max(vals),
             len(set(round(v, 6) for v in vals))))

    # GRPO learns ONLY from variance WITHIN a group of rollouts on the SAME prompt, so
    # the pooled sd is the wrong statistic — measure per prompt id.
    import re
    all_decks = sorted(glob.glob(os.path.join(here, "gallery", "deck_*.pptx")))
    prompts, failures = {}, 0
    for p in all_decks:
        m = re.search(r"deck_(p\d+)_", os.path.basename(p))
        if not m:
            continue
        try:
            v = score_conformance(p, template_path=tpl)["conformance"]
        except Exception:
            failures += 1
            continue
        prompts.setdefault(m.group(1), []).append(v)

    print()
    print("  WITHIN-PROMPT variance (%d decks, %d prompts, %d scoring failures):"
          % (sum(len(v) for v in prompts.values()), len(prompts), failures))
    print("    %-8s %4s %8s %8s %8s %8s %9s"
          % ("prompt", "n", "mean", "sd", "min", "max", "distinct"))
    sds = []
    for name in sorted(prompts):
        v = prompts[name]
        s = _sd(v)
        sds.append(s)
        print("    %-8s %4d %8.4f %8.4f %8.4f %8.4f %9d"
              % (name, len(v), _mean(v), s, min(v), max(v),
                 len(set(round(x, 6) for x in v))))

    # Simulate the real training signal: chunks of 8 rollouts from one prompt.
    flat, dead = 0, 0
    for v in prompts.values():
        for i in range(0, len(v) - 7, 8):
            chunk = v[i:i + 8]
            flat += 1
            if _sd(chunk) < 1e-9:
                dead += 1
    mean_sd = _mean(sds)
    print()
    print("    mean within-prompt sd = %.4f  <- what GRPO turns into gradient" % mean_sd)
    print("    simulated groups of 8: %d total, %d flat (%.0f%% would give NO gradient)"
          % (flat, dead, (100.0 * dead / flat) if flat else 0.0))
    print()
    if mean_sd > 0.05 and dead == 0:
        print("  VERDICT: VARIES — real gradient today, not constant-at-zero.")
    elif mean_sd > 1e-6:
        print("  VERDICT: varies weakly — check flat-group rate before weighting it.")
    else:
        print("  VERDICT: CONSTANT — contributes zero gradient, do not weight it.")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        for path in sys.argv[1:]:
            try:
                print(path, score_conformance(path))
            except Exception as exc:  # noqa: BLE001
                print(path, "ERROR:", type(exc).__name__, exc)
    else:
        _self_test()
