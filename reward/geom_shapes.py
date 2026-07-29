"""Deterministic slide-layout geometry scoring for .pptx decks (python-pptx only).

No rendering, no images, no LLM: every number here comes from the shape
geometry stored in the file (EMU offsets/extents) and the slide dimensions.

Public API
----------
    score_shapes(pptx_path) -> {'collision': float, 'overflow': float, 'imbalance': float}

SCORE CONVENTION: every value is a float in [0.0, 1.0] where 1.0 = GOOD.

Metrics (each computed per slide, then averaged over the slides of the deck):

  collision  Pairwise bounding-box overlap between the visible shapes of a
             slide.  overlap_ratio = sum(pairwise intersection areas) /
             sum(shape areas); score = 1 - min(1, overlap_ratio).
             Nothing overlapping -> 1.0.

  overflow   Area escaping the slide canvas.  overflow_ratio =
             sum(area of each shape lying outside the slide rectangle) /
             sum(shape areas); score = 1 - min(1, overflow_ratio).
             Everything on-canvas -> 1.0.

  imbalance  Offset of the area-weighted centroid of the shapes from the slide
             centre, using the ellipse normalisation of the AeSlides paper
             (horizontal drift is penalised 3x harder than vertical):
                 d = sqrt(((cx-0.5)/0.05)**2 + ((cy-0.5)/0.15)**2)
             with cx, cy the centroid normalised to 0-1; score = 1/(1+d).
             Perfectly centred composition -> ~1.0.

IMPORTANT - empty slides score 0.0, not 1.0
-------------------------------------------
A slide with no visible shapes is trivially collision-free, overflow-free and
(vacuously) "balanced".  Returning 1.0 there would make a blank deck look
perfect and hand a policy a free reward-hacking hole, so a slide with zero
visible shapes returns 0.0 for all three metrics.  Same for a deck with zero
slides.  A slide whose only content is *unfilled layout placeholders* counts as
empty as well (an empty placeholder renders as nothing), which closes the
related "emit slides straight from a layout and touch nothing" hole.

What counts as a "visible shape"
--------------------------------
  * left / top / width / height must all be present (None geometry is skipped;
    e.g. some inherited placeholders carry no explicit xfrm).
  * area must be > 0 (zero-width/height connectors etc. are skipped).
  * placeholders whose text frame is blank are skipped (nothing renders).
  * group shapes are scored atomically by their own bounding box; the code does
    NOT recurse into a group, because counting both the group box and its
    children would double-count area and report a fake collision.

Robustness: never raises on a valid .pptx.  If the file cannot be opened at all
the underlying exception propagates - the caller owns the validity gate.
"""

from __future__ import annotations

import math
from typing import Dict, List, Optional, Tuple

from pptx import Presentation

__all__ = ["score_shapes"]

# Ellipse tolerances from the AeSlides paper: horizontal offset hurts more.
X_TOL = 0.05
Y_TOL = 0.15

# Fallback canvas if a deck somehow reports no slide dimensions (10 x 7.5 in).
DEFAULT_SLIDE_W = 9144000
DEFAULT_SLIDE_H = 6858000

METRIC_KEYS = ("collision", "overflow", "imbalance", "textfit", "contrast", "picfit", "alignment")

# Fraction of each edge treated as "too close to the edge". A shape jammed flush against the
# canvas boundary reads as cut off even when technically inside, so overflow uses this inset.
SAFE_MARGIN = 0.02

# A shape covering more than this fraction of the slide is treated as a background.
BACKGROUND_COVER = 0.85
_CANVAS_AREA = []

# --- text-fit estimation (no rendering) -------------------------------------
# We estimate rendered text height from character count, font size and box width:
#   chars_per_line ~= box_width / (font_size * CHAR_W_RATIO)
#   lines          ~= ceil(chars / chars_per_line)  (plus explicit paragraph breaks)
#   text_height    ~= lines * font_size * LINE_H_RATIO
# then compare to the box height. Coarse, but it reliably catches the real failure mode:
# a large font in a small box, or a paragraph crammed into a one-line box.
EMU_PER_PT = 12700.0
DEFAULT_FONT_PT = 18.0       # python-pptx default when a run sets no size
CHAR_W_RATIO = 0.50          # avg glyph advance as a fraction of point size (sans-serif)
LINE_H_RATIO = 1.20          # line height as a multiple of point size

# (x0, y0, x1, y1) in EMU, x1 >= x0 and y1 >= y0
Rect = Tuple[float, float, float, float]


# --------------------------------------------------------------------------- #
# geometry helpers
# --------------------------------------------------------------------------- #
def _rect_of(shape) -> Optional[Rect]:
    """Bounding box of a shape as (x0, y0, x1, y1), or None if unusable."""
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return None
    if left is None or top is None or width is None or height is None:
        return None
    try:
        left = float(left)
        top = float(top)
        width = float(width)
        height = float(height)
    except (TypeError, ValueError):
        return None
    if not all(math.isfinite(v) for v in (left, top, width, height)):
        return None
    # Tolerate negative extents by normalising the corners.
    x0, x1 = (left, left + width) if width >= 0 else (left + width, left)
    y0, y1 = (top, top + height) if height >= 0 else (top + height, top)
    if (x1 - x0) <= 0 or (y1 - y0) <= 0:
        return None  # zero-area shape contributes nothing
    return (x0, y0, x1, y1)


def _area(r: Rect) -> float:
    return (r[2] - r[0]) * (r[3] - r[1])


def _intersection_area(a: Rect, b: Rect) -> float:
    dx = min(a[2], b[2]) - max(a[0], b[0])
    dy = min(a[3], b[3]) - max(a[1], b[1])
    if dx <= 0 or dy <= 0:
        return 0.0
    return dx * dy


def _is_visible(shape) -> bool:
    """False for things that occupy geometry but render as nothing."""
    try:
        if shape.is_placeholder:
            if shape.has_text_frame and not shape.text_frame.text.strip():
                # An unfilled title/body placeholder draws nothing at render time.
                return False
    except Exception:
        # Malformed shape: fall through and keep it rather than crashing.
        return True
    return True


def _text_extent_rect(shape, box: Rect) -> Optional[Rect]:
    """The rectangle the RENDERED TEXT actually occupies, or None to keep the box.

    Humans see glyphs, not shape boxes. Scoring the box produces errors in BOTH directions:
      * text WIDER than its box (wrap=none) renders past the box and is visually clipped,
        yet the box is on-canvas so overflow says 1.00   (blind study: B_hl15, grader 0.91
        vs human 0.30)
      * text NARROWER than its box makes the box overlap a neighbour although the glyphs do
        not, inventing a collision                        (blind study: C_v31, grader 0.39
        vs human 0.72)

    Only applied to shapes with NO visible fill/line: a filled panel really does occupy its
    whole box, so its box is the honest rectangle.
    """
    try:
        if not shape.has_text_frame:
            return None
        tf = shape.text_frame
        txt = (tf.text or "").strip()
        if not txt:
            return None
        # a shape with a visible fill or outline occupies its full box
        try:
            ft = shape.fill.type
            if ft is not None and "BACKGROUND" not in str(ft):
                return None
        except Exception:
            pass
        import text_metrics as _tm
        x0, y0, x1, y1 = box
        box_w_pt = (x1 - x0) / EMU_PER_PT
        wrap = True
        try:
            wrap = tf.word_wrap is not False
        except Exception:
            pass
        widest = 0.0
        lines = 0
        for para in tf.paragraphs:
            t = "".join(r.text or "" for r in para.runs)
            if not t.strip():
                lines += 1
                continue
            size = 0.0
            try:
                if para.font.size:
                    size = para.font.size.pt
            except Exception:
                pass
            for r in para.runs:
                try:
                    if r.font.size:
                        size = max(size, r.font.size.pt)
                except Exception:
                    pass
            size = size or DEFAULT_FONT_PT
            w = _tm.measure_text_width(t, size)
            if wrap and w > box_w_pt:
                widest = max(widest, box_w_pt)
                lines += max(1, int(w / max(1e-6, box_w_pt)) + 1)
            else:
                widest = max(widest, w)
                lines += 1
        if widest <= 0 or lines <= 0:
            return None
        big = 0.0
        for para in tf.paragraphs:
            try:
                if para.font.size:
                    big = max(big, para.font.size.pt)
            except Exception:
                pass
            for r in para.runs:
                try:
                    if r.font.size:
                        big = max(big, r.font.size.pt)
                except Exception:
                    pass
        big = big or DEFAULT_FONT_PT
        tw = widest * EMU_PER_PT
        th = lines * big * LINE_H_RATIO * EMU_PER_PT
        cx = (x0 + x1) / 2.0
        # unwrapped centred text spills symmetrically; that spill is what gets clipped
        nx0, nx1 = cx - tw / 2.0, cx + tw / 2.0
        return (nx0, y0, nx1, y0 + th)
    except Exception:
        return None


def _slide_rects(slide) -> List[Rect]:
    rects: List[Rect] = []
    _TEXT_FLAGS.clear()
    _CONTENT_FLAGS.clear()
    try:
        shapes = list(slide.shapes)
    except Exception:
        return rects
    for shape in shapes:
        try:
            if not _is_visible(shape):
                continue
            r = _rect_of(shape)
        except Exception:
            continue
        if r is not None:
            tr = _text_extent_rect(shape, r)
            if tr is not None:
                r = tr
            try:
                has_txt = bool(shape.has_text_frame and (shape.text_frame.text or "").strip())
            except Exception:
                has_txt = False
            try:
                is_pic = hasattr(shape, "image")
            except Exception:
                is_pic = False
            _TEXT_FLAGS[len(rects)] = has_txt
            # CONTENT = information a viewer can lose. A decorative fill carries none.
            _CONTENT_FLAGS[len(rects)] = bool(has_txt or is_pic)
            rects.append(r)
    return rects


# --------------------------------------------------------------------------- #
# per-slide metrics
# --------------------------------------------------------------------------- #
_TEXT_FLAGS = {}
_CONTENT_FLAGS = {}
# slack when testing containment, so a card drawn a hair smaller than its label still counts
_CONTAIN_TOL = 45720  # 0.05 in


def _contains(outer, inner, tol=_CONTAIN_TOL):
    """True if `outer` fully encloses `inner` (within tol)."""
    return (outer[0] <= inner[0] + tol and outer[1] <= inner[1] + tol
            and outer[2] >= inner[2] - tol and outer[3] >= inner[3] - tol)


def _soft_penalty(fracs):
    """Combine per-shape penalties WITHOUT saturating.

    The old code summed fracs and capped at 1.0, so any slide with 2-3 flagged shapes scored
    EXACTLY 0.000 -- the metric went binary and structurally punished shape-rich slides.
    Measured on real decks: 100% of 20+-shape decks sat at 0.000 vs 27% of 0-5-shape decks,
    and our own model drifted 4.2 -> 3.5 shapes/slide learning to dodge it.

    Same worst+mean blend `_textfit_score` already uses: one badly broken shape still
    dominates the score, but a long tail of clean shapes is no longer ignored.
    """
    if not fracs:
        return 0.0
    worst = max(fracs)
    mean = sum(fracs) / len(fracs)
    return min(1.0, 0.65 * worst + 0.35 * mean)


def _collision_score(rects: List[Rect], total_area: float = 0.0) -> float:
    """Per-shape overlap, relative to EACH shape's own area (not the slide total), so a
    localized bad overlap isn't diluted by the rest of the slide. For every shape we sum how
    much of IT is covered by others (capped at 1), then sum those coverages. A single image
    sitting on a text box, or two text boxes overlapping, now registers at full strength."""
    n = len(rects)
    if n < 2:
        return 1.0
    # A shape covering most of the canvas is a BACKGROUND (full-bleed image / colour panel).
    # Putting text on a background is correct design, not a collision -- but the naive
    # pairwise test called it a 100% overlap and scored such slides 0.00, actively teaching
    # the model to avoid backgrounds. Backgrounds are excluded from BOTH sides of the test.
    # (hand-audit ID 9f2013)
    canvas = _CANVAS_AREA[0] if _CANVAS_AREA else 0.0
    is_bg = [canvas > 0 and _area(r) / canvas > BACKGROUND_COVER for r in rects]
    fracs = []
    for i in range(n):
        if is_bg[i]:
            continue
        ai = _area(rects[i])
        if ai <= 0:
            continue
        covered = 0.0
        for j in range(n):
            if i == j or is_bg[j]:
                continue
            # LAYERING IS NOT COLLISION. A shape that fully CONTAINS this one and carries no
            # text of its own is a deliberate backdrop: a label on a card, a caption on a
            # photo, a chip inside a panel. Before this, only shapes >85% of canvas were
            # exempt, so a 31%-of-canvas card turned every label on it into a 100% overlap
            # and zeroed 19% of the reward. Verified: identical text scored collision 1.000
            # bare and 0.000 on a card. The container must not itself hold text, so genuine
            # text-on-text still registers.
            if not _TEXT_FLAGS.get(j) and _contains(rects[j], rects[i]):
                continue
            # ...and symmetrically: if THIS shape contains the other and carries no text of
            # its own, it is a backdrop being covered by its own content -- a frame around a
            # photo, a card behind a chart. Being covered is the whole point of a frame.
            # (social_media/slide_3: a picture filled 92% of its own frame and cost 19%.)
            if not _TEXT_FLAGS.get(i) and _contains(rects[i], rects[j]):
                continue
            covered += _intersection_area(rects[i], rects[j])
        frac = min(1.0, covered / ai)
        # Covering TEXT is visually fatal even at small area (the words become unreadable),
        # so text shapes get a steep sqrt response: 10% covered -> 32% penalty, not 10%.
        # Hand-audit ID b065a9: image over ~10% of a text line still scored 0.89.
        if _TEXT_FLAGS.get(i):
            frac = math.sqrt(frac)
        fracs.append(frac)
    return 1.0 - _soft_penalty(fracs)


def _text_offcanvas_penalty(slide, sw: float, sh: float) -> float:
    """Text can overflow its BOX and then run off the SLIDE while the box itself is legally
    on-canvas -- so the box-based overflow test says 1.00 while the reader sees text cut off
    (hand-audit d12f2f, 2fe501). We estimate rendered text height and check whether it spills
    past the bottom edge from where the box starts."""
    pen = 0.0
    try:
        for sh_ in slide.shapes:
            try:
                if not sh_.has_text_frame or sh_.top is None or sh_.height is None:
                    continue
                tf = sh_.text_frame
                txt = tf.text or ""
                if not txt.strip():
                    continue
                size = 0.0
                for para in tf.paragraphs:
                    try:
                        if para.font.size is not None:
                            size = max(size, para.font.size.pt)
                    except Exception:
                        pass
                    for r in para.runs:
                        try:
                            if r.font.size is not None:
                                size = max(size, r.font.size.pt)
                        except Exception:
                            pass
                size = size or DEFAULT_FONT_PT
                bw = float(sh_.width) / EMU_PER_PT if sh_.width else 1.0
                cpl = max(1.0, bw / max(1e-6, size * CHAR_W_RATIO))
                lines = 0.0
                for para in tf.paragraphs:
                    t = "".join(r.text or "" for r in para.runs)
                    lines += max(1.0, math.ceil(len(t) / cpl)) if t.strip() else 1.0
                need = lines * size * LINE_H_RATIO * EMU_PER_PT
                bottom = float(sh_.top) + need
                if bottom > sh:
                    pen += min(1.0, (bottom - sh) / max(1.0, need))
            except Exception:
                continue
    except Exception:
        return 0.0
    return pen


def _overflow_score(rects: List[Rect], total_area: float, sw: float, sh: float) -> float:
    """Per-shape off-canvas fraction, relative to EACH shape's own area. A small safe margin
    (SAFE_MARGIN of each edge) means a shape jammed flush against the very edge also counts as
    partly spilling, since that reads as 'cut off' even when technically inside."""
    m = SAFE_MARGIN
    safe: Rect = (sw * m, sh * m, sw * (1.0 - m), sh * (1.0 - m))
    canvas = _CANVAS_AREA[0] if _CANVAS_AREA else 0.0
    fracs = []
    for idx, r in enumerate(rects):
        a = _area(r)
        if a <= 0:
            continue
        # A decorative fill running off the edge loses NOTHING -- edge bleed is a standard
        # accent technique (corner circles, side rules), and we were scoring it identically
        # to a sentence falling off the slide. Only text and pictures carry information that
        # can actually be lost. Verified: three accent circles scored overflow 0.000.
        if not _CONTENT_FLAGS.get(idx, True):
            continue
        # a full-bleed background image is MEANT to reach the edges
        if canvas > 0 and a / canvas > BACKGROUND_COVER:
            continue
        outside = a - _intersection_area(r, safe)
        fracs.append(min(1.0, max(0.0, outside) / a))
    return 1.0 - _soft_penalty(fracs)


def _textfit_shape(shape) -> float:
    """Estimate whether a shape's text fits inside its box. 1.0 = fits, 0.0 = badly overflowing.

    No rendering: we approximate the laid-out text height from character count, font size and
    box width (see constants above) and compare it to the box height. Catches the 'font too big
    / box too small' failure where text visibly spills out of its container.
    """
    try:
        if not shape.has_text_frame:
            return None
        tf = shape.text_frame
        text = tf.text or ""
        if not text.strip():
            return None
        w, h = shape.width, shape.height
        if not w or not h or w <= 0 or h <= 0:
            return None

        # largest font in the box drives overflow
        size_pt = 0.0
        for para in tf.paragraphs:
            try:
                if para.font.size is not None:
                    size_pt = max(size_pt, para.font.size.pt)
            except Exception:
                pass
            for run in para.runs:
                try:
                    if run.font.size is not None:
                        size_pt = max(size_pt, run.font.size.pt)
                except Exception:
                    pass
        if size_pt <= 0:
            size_pt = DEFAULT_FONT_PT

        box_w_pt = float(w) / EMU_PER_PT
        box_h_pt = float(h) / EMU_PER_PT
        cpl = max(1.0, box_w_pt / max(1e-6, size_pt * CHAR_W_RATIO))

        no_wrap = False
        try:
            no_wrap = (tf.word_wrap is False)
        except Exception:
            pass
        lines = 0.0
        widest = 0.0
        for para in tf.paragraphs:
            ptxt = "".join(r.text or "" for r in para.runs)
            if not ptxt.strip():
                lines += 1.0
                continue
            widest = max(widest, (len(ptxt) * size_pt * CHAR_W_RATIO) / max(1e-6, box_w_pt))
            lines += 1.0 if no_wrap else max(1.0, math.ceil(len(ptxt) / cpl))
        if no_wrap and widest > 1.0:
            return max(0.0, min(1.0, 1.0 - (widest - 1.0)))
        need_h = lines * size_pt * LINE_H_RATIO

        ratio = need_h / max(1e-6, box_h_pt)
        if ratio <= 1.0:
            return 1.0
        # 1x = fits, 2x or more of the box height = fully overflowing
        return max(0.0, min(1.0, 1.0 - (ratio - 1.0)))
    except Exception:
        return None


def _textfit_score(slide) -> float:
    """Real-font-metric text fitting (text_metrics.py, Pillow FreeTypeFont).

    The previous character-count estimate (0.5 x font_size per glyph) had a mean error of
    22.8% and up to 88.8% on wide glyphs -- it scored "Community" at 57.6pt as fitting a
    one-line box when it genuinely wraps to two. Real metrics are both more accurate AND
    give more within-group variance (0.0844 -> 0.0928), i.e. more gradient for GRPO.
    Falls back to the old estimator if no TTF is available.
    """
    try:
        import text_metrics as _tm
        vals = []
        for shp in slide.shapes:
            try:
                if not shp.has_text_frame or not (shp.text_frame.text or "").strip():
                    continue
                vals.append(_tm.textfit_shape(shp))
            except Exception:
                continue
        if vals:
            # AGGREGATION FIX (blind LLM-vs-grader study): averaging hid real defects. A deck
            # with ONE badly clipped text box among seven good ones scored 0.985 while a human
            # graded it 0.30 -- the single worst source of grader/human disagreement.
            # Clipped text is unreadable text: severity should follow the WORST box, not the
            # mean. We blend worst and mean (70/30) so the metric still varies smoothly
            # (pure min collapses group variance, which GRPO needs).
            worst = min(vals)
            mean = sum(vals) / len(vals)
            return 0.70 * worst + 0.30 * mean
        return 1.0
    except Exception:
        pass
    return _textfit_score_legacy(slide)


def _textfit_score_legacy(slide) -> float:
    """Mean text-fit across the slide's text shapes; 1.0 when there is no text to check."""
    vals = []
    try:
        for shape in slide.shapes:
            try:
                if not _is_visible(shape):
                    continue
            except Exception:
                pass
            v = _textfit_shape(shape)
            if v is not None:
                vals.append(v)
    except Exception:
        return 1.0
    return sum(vals) / len(vals) if vals else 1.0


def _alignment_score(rects: List[Rect], sw: float, sh: float) -> float:
    """Reference-free ALIGNMENT (LayoutGAN / Kikuchi et al.): well-designed layouts share
    edges and centre-lines. For each element we take the smallest misalignment to ANY other
    element across the six axes (left, centre-x, right, top, centre-y, bottom), normalise by
    the canvas, and average. This is what makes a slide read as 'tidy/symmetrical' vs 'ugly'
    even when nothing overlaps or spills -- the gap our overlap/off-canvas metrics missed.
    """
    n = len(rects)
    if n < 2:
        return 1.0
    xs = [(r[0], (r[0] + r[2]) / 2.0, r[2]) for r in rects]
    ys = [(r[1], (r[1] + r[3]) / 2.0, r[3]) for r in rects]
    total = 0.0
    for i in range(n):
        best = None
        for j in range(n):
            if i == j:
                continue
            dx = min(abs(xs[i][k] - xs[j][k]) for k in range(3)) / max(1.0, sw)
            dy = min(abs(ys[i][k] - ys[j][k]) for k in range(3)) / max(1.0, sh)
            d = min(dx, dy)
            best = d if best is None else min(best, d)
        total += best if best is not None else 0.0
    mean_mis = total / n
    # 0 misalignment -> 1.0 ; 5% of canvas average misalignment -> ~0.0
    return max(0.0, min(1.0, 1.0 - mean_mis / 0.05))


def _authored_boxes(slide) -> List[Rect]:
    """The shape boxes as the author placed them, with no ink-extent substitution.
    Used by alignment, which scores the layout grid rather than glyph placement."""
    out = []
    for shape in slide.shapes:
        try:
            l, t = shape.left, shape.top
            w, h = shape.width, shape.height
            if None in (l, t, w, h):
                continue
            out.append((float(l), float(t), float(l + w), float(t + h)))
        except Exception:
            continue
    return out


def _picture_distortion_score(slide) -> float:
    """Images stretched away from their native aspect ratio look bad. 1.0 = undistorted."""
    vals = []
    try:
        for sh in slide.shapes:
            try:
                if not hasattr(sh, "image"):
                    continue
                iw, ih = sh.image.size
                if not iw or not ih or not sh.width or not sh.height:
                    continue
                native = float(iw) / float(ih)
                shown = float(sh.width) / float(sh.height)
                r = max(native, shown) / max(1e-6, min(native, shown))
                vals.append(max(0.0, min(1.0, 1.0 - (r - 1.0))))
            except Exception:
                continue
    except Exception:
        return 1.0
    return sum(vals) / len(vals) if vals else 1.0


# --- colour contrast (render-free) ------------------------------------------
# Text the same colour as what is behind it renders as a BLANK SLIDE to a human even
# though the geometry is perfect (hand-audit ID 2fe501). WCAG relative luminance gives a
# contrast ratio; 4.5:1 (AA for body text) maps to 1.0. Only runs that set an explicit RGB
# are judged -- theme-inherited colour is unknowable render-free and scores neutral 1.0,
# so we never punish a deck for something we cannot see.
def _rel_lum(rgb):
    def f(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * f(rgb[0]) + 0.7152 * f(rgb[1]) + 0.0722 * f(rgb[2])


def _contrast_ratio(a, b):
    la, lb = _rel_lum(a), _rel_lum(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def _solid_rgb(fmt):
    try:
        if fmt.type is not None and str(fmt.type).startswith("SOLID"):
            c = fmt.fore_color
            if c.type is not None and str(c.type).startswith("RGB"):
                return tuple(c.rgb)
    except Exception:
        pass
    return None


def _contrast_score(slide, sw, sh):
    """1.0 = readable (>=4.5:1). Neutral 1.0 when no explicit colours are set."""
    vals = []
    try:
        bg = (255, 255, 255)
        pic_bg = False
        for shp in slide.shapes:
            try:
                if shp.width and shp.height and (shp.width * shp.height) / float(sw * sh) > BACKGROUND_COVER:
                    if "PICTURE" in str(shp.shape_type):
                        # A full-bleed PHOTO background: its pixel colours are unknowable
                        # render-free, so judging text contrast against an assumed white
                        # canvas would wrongly flag correct white-on-photo design and push
                        # the model away from background images (hand-audit ID 9f2013).
                        pic_bg = True
                        continue
                    c = _solid_rgb(shp.fill)
                    if c:
                        bg = c
            except Exception:
                pass
        if pic_bg:
            return 1.0  # neutral: cannot judge, so never punish
        for shp in slide.shapes:
            try:
                if not shp.has_text_frame or not (shp.text_frame.text or "").strip():
                    continue
                local = _solid_rgb(shp.fill) or bg
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        if not (run.text or "").strip():
                            continue
                        try:
                            fc = run.font.color
                            rgb = tuple(fc.rgb) if (fc and fc.type is not None
                                                    and str(fc.type).startswith("RGB")) else None
                        except Exception:
                            rgb = None
                        if rgb is None:
                            continue
                        vals.append(min(1.0, max(0.0, (_contrast_ratio(rgb, local) - 1.0) / 3.5)))
            except Exception:
                continue
    except Exception:
        return 1.0
    return sum(vals) / len(vals) if vals else 1.0


def _blank_table_rows(slide):
    """Fraction of table rows that are entirely empty (hand-audit ID 3c1ba8). Folded into
    density as wasted space rather than given its own weight: tables appear on ~1% of decks,
    so a standalone metric would be constant -> zero gradient -> pure dilution."""
    tot = blank = 0
    try:
        for shp in slide.shapes:
            if not getattr(shp, "has_table", False):
                continue
            for row in shp.table.rows:
                tot += 1
                if all(not c.text.strip() for c in row.cells):
                    blank += 1
    except Exception:
        return 0.0
    return (blank / float(tot)) if tot else 0.0


def _imbalance_score(rects: List[Rect], total_area: float, sw: float, sh: float) -> float:
    wx = 0.0
    wy = 0.0
    for r in rects:
        a = _area(r)
        wx += a * (r[0] + r[2]) / 2.0
        wy += a * (r[1] + r[3]) / 2.0
    cx = (wx / total_area) / sw
    cy = (wy / total_area) / sh
    d = math.hypot((cx - 0.5) / X_TOL, (cy - 0.5) / Y_TOL)
    return 1.0 / (1.0 + d)


def _score_slide(slide, sw: float, sh: float) -> Dict[str, float]:
    rects = _slide_rects(slide)
    _CANVAS_AREA[:] = [float(sw) * float(sh)]
    total_area = sum(_area(r) for r in rects)
    if not rects or total_area <= 0:
        # Deliberate: a blank slide is NOT a perfect slide. See module docstring.
        return {k: 0.0 for k in METRIC_KEYS}
    return {
        "collision": _collision_score(rects, total_area),
        "overflow": max(0.0, min(1.0, _overflow_score(rects, total_area, sw, sh)
                                 - _text_offcanvas_penalty(slide, sw, sh))),
        "imbalance": _imbalance_score(rects, total_area, sw, sh),
        "textfit": _textfit_score(slide),
        "contrast": _contrast_score(slide, sw, sh),
        "picfit": _picture_distortion_score(slide),
        # ALIGNMENT must see the AUTHORED boxes, not the ink rects. It measures the design
        # grid -- whether the author put shapes on shared edges/centre-lines. Feeding it
        # _text_extent_rect output measures the ragged edges of rendered glyphs instead,
        # which destroys the signal: a textbook two-column layout (shared lefts 0.5/0.5,
        # shared tops 2.0/2.0) scored 1.000 on boxes and 0.000 on ink rects.
        # Found by hand-audit: technology/slide_7, alignment 0.00 on a tidy layout.
        # collision/overflow keep the ink rects -- for those, where the glyphs land IS the
        # right question.
        "alignment": _alignment_score(_authored_boxes(slide), sw, sh),
    }


# --------------------------------------------------------------------------- #
# public API
# --------------------------------------------------------------------------- #
def score_shapes(pptx_path: str) -> dict:
    """Score slide-layout geometry of a .pptx file.

    Returns {'collision': float, 'overflow': float, 'imbalance': float},
    each in [0.0, 1.0] where 1.0 = good.  Averaged over all slides; a deck with
    no slides (or only blank slides) scores 0.0 on every metric.

    Raises whatever python-pptx raises if the file cannot be opened.
    """
    prs = Presentation(pptx_path)  # let a bad file raise: caller owns the gate

    sw = prs.slide_width
    sh = prs.slide_height
    sw = float(sw) if sw else float(DEFAULT_SLIDE_W)
    sh = float(sh) if sh else float(DEFAULT_SLIDE_H)
    if sw <= 0:
        sw = float(DEFAULT_SLIDE_W)
    if sh <= 0:
        sh = float(DEFAULT_SLIDE_H)

    try:
        slides = list(prs.slides)
    except Exception:
        slides = []

    if not slides:
        return {k: 0.0 for k in METRIC_KEYS}

    totals = {k: 0.0 for k in METRIC_KEYS}
    for slide in slides:
        try:
            s = _score_slide(slide, sw, sh)
        except Exception:
            s = {k: 0.0 for k in METRIC_KEYS}  # never raise on a valid pptx
        for k in METRIC_KEYS:
            totals[k] += s[k]

    n = float(len(slides))
    return {k: max(0.0, min(1.0, totals[k] / n)) for k in METRIC_KEYS}


if __name__ == "__main__":  # tiny CLI: python geom_shapes.py deck1.pptx deck2.pptx
    import sys

    for path in sys.argv[1:]:
        try:
            print(path, score_shapes(path))
        except Exception as exc:  # noqa: BLE001
            print(path, "ERROR:", type(exc).__name__, exc)
