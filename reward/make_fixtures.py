"""Adversarial .pptx fixtures for gating the agentic reward function.

Builds three 16:9 decks and rasterizes each to PNG:

    good   -- tidy, well-composed, moderate text, nothing overlapping or
              off-slide. This is what a *correct* reward function must rank #1.
    padded -- THE REWARD-HACK. Contains every single string that "good"
              contains (so any text-coverage / keyword-recall metric scores it
              >= good) but buries it under rambling paragraphs, extra
              textboxes, heavy overlap, and boxes that spill off the slide.
              A layout-blind reward will prefer this deck; a good one must not.
    empty  -- structurally valid deck whose slides carry no meaningful content
              (one truly blank slide, one holding a single empty textbox).
              Guards the floor of the metric.

Rendering reuses the exact mechanism already used by
/home/ubuntu/powerbench/agentic/pptx_tools.py (PptxRenderSlides):
    soffice --headless --convert-to pdf --outdir <dir> <deck>
    pdftoppm -png -r <dpi> <pdf> <dir>/slide
so fixture images look identical to what the policy model sees at train time.
This module only READS that file's approach; it imports and modifies nothing.

Usage:
    from make_fixtures import build_fixtures
    fx = build_fixtures()            # -> {"good": {...}, "padded": {...}, "empty": {...}}
    fx["padded"]["pptx"]             # path to the .pptx
    fx["padded"]["pngs"]             # list of rendered PNG paths (may be [] if
                                     #   LibreOffice/poppler are unavailable)
"""

from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- constants --
# 16:9 at 13.333in x 7.5in, the standard widescreen PowerPoint canvas.
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000

RENDER_DPI = int(os.environ.get("FIXTURE_RENDER_DPI", "80"))  # match pptx_tools default
RENDER_TIMEOUT = int(os.environ.get("FIXTURE_RENDER_TIMEOUT", "180"))

# The shared corpus. "good" says exactly this and no more; "padded" says all of
# it too, plus a landfill of extra prose. Keeping one source of truth is what
# makes the text-coverage trap airtight.
GOOD_SLIDES = [
    {
        "title": "Quarterly Platform Review",
        "subtitle": "Engineering summary - Q3 2026",
        "bullets": [
            "Latency down 18% after the caching rollout",
            "Error budget spent: 41% of the quarterly allowance",
            "Two regions migrated to the new scheduler",
            "Headcount steady, with one open SRE role",
        ],
    },
    {
        "title": "Next Steps",
        "subtitle": "Owners and dates on the following page",
        "bullets": [
            "Finish the scheduler migration in EU-West",
            "Cut p99 latency by a further 10%",
            "Close the open SRE hire by mid-August",
            "Publish the postmortem for the June incident",
        ],
    },
]

# Filler used only by "padded": long, low-information, rambling.
FILLER = [
    "In terms of the overall strategic direction that we have been pursuing over "
    "the course of the last several quarters, it is worth noting that the platform "
    "team has continued to iterate on a broad range of initiatives, many of which "
    "are ongoing and some of which have already begun to show early indications of "
    "the kind of progress that leadership has been asking about in recent reviews.",
    "Additionally, and this is something that came up repeatedly in the working "
    "sessions, there is a shared understanding across the organization that the "
    "measurement methodology itself deserves further scrutiny, because a number of "
    "the dashboards currently in use were assembled at different times by different "
    "people using slightly different definitions of the underlying quantities.",
    "It should also be mentioned, for completeness, that the migration work streams "
    "remain coupled to the scheduler rewrite, and therefore any slippage in one is "
    "likely to propagate to the other, which is a risk we are tracking closely and "
    "will continue to track closely through the remainder of the fiscal year and "
    "quite possibly into the one that follows it as well.",
    "Finally, we want to reiterate the points raised above regarding latency, the "
    "error budget, the regional migrations, and headcount, all of which are covered "
    "in more detail elsewhere in this document and in the appendix materials that "
    "accompany it, and all of which will be revisited at the next review cycle.",
]


# ------------------------------------------------------------------ helpers --
def _new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    return prs


def _blank_slide(prs: Presentation):
    """Layout 6 is 'Blank' -- no placeholders, so shape counts stay honest."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, lines, size=18, bold=False,
             color=(0x20, 0x20, 0x20), wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
    return box


# ------------------------------------------------------------ deck builders --
def build_good(path: str) -> str:
    """Tidy 2-slide deck: title up top, one body block, generous margins."""
    # v2: the old "good" fixture was title + subtitle + bullets on white, which measures
    # detail_cov 0.101 -- SPARSER than our own model's average of 0.150 and far below the
    # 0.264 of real human decks. It was built to satisfy a geometry-only grader, so under a
    # content-aware reward it scored 0.571 and failed its own "good stays high" check.
    # A fixture that anchors "good" must look like what we actually want the model to make:
    # a two-column layout with a supporting visual and a filled callout.
    from pptx.enum.shapes import MSO_SHAPE
    prs = _new_deck()
    for spec in GOOD_SLIDES:
        s = _blank_slide(prs)
        _textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.0),
                 [spec["title"]], size=36, bold=True, color=(0x11, 0x2B, 0x4A))
        _textbox(s, Inches(0.9), Inches(1.62), Inches(11.5), Inches(0.45),
                 [spec["subtitle"]], size=17, color=(0x5A, 0x66, 0x72))
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.20),
                                  Inches(2.4), Inches(0.06))
        rule.fill.solid(); rule.fill.fore_color.rgb = RGBColor(0x1F, 0x6F, 0xEB)
        rule.line.fill.background()
        _textbox(s, Inches(0.9), Inches(2.55), Inches(6.6), Inches(3.9),
                 ["•  " + b for b in spec["bullets"]], size=19)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.95), Inches(2.55),
                                  Inches(4.45), Inches(3.9))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xEF, 0xF3, 0xF8)
        card.line.fill.background()
        _textbox(s, Inches(8.30), Inches(2.85), Inches(3.75), Inches(0.5),
                 ["In practice"], size=18, bold=True, color=(0x11, 0x2B, 0x4A))
        _textbox(s, Inches(8.30), Inches(3.45), Inches(3.75), Inches(2.8),
                 [spec["subtitle"], "", "•  " + spec["bullets"][0],
                  "•  " + spec["bullets"][-1]], size=15, color=(0x30, 0x3A, 0x46))
    prs.save(path)
    return path


def build_padded(path: str) -> str:
    """The adversarial deck: a superset of good's text, laid out horribly.

    Every string from GOOD_SLIDES appears verbatim, so token/keyword coverage
    can only go UP relative to good. What degrades is purely geometric --
    overlap, off-canvas bleed, edge-to-edge crowding, text volume.
    """
    prs = _new_deck()
    for spec in GOOD_SLIDES:
        s = _blank_slide(prs)

        # 1. All of good's text, but jammed against the very top-left corner.
        _textbox(s, Inches(0.02), Inches(0.01), Inches(13.3), Inches(0.9),
                 [spec["title"] + " " + spec["subtitle"]], size=34, bold=True)
        _textbox(s, Inches(0.02), Inches(0.75), Inches(13.4), Inches(2.2),
                 ["• " + b for b in spec["bullets"]], size=19)

        # 2. Rambling paragraphs stacked on top of the bullets (overlap).
        _textbox(s, Inches(0.05), Inches(1.30), Inches(7.2), Inches(3.2),
                 FILLER[:2], size=14)
        _textbox(s, Inches(5.90), Inches(1.60), Inches(7.2), Inches(3.2),
                 FILLER[2:], size=14)

        # 3. Boxes that spill off the right and bottom edges of the canvas.
        _textbox(s, Inches(9.4), Inches(4.3), Inches(6.5), Inches(2.0),
                 [FILLER[0]], size=13, wrap=True)
        _textbox(s, Inches(0.6), Inches(6.6), Inches(9.0), Inches(2.4),
                 [FILLER[1]], size=13)
        _textbox(s, Inches(-1.4), Inches(3.1), Inches(5.5), Inches(1.8),
                 [FILLER[3]], size=13)

        # 4. A pile of small crammed boxes, deliberately on top of each other,
        #    each repeating the bullets so text coverage climbs further.
        for i, b in enumerate(spec["bullets"]):
            _textbox(s,
                     Inches(1.1 + 1.55 * i), Inches(3.9 + 0.42 * i),
                     Inches(4.2), Inches(0.9),
                     [b, b.upper()], size=16, wrap=False)

        # 5. Edge-to-edge banner clipped by the bottom of the slide.
        _textbox(s, Inches(0.0), Inches(7.15), Inches(13.4), Inches(1.2),
                 [" / ".join(spec["bullets"])], size=15, wrap=False)

        # 6. Dead-centre pile: three more boxes straight over the middle of
        #    everything already placed, plus one hanging off the top edge.
        _textbox(s, Inches(3.6), Inches(2.55), Inches(6.6), Inches(1.7),
                 [spec["title"], FILLER[1]], size=15)
        _textbox(s, Inches(4.4), Inches(3.05), Inches(6.6), Inches(1.7),
                 [spec["subtitle"], FILLER[2]], size=15)
        _textbox(s, Inches(2.9), Inches(5.35), Inches(8.4), Inches(1.6),
                 [FILLER[3], " ".join(spec["bullets"])], size=14)
        _textbox(s, Inches(7.8), Inches(-0.55), Inches(5.4), Inches(1.5),
                 [FILLER[0]], size=13)

    prs.save(path)
    return path


def build_empty(path: str) -> str:
    """Valid deck, zero meaningful content: one blank slide, one empty textbox."""
    prs = _new_deck()
    _blank_slide(prs)  # slide 1: genuinely nothing on it
    s2 = _blank_slide(prs)
    box = s2.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(8.0), Inches(1.5))
    box.text_frame.text = ""  # slide 2: a single, empty textbox
    prs.save(path)
    return path


# --------------------------------------------------------------- rendering --
def _have_renderer() -> bool:
    return bool(shutil.which("soffice")) and bool(shutil.which("pdftoppm"))


def render_pptx(pptx_path: str, outdir: str) -> List[str]:
    """pptx -> PNG per slide, using the same commands as pptx_tools.PptxRenderSlides.

    Returns [] (never raises) if the toolchain is missing or conversion fails.
    A private LibreOffice profile dir is used so concurrent soffice runs by
    other agents on this box cannot collide on the shared ~/.config profile
    lock; the binary and conversion flags are otherwise identical.
    """
    if not _have_renderer():
        return []
    os.makedirs(outdir, exist_ok=True)
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    for stale in os.listdir(outdir):
        if stale.endswith(".png") or stale.endswith(".pdf"):
            os.remove(os.path.join(outdir, stale))
    profile = tempfile.mkdtemp(prefix="lo_profile_")
    try:
        subprocess.run(
            ["soffice", f"-env:UserInstallation=file://{profile}",
             "--headless", "--convert-to", "pdf", "--outdir", outdir, pptx_path],
            capture_output=True, text=True, timeout=RENDER_TIMEOUT, check=True,
        )
        pdf = os.path.join(outdir, stem + ".pdf")
        if not os.path.isfile(pdf):
            return []
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(RENDER_DPI), pdf,
             os.path.join(outdir, "slide")],
            capture_output=True, text=True, timeout=RENDER_TIMEOUT, check=True,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
        return []
    finally:
        shutil.rmtree(profile, ignore_errors=True)
    return sorted(
        os.path.join(outdir, p) for p in os.listdir(outdir)
        if p.startswith("slide") and p.endswith(".png")
    )


# -------------------------------------------------------------- public API --
BUILDERS = {"good": build_good, "padded": build_padded, "empty": build_empty}


def build_fixtures(outdir: str = "/tmp/reward_fixtures") -> Dict[str, dict]:
    """Build the 3 fixture decks and render them.

    Returns {name: {"pptx": <path>, "pngs": [<path>, ...]}} for the names
    "good", "padded", "empty". "pngs" is [] when rendering is unavailable.
    """
    os.makedirs(outdir, exist_ok=True)
    out: Dict[str, dict] = {}
    for name, builder in BUILDERS.items():
        pptx_path = os.path.join(outdir, f"{name}.pptx")
        builder(pptx_path)
        pngs = render_pptx(pptx_path, os.path.join(outdir, f"{name}_png"))
        out[name] = {"pptx": pptx_path, "pngs": pngs}
    return out


def deck_stats(pptx_path: str) -> dict:
    """Slide/shape/character counts -- used by the self-test and handy for
    sanity-checking that a metric is looking at the deck we think it is."""
    prs = Presentation(pptx_path)
    slides = len(prs.slides)
    shapes = chars = 0
    for s in prs.slides:
        for sh in s.shapes:
            shapes += 1
            if sh.has_text_frame:
                chars += len(sh.text_frame.text)
    return {"slides": slides, "shapes": shapes, "chars": chars}


# ------------------------------------------------------------- self-test --
if __name__ == "__main__":
    print(f"renderer available: soffice={bool(shutil.which('soffice'))} "
          f"pdftoppm={bool(shutil.which('pdftoppm'))} dpi={RENDER_DPI}")
    fx = build_fixtures()
    print()
    hdr = f"{'deck':<8}{'pptx bytes':>12}{'slides':>8}{'shapes':>8}{'chars':>8}{'pngs':>6}"
    print(hdr)
    print("-" * len(hdr))
    stats = {}
    for name in ("good", "padded", "empty"):
        p = fx[name]["pptx"]
        assert os.path.isfile(p), f"MISSING {p}"
        st = deck_stats(p)
        stats[name] = st
        print(f"{name:<8}{os.path.getsize(p):>12}{st['slides']:>8}"
              f"{st['shapes']:>8}{st['chars']:>8}{len(fx[name]['pngs']):>6}")

    print()
    for name in ("good", "padded", "empty"):
        for png in fx[name]["pngs"]:
            print(f"  {name}: {png}  ({os.path.getsize(png)} bytes)")

    print()
    g, pd_, e = stats["good"], stats["padded"], stats["empty"]
    print(f"padded/good shape ratio : {pd_['shapes'] / max(1, g['shapes']):.1f}x")
    print(f"padded/good char  ratio : {pd_['chars'] / max(1, g['chars']):.1f}x")
    print(f"empty chars             : {e['chars']} (expect 0)")

    # The trap must actually be a trap: padded has to be a strict text superset
    # of good, otherwise a coverage metric could separate them for free.
    good_text = " ".join(
        sh.text_frame.text for s in Presentation(fx["good"]["pptx"]).slides
        for sh in s.shapes if sh.has_text_frame)
    padded_text = " ".join(
        sh.text_frame.text for s in Presentation(fx["padded"]["pptx"]).slides
        for sh in s.shapes if sh.has_text_frame)
    missing = [b for spec in GOOD_SLIDES
               for b in [spec["title"], spec["subtitle"], *spec["bullets"]]
               if b not in padded_text]
    gw, pw = set(good_text.split()), set(padded_text.split())
    print(f"good strings missing from padded : {len(missing)} (expect 0)")
    print(f"good vocab covered by padded     : {len(gw & pw)}/{len(gw)} "
          f"({100 * len(gw & pw) / max(1, len(gw)):.0f}%)")
    assert not missing, missing
    assert pd_["shapes"] >= 4 * g["shapes"], "padded is not crowded enough"
    assert pd_["chars"] >= 5 * g["chars"], "padded is not text-heavy enough"
    assert e["chars"] == 0, "empty deck has text"
    print("\nSELF-TEST PASSED")
