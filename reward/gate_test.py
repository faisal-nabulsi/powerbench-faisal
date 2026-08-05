"""
THE ADVERSARIAL GATE.

Go/no-go before a single GPU-hour is spent on the next training run.

REQUIREMENT: the padded deck (text-rich but visually terrible — exactly the artefact our
last policy learned to produce) must score STRICTLY WORSE than the good deck.

We also score both decks with the OLD text-coverage proxy, to show the contrast: the old
reward rates the padded deck >= the good deck (which is precisely why the model padded),
while the new geometric reward must rate it lower.

Exit code 0 = gate passed, 1 = gate failed.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation  # noqa: E402

from geometric_reward import score_deck  # noqa: E402
from make_fixtures import build_fixtures  # noqa: E402

def deck_text(pptx_path):
    """All visible text in a deck, lowercased."""
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return ""
    return " ".join(
        run.text
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
        for para in sh.text_frame.paragraphs
        for run in para.runs
    ).lower()


def derive_required(good_pptx):
    """Use the GOOD deck's own content words as the 'required content'.

    Word-level, not phrase-level: phrase matching is brittle to whitespace/bullet
    differences, and word coverage is a faithful reproduction of how a text-coverage
    reward actually behaves. Derived from the fixture rather than hardcoded so this stays
    correct whatever wording the fixture builder chose.
    """
    words = [w.strip(".,:;•-()").lower() for w in deck_text(good_pptx).split()]
    return sorted({w for w in words if len(w) >= 4})


def old_coverage_reward(pptx_path, required):
    """Reproduction of the OLD, gameable reward: fraction of required strings present.

    This is presence-based: it can only go UP as more text is added. It has no way to
    notice that the slide has become an unreadable mess.
    """
    if not required:
        return 0.0
    words = {w.strip(".,:;•-()").lower() for w in deck_text(pptx_path).split()}
    hits = sum(1 for q in required if q in words)
    return hits / float(len(required))


def main():
    print("=" * 74)
    print("ADVERSARIAL REWARD GATE")
    print("=" * 74)

    print("\nBuilding fixtures ...")
    fx = build_fixtures()

    if "good" not in fx:
        print("MISSING FIXTURE: good")
        return 1
    required = derive_required(fx["good"].get("pptx"))
    print("Derived %d required phrases from the good deck (for the old-reward contrast)."
          % len(required))

    rows = []
    for name in ("good", "padded", "empty", "empty_chart", "sparse"):
        if name not in fx:
            print("MISSING FIXTURE: %s" % name)
            return 1
        pptx = fx[name].get("pptx")
        pngs = fx[name].get("pngs") or []
        res = score_deck(pptx, pngs)
        old = old_coverage_reward(pptx, required)
        nwords = len(deck_text(pptx).split())
        rows.append((name, res, old, len(pngs), nwords))

    # ---- report ---------------------------------------------------------
    print("\n%-8s | %-6s | %-9s | %-9s | %s"
          % ("deck", "words", "NEW geo", "OLD cover", "per-metric (new)"))
    print("-" * 84)
    for name, res, old, npng, nwords in rows:
        m = res["metrics"]
        mstr = " ".join("%s=%.2f" % (k[:4], v) for k, v in sorted(m.items()))
        print("%-8s | %6d | %9.4f | %9.4f | %s%s"
              % (name, nwords, res["score"], old, mstr, "" if npng else "  (no renders)"))

    scores = {name: res["score"] for name, res, _, _, _ in rows}
    old_scores = {name: old for name, _, old, _, _ in rows}

    # ---- the gate -------------------------------------------------------
    print("\n" + "=" * 74)
    print("VERDICT")
    print("=" * 74)

    ok = True

    # 1. THE gate: padded must be strictly worse than good.
    if scores["padded"] < scores["good"]:
        print("PASS  padded (%.4f) < good (%.4f)  -- padding is punished, margin %.4f"
              % (scores["padded"], scores["good"], scores["good"] - scores["padded"]))
    else:
        ok = False
        print("FAIL  padded (%.4f) >= good (%.4f)  -- reward is STILL gameable by padding"
              % (scores["padded"], scores["good"]))

    # 2. Empty must not look good (blank-deck hole).
    if scores["empty"] < scores["good"]:
        print("PASS  empty (%.4f) < good (%.4f)  -- blank decks are not a free win"
              % (scores["empty"], scores["good"]))
    else:
        ok = False
        print("FAIL  empty (%.4f) >= good (%.4f)  -- degenerate blank deck scores well"
              % (scores["empty"], scores["good"]))

    # 2a. Empty-chart must not look good (the detail_cov trap from item_30).
    # A chart's axes/gridlines/legend are structured ink; if the content term can be
    # farmed by chart chrome with no data, this passes the gate but fails a viewer.
    if scores["empty_chart"] < scores["good"]:
        print("PASS  empty_chart (%.4f) < good (%.4f)  -- chart chrome is not free content"
              % (scores["empty_chart"], scores["good"]))
    else:
        ok = False
        print("FAIL  empty_chart (%.4f) >= good (%.4f)  -- detail_cov fooled by empty-chart "
              "chrome; content term needs a real data/fill signal"
              % (scores["empty_chart"], scores["good"]))

    # 2b. Sparse monochrome slide must land below good (item_25/35 blind spot).
    if scores["sparse"] < scores["good"]:
        print("PASS  sparse (%.4f) < good (%.4f)  -- a near-empty slide is not a free win"
              % (scores["sparse"], scores["good"]))
    else:
        ok = False
        print("FAIL  sparse (%.4f) >= good (%.4f)  -- near-empty monochrome slide scores well"
              % (scores["sparse"], scores["good"]))

    # 2c. Instruction-adherence: an on-topic deck must beat an off-topic one of IDENTICAL form.
    # This is the hole a form-only reward cannot see -- both decks are well-composed; they
    # differ only in whether the title carries the task's required text.
    REQ = ["Quarterly Platform Review"]
    on_res = score_deck(fx["ontopic"].get("pptx"), fx["ontopic"].get("pngs") or [], required_texts=REQ)
    off_res = score_deck(fx["offtopic"].get("pptx"), fx["offtopic"].get("pngs") or [], required_texts=REQ)
    on_s, off_s = on_res["score"], off_res["score"]
    off_adh = off_res["metrics"].get("adherence")
    print("\nAdherence check (required_texts=%s):" % REQ)
    print("  on-topic  score %.4f  adherence %.2f" % (on_s, on_res["metrics"].get("adherence", -1)))
    print("  off-topic score %.4f  adherence %.2f" % (off_s, off_adh if off_adh is not None else -1))
    if on_s > off_s and off_adh == 0.0:
        print("PASS  on-topic (%.4f) > off-topic (%.4f), off-topic adherence 0.0  -- the reward "
              "now sees instruction-adherence, margin %.4f" % (on_s, off_s, on_s - off_s))
    else:
        ok = False
        print("FAIL  on-topic (%.4f) vs off-topic (%.4f), off adherence %s  -- adherence signal "
              "not working" % (on_s, off_s, off_adh))

    # 3. No negative rewards anywhere (the collapse cause).
    if all(s >= 0.0 for s in scores.values()):
        print("PASS  all rewards >= 0 (invalid scores 0, never negative)")
    else:
        ok = False
        print("FAIL  a negative reward exists -- this is what collapsed the last run")

    # 4. Contrast with the old reward (diagnostic, not a gate).
    print("\nContrast — why the old reward failed:")
    print("  OLD coverage:  good=%.4f  padded=%.4f   -> padding rated %s"
          % (old_scores["good"], old_scores["padded"],
             "BETTER or EQUAL (gameable!)" if old_scores["padded"] >= old_scores["good"]
             else "worse"))
    print("  NEW geometric: good=%.4f  padded=%.4f   -> padding rated %s"
          % (scores["good"], scores["padded"],
             "worse (correct)" if scores["padded"] < scores["good"] else "BETTER (broken!)"))

    print("\n%s" % ("GATE PASSED — safe to spend GPU-hours." if ok
                    else "GATE FAILED — do NOT train against this reward yet."))
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
