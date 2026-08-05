"""
Single-turn PPTX reward: model emits python-pptx code -> we execute it -> render it ->
score it GEOMETRICALLY.

This is the verl reward hook for the single-turn hill-climb. It deliberately contains NO
text-coverage component. Reward comes only from the geometry of the produced slide, so the
padding strategy that collapsed the previous run cannot pay off here.

Reward scale (all in [0,1], higher is better):
    0.0            code missing / doesn't run / no deck produced   <- NEVER negative
    score_deck()   otherwise: mean of collision, overflow, imbalance (+ whitespace if the
                   deck rendered)

Why invalid == 0.0 and not a negative penalty: a negative floor makes every rollout in a
GRPO group score the same, which gives zero relative advantage and therefore zero gradient
(arXiv 2601.03525). That is precisely how the previous run died.
"""

import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import uuid

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from geometric_reward import score_deck, _parse_required  # noqa: E402

# Persist a sample of real rollout decks + renders + scores so we can SHOW what the model
# actually produced (the rollout tmpdir is otherwise deleted after scoring).
# NOTE: this reward runs inside Ray training workers that do NOT inherit the launcher's env,
# so we DEFAULT to a fixed on-box path rather than relying on GEO_GALLERY being propagated.
GALLERY_DIR = os.environ.get("GEO_GALLERY", "") or "/home/ubuntu/powerbench/agentic/gallery"
GALLERY_CAP = int(os.environ.get("GEO_GALLERY_CAP", "24"))    # max saved per worker process
GALLERY_EVERY = int(os.environ.get("GEO_GALLERY_EVERY", "6"))  # sample 1 in N (span the run)
_gallery_saved = [0]
_rollout_seen = [0]


def _save_gallery(deck, res, solution_str="", extra_info=None):
    """Persist a per-slide record for human review: the .pptx, the exact code the model
    wrote, and rich meta (reward, per-metric, validity, response token estimate, task).
    Sampled 1-in-N and capped so it spans the run without filling the disk. Never raises."""
    _rollout_seen[0] += 1
    if not GALLERY_DIR or _gallery_saved[0] >= GALLERY_CAP:
        return
    if _rollout_seen[0] % GALLERY_EVERY != 0:
        return
    try:
        os.makedirs(GALLERY_DIR, exist_ok=True)
        ei = extra_info or {}
        code = extract_code(solution_str or "")
        tag = "p%d_n%03d_s%03d_%s" % (
            os.getpid() % 100000, _gallery_saved[0],
            int(round(float(res.get("score", 0.0)) * 100)), uuid.uuid4().hex[:6])
        shutil.copy(deck, os.path.join(GALLERY_DIR, "deck_%s.pptx" % tag))
        with open(os.path.join(GALLERY_DIR, "code_%s.py" % tag), "w") as f:
            f.write(code)
        with open(os.path.join(GALLERY_DIR, "meta_%s.json" % tag), "w") as f:
            json.dump({
                "score": res.get("score"), "metrics": res.get("metrics"),
                "valid": res.get("valid"),
                "resp_chars": len(solution_str or ""),
                "resp_tokens_est": round(len(solution_str or "") / 4),
                "code_chars": len(code),
                "task": ei.get("task"), "index": ei.get("index"),
                "n_required": ei.get("n_required"),
                "seen": _rollout_seen[0],
            }, f)
        _gallery_saved[0] += 1
    except Exception:
        pass

EXEC_TIMEOUT = 12        # seconds to run the model's code. Slide code runs in <5s; a long
                         # cap just lets pathological/looping generated code eat minutes
                         # (64 rollouts x 60s was a big chunk of the 83-min steps).
RENDER = False           # render-free: density comes from python-pptx, no soffice storm.
PY = sys.executable or "/home/ubuntu/powerbench/.venv/bin/python"

_CODE_BLOCK = re.compile(r"```(?:python|py)?\s*\n(.*?)```", re.DOTALL)


def extract_code(solution_str):
    """Pull the python source out of a model response.

    The model is a reasoning model: it emits <think>...</think> and then the answer, so we
    only look AFTER the closing think tag. We take the LAST fenced block, which is the
    final answer when the model shows intermediate drafts.
    """
    if not solution_str:
        return ""
    text = solution_str
    if "</think>" in text:
        text = text.rsplit("</think>", 1)[1]
    blocks = _CODE_BLOCK.findall(text)
    if blocks:
        return blocks[-1].strip()
    # Fallback: unfenced code that at least imports pptx
    if "pptx" in text and ("import" in text or "Presentation(" in text):
        return text.strip()
    return ""


ASSETS_DIR = os.environ.get("GEO_ASSETS", "") or "/home/ubuntu/powerbench/agentic/assets"
# The public SlidesBench tasks reference images ("insert an image of X") but ship NO image
# files. We drop ONE neutral placeholder into the rollout dir under every filename the model
# is likely to use, so add_picture(...) resolves and the requested image-LAYOUT gets scored
# instead of crashing the script to reward 0. The system prompt points the model at 'image.png'.
_PLACEHOLDER_NAMES = ("image.png", "placeholder.png", "picture.png", "img.png",
                      "photo.png", "logo.png", "background.png", "bg.png", "image.jpg")


def _provision_assets(workdir):
    src = os.path.join(ASSETS_DIR, "image.png")
    if not os.path.isfile(src):
        return
    for nm in _PLACEHOLDER_NAMES:
        try:
            shutil.copy(src, os.path.join(workdir, nm))
        except Exception:
            pass


def run_code_to_deck(code, workdir):
    """Execute the model's code in workdir. Returns path to the produced .pptx or None."""
    if not code.strip():
        return None
    _provision_assets(workdir)  # placeholder image(s) so add_picture(...) resolves, not crashes
    script = os.path.join(workdir, "gen.py")
    with open(script, "w") as f:
        f.write(code)
    try:
        subprocess.run([PY, script], cwd=workdir, timeout=EXEC_TIMEOUT,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except (subprocess.TimeoutExpired, OSError):
        return None
    # Accept whatever .pptx the code produced (prefer output.pptx by convention).
    preferred = os.path.join(workdir, "output.pptx")
    if os.path.isfile(preferred):
        return preferred
    for fn in sorted(os.listdir(workdir)):
        if fn.lower().endswith(".pptx"):
            return os.path.join(workdir, fn)
    return None


def score_solution(solution_str, extra_info=None, keep_dir=None, keep_name=None, required_texts=None):
    """Full pipeline: code -> deck -> geometric score. Returns the score dict.

    keep_dir: if set, EVERY deck is copied there (not the ~1-in-6 gallery sample) and the
    training gallery is left untouched. Used by the frontier baselines so an external model's
    decks never contaminate the corpus we analyse as our own model's output -- 7 Fable decks
    silently landed in gallery/ before this existed.
    """
    code = extract_code(solution_str)
    if not code:
        return {"score": 0.0, "metrics": {}, "valid": False, "reason": "no code emitted"}

    with tempfile.TemporaryDirectory(prefix="stgeo_") as wd:
        deck = run_code_to_deck(code, wd)
        if not deck:
            return {"score": 0.0, "metrics": {}, "valid": False,
                    "reason": "code failed to produce a deck"}

        res = score_deck(deck, required_texts=required_texts)  # + instruction-adherence
        if keep_dir:
            import re as _re
            os.makedirs(keep_dir, exist_ok=True)
            base = _re.sub(r"\W+", "_", str(keep_name or (extra_info or {}).get("task") or "task"))
            shutil.copy(deck, os.path.join(keep_dir, "deck_%s.pptx" % base))
            with open(os.path.join(keep_dir, "code_%s.py" % base), "w") as _f:
                _f.write(extract_code(solution_str or ""))
            with open(os.path.join(keep_dir, "meta_%s.json" % base), "w") as _f:
                json.dump({"score": res.get("score"), "metrics": res.get("metrics"),
                           "valid": res.get("valid"), "task": keep_name or (extra_info or {}).get("task"),
                           "resp_chars": len(solution_str or "")}, _f)
        else:
            _save_gallery(deck, res, solution_str, extra_info)  # per-slide record for review
        return res


# --------------------------------------------------------------------------
# verl entry point
# --------------------------------------------------------------------------
# verl collects the extra keys from every rollout and stacks them into one array per key:
#   non_tensor_batch[key] = np.array([info[key] for info in reward_extra_infos])
# That requires EVERY sample to return EXACTLY the same set of keys. If an invalid deck
# returns fewer keys than a valid one, verl KeyErrors. So we always return this fixed
# schema, filling missing metrics with 0.0.
# v2 additions (content, clipping, chart_ok) are logged so W&B shows them per step: content
# and chart_ok are the two metrics that close the empty-deck / empty-chart holes, so a run
# that starts farming either will show it here BEFORE the held-out score reacts. chart_ok
# defaults to 1.0 (no chart present) when a deck has no metrics; the others default 0.0.
_METRIC_KEYS = ("collision", "overflow", "imbalance", "density", "aspect", "textfit",
                "contrast", "picfit", "alignment", "content", "clipping", "chart_ok",
                "adherence")
_METRIC_DEFAULT = {"chart_ok": 1.0, "adherence": 1.0}


def compute_score(data_source, solution_str, ground_truth, extra_info=None, **kwargs):
    try:
        res = score_solution(solution_str, extra_info, required_texts=_parse_required(ground_truth))
        score = float(res.get("score", 0.0))
        valid = 1.0 if res.get("valid") else 0.0
        metrics = res.get("metrics", {}) or {}
    except Exception:
        score, valid, metrics = 0.0, 0.0, {}

    out = {"score": score, "valid": valid}
    for k in _METRIC_KEYS:  # always present, consistent schema across the batch
        out["m_" + k] = float(metrics.get(k, _METRIC_DEFAULT.get(k, 0.0)))
    return out


if __name__ == "__main__":
    # Self-test: a good deck, a padded deck, and broken code.
    GOOD = '''```python
from pptx import Presentation
from pptx.util import Emu, Pt
prs = Presentation(); prs.slide_width=Emu(12192000); prs.slide_height=Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Emu(1200000), Emu(900000), Emu(9700000), Emu(1200000))
tb.text_frame.text = "Quarterly Platform Review"
body = s.shapes.add_textbox(Emu(1200000), Emu(2600000), Emu(9700000), Emu(2600000))
tf = body.text_frame
for line in ["Latency down 18%", "Error budget healthy", "Three regions live"]:
    tf.add_paragraph().text = line
prs.save("output.pptx")
```'''
    PADDED = '''```python
from pptx import Presentation
from pptx.util import Emu
prs = Presentation(); prs.slide_width=Emu(12192000); prs.slide_height=Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
blob = ("Quarterly Platform Review Latency down 18% Error budget healthy " * 60)
for i in range(14):
    tb = s.shapes.add_textbox(Emu(300000 + i*200000), Emu(200000 + i*180000),
                              Emu(9000000), Emu(2500000))
    tb.text_frame.text = blob
prs.save("output.pptx")
```'''
    BROKEN = "```python\nthis is not valid python(((\n```"
    for name, sol in (("good", GOOD), ("padded", PADDED), ("broken", BROKEN)):
        r = score_solution(sol)
        print("%-7s score=%.4f valid=%s metrics=%s reason=%s"
              % (name, r["score"], r["valid"],
                 {k: round(v, 3) for k, v in r.get("metrics", {}).items()}, r["reason"]))
